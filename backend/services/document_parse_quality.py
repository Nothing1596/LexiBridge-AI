"""Document parse quality classification and lightweight parsing helpers."""

from __future__ import annotations

import json
import os
import re
import tempfile
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from services.formula_detection import contains_formula_text, detect_pdf_formula_regions
from services.layout_analysis import (
    SKIPPED_TEXT_LAYOUT_TYPES,
    layout_blocks_to_text,
    parse_pdf_layout,
)
from services.ocr import get_ocr_provider, join_tesseract_blocks_text


QUALITY_STATUSES = {
    "native_text_ok",
    "ocr_text_ok",
    "partial_text",
    "empty_text",
    "ocr_required",
    "ocr_unavailable",
    "ocr_low_confidence",
    "formula_detected",
    "formula_ocr_required",
    "formula_ocr_unavailable",
    "unsupported_file_type",
    "parse_failed",
    "mixed_quality",
}
PARSE_STATUSES = {"success", "partial", "failed"}
SUPPORTED_TEXT_TYPES = {"txt", "md", "markdown"}
SUPPORTED_OFFICE_TYPES = {"docx", "pptx"}
SUPPORTED_IMAGE_TYPES = {"jpg", "jpeg", "png"}
SUPPORTED_TYPES = {"pdf", *SUPPORTED_TEXT_TYPES, *SUPPORTED_OFFICE_TYPES, *SUPPORTED_IMAGE_TYPES}
BLOCKING_QUALITY_STATUSES = {
    "empty_text",
    "ocr_required",
    "ocr_unavailable",
    "parse_failed",
    "unsupported_file_type",
}
OCR_PDF_DEFAULT_DPI = 220
OCR_PDF_MAX_PAGES = 50
OCR_PDF_MAX_PIXELS_PER_PAGE = 30_000_000
OCR_PDF_MAX_TOTAL_PIXELS = 160_000_000


@dataclass
class DocumentParseResult:
    parse_record_data: dict[str, Any]
    blocks: list[dict[str, Any]] = field(default_factory=list)
    raw_text: str = ""
    formula_regions: list[Any] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)


def _json_loads(value: Any, fallback: Any) -> Any:
    if value in (None, ""):
        return fallback
    if isinstance(value, (list, dict)):
        return value
    try:
        return json.loads(value)
    except (TypeError, ValueError):
        return fallback


def _json_dumps(value: Any) -> str:
    return json.dumps(value if value is not None else [], ensure_ascii=False, sort_keys=True)


def _clean_text(text: Any) -> str:
    text = str(text or "").replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _now_text(now_fn=None) -> str:
    return now_fn() if now_fn else ""


def _file_size(path: str | os.PathLike[str] | None) -> int | None:
    if not path:
        return None


def _env_int(name: str, default: int) -> int:
    try:
        return int(os.environ.get(name, str(default)))
    except (TypeError, ValueError):
        return default
    try:
        return os.path.getsize(path)
    except OSError:
        return None


def detect_file_type(filename: str, mime_type: str | None = None) -> str:
    suffix = Path(str(filename or "")).suffix.lower().lstrip(".")
    if suffix in {"txt", "md", "markdown"}:
        return suffix
    if suffix in {"pdf", "docx", "pptx", "jpg", "jpeg", "png"}:
        return "image" if suffix in SUPPORTED_IMAGE_TYPES else suffix
    mime = str(mime_type or "").lower()
    if "pdf" in mime:
        return "pdf"
    if "wordprocessingml" in mime or "msword" in mime:
        return "docx"
    if "presentationml" in mime or "powerpoint" in mime:
        return "pptx"
    if mime.startswith("text/"):
        return "txt"
    if mime.startswith("image/"):
        return "image"
    return "unknown"


def normalize_quality_flags(flags: list[Any] | tuple[Any, ...] | set[Any] | None) -> list[str]:
    normalized = []
    seen = set()
    for flag in flags or []:
        text = str(flag or "").strip()
        if not text or text in seen:
            continue
        normalized.append(text)
        seen.add(text)
    return sorted(normalized)


def _ocr_available(ocr_provider_name: str | None = None) -> bool:
    provider = get_ocr_provider(ocr_provider_name or os.environ.get("OCR_PROVIDER", "none"))
    try:
        return bool(provider.is_available())
    except Exception:
        return False


def _text_blocks_from_text(text: str, *, page_number=None, slide_number=None, parser_type="native", source_locator="") -> list[dict[str, Any]]:
    text = _clean_text(text)
    if not text:
        return []
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]
    if not paragraphs:
        paragraphs = [text]
    blocks = []
    for index, paragraph in enumerate(paragraphs, start=1):
        blocks.append({
            "block_uid": str(uuid.uuid4()),
            "page_number": page_number,
            "slide_number": slide_number,
            "block_index": index,
            "block_type": "text",
            "text": paragraph,
            "confidence": 1.0,
            "parser_type": parser_type,
            "source_locator": source_locator or f"block:{index}",
            "quality_flags": [],
        })
    return blocks


def _read_text_file(path: str) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            with open(path, "r", encoding=encoding) as file:
                return _clean_text(file.read())
        except UnicodeDecodeError:
            continue
    raise ValueError("Unable to decode text file.")


def _parse_pdf_native(path: str) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is not available for PDF parsing.") from exc

    blocks = []
    text_parts = []
    page_count = 0
    image_only_pages = 0
    image_only_page_numbers = []
    with fitz.open(path) as doc:
        page_count = len(doc)
        for page_index, page in enumerate(doc, start=1):
            page_text = _clean_text(page.get_text("text") or "")
            if page_text:
                text_parts.append(f"[Page {page_index}]\n{page_text}")
                for block in _text_blocks_from_text(
                    page_text,
                    page_number=page_index,
                    parser_type="native",
                    source_locator=f"page:{page_index}",
                ):
                    block["block_index"] = len(blocks) + 1
                    blocks.append(block)
            else:
                image_only_pages += 1
                image_only_page_numbers.append(page_index)
    return "\n\n".join(text_parts), blocks, {
        "page_count": page_count,
        "image_only_suspected": page_count > 0 and image_only_pages == page_count,
        "partial_text": 0 < image_only_pages < page_count,
        "image_only_pages": image_only_page_numbers,
    }


def _parse_pdf_with_layout(path: str) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    """Parse a PDF through the layout analysis pipeline.

    Same return contract as ``_parse_pdf_native``. Raises on layout failure so
    the caller can fall back to the legacy per-page extraction; layout analysis
    must never be the reason a parse fails.
    """
    result = parse_pdf_layout(path)
    if not result.ok:
        raise RuntimeError(result.error or f"Layout analysis failed with status {result.status}.")

    raw_text = _clean_text(layout_blocks_to_text(result.blocks))
    blocks = []
    text_pages = set()
    for block in result.blocks:
        if block.layout_type in SKIPPED_TEXT_LAYOUT_TYPES or not block.text.strip():
            continue
        text_pages.add(block.page_number)
        bbox = block.bbox
        locator = f"page:{block.page_number};bbox:{round(bbox.x0, 2)},{round(bbox.y0, 2)},{round(bbox.x1, 2)},{round(bbox.y1, 2)}"
        blocks.append({
            "block_uid": str(uuid.uuid4()),
            "page_number": block.page_number,
            "slide_number": None,
            "block_index": len(blocks) + 1,
            "block_type": block.layout_type,
            "text": _clean_text(block.text),
            "confidence": block.confidence,
            "parser_type": f"layout_{result.provider}",
            "source_locator": locator[:160],
            "quality_flags": normalize_quality_flags([
                "layout",
                f"layout_provider_{result.provider}",
                f"layout_type_{block.layout_type}",
            ]),
        })

    page_count = result.page_count
    image_only_page_numbers = [page for page in range(1, page_count + 1) if page not in text_pages]
    return raw_text, blocks, {
        "page_count": page_count,
        "image_only_suspected": page_count > 0 and len(image_only_page_numbers) == page_count,
        "partial_text": 0 < len(image_only_page_numbers) < page_count,
        "image_only_pages": image_only_page_numbers,
        "layout_provider": result.provider,
        "layout_warnings": list(result.warnings or ()),
    }


def _aggregate_bbox(words: list[dict[str, Any]]) -> tuple[int, int, int, int] | None:
    boxes = []
    for word in words:
        left = word.get("left")
        top = word.get("top")
        width = word.get("width")
        height = word.get("height")
        if None in {left, top, width, height}:
            continue
        boxes.append((int(left), int(top), int(left) + int(width), int(top) + int(height)))
    if not boxes:
        return None
    return (
        min(item[0] for item in boxes),
        min(item[1] for item in boxes),
        max(item[2] for item in boxes),
        max(item[3] for item in boxes),
    )


def _ocr_blocks_to_parse_blocks(
    ocr_result,
    *,
    page_number: int,
    block_index_start: int,
    render_config: dict[str, Any],
) -> list[dict[str, Any]]:
    word_blocks = list(getattr(ocr_result, "blocks", []) or [])
    if not word_blocks:
        return _text_blocks_from_text(
            getattr(ocr_result, "text", ""),
            page_number=page_number,
            parser_type="ocr",
            source_locator=f"page:{page_number};ocr_text",
        )
    groups: dict[tuple[int, int, int], list[dict[str, Any]]] = {}
    order: list[tuple[int, int, int]] = []
    for word in word_blocks:
        key = (
            int(word.get("block_number") or 0),
            int(word.get("paragraph_number") or 0),
            int(word.get("line_number") or 0),
        )
        if key not in groups:
            groups[key] = []
            order.append(key)
        item = dict(word)
        item["page_number"] = page_number
        groups[key].append(item)

    parse_blocks = []
    for offset, key in enumerate(order, start=0):
        words = groups[key]
        text = _clean_text(join_tesseract_blocks_text(words))
        if not text:
            continue
        bbox = _aggregate_bbox(words)
        confidences = [float(word["confidence"]) for word in words if word.get("confidence") is not None]
        confidence = sum(confidences) / len(confidences) / 100 if confidences else None
        locator = f"page:{page_number};ocr:{key[0]}.{key[1]}.{key[2]}"
        if bbox is not None:
            locator = f"{locator};bbox:{bbox[0]},{bbox[1]},{bbox[2]},{bbox[3]}"
        flags = [
            "ocr",
            "ocr_tsv_provenance",
            f"ocr_provider_{getattr(ocr_result, 'provider', '') or 'unknown'}",
            f"ocr_language_{getattr(ocr_result, 'language', '') or 'unknown'}",
            f"ocr_dpi_{render_config.get('dpi', OCR_PDF_DEFAULT_DPI)}",
        ]
        if getattr(ocr_result, "engine_version", ""):
            flags.append("ocr_engine_version_recorded")
        flags.extend(getattr(ocr_result, "quality_flags", []) or [])
        parse_blocks.append({
            "block_uid": str(uuid.uuid4()),
            "page_number": page_number,
            "slide_number": None,
            "block_index": block_index_start + offset,
            "block_type": "text",
            "text": text,
            "confidence": None if confidence is None else max(0, min(confidence, 1)),
            "parser_type": "ocr",
            "source_locator": locator[:160],
            "quality_flags": normalize_quality_flags(flags),
        })
    return parse_blocks


def _render_pdf_page_for_ocr(doc, page_number: int, temp_dir: str, *, dpi: int) -> tuple[str, dict[str, Any]]:
    page = doc[page_number - 1]
    zoom = max(1, dpi) / 72
    width = int(page.rect.width * zoom)
    height = int(page.rect.height * zoom)
    pixels = width * height
    if pixels > _env_int("OCR_PDF_MAX_PIXELS_PER_PAGE", OCR_PDF_MAX_PIXELS_PER_PAGE):
        raise RuntimeError("Rendered OCR page exceeds the safe pixel limit.")
    output = Path(temp_dir) / f"page-{page_number}.png"
    pixmap = page.get_pixmap(matrix=__import__("fitz").Matrix(zoom, zoom), alpha=False)
    pixmap.save(str(output))
    return str(output), {
        "page_number": page_number,
        "dpi": dpi,
        "width": pixmap.width,
        "height": pixmap.height,
        "pixels": pixmap.width * pixmap.height,
    }


def _ocr_pdf_image_pages(path: str, page_numbers: list[int], provider, *, language_hint: str = "bilingual") -> dict[str, Any]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is not available for PDF OCR rendering.") from exc

    dpi = max(72, min(_env_int("OCR_PDF_RENDER_DPI", OCR_PDF_DEFAULT_DPI), 300))
    max_pages = max(1, _env_int("OCR_PDF_MAX_PAGES", OCR_PDF_MAX_PAGES))
    max_total_pixels = max(1, _env_int("OCR_PDF_MAX_TOTAL_PIXELS", OCR_PDF_MAX_TOTAL_PIXELS))
    pages = [int(page) for page in page_numbers if int(page) > 0]
    if len(pages) > max_pages:
        raise RuntimeError("PDF OCR page count exceeds the safe limit.")

    raw_parts: list[str] = []
    blocks: list[dict[str, Any]] = []
    warnings: list[str] = []
    errors: list[str] = []
    total_pixels = 0
    executed = False
    with fitz.open(path) as doc, tempfile.TemporaryDirectory(prefix="lexibridge-ocr-") as temp_dir:
        for page_number in pages:
            if page_number > len(doc):
                continue
            image_path, render_config = _render_pdf_page_for_ocr(doc, page_number, temp_dir, dpi=dpi)
            total_pixels += int(render_config["pixels"])
            if total_pixels > max_total_pixels:
                raise RuntimeError("PDF OCR total rendered pixels exceed the safe limit.")
            result = provider.recognize_image(image_path, language=language_hint or "bilingual")
            executed = True
            if result.status in {"ok", "low_confidence"} and _clean_text(result.text):
                raw_parts.append(f"[Page {page_number} OCR]\n{_clean_text(result.text)}")
                blocks.extend(_ocr_blocks_to_parse_blocks(
                    result,
                    page_number=page_number,
                    block_index_start=len(blocks) + 1,
                    render_config=render_config,
                ))
                if result.status == "low_confidence":
                    warnings.append("OCR returned low confidence text.")
                warnings.extend(result.quality_flags or [])
            elif result.status == "empty_result":
                warnings.append(f"OCR produced no text for page {page_number}.")
                warnings.extend(result.quality_flags or [])
            else:
                errors.append(result.error or result.status)
                warnings.extend(result.quality_flags or [])
    return {
        "raw_text": "\n\n".join(raw_parts),
        "blocks": blocks,
        "warnings": normalize_quality_flags(warnings),
        "errors": errors,
        "executed": executed,
        "completed": bool(_clean_text("\n\n".join(raw_parts)) and blocks),
    }


def _parse_docx_native(path: str) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is not available for DOCX parsing.") from exc
    document = Document(path)
    parts = []
    for paragraph in document.paragraphs:
        content = paragraph.text.strip()
        if content:
            parts.append(content)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = _clean_text("\n\n".join(parts))
    return text, _text_blocks_from_text(text), {"page_count": None}


def _parse_pptx_native(path: str) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not available for PPTX parsing.") from exc
    presentation = Presentation(path)
    blocks = []
    text_parts = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    slide_parts.append(content)
        slide_text = _clean_text("\n".join(slide_parts))
        if slide_text:
            text_parts.append(f"[Slide {slide_index}]\n{slide_text}")
            for block in _text_blocks_from_text(
                slide_text,
                slide_number=slide_index,
                parser_type="native",
                source_locator=f"slide:{slide_index}",
            ):
                block["block_index"] = len(blocks) + 1
                blocks.append(block)
    return "\n\n".join(text_parts), blocks, {"page_count": len(presentation.slides)}


def classify_parse_quality(parse_result: dict[str, Any]) -> dict[str, Any]:
    file_type = parse_result.get("file_type", "unknown")
    raw_text = _clean_text(parse_result.get("raw_text", ""))
    blocks = parse_result.get("blocks", []) or []
    warnings = list(parse_result.get("warnings", []) or [])
    errors = list(parse_result.get("errors", []) or [])
    flags = list(parse_result.get("quality_flags", []) or [])
    ocr_required = bool(parse_result.get("ocr_required", False))
    ocr_available = bool(parse_result.get("ocr_available", False))
    ocr_completed = bool(parse_result.get("ocr_completed", False))
    formula_detected = bool(parse_result.get("formula_detected", False))
    image_only_suspected = bool(parse_result.get("image_only_suspected", False))

    if file_type == "unknown":
        quality_status = "unsupported_file_type"
        parse_status = "failed"
        flags.append("unsupported_file_type")
    elif parse_result.get("exception"):
        quality_status = "parse_failed"
        parse_status = "failed"
        flags.append("parse_failed")
    elif ocr_required and not raw_text and not blocks:
        quality_status = "ocr_unavailable" if not ocr_available else "ocr_required"
        parse_status = "failed"
        flags.append("ocr_required")
        if not ocr_available:
            flags.append("ocr_unavailable")
    elif not raw_text and not blocks:
        quality_status = "empty_text"
        parse_status = "failed"
        flags.append("empty_text")
    elif parse_result.get("partial_text") or image_only_suspected or ocr_required:
        quality_status = "ocr_text_ok" if ocr_completed else ("mixed_quality" if raw_text and ocr_required else "partial_text")
        parse_status = "success" if ocr_completed else "partial"
        if ocr_required:
            flags.append("ocr_required")
        if ocr_completed:
            flags.append("ocr_completed")
        if image_only_suspected:
            flags.append("image_only_suspected")
    else:
        quality_status = "native_text_ok"
        parse_status = "success"
        flags.append("native_text_ok")

    if formula_detected or contains_formula_text(raw_text):
        flags.append("formula_detected")
        flags.append("formula_ocr_required")
    if "formula_ocr_unavailable" in flags:
        formula_detected = True
    if "ocr_low_confidence" in flags:
        quality_status = "ocr_low_confidence"
        parse_status = "partial"

    return {
        "parse_status": parse_status,
        "quality_status": quality_status,
        "quality_flags": normalize_quality_flags(flags),
        "warnings": normalize_quality_flags(warnings),
        "errors": errors,
        "ocr_required": ocr_required,
        "ocr_available": ocr_available,
        "formula_detected": formula_detected or contains_formula_text(raw_text),
        "image_only_suspected": image_only_suspected,
    }


def build_parse_record_from_result(
    *,
    source_filename: str,
    stored_path: str = "",
    mime_type: str = "",
    file_size_bytes: int | None = None,
    parser_name: str,
    parser_version: str = "parse_quality_v1",
    file_type: str,
    raw_text: str,
    blocks: list[dict[str, Any]],
    quality: dict[str, Any],
    page_count: int | None = None,
    error_code: str = "",
    error_message: str = "",
    now_fn=None,
) -> dict[str, Any]:
    now = _now_text(now_fn)
    return {
        "parse_uid": str(uuid.uuid4()),
        "source_filename": source_filename,
        "stored_path": stored_path,
        "file_type": file_type,
        "mime_type": mime_type,
        "file_size_bytes": file_size_bytes,
        "parser_name": parser_name,
        "parser_version": parser_version,
        "parse_status": quality["parse_status"],
        "quality_status": quality["quality_status"],
        "quality_flags": quality["quality_flags"],
        "page_count": page_count,
        "block_count": len(blocks),
        "extracted_text_chars": len(_clean_text(raw_text)),
        "ocr_required": quality["ocr_required"],
        "ocr_available": quality["ocr_available"],
        "formula_detected": quality["formula_detected"],
        "table_detected": False,
        "image_only_suspected": quality["image_only_suspected"],
        "error_code": error_code,
        "error_message": error_message,
        "warnings": quality["warnings"],
        "created_at": now,
        "updated_at": now,
    }


def parse_document_with_quality(
    file_path: str,
    filename: str | None = None,
    mime_type: str | None = None,
    *,
    ocr_provider_name: str | None = None,
    language_hint: str | None = None,
    now_fn=None,
) -> DocumentParseResult:
    source_filename = filename or os.path.basename(str(file_path or ""))
    file_type = detect_file_type(source_filename, mime_type)
    raw_text = ""
    blocks: list[dict[str, Any]] = []
    formula_regions: list[Any] = []
    warnings: list[str] = []
    errors: list[str] = []
    parser_name = "document_parse_quality_v1"
    page_count = None
    image_only_suspected = False
    partial_text = False
    ocr_required = False
    ocr_available = _ocr_available(ocr_provider_name)
    ocr_language_hint = language_hint or "bilingual"
    ocr_completed = False
    exception = None
    layout_quality_flags: list[str] = []

    try:
        if file_type == "unknown":
            errors.append("Unsupported file type.")
        elif file_type in SUPPORTED_TEXT_TYPES:
            raw_text = _read_text_file(file_path)
            blocks = _text_blocks_from_text(raw_text)
            parser_name = "native_text"
        elif file_type == "pdf":
            layout_provider = ""
            if os.environ.get("LAYOUT_PROVIDER", "").strip():
                try:
                    raw_text, blocks, meta = _parse_pdf_with_layout(file_path)
                except Exception:
                    warnings.append("layout_fallback_native")
                    raw_text, blocks, meta = _parse_pdf_native(file_path)
                else:
                    layout_provider = str(meta.get("layout_provider") or "")
            else:
                raw_text, blocks, meta = _parse_pdf_native(file_path)
            if layout_provider:
                warnings.extend(meta.get("layout_warnings") or [])
                layout_quality_flags = ["layout_applied", f"layout_provider_{layout_provider}"]
            page_count = meta.get("page_count")
            image_only_suspected = bool(meta.get("image_only_suspected"))
            partial_text = bool(meta.get("partial_text"))
            ocr_required = image_only_suspected or partial_text
            try:
                formula_regions = detect_pdf_formula_regions(file_path)
            except Exception:
                warnings.append("formula_region_detection_failed")
            if ocr_required and ocr_available:
                provider = get_ocr_provider(ocr_provider_name or os.environ.get("OCR_PROVIDER", "none"))
                image_only_pages = list(meta.get("image_only_pages") or [])
                if not image_only_pages and page_count:
                    image_only_pages = list(range(1, int(page_count) + 1))
                ocr_data = _ocr_pdf_image_pages(file_path, image_only_pages, provider, language_hint=ocr_language_hint)
                if _clean_text(ocr_data.get("raw_text", "")):
                    raw_text = _clean_text("\n\n".join(part for part in [raw_text, ocr_data["raw_text"]] if part))
                    existing_count = len(blocks)
                    for index, block in enumerate(ocr_data.get("blocks", []), start=1):
                        block["block_index"] = existing_count + index
                        blocks.append(block)
                    ocr_completed = True
                warnings.extend(ocr_data.get("warnings", []) or [])
                errors.extend(ocr_data.get("errors", []) or [])
            elif ocr_required and not ocr_available:
                warnings.append("OCR is required but no OCR provider is available.")
            parser_base = f"pymupdf_layout_{layout_provider}" if layout_provider else "pymupdf_native"
            parser_name = f"{parser_base}_tesseract_ocr" if ocr_completed else parser_base
        elif file_type == "docx":
            raw_text, blocks, meta = _parse_docx_native(file_path)
            page_count = meta.get("page_count")
            parser_name = "python_docx_native"
        elif file_type == "pptx":
            raw_text, blocks, meta = _parse_pptx_native(file_path)
            page_count = meta.get("page_count")
            parser_name = "python_pptx_native"
        elif file_type == "image":
            ocr_required = True
            if not ocr_available:
                warnings.append("Image text requires OCR, but no OCR provider is available.")
            else:
                provider = get_ocr_provider(ocr_provider_name or os.environ.get("OCR_PROVIDER", "none"))
                result = provider.recognize_image(file_path, language=ocr_language_hint)
                if result.status in {"ok", "low_confidence"} and _clean_text(result.text):
                    raw_text = _clean_text(result.text)
                    confidence = max(0, min(float(result.confidence or 0) / 100, 1))
                    blocks = _ocr_blocks_to_parse_blocks(
                        result,
                        page_number=1,
                        block_index_start=1,
                        render_config={"dpi": 0},
                    ) or _text_blocks_from_text(raw_text, parser_type="ocr", source_locator="image:1")
                    for block in blocks:
                        if block.get("confidence") is None:
                            block["confidence"] = confidence
                        block["parser_type"] = "ocr"
                    ocr_completed = True
                    if result.status == "low_confidence":
                        warnings.append("OCR returned low confidence text.")
                else:
                    errors.append(result.error or result.status)
                    warnings.extend(result.quality_flags or [])
            parser_name = "image_ocr_gate"
    except Exception as exc:
        exception = exc
        errors.append(str(exc))

    parse_input = {
        "file_type": file_type,
        "raw_text": raw_text,
        "blocks": blocks,
        "warnings": warnings,
        "errors": errors,
        "quality_flags": layout_quality_flags,
        "ocr_required": ocr_required,
        "ocr_available": ocr_available,
        "formula_detected": contains_formula_text(raw_text) or bool(formula_regions),
        "image_only_suspected": image_only_suspected,
        "partial_text": partial_text,
        "ocr_completed": ocr_completed,
        "exception": exception,
    }
    quality = classify_parse_quality(parse_input)
    if formula_regions:
        quality["quality_flags"] = normalize_quality_flags(
            list(quality.get("quality_flags", []) or [])
            + ["formula_region_detected", "formula_recognizer_unavailable"]
        )
    error_code = ""
    error_message = ""
    if file_type == "unknown":
        error_code = "unsupported_file_type"
        error_message = "Unsupported file type."
    elif exception:
        error_code = "parse_failed"
        error_message = str(exception)
    elif quality["quality_status"] in {"ocr_required", "ocr_unavailable"}:
        error_code = quality["quality_status"]
        error_message = "OCR is required before reliable text extraction."
    elif quality["quality_status"] == "empty_text":
        error_code = "empty_text"
        error_message = "No effective text was extracted."

    record_data = build_parse_record_from_result(
        source_filename=source_filename,
        stored_path=str(file_path or ""),
        mime_type=mime_type or "",
        file_size_bytes=_file_size(file_path),
        parser_name=parser_name,
        file_type=file_type,
        raw_text=raw_text,
        blocks=blocks,
        quality=quality,
        page_count=page_count,
        error_code=error_code,
        error_message=error_message,
        now_fn=now_fn,
    )
    return DocumentParseResult(
        parse_record_data=record_data,
        blocks=blocks,
        raw_text=raw_text,
        formula_regions=formula_regions,
        warnings=quality["warnings"],
        errors=errors,
    )


def should_allow_term_extraction(parse_record: Any) -> bool:
    quality_status = getattr(parse_record, "quality_status", None)
    if isinstance(parse_record, dict):
        quality_status = parse_record.get("quality_status")
        text_chars = int(parse_record.get("extracted_text_chars") or 0)
    else:
        text_chars = int(getattr(parse_record, "extracted_text_chars", 0) or 0)
    if quality_status in {"native_text_ok", "ocr_text_ok", "partial_text", "ocr_low_confidence"}:
        return text_chars > 0
    return False


def serialize_parse_record(record: Any) -> dict[str, Any]:
    return {
        "id": getattr(record, "id", None),
        "parse_uid": getattr(record, "parse_uid", ""),
        "source_filename": getattr(record, "source_filename", ""),
        "stored_path": getattr(record, "stored_path", ""),
        "file_type": getattr(record, "file_type", ""),
        "mime_type": getattr(record, "mime_type", ""),
        "file_size_bytes": getattr(record, "file_size_bytes", None),
        "parser_name": getattr(record, "parser_name", ""),
        "parser_version": getattr(record, "parser_version", ""),
        "parse_status": getattr(record, "parse_status", ""),
        "quality_status": getattr(record, "quality_status", ""),
        "quality_flags": _json_loads(getattr(record, "quality_flags", "[]"), []),
        "page_count": getattr(record, "page_count", None),
        "block_count": getattr(record, "block_count", 0),
        "extracted_text_chars": getattr(record, "extracted_text_chars", 0),
        "ocr_required": bool(getattr(record, "ocr_required", False)),
        "ocr_available": bool(getattr(record, "ocr_available", False)),
        "formula_detected": bool(getattr(record, "formula_detected", False)),
        "table_detected": bool(getattr(record, "table_detected", False)),
        "image_only_suspected": bool(getattr(record, "image_only_suspected", False)),
        "error_code": getattr(record, "error_code", ""),
        "error_message": getattr(record, "error_message", ""),
        "warnings": _json_loads(getattr(record, "warnings", "[]"), []),
        "allow_term_extraction": should_allow_term_extraction(record),
        "created_at": getattr(record, "created_at", ""),
        "updated_at": getattr(record, "updated_at", ""),
    }


def serialize_parse_block(block: Any) -> dict[str, Any]:
    return {
        "id": getattr(block, "id", None),
        "block_uid": getattr(block, "block_uid", ""),
        "parse_uid": getattr(block, "parse_uid", ""),
        "page_number": getattr(block, "page_number", None),
        "slide_number": getattr(block, "slide_number", None),
        "block_index": getattr(block, "block_index", 0),
        "block_type": getattr(block, "block_type", ""),
        "text": getattr(block, "text", ""),
        "confidence": getattr(block, "confidence", None),
        "parser_type": getattr(block, "parser_type", ""),
        "source_locator": getattr(block, "source_locator", ""),
        "quality_flags": _json_loads(getattr(block, "quality_flags", "[]"), []),
        "created_at": getattr(block, "created_at", ""),
    }


def json_dumps(value: Any) -> str:
    return _json_dumps(value)
