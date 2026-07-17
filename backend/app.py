from flask import Flask, Response, g, jsonify, request, send_file
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from werkzeug.exceptions import RequestEntityTooLarge
from werkzeug.utils import secure_filename

import json
import hashlib
import io
import math
import os
import re
import secrets
import tempfile
import urllib.error
import urllib.request
import uuid
import zipfile
from datetime import datetime, timedelta
from pathlib import Path
from xml.etree import ElementTree as ET

from sqlalchemy import event
from sqlalchemy.orm import validates
from werkzeug.security import check_password_hash, generate_password_hash
from services.ai_providers import get_ai_provider
from services.ocr import get_ocr_provider
from services.formula_detection import contains_formula_text, looks_like_formula_image
from services.formula_ocr import get_formula_ocr_provider
from services import retrieval as retrieval_service
from services.scoring import score_knowledge_chunk as score_evidence_chunk
from services.alignment import (
    finalize_alignment_decision,
    validate_card_status_transition,
)
from services.evaluation import (
    evaluate_single_item as evaluate_single_item_result,
    normalize_evaluation_record,
    read_evaluation_jsonl,
)
from services.evaluation_metrics import compute_evaluation_metrics
from services.evaluation_report import generate_evaluation_report
from services.pilot_feedback import (
    FEEDBACK_SOURCES,
    FEEDBACK_STATUSES,
    FEEDBACK_TYPES,
    CLASSIFICATIONS,
    ROOT_CAUSES,
    SEVERITIES,
    anonymize_user,
    classify_feedback,
    normalize_choice,
    should_escalate_card,
)
from services.iteration_backlog import (
    BACKLOG_STATUSES,
    default_acceptance_criteria,
    map_feedback_to_category,
    map_feedback_to_priority,
)
from services.pilot_report import generate_pilot_report_markdown
from services.storage import StorageService, safe_filename as storage_safe_filename
from services.ai_registry import (
    env_provider_selection,
    validate_ai_config as validate_ai_environment_config,
    can_default_provider,
    is_placeholder_secret,
)
from services.legacy_provider_registry_seed import (
    LegacyProviderRegistrySeedModels,
    ensure_legacy_provider_registry_seed,
)
from services.legacy_provider_prompt_mutation import (
    LegacyPromptMutationDependencies,
    LegacyPromptMutationRequest,
    execute_legacy_prompt_mutation,
)
from services.ai_provider import provider_from_selection
from services.prompt_registry import (
    DEFAULT_PROMPTS,
    ALIGNMENT_STATUS_ENUM,
    default_prompt_lookup,
    validate_ai_json,
)
from services.ai_call_log import (
    estimate_tokens as estimate_ai_tokens,
    hash_payload as hash_ai_payload,
    preview_payload as preview_ai_payload,
)
from services.ai_cost import (
    AI_EVENT_BY_TASK,
    check_ai_quota,
    estimate_ai_cost,
    summarize_ai_calls,
)
from services.legacy_provider_local_readiness import evaluate_legacy_provider_local_readiness
from services.chunk_dedup import compute_content_hash, find_duplicate_chunk, mark_duplicate_chunk, normalize_chunk_text
from services.knowledge_health import summarize_health as summarize_kb_health
from services.knowledge_indexing import (
    INDEX_BACKEND,
    INDEX_VERSION,
    RETRIEVAL_VERSION as KB_RETRIEVAL_VERSION,
    build_knowledge_chunk_fields,
)
from services.knowledge_versioning import can_publish_version, default_version_name, next_version_number
from services.retrieval_regression import evaluate_retrieval_cases
from services.source_governance import can_source_generate_public_evidence, source_quality_from_governance, source_status_flags
from services.embedding_provider import get_embedding_provider
from services.vector_index import get_vector_index_backend
from services.retrieval_backends import VALID_RETRIEVAL_BACKENDS, get_retrieval_backend
from services.retrieval_experiments import (
    BACKENDS_TO_COMPARE,
    evaluate_backend_cases,
    markdown_report as retrieval_experiment_markdown,
    recommend_backend,
)
from services import concept_alignment_cards as concept_card_service
from services import concept_card_review as concept_card_review_service
from services import course_review_policy as course_review_policy_service
from services import audit_records as audit_record_service
from services import audit_context as audit_context_service
from services import document_parse_quality as document_parse_quality_service
from services import parse_quality_risk as parse_quality_risk_service
from services import knowledge_governance as knowledge_governance_service
from services import knowledge_ingestion as knowledge_ingestion_service
from services import evidence_retrieval as evidence_retrieval_service
from services import bilingual_evidence_workflow as bilingual_evidence_service
from services import concept_card_drafts as concept_card_draft_service
from services import chinese_term_candidates as chinese_term_candidate_service
from services import alignment_verification as alignment_verification_service
from services import alignment_verification_execution as alignment_verification_execution_service
from services import alignment_providers as alignment_provider_service
from services import provider_governance as provider_governance_service
from services import provider_preflight as provider_preflight_service
from services import student_course_access as student_course_access_service
from services import student_learning_progress as student_learning_progress_service
from services import concept_card_feedback as concept_card_feedback_service
from routes.admin_alignment_runs import AdminAlignmentRunModels, register_admin_alignment_run_routes
from routes.alignment_verification import register_alignment_verification_routes
from routes.concept_card_feedback import ConceptCardFeedbackModels, register_concept_card_feedback_routes
from routes.concept_card_review import ConceptCardReviewModels, register_concept_card_review_routes
from routes.legacy_provider_admin_observability import (
    LegacyProviderAdminObservabilityModels,
    LegacyProviderAdminObservabilitySerializers,
    register_legacy_provider_admin_observability_routes,
)
from routes.legacy_provider_admin_configuration import (
    LegacyProviderAdminConfigurationModels,
    LegacyProviderAdminConfigurationSerializers,
    register_legacy_provider_admin_configuration_routes,
)
from routes.legacy_provider_admin_healthcheck import (
    LegacyProviderAdminHealthcheckModels,
    LegacyProviderAdminHealthcheckSerializers,
    register_legacy_provider_admin_healthcheck_routes,
)
from routes.provider_governance import ProviderGovernanceModels, register_provider_governance_routes
from routes.provider_policy import ProviderPolicyModels, register_provider_policy_routes
from routes.provider_preflight import ProviderPreflightModels, register_provider_preflight_routes
from routes.shared import RouteCoreDependencies
from routes.student_concept_cards import StudentConceptCardModels, register_student_concept_card_routes
from routes.teacher_learning_analytics import TeacherLearningAnalyticsModels, register_teacher_learning_analytics_routes


# ============================================================
# LexiBridge AI local product prototype
# AI 双语课程知识对齐平台
#
# 当前版本定位：
# 1. 不做普通术语直译工具
# 2. 上传英文课件后抽取候选术语
# 3. 先检索英文教材证据，再生成中文候选译名
# 4. 再检索中文教材证据并给出对齐理由
# 5. 当前阶段只保证本地 Flask + SQLite 可演示，API 失败时使用 mock fallback
# ============================================================


def load_env_file(path):
    """
    Load simple KEY=VALUE pairs from a local .env file.
    This keeps the MVP dependency-light while still supporting secret config.
    """
    if not os.path.exists(path):
        return

    with open(path, "r", encoding="utf-8") as env_file:
        for raw_line in env_file:
            line = raw_line.strip()

            if not line or line.startswith("#") or "=" not in line:
                continue

            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip().strip('"').strip("'")

            if key and key not in os.environ:
                os.environ[key] = value


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE_DIR)

load_env_file(os.path.join(PROJECT_ROOT, ".env"))
load_env_file(os.path.join(BASE_DIR, ".env"))

app = Flask(__name__)

FRONTEND_ORIGIN = os.environ.get("FRONTEND_ORIGIN", "*").strip() or "*"
CORS(app, origins="*" if FRONTEND_ORIGIN == "*" else [FRONTEND_ORIGIN])


# ============================================================
# 路径与数据库配置
# ============================================================

# 上传目录放在用户目录，避免 Windows 中误把 uploads 建成文件后报错
UPLOAD_FOLDER = os.environ.get(
    "UPLOAD_FOLDER",
    os.path.join(os.path.expanduser("~"), "LexiBridge-AI-uploads")
)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
STORAGE_BACKEND = os.environ.get("STORAGE_BACKEND", "local").strip().lower() or "local"
LOCAL_STORAGE_ROOT = os.environ.get("LOCAL_STORAGE_ROOT", UPLOAD_FOLDER)
DATABASE_ENGINE = os.environ.get("DATABASE_ENGINE", "").strip().lower()

DATABASE_URL = os.environ.get("DATABASE_URL", "").strip()

if DATABASE_URL:
    app.config["SQLALCHEMY_DATABASE_URI"] = DATABASE_URL
else:
    DATABASE_FOLDER = os.path.join(os.path.expanduser("~"), "LexiBridge-AI-data")
    os.makedirs(DATABASE_FOLDER, exist_ok=True)
    DATABASE_PATH = os.path.join(DATABASE_FOLDER, "lexibridge.db")
    app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///" + DATABASE_PATH
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", "lexibridge-local-dev-secret")
if not DATABASE_ENGINE:
    DATABASE_ENGINE = "sqlite" if app.config["SQLALCHEMY_DATABASE_URI"].startswith("sqlite") else "postgresql"

db = SQLAlchemy(app)

ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "jpg", "jpeg", "png", "txt", "md", "markdown"}
DANGEROUS_EXTENSIONS = {"exe", "bat", "sh", "js", "html", "php", "docm", "xlsm", "zip"}
KNOWLEDGE_BASE_TYPES = {"en_course_kb", "zh_course_kb", "student_personal_kb"}
KB_SCOPES = {"global", "course", "personal"}

AI_PROVIDER = os.environ.get("AI_PROVIDER", "none").strip().lower()
AI_MODEL = os.environ.get("AI_MODEL", "").strip()
AI_PROVIDER_MODE = os.environ.get("AI_PROVIDER_MODE", "").strip().lower()
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "").strip()
DEEPSEEK_BASE_URL = os.environ.get("DEEPSEEK_BASE_URL", "https://api.deepseek.com").rstrip("/")
DEEPSEEK_MODEL = os.environ.get("DEEPSEEK_MODEL", os.environ.get("AI_MODEL", "deepseek-chat")).strip()
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "").strip()
OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
OPENAI_MODEL = os.environ.get("OPENAI_MODEL", "").strip()
AUTH_REQUIRED = os.environ.get("AUTH_REQUIRED", "true").strip().lower() == "true"
ALLOW_MOCK_AI = os.environ.get("ALLOW_MOCK_AI", "true").strip().lower() == "true"
ALLOW_LOCAL_HEURISTIC_AI = os.environ.get("ALLOW_LOCAL_HEURISTIC_AI", "true").strip().lower() == "true"
AI_REQUEST_TIMEOUT_SECONDS = int(os.environ.get("AI_REQUEST_TIMEOUT_SECONDS", os.environ.get("DEEPSEEK_TIMEOUT_SECONDS", "45")) or 45)
AI_MAX_RETRIES = int(os.environ.get("AI_MAX_RETRIES", "2") or 2)
AI_RETRY_BACKOFF_SECONDS = int(os.environ.get("AI_RETRY_BACKOFF_SECONDS", "2") or 2)
AI_DAILY_CALL_LIMIT_PER_USER = int(os.environ.get("AI_DAILY_CALL_LIMIT_PER_USER", "100") or 100)
AI_MONTHLY_CALL_LIMIT_PER_USER = int(os.environ.get("AI_MONTHLY_CALL_LIMIT_PER_USER", "1000") or 1000)
AI_DAILY_COST_LIMIT_PER_USER = float(os.environ.get("AI_DAILY_COST_LIMIT_PER_USER", "5.00") or 5.0)
AI_PROVIDER_HEALTHCHECK_ENABLED = os.environ.get("AI_PROVIDER_HEALTHCHECK_ENABLED", "true").strip().lower() == "true"
AI_LOG_PROMPT_FULL = os.environ.get("AI_LOG_PROMPT_FULL", "false").strip().lower() == "true"
AI_LOG_RESPONSE_FULL = os.environ.get("AI_LOG_RESPONSE_FULL", "false").strip().lower() == "true"
AI_LOG_REDACT_SECRETS = os.environ.get("AI_LOG_REDACT_SECRETS", os.environ.get("LOG_REDACT_SECRETS", "true")).strip().lower() == "true"
OCR_PROVIDER = os.environ.get("OCR_PROVIDER", "auto").strip().lower()
OCR_LANGS = os.environ.get("OCR_LANGS", "eng+chi_sim").strip() or "eng+chi_sim"
OCR_MIN_CONFIDENCE = int(os.environ.get("OCR_MIN_CONFIDENCE", "60"))
OCR_ENABLE_REGION_EXTRACTION = os.environ.get("OCR_ENABLE_REGION_EXTRACTION", "true").strip().lower() == "true"
PDF_OCR_DPI = int(os.environ.get("PDF_OCR_DPI", "300"))
PDF_IMAGE_MIN_WIDTH = int(os.environ.get("PDF_IMAGE_MIN_WIDTH", "80"))
PDF_IMAGE_MIN_HEIGHT = int(os.environ.get("PDF_IMAGE_MIN_HEIGHT", "40"))
PDF_MIXED_PAGE_IMAGE_OCR = os.environ.get("PDF_MIXED_PAGE_IMAGE_OCR", "true").strip().lower() == "true"
FORMULA_OCR_PROVIDER = os.environ.get("FORMULA_OCR_PROVIDER", "none").strip().lower()
FORMULA_OCR_MIN_CONFIDENCE = int(os.environ.get("FORMULA_OCR_MIN_CONFIDENCE", "60"))
FORMULA_DETECTION_MODE = os.environ.get("FORMULA_DETECTION_MODE", "heuristic").strip().lower()
MOCK_EMAIL_ENABLED = os.environ.get("MOCK_EMAIL_ENABLED", "true").strip().lower() == "true"
LOCAL_EMBEDDING_DIM = int(os.environ.get("LOCAL_EMBEDDING_DIM", "256"))
MAX_UPLOAD_SIZE_MB = int(os.environ.get("MAX_UPLOAD_SIZE_MB", "50"))
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_SIZE_MB * 1024 * 1024
RETRIEVAL_VERSION = os.environ.get("RETRIEVAL_VERSION", "local_lexical_v1").strip() or "local_lexical_v1"
RETRIEVAL_BACKEND = os.environ.get("RETRIEVAL_BACKEND", "lexical").strip().lower() or "lexical"
ENABLE_VECTOR_SEARCH = os.environ.get("ENABLE_VECTOR_SEARCH", "false").strip().lower() == "true"
VECTOR_INDEX_BACKEND = os.environ.get("VECTOR_INDEX_BACKEND", "none").strip().lower() or "none"
VECTOR_INDEX_DIR = os.environ.get("VECTOR_INDEX_DIR", "data/vector_indexes").strip() or "data/vector_indexes"
EMBEDDING_PROVIDER = os.environ.get("EMBEDDING_PROVIDER", "none").strip().lower() or "none"
EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL", "").strip()
ENABLE_HYBRID_SEARCH = os.environ.get("ENABLE_HYBRID_SEARCH", "false").strip().lower() == "true"
ENABLE_RERANKER = os.environ.get("ENABLE_RERANKER", "false").strip().lower() == "true"
RERANKER_PROVIDER = os.environ.get("RERANKER_PROVIDER", "none").strip().lower() or "none"
TERM_EXTRACTION_PROMPT_VERSION = os.environ.get("TERM_EXTRACTION_PROMPT_VERSION", "term_extraction_v1").strip() or "term_extraction_v1"
ALIGNMENT_PROMPT_VERSION = os.environ.get("ALIGNMENT_PROMPT_VERSION", "alignment_v1").strip() or "alignment_v1"
TOKEN_HASH_SECRET = os.environ.get("TOKEN_HASH_SECRET", app.config["SECRET_KEY"])
JOB_WORKER_ID = os.environ.get("JOB_WORKER_ID", "local-worker").strip() or "local-worker"
JOB_MAX_ATTEMPTS = int(os.environ.get("JOB_MAX_ATTEMPTS", "3"))

JOB_TYPES = {"document_ingestion", "alignment_run", "evaluation_run"}
JOB_STATUSES = {"queued", "running", "completed", "failed", "canceled", "retrying"}
TERMINAL_JOB_STATUSES = {"completed", "failed", "canceled"}

ERROR_CODES = {
    "AUTH_REQUIRED": 401,
    "TOKEN_EXPIRED": 401,
    "PERMISSION_DENIED": 403,
    "RESOURCE_NOT_FOUND": 404,
    "VALIDATION_ERROR": 400,
    "FILE_TOO_LARGE": 413,
    "UNSUPPORTED_FILE_TYPE": 415,
    "OCR_UNAVAILABLE": 422,
    "FORMULA_OCR_UNAVAILABLE": 422,
    "PARSING_FAILED": 422,
    "QUOTA_EXCEEDED": 402,
    "AI_PROVIDER_FAILED": 502,
    "AI_PROVIDER_NOT_CONFIGURED": 422,
    "AI_INVALID_RESPONSE": 502,
    "PDF_FONT_UNAVAILABLE": 422,
    "TOO_MANY_REQUESTS": 429,
    "INTERNAL_ERROR": 500,
}


# ============================================================
# 数据表
# ============================================================

class User(db.Model):
    """
    Local account table with password hash, email verification, reset tokens, and roles.
    """
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(160), default="")
    password_hash = db.Column(db.String(260), nullable=False)
    role = db.Column(db.String(30), nullable=False, default="student")
    display_name = db.Column(db.String(120), default="")
    is_verified = db.Column(db.Boolean, default=False)
    verification_token = db.Column(db.String(160), default="")
    verification_token_expires_at = db.Column(db.String(40), default="")
    reset_token = db.Column(db.String(160), default="")
    reset_token_expires_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    last_login_at = db.Column(db.String(40), default="")


class AuthToken(db.Model):
    """
    Bearer token table for local session persistence.
    """
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False)
    token = db.Column(db.String(120), unique=True, nullable=False)
    token_hash = db.Column(db.String(128), unique=True, default="")
    created_at = db.Column(db.String(40), default="")
    expires_at = db.Column(db.String(40), default="")
    last_used_at = db.Column(db.String(40), default="")
    revoked = db.Column(db.Boolean, default=False)


class Course(db.Model):
    """
    课程表：用于最终课程隔离。
    """
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(120), unique=True, nullable=False)
    course_code = db.Column(db.String(80), default="")
    semester = db.Column(db.String(80), default="")
    description = db.Column(db.Text, default="")
    language_mode = db.Column(db.String(60), default="bilingual")
    teacher_id = db.Column(db.Integer, default=0)
    status = db.Column(db.String(40), default="active")
    deleted_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")


class CourseMember(db.Model):
    """
    课程成员关系：学生/教师与课程绑定。
    """
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, nullable=False)
    user_id = db.Column(db.Integer, nullable=False)
    role = db.Column(db.String(30), default="student")
    role_in_course = db.Column(db.String(30), default="student")
    status = db.Column(db.String(40), default="active")
    created_at = db.Column(db.String(40), default="")
    joined_at = db.Column(db.String(40), default="")


class StorageObject(db.Model):
    """
    Storage metadata for uploaded and derived files.
    Local MVP uses LocalStorageBackend; staging/production can map storage_key
    to S3-compatible object storage without storing binaries in the database.
    """
    id = db.Column(db.Integer, primary_key=True)
    storage_backend = db.Column(db.String(40), default="local")
    bucket = db.Column(db.String(160), default="")
    storage_key = db.Column(db.String(600), default="", unique=True)
    original_filename = db.Column(db.String(260), default="")
    content_type = db.Column(db.String(160), default="")
    size_bytes = db.Column(db.Integer, default=0)
    sha256 = db.Column(db.String(64), default="")
    owner_user_id = db.Column(db.Integer, nullable=True)
    course_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    visibility = db.Column(db.String(40), default="private")
    purpose = db.Column(db.String(80), default="uploaded_document")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")
    status = db.Column(db.String(40), default="active")


class Term(db.Model):
    """
    术语表：
    pending  = 候选术语，等待质量控制
    approved = 质量控制通过，学生端可见
    """
    id = db.Column(db.Integer, primary_key=True)

    course = db.Column(db.String(120), nullable=False, default="Data Structures and Algorithms")
    chapter = db.Column(db.String(120), nullable=False, default="Chapter 4 - Hashing")

    english_term = db.Column(db.String(200), nullable=False)
    chinese_term = db.Column(db.String(200), default="待质量控制")
    explanation = db.Column(db.Text, default="待质量控制：系统已抽取候选术语，等待证据对齐结果确认。")
    context = db.Column(db.Text, default="")
    courseware_sentence = db.Column(db.Text, default="")
    english_kb_evidence = db.Column(db.Text, default="")
    ai_translation_candidate = db.Column(db.String(200), default="")
    chinese_kb_evidence = db.Column(db.Text, default="")
    final_chinese_term = db.Column(db.String(200), default="")
    alignment_reason = db.Column(db.Text, default="")
    review_status = db.Column(db.String(40), default="pending")

    confidence = db.Column(db.Integer, default=60)
    status = db.Column(db.String(30), default="pending")
    ai_status = db.Column(db.String(40), default="pending")
    english_evidence = db.Column(db.Text, default="")
    chinese_evidence = db.Column(db.Text, default="")
    risk_note = db.Column(db.Text, default="")
    parse_uid = db.Column(db.String(64), default="")
    parse_block_uid = db.Column(db.String(64), default="")
    parse_quality_status = db.Column(db.String(80), default="")
    parse_quality_flags = db.Column(db.Text, default="[]")
    input_risk_labels = db.Column(db.Text, default="[]")
    source_uid = db.Column(db.String(64), default="")
    chunk_uid = db.Column(db.String(64), default="")
    ai_model = db.Column(db.String(80), default="")

class Feedback(db.Model):
    """
    学生反馈表：
    用于记录学生对已发布术语提出的问题。
    """
    id = db.Column(db.Integer, primary_key=True)
    feedback_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))

    term_id = db.Column(db.Integer, nullable=False, default=0)
    user_id = db.Column(db.Integer, default=0)
    course_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    terminology_card_id = db.Column(db.Integer, nullable=True)
    formula_block_id = db.Column(db.Integer, nullable=True)
    job_id = db.Column(db.Integer, nullable=True)
    alignment_run_id = db.Column(db.Integer, nullable=True)
    evaluation_run_id = db.Column(db.Integer, nullable=True)
    user_role = db.Column(db.String(40), default="")

    course = db.Column(db.String(120), default="")
    chapter = db.Column(db.String(120), default="")
    card_uid = db.Column(db.String(64), default="")

    english_term = db.Column(db.String(200), default="")
    chinese_term = db.Column(db.String(200), default="")

    feedback_type = db.Column(db.String(80), default="其他问题")
    feedback_source = db.Column(db.String(80), default="student_card_detail")
    severity = db.Column(db.String(40), default="normal")
    priority = db.Column(db.String(20), default="P2")
    message = db.Column(db.Text, default="")
    suggested_chinese_term = db.Column(db.String(220), default="")
    feedback_content = db.Column(db.Text, default="")
    reported_issue = db.Column(db.Text, default="")
    expected_result = db.Column(db.Text, default="")
    actual_result = db.Column(db.Text, default="")
    evidence_comment = db.Column(db.Text, default="")
    screenshot_path = db.Column(db.String(500), default="")
    classification = db.Column(db.String(80), default="")
    root_cause = db.Column(db.String(80), default="")
    resolution_action = db.Column(db.String(80), default="")
    resolution_note = db.Column(db.Text, default="")
    handled_by = db.Column(db.Integer, nullable=True)
    handled_at = db.Column(db.String(40), default="")
    handler_role = db.Column(db.String(40), default="")
    teacher_note = db.Column(db.Text, default="")
    linked_review_uid = db.Column(db.String(64), default="")
    linked_card_uid = db.Column(db.String(64), default="")
    converted_to_evaluation_item_id = db.Column(db.Integer, nullable=True)
    linked_backlog_item_id = db.Column(db.Integer, nullable=True)

    status = db.Column(db.String(30), default="open")
    resolved_by = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")
    resolved_at = db.Column(db.String(40), default="")


@event.listens_for(Feedback, "before_insert")
def before_insert_feedback(mapper, connection, target):
    now = current_time_text()
    target.feedback_uid = target.feedback_uid or str(uuid.uuid4())
    target.message = target.message or target.reported_issue or target.feedback_content or ""
    target.card_uid = target.card_uid or (target.actual_result if target.feedback_source == "student_concept_card" else "")
    target.linked_card_uid = target.linked_card_uid or target.card_uid
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now


@event.listens_for(Feedback, "before_update")
def before_update_feedback(mapper, connection, target):
    target.message = target.message or target.reported_issue or target.feedback_content or ""
    target.card_uid = target.card_uid or (target.actual_result if target.feedback_source == "student_concept_card" else "")
    target.linked_card_uid = target.linked_card_uid or target.card_uid
    target.updated_at = current_time_text()


class ConceptCardFeedbackTriageRecord(db.Model):
    __tablename__ = "concept_card_feedback_triage_record"

    id = db.Column(db.Integer, primary_key=True)
    triage_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    feedback_uid = db.Column(db.String(64), nullable=False, default="")
    card_uid = db.Column(db.String(64), default="")
    course = db.Column(db.String(160), default="")
    action = db.Column(db.String(80), default="")
    previous_status = db.Column(db.String(40), default="")
    new_status = db.Column(db.String(40), default="")
    handled_by = db.Column(db.Integer, nullable=True)
    handler_role = db.Column(db.String(40), default="")
    reason_code = db.Column(db.String(120), default="")
    teacher_note = db.Column(db.Text, default="")
    linked_review_uid = db.Column(db.String(64), default="")
    created_at = db.Column(db.String(40), default="")


@event.listens_for(ConceptCardFeedbackTriageRecord, "before_insert")
def before_insert_concept_card_feedback_triage_record(mapper, connection, target):
    target.triage_uid = target.triage_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()


class IterationBacklogItem(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(220), default="")
    description = db.Column(db.Text, default="")
    source_type = db.Column(db.String(80), default="feedback")
    source_feedback_id = db.Column(db.Integer, nullable=True)
    course_id = db.Column(db.Integer, nullable=True)
    severity = db.Column(db.String(40), default="medium")
    priority = db.Column(db.String(20), default="P2")
    category = db.Column(db.String(80), default="documentation")
    status = db.Column(db.String(40), default="open")
    owner = db.Column(db.String(120), default="")
    target_pr = db.Column(db.String(80), default="")
    acceptance_criteria = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")
    closed_at = db.Column(db.String(40), default="")

class KnowledgeDocument(db.Model):
    """
    课程知识库文档表：
    用于记录教师上传的中文教材、中文课件、中文参考资料。
    """
    id = db.Column(db.Integer, primary_key=True)

    course = db.Column(db.String(120), nullable=False, default="")
    title = db.Column(db.String(200), nullable=False, default="")
    filename = db.Column(db.String(260), default="")
    saved_filename = db.Column(db.String(260), default="")
    file_type = db.Column(db.String(30), default="")
    parse_uid = db.Column(db.String(64), default="")

    language = db.Column(db.String(30), default="zh")
    source_type = db.Column(db.String(80), default="教师上传资料")
    knowledge_base_type = db.Column(db.String(40), default="zh_course_kb")
    owner_user_id = db.Column(db.String(80), default="")
    visibility = db.Column(db.String(30), default="course")

    text_length = db.Column(db.Integer, default=0)
    chunk_count = db.Column(db.Integer, default=0)

    created_at = db.Column(db.String(40), default="")


class KnowledgeChunk(db.Model):
    """
    课程知识库片段表：
    每个文档会被切分为多个 chunk，后续用于检索和 RAG。
    """
    id = db.Column(db.Integer, primary_key=True)

    chunk_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    source_uid = db.Column(db.String(64), default="")
    document_id = db.Column(db.Integer, nullable=False)
    source_id = db.Column(db.Integer, nullable=True)
    knowledge_source_id = db.Column(db.Integer, nullable=True)
    knowledge_base_version_id = db.Column(db.Integer, nullable=True)
    parse_uid = db.Column(db.String(64), default="")
    parse_block_uid = db.Column(db.String(64), default="")
    course_id = db.Column(db.Integer, nullable=True)
    scope_type = db.Column(db.String(30), default="course")

    course = db.Column(db.String(120), nullable=False, default="")
    title = db.Column(db.String(200), default="")
    discipline = db.Column(db.String(120), default="")
    chapter = db.Column(db.String(120), default="")

    chunk_index = db.Column(db.Integer, default=0)
    content = db.Column(db.Text, nullable=False, default="")
    normalized_text = db.Column(db.Text, default="")
    content_hash = db.Column(db.String(64), default="")

    source_page = db.Column(db.String(80), default="")
    source_slide = db.Column(db.String(80), default="")
    source_section = db.Column(db.String(160), default="")
    source_locator = db.Column(db.String(160), default="")
    page_number = db.Column(db.Integer, nullable=True)
    slide_number = db.Column(db.Integer, nullable=True)
    block_type = db.Column(db.String(40), default="text")
    token_count = db.Column(db.Integer, nullable=True)
    char_count = db.Column(db.Integer, nullable=True)
    formula_block_ids_json = db.Column(db.Text, default="[]")
    keywords = db.Column(db.Text, default="")
    source_citation = db.Column(db.Text, default="")
    embedding_id = db.Column(db.String(120), default="")
    language = db.Column(db.String(30), default="")
    knowledge_base_type = db.Column(db.String(40), default="zh_course_kb")
    owner_user_id = db.Column(db.String(80), default="")
    visibility = db.Column(db.String(30), default="course")
    index_status = db.Column(db.String(40), default="indexed")
    quality_status = db.Column(db.String(80), default="")
    quality_flags = db.Column(db.Text, default="[]")
    trust_level = db.Column(db.String(60), default="unknown")
    status = db.Column(db.String(40), default="active")
    embedding_status = db.Column(db.String(40), default="not_started")
    is_duplicate = db.Column(db.Boolean, default=False)
    duplicate_of_chunk_id = db.Column(db.Integer, nullable=True)
    is_active = db.Column(db.Boolean, default=True)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class CoursewareUpload(db.Model):
    """
    课件上传记录：保存上传源、课程章节和解析文本。
    """
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(260), default="")
    saved_filename = db.Column(db.String(260), default="")
    parse_uid = db.Column(db.String(64), default="")
    course = db.Column(db.String(120), default="")
    chapter = db.Column(db.String(120), default="")
    uploaded_by = db.Column(db.String(120), default="demo_teacher")
    upload_time = db.Column(db.String(40), default="")
    parsed_text = db.Column(db.Text, default="")


class KnowledgeBaseEntry(db.Model):
    """
    本地演示用结构化教材知识库条目。
    """
    id = db.Column(db.Integer, primary_key=True)
    language = db.Column(db.String(30), default="en")
    course = db.Column(db.String(120), default="")
    chapter = db.Column(db.String(120), default="")
    term = db.Column(db.String(200), default="")
    definition = db.Column(db.Text, default="")
    source = db.Column(db.String(260), default="")
    keywords = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


class EmbeddingRecord(db.Model):
    """
    本地 embedding 记录。
    当前演示版使用 deterministic hashing embedding 存 JSON，后续可替换为向量数据库。
    """
    id = db.Column(db.Integer, primary_key=True)
    chunk_id = db.Column(db.Integer, nullable=False)
    provider = db.Column(db.String(40), default="local_hashing")
    dim = db.Column(db.Integer, default=LOCAL_EMBEDDING_DIM)
    vector_json = db.Column(db.Text, nullable=False, default="[]")
    created_at = db.Column(db.String(40), default="")


class TaskJob(db.Model):
    """
    本地演示任务状态表。
    当前阶段不运行独立 worker，仅保留接口扩展空间。
    """
    id = db.Column(db.Integer, primary_key=True)
    job_type = db.Column(db.String(80), nullable=False)
    status = db.Column(db.String(30), default="queued")
    payload_json = db.Column(db.Text, default="{}")
    result_json = db.Column(db.Text, default="{}")
    error = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class Document(db.Model):
    """
    Uploaded source document, scoped to global/course/personal knowledge.
    """
    id = db.Column(db.Integer, primary_key=True)
    owner_user_id = db.Column(db.Integer, default=0)
    course_id = db.Column(db.Integer, nullable=True)
    scope_type = db.Column(db.String(30), default="course")
    filename = db.Column(db.String(260), default="")
    saved_filename = db.Column(db.String(260), default="")
    file_sha256 = db.Column(db.String(64), default="")
    storage_object_id = db.Column(db.Integer, nullable=True)
    storage_backend = db.Column(db.String(40), default="")
    storage_key = db.Column(db.String(600), default="")
    original_filename = db.Column(db.String(260), default="")
    content_type = db.Column(db.String(160), default="")
    size_bytes = db.Column(db.Integer, default=0)
    sha256 = db.Column(db.String(64), default="")
    file_type = db.Column(db.String(30), default="")
    language = db.Column(db.String(30), default="")
    upload_time = db.Column(db.String(40), default="")
    parsing_status = db.Column(db.String(40), default="pending")
    parse_uid = db.Column(db.String(64), default="")
    ocr_required = db.Column(db.Boolean, default=False)
    ocr_provider = db.Column(db.String(40), default="")
    ocr_status = db.Column(db.String(60), default="")
    ocr_error = db.Column(db.Text, default="")
    quality_flags_json = db.Column(db.Text, default="[]")
    source_type = db.Column(db.String(80), default="teacher_upload")
    parsed_text = db.Column(db.Text, default="")
    error_message = db.Column(db.Text, default="")
    deleted_at = db.Column(db.String(40), default="")


class DocumentChunk(db.Model):
    """
    Structured parsed text block. AI receives chunks instead of raw files.
    """
    id = db.Column(db.Integer, primary_key=True)
    document_id = db.Column(db.Integer, nullable=False)
    course_id = db.Column(db.Integer, nullable=True)
    user_id = db.Column(db.Integer, default=0)
    owner_user_id = db.Column(db.Integer, default=0)
    chunk_index = db.Column(db.Integer, default=0)
    parse_uid = db.Column(db.String(64), default="")
    parse_block_uid = db.Column(db.String(64), default="")
    language = db.Column(db.String(30), default="")
    page_number = db.Column(db.Integer, nullable=True)
    slide_number = db.Column(db.Integer, nullable=True)
    section_title = db.Column(db.String(220), default="")
    content = db.Column(db.Text, default="")
    source_type = db.Column(db.String(80), default="")
    source_location = db.Column(db.String(160), default="")
    ocr_confidence = db.Column(db.Integer, default=100)
    ocr_provider = db.Column(db.String(40), default="")
    ocr_status = db.Column(db.String(60), default="")
    ocr_error = db.Column(db.Text, default="")
    quality_flags_json = db.Column(db.Text, default="[]")
    created_at = db.Column(db.String(40), default="")


class DocumentParseRecord(db.Model):
    __tablename__ = "document_parse_record"

    id = db.Column(db.Integer, primary_key=True)
    parse_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    source_filename = db.Column(db.String(260), default="")
    stored_path = db.Column(db.String(600), default="")
    file_type = db.Column(db.String(40), default="unknown")
    mime_type = db.Column(db.String(160), default="")
    file_size_bytes = db.Column(db.Integer, nullable=True)
    parser_name = db.Column(db.String(120), default="")
    parser_version = db.Column(db.String(80), default="")
    parse_status = db.Column(db.String(40), default="failed")
    quality_status = db.Column(db.String(80), default="parse_failed")
    quality_flags = db.Column(db.Text, default="[]")
    page_count = db.Column(db.Integer, nullable=True)
    block_count = db.Column(db.Integer, default=0)
    extracted_text_chars = db.Column(db.Integer, default=0)
    ocr_required = db.Column(db.Boolean, default=False)
    ocr_available = db.Column(db.Boolean, default=False)
    formula_detected = db.Column(db.Boolean, default=False)
    table_detected = db.Column(db.Boolean, default=False)
    image_only_suspected = db.Column(db.Boolean, default=False)
    error_code = db.Column(db.String(120), default="")
    error_message = db.Column(db.Text, default="")
    warnings = db.Column(db.Text, default="[]")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class DocumentParseBlock(db.Model):
    __tablename__ = "document_parse_block"

    id = db.Column(db.Integer, primary_key=True)
    block_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    parse_uid = db.Column(db.String(64), nullable=False, default="")
    page_number = db.Column(db.Integer, nullable=True)
    slide_number = db.Column(db.Integer, nullable=True)
    block_index = db.Column(db.Integer, default=0)
    block_type = db.Column(db.String(40), default="text")
    text = db.Column(db.Text, default="")
    confidence = db.Column(db.Float, nullable=True)
    parser_type = db.Column(db.String(40), default="native")
    source_locator = db.Column(db.String(160), default="")
    quality_flags = db.Column(db.Text, default="[]")
    created_at = db.Column(db.String(40), default="")


@event.listens_for(DocumentParseRecord, "before_insert")
def before_insert_document_parse_record(mapper, connection, target):
    now = current_time_text()
    target.parse_uid = target.parse_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    target.quality_flags = document_parse_quality_service.json_dumps(
        document_parse_quality_service.normalize_quality_flags(
            document_parse_quality_service._json_loads(target.quality_flags, [])
        )
    )
    target.warnings = document_parse_quality_service.json_dumps(
        document_parse_quality_service.normalize_quality_flags(
            document_parse_quality_service._json_loads(target.warnings, [])
        )
    )


@event.listens_for(DocumentParseRecord, "before_update")
def before_update_document_parse_record(mapper, connection, target):
    target.updated_at = current_time_text()


@event.listens_for(DocumentParseBlock, "before_insert")
def before_insert_document_parse_block(mapper, connection, target):
    target.block_uid = target.block_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    target.quality_flags = document_parse_quality_service.json_dumps(
        document_parse_quality_service.normalize_quality_flags(
            document_parse_quality_service._json_loads(target.quality_flags, [])
        )
    )


class FormulaBlock(db.Model):
    """
    Image-based formula recognition result.
    Formula blocks are stored separately from text chunks so LaTeX/noisy symbols do
    not enter terminology extraction as English terms.
    """
    id = db.Column(db.Integer, primary_key=True)
    document_id = db.Column(db.Integer, nullable=False)
    course_id = db.Column(db.Integer, nullable=True)
    owner_user_id = db.Column(db.Integer, default=0)
    scope_type = db.Column(db.String(30), default="course")
    page_number = db.Column(db.Integer, nullable=True)
    slide_number = db.Column(db.Integer, nullable=True)
    bbox_json = db.Column(db.Text, default="{}")
    image_path = db.Column(db.String(500), default="")
    image_storage_object_id = db.Column(db.Integer, nullable=True)
    image_storage_key = db.Column(db.String(600), default="")
    image_content_type = db.Column(db.String(160), default="")
    image_sha256 = db.Column(db.String(64), default="")
    latex = db.Column(db.Text, default="")
    plain_text = db.Column(db.Text, default="")
    provider = db.Column(db.String(80), default="")
    confidence = db.Column(db.Float, default=0)
    status = db.Column(db.String(80), default="")
    error = db.Column(db.Text, default="")
    quality_flags_json = db.Column(db.Text, default="[]")
    created_at = db.Column(db.String(40), default="")


class KnowledgeSource(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    source_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    title = db.Column(db.String(220), default="")
    course = db.Column(db.String(160), default="")
    chapter = db.Column(db.String(160), default="")
    course_id = db.Column(db.Integer, nullable=True)
    scope_type = db.Column(db.String(30), default="course")
    owner_user_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    name = db.Column(db.String(220), nullable=False, default="")
    source_title = db.Column(db.String(220), default="")
    language = db.Column(db.String(30), default="")
    discipline = db.Column(db.String(120), default="")
    source_type = db.Column(db.String(80), default="teacher_upload")
    source_role = db.Column(db.String(80), default="unknown")
    owner_type = db.Column(db.String(40), default="unknown")
    owner_id = db.Column(db.String(80), default="")
    visibility = db.Column(db.String(40), default="course")
    trust_level = db.Column(db.String(60), default="unknown")
    parse_uid = db.Column(db.String(64), default="")
    source_filename = db.Column(db.String(260), default="")
    file_type = db.Column(db.String(40), default="unknown")
    content_hash = db.Column(db.String(64), default="")
    version = db.Column(db.Integer, default=1)
    license_note = db.Column(db.Text, default="")
    quality_status = db.Column(db.String(80), default="")
    quality_flags = db.Column(db.Text, default="[]")
    knowledge_base_type = db.Column(db.String(40), default="")
    access_method = db.Column(db.String(80), default="manual_upload")
    license_status = db.Column(db.String(80), default="unknown")
    license_type = db.Column(db.String(80), default="unknown")
    authorization_status = db.Column(db.String(80), default="unknown")
    source_quality = db.Column(db.Float, default=0.4)
    version_introduced_id = db.Column(db.Integer, nullable=True)
    version_removed_id = db.Column(db.Integer, nullable=True)
    status = db.Column(db.String(40), default="active")
    effective_from = db.Column(db.String(40), default="")
    effective_to = db.Column(db.String(40), default="")
    update_frequency = db.Column(db.String(80), default="manual")
    allow_full_text_indexing = db.Column(db.Boolean, default=False)
    allow_student_search = db.Column(db.Boolean, default=False)
    allow_derivative_cards = db.Column(db.Boolean, default=False)
    created_by = db.Column(db.Integer, default=0)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class KnowledgeVersion(db.Model):
    __tablename__ = "knowledge_version"

    id = db.Column(db.Integer, primary_key=True)
    version_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    source_uid = db.Column(db.String(64), nullable=False, default="")
    version_number = db.Column(db.Integer, default=1)
    change_type = db.Column(db.String(40), default="created")
    previous_content_hash = db.Column(db.String(64), default="")
    new_content_hash = db.Column(db.String(64), default="")
    parse_uid = db.Column(db.String(64), default="")
    changed_by = db.Column(db.Integer, nullable=True)
    change_note = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


class KnowledgePermission(db.Model):
    __tablename__ = "knowledge_permission"

    id = db.Column(db.Integer, primary_key=True)
    permission_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    source_uid = db.Column(db.String(64), nullable=False, default="")
    principal_type = db.Column(db.String(40), default="system")
    principal_id = db.Column(db.String(120), default="")
    access_level = db.Column(db.String(40), default="read")
    created_at = db.Column(db.String(40), default="")


@event.listens_for(KnowledgeSource, "before_insert")
def before_insert_knowledge_source(mapper, connection, target):
    now = current_time_text()
    target.source_uid = target.source_uid or str(uuid.uuid4())
    target.title = target.title or target.source_title or target.name
    target.source_title = target.source_title or target.title or target.name
    target.name = target.name or target.title or target.source_title
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    target.quality_flags = knowledge_governance_service.dumps_json_list(target.quality_flags)


@event.listens_for(KnowledgeSource, "before_update")
def before_update_knowledge_source(mapper, connection, target):
    target.updated_at = current_time_text()
    target.quality_flags = knowledge_governance_service.dumps_json_list(target.quality_flags)


@event.listens_for(KnowledgeChunk, "before_insert")
def before_insert_knowledge_chunk(mapper, connection, target):
    now = current_time_text()
    target.chunk_uid = target.chunk_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    target.source_locator = target.source_locator or target.source_section or target.source_page
    target.char_count = target.char_count if target.char_count is not None else len(target.content or "")
    target.token_count = target.token_count if target.token_count is not None else len((target.normalized_text or target.content or "").split())
    target.quality_flags = knowledge_governance_service.dumps_json_list(target.quality_flags)


@event.listens_for(KnowledgeChunk, "before_update")
def before_update_knowledge_chunk(mapper, connection, target):
    target.updated_at = current_time_text()
    target.quality_flags = knowledge_governance_service.dumps_json_list(target.quality_flags)


@event.listens_for(KnowledgeVersion, "before_insert")
def before_insert_knowledge_version(mapper, connection, target):
    target.version_uid = target.version_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()


@event.listens_for(KnowledgePermission, "before_insert")
def before_insert_knowledge_permission(mapper, connection, target):
    target.permission_uid = target.permission_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()


class KnowledgeBaseVersion(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    kb_scope = db.Column(db.String(30), default="course")
    scope_type = db.Column(db.String(30), default="course")
    course_id = db.Column(db.Integer, nullable=True)
    owner_user_id = db.Column(db.Integer, nullable=True)
    version_name = db.Column(db.String(120), default="")
    version_number = db.Column(db.Integer, default=1)
    status = db.Column(db.String(40), default="draft")
    description = db.Column(db.Text, default="")
    source_count = db.Column(db.Integer, default=0)
    chunk_count = db.Column(db.Integer, default=0)
    formula_block_count = db.Column(db.Integer, default=0)
    deduped_chunk_count = db.Column(db.Integer, default=0)
    index_backend = db.Column(db.String(80), default="local_lexical")
    index_version = db.Column(db.String(80), default="local_lexical_v1")
    retrieval_version = db.Column(db.String(80), default="local_lexical_v1")
    embedding_provider = db.Column(db.String(80), default="")
    embedding_model = db.Column(db.String(160), default="")
    embedding_dimension = db.Column(db.Integer, default=0)
    vector_index_status = db.Column(db.String(40), default="")
    vector_index_updated_at = db.Column(db.String(40), default="")
    evaluation_run_id = db.Column(db.Integer, nullable=True)
    quality_gate_status = db.Column(db.String(40), default="")
    manifest_json = db.Column(db.Text, default="{}")
    created_at = db.Column(db.String(40), default="")
    created_by = db.Column(db.Integer, default=0)
    published_at = db.Column(db.String(40), default="")
    archived_at = db.Column(db.String(40), default="")
    parent_version_id = db.Column(db.Integer, nullable=True)
    is_active = db.Column(db.Boolean, default=True)


class RetrievalRun(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    query = db.Column(db.Text, default="")
    course_id = db.Column(db.Integer, nullable=True)
    scope_type = db.Column(db.String(30), default="course")
    owner_user_id = db.Column(db.Integer, nullable=True)
    knowledge_base_version_id = db.Column(db.Integer, nullable=True)
    retrieval_version = db.Column(db.String(80), default="local_lexical_v1")
    index_version = db.Column(db.String(80), default="local_lexical_v1")
    result_count = db.Column(db.Integer, default=0)
    top_score = db.Column(db.Float, default=0.0)
    status = db.Column(db.String(40), default="completed")
    created_at = db.Column(db.String(40), default="")
    metadata_json = db.Column(db.Text, default="{}")


class RetrievalExperimentRun(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, nullable=True)
    evaluation_set_id = db.Column(db.Integer, nullable=True)
    kb_version_id = db.Column(db.Integer, nullable=True)
    experiment_name = db.Column(db.String(160), default="")
    backends_tested_json = db.Column(db.Text, default="[]")
    best_backend = db.Column(db.String(80), default="")
    recommendation = db.Column(db.Text, default="")
    metrics_json = db.Column(db.Text, default="{}")
    created_by = db.Column(db.Integer, default=0)
    created_at = db.Column(db.String(40), default="")
    finished_at = db.Column(db.String(40), default="")
    status = db.Column(db.String(40), default="completed")
    report_markdown = db.Column(db.Text, default="")


CONCEPT_ALIGNMENT_CARD_STATUSES = {"draft", "needs_review", "approved", "rejected", "deprecated"}


def concept_card_json_dumps(value):
    return json.dumps(value, ensure_ascii=False, sort_keys=True)


def concept_card_json_loads(value, fallback):
    if value in (None, ""):
        return fallback
    try:
        return json.loads(value)
    except (TypeError, ValueError):
        return fallback


def concept_card_serialize_list(value, field_name):
    if value in (None, ""):
        return "[]"
    if isinstance(value, str):
        parsed = concept_card_json_loads(value, None)
        if parsed in (None, ""):
            raise ValueError(f"{field_name} must be a JSON list.")
        value = parsed
    if not isinstance(value, list):
        raise ValueError(f"{field_name} must be a list.")
    return concept_card_json_dumps(value)


def concept_card_serialize_json(value):
    if value in (None, ""):
        return "[]"
    if isinstance(value, (list, dict)):
        return concept_card_json_dumps(value)
    if isinstance(value, str):
        parsed = concept_card_json_loads(value, None)
        if isinstance(parsed, (list, dict)):
            return concept_card_json_dumps(parsed)
        return value
    raise ValueError("Evidence fields must be JSON-serializable list/dict values or text.")


def concept_card_evidence_has_content(value):
    if value in (None, ""):
        return False
    parsed = concept_card_json_loads(value, None)
    if isinstance(parsed, list):
        return any(bool(item) for item in parsed)
    if isinstance(parsed, dict):
        return bool(parsed)
    return bool(str(value).strip())


def validate_concept_alignment_card(card):
    if not str(card.english_term or "").strip():
        raise ValueError("english_term is required.")
    if not str(card.course or "").strip():
        raise ValueError("course is required.")
    if card.status not in CONCEPT_ALIGNMENT_CARD_STATUSES:
        raise ValueError(f"status must be one of {sorted(CONCEPT_ALIGNMENT_CARD_STATUSES)}.")
    if card.confidence_score is not None:
        score = float(card.confidence_score)
        if score < 0 or score > 1:
            raise ValueError("confidence_score must be between 0 and 1.")
    risk_labels = concept_card_json_loads(card.risk_labels, [])
    if not isinstance(risk_labels, list):
        raise ValueError("risk_labels must serialize to a JSON list.")
    if (
        card.status == "approved"
        and set(risk_labels) & parse_quality_risk_service.FORCE_REVIEW_RISK_LABELS
        and not getattr(card, "reviewed_by", None)
    ):
        raise ValueError("ConceptAlignmentCard with input parse quality risk cannot be approved.")
    if card.status == "approved" and not (
        concept_card_evidence_has_content(card.english_evidence)
        or concept_card_evidence_has_content(card.chinese_evidence)
    ):
        raise ValueError("approved ConceptAlignmentCard requires English or Chinese evidence.")
    return True


class TerminologyCard(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    scope_type = db.Column(db.String(30), default="course")
    course_id = db.Column(db.Integer, nullable=True)
    owner_user_id = db.Column(db.Integer, nullable=True)
    english_term = db.Column(db.String(220), nullable=False, default="")
    normalized_english_term = db.Column(db.String(255), default="")
    final_chinese_term = db.Column(db.String(220), default="")
    normalized_chinese_term = db.Column(db.String(255), default="")
    ai_translation_candidate = db.Column(db.String(220), default="")
    courseware_sentence = db.Column(db.Text, default="")
    english_kb_evidence = db.Column(db.Text, default="")
    chinese_kb_evidence = db.Column(db.Text, default="")
    english_evidence_snapshot = db.Column(db.Text, default="")
    chinese_evidence_snapshot = db.Column(db.Text, default="")
    english_evidence_score = db.Column(db.Float, default=0.0)
    chinese_evidence_score = db.Column(db.Float, default=0.0)
    concept_explanation = db.Column(db.Text, default="")
    alignment_reason = db.Column(db.Text, default="")
    alignment_status = db.Column(db.String(64), default="unverified_translation")
    score_breakdown_json = db.Column(db.Text, default="{}")
    quality_flags_json = db.Column(db.Text, default="[]")
    parse_uid = db.Column(db.String(64), default="")
    parse_block_uid = db.Column(db.String(64), default="")
    parse_quality_status = db.Column(db.String(80), default="")
    parse_quality_flags = db.Column(db.Text, default="[]")
    input_risk_labels = db.Column(db.Text, default="[]")
    confidence_score = db.Column(db.Integer, default=0)
    status = db.Column(db.String(50), default="pending_quality_control")
    ai_provider = db.Column(db.String(80), default="")
    ai_provider_mode = db.Column(db.String(40), default="")
    ai_model = db.Column(db.String(80), default="")
    prompt_key = db.Column(db.String(80), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    knowledge_base_version_id = db.Column(db.Integer, nullable=True)
    english_kb_version_id = db.Column(db.Integer, nullable=True)
    chinese_kb_version_id = db.Column(db.Integer, nullable=True)
    retrieval_run_id = db.Column(db.Integer, nullable=True)
    index_version = db.Column(db.String(80), default="")
    evidence_content_hashes_json = db.Column(db.Text, default="[]")
    evidence_status = db.Column(db.String(40), default="evidence_current")
    ai_call_log_id = db.Column(db.Integer, nullable=True)
    alignment_run_id = db.Column(db.Integer, nullable=True)
    source_alignment_run_id = db.Column(db.Integer, nullable=True)
    risk_note = db.Column(db.Text, default="")
    source_document_id = db.Column(db.Integer, nullable=True)
    english_evidence_chunk_id = db.Column(db.Integer, nullable=True)
    chinese_evidence_chunk_id = db.Column(db.Integer, nullable=True)
    approved_by = db.Column(db.Integer, nullable=True)
    approved_at = db.Column(db.String(40), default="")
    rejected_reason = db.Column(db.Text, default="")
    reviewer_note = db.Column(db.Text, default="")
    feedback_count = db.Column(db.Integer, default=0)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class ConceptAlignmentCard(db.Model):
    """
    Core bilingual concept object for future course-context alignment workflows.

    This table does not replace legacy Term or TerminologyCard records. It provides
    a stable, auditable data model that later migration steps can map legacy terms
    into without deleting existing data.
    """
    __tablename__ = "concept_alignment_card"

    id = db.Column(db.Integer, primary_key=True)
    card_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    english_term = db.Column(db.String(220), nullable=False)
    chinese_term = db.Column(db.String(220), default="")
    course = db.Column(db.String(160), nullable=False)
    chapter = db.Column(db.String(160), default="")
    concept_scope = db.Column(db.Text, default="")
    english_explanation = db.Column(db.Text, default="")
    chinese_explanation = db.Column(db.Text, default="")
    english_evidence = db.Column(db.Text, default="[]")
    chinese_evidence = db.Column(db.Text, default="[]")
    alignment_reason = db.Column(db.Text, default="")
    confidence_score = db.Column(db.Float, nullable=True)
    risk_labels = db.Column(db.Text, default="[]")
    parse_uid = db.Column(db.String(64), default="")
    parse_block_uid = db.Column(db.String(64), default="")
    parse_quality_status = db.Column(db.String(80), default="")
    parse_quality_flags = db.Column(db.Text, default="[]")
    input_risk_labels = db.Column(db.Text, default="[]")
    status = db.Column(db.String(40), nullable=False, default="draft")
    source_document_id = db.Column(db.Integer, nullable=True)
    source_chunk_id = db.Column(db.Integer, nullable=True)
    created_by = db.Column(db.Integer, nullable=True)
    reviewed_by = db.Column(db.Integer, nullable=True)
    reviewed_at = db.Column(db.String(40), default="")
    model_name = db.Column(db.String(120), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    version = db.Column(db.Integer, default=1)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")

    @validates("english_term")
    def validate_english_term(self, key, value):
        value = str(value or "").strip()
        if not value:
            raise ValueError("english_term is required.")
        return value

    @validates("course")
    def validate_course(self, key, value):
        value = str(value or "").strip()
        if not value:
            raise ValueError("course is required.")
        return value

    @validates("status")
    def validate_status(self, key, value):
        value = str(value or "").strip()
        if value not in CONCEPT_ALIGNMENT_CARD_STATUSES:
            raise ValueError(f"status must be one of {sorted(CONCEPT_ALIGNMENT_CARD_STATUSES)}.")
        return value

    @validates("confidence_score")
    def validate_confidence_score(self, key, value):
        if value is None:
            return None
        score = float(value)
        if score < 0 or score > 1:
            raise ValueError("confidence_score must be between 0 and 1.")
        return score

    @validates("risk_labels")
    def validate_risk_labels(self, key, value):
        return concept_card_serialize_list(value, "risk_labels")

    @validates("english_evidence", "chinese_evidence")
    def validate_evidence(self, key, value):
        return concept_card_serialize_json(value)

    def set_risk_labels(self, labels):
        self.risk_labels = labels

    def get_risk_labels(self):
        return concept_card_json_loads(self.risk_labels, [])

    def set_english_evidence(self, evidence):
        self.english_evidence = evidence

    def get_english_evidence(self):
        return concept_card_json_loads(self.english_evidence, self.english_evidence or [])

    def set_chinese_evidence(self, evidence):
        self.chinese_evidence = evidence

    def get_chinese_evidence(self):
        return concept_card_json_loads(self.chinese_evidence, self.chinese_evidence or [])


@event.listens_for(ConceptAlignmentCard, "before_insert")
def before_insert_concept_alignment_card(mapper, connection, target):
    now = current_time_text()
    target.card_uid = target.card_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    target.risk_labels = concept_card_serialize_list(concept_card_json_loads(target.risk_labels, []), "risk_labels")
    target.parse_quality_flags = concept_card_serialize_list(concept_card_json_loads(target.parse_quality_flags, []), "parse_quality_flags")
    target.input_risk_labels = concept_card_serialize_list(concept_card_json_loads(target.input_risk_labels, []), "input_risk_labels")
    validate_concept_alignment_card(target)


@event.listens_for(ConceptAlignmentCard, "before_update")
def before_update_concept_alignment_card(mapper, connection, target):
    target.updated_at = current_time_text()
    target.risk_labels = concept_card_serialize_list(concept_card_json_loads(target.risk_labels, []), "risk_labels")
    target.parse_quality_flags = concept_card_serialize_list(concept_card_json_loads(target.parse_quality_flags, []), "parse_quality_flags")
    target.input_risk_labels = concept_card_serialize_list(concept_card_json_loads(target.input_risk_labels, []), "input_risk_labels")
    validate_concept_alignment_card(target)


class ConceptCardReviewRecord(db.Model):
    __tablename__ = "concept_card_review_record"

    id = db.Column(db.Integer, primary_key=True)
    review_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    card_uid = db.Column(db.String(64), nullable=False, default="")
    reviewer_id = db.Column(db.Integer, nullable=True)
    reviewer_role = db.Column(db.String(40), default="")
    reviewer_name = db.Column(db.String(160), default="")
    action = db.Column(db.String(80), nullable=False, default="")
    previous_status = db.Column(db.String(40), default="")
    new_status = db.Column(db.String(40), default="")
    decision = db.Column(db.String(80), default="")
    reason_code = db.Column(db.String(120), default="")
    review_comment = db.Column(db.Text, default="")
    evidence_assessment = db.Column(db.Text, default="{}")
    term_assessment = db.Column(db.Text, default="{}")
    risk_assessment = db.Column(db.Text, default="{}")
    required_changes = db.Column(db.Text, default="[]")
    resolved_risk_labels = db.Column(db.Text, default="[]")
    remaining_risk_labels = db.Column(db.Text, default="[]")
    verification_run_uid = db.Column(db.String(64), default="")
    request_id = db.Column(db.String(120), default="")
    created_at = db.Column(db.String(40), default="")


def validate_concept_card_review_record(record):
    if not str(record.card_uid or "").strip():
        raise ValueError("card_uid is required.")
    if record.action not in concept_card_review_service.REVIEW_ACTIONS:
        raise ValueError(f"action must be one of {sorted(concept_card_review_service.REVIEW_ACTIONS)}.")
    if record.decision and record.decision not in concept_card_review_service.REVIEW_DECISIONS:
        raise ValueError(f"decision must be one of {sorted(concept_card_review_service.REVIEW_DECISIONS)}.")
    if record.reason_code and record.reason_code not in concept_card_review_service.REASON_CODES:
        raise ValueError(f"reason_code must be one of {sorted(concept_card_review_service.REASON_CODES)}.")
    record.evidence_assessment = concept_card_json_dumps(concept_card_json_loads(record.evidence_assessment, {}))
    record.term_assessment = concept_card_json_dumps(concept_card_json_loads(record.term_assessment, {}))
    record.risk_assessment = concept_card_json_dumps(concept_card_json_loads(record.risk_assessment, {}))
    record.required_changes = concept_card_serialize_list(concept_card_json_loads(record.required_changes, []), "required_changes")
    record.resolved_risk_labels = concept_card_serialize_list(concept_card_json_loads(record.resolved_risk_labels, []), "resolved_risk_labels")
    record.remaining_risk_labels = concept_card_serialize_list(concept_card_json_loads(record.remaining_risk_labels, []), "remaining_risk_labels")
    return True


@event.listens_for(ConceptCardReviewRecord, "before_insert")
def before_insert_concept_card_review_record(mapper, connection, target):
    target.review_uid = target.review_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    validate_concept_card_review_record(target)


@event.listens_for(ConceptCardReviewRecord, "before_update")
def before_update_concept_card_review_record(mapper, connection, target):
    validate_concept_card_review_record(target)


class ConceptCardReviewAssignment(db.Model):
    __tablename__ = "concept_card_review_assignment"

    id = db.Column(db.Integer, primary_key=True)
    assignment_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    card_uid = db.Column(db.String(64), nullable=False, default="")
    assigned_to = db.Column(db.String(120), nullable=False, default="")
    assigned_by = db.Column(db.Integer, nullable=True)
    assignment_status = db.Column(db.String(40), default="active")
    due_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


def validate_concept_card_review_assignment(assignment):
    if not str(assignment.card_uid or "").strip():
        raise ValueError("card_uid is required.")
    if not str(assignment.assigned_to or "").strip():
        raise ValueError("assigned_to is required.")
    if assignment.assignment_status not in concept_card_review_service.ASSIGNMENT_STATUSES:
        assignment.assignment_status = "active"
    return True


@event.listens_for(ConceptCardReviewAssignment, "before_insert")
def before_insert_concept_card_review_assignment(mapper, connection, target):
    now = current_time_text()
    target.assignment_uid = target.assignment_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    validate_concept_card_review_assignment(target)


@event.listens_for(ConceptCardReviewAssignment, "before_update")
def before_update_concept_card_review_assignment(mapper, connection, target):
    target.updated_at = current_time_text()
    validate_concept_card_review_assignment(target)


class CourseReviewPolicy(db.Model):
    __tablename__ = "course_review_policy"

    id = db.Column(db.Integer, primary_key=True)
    policy_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    course = db.Column(db.String(160), nullable=False, default="")
    chapter = db.Column(db.String(160), default="")
    require_human_review = db.Column(db.Boolean, default=True)
    require_two_step_review = db.Column(db.Boolean, default=False)
    require_admin_for_override = db.Column(db.Boolean, default=True)
    allow_teacher_override = db.Column(db.Boolean, default=False)
    allow_approve_with_unverified_alignment = db.Column(db.Boolean, default=False)
    allow_approve_with_partial_text = db.Column(db.Boolean, default=False)
    allow_approve_with_missing_chinese_evidence = db.Column(db.Boolean, default=False)
    allow_approve_with_missing_english_evidence = db.Column(db.Boolean, default=False)
    blocking_risk_labels = db.Column(db.Text, default="[]")
    override_allowed_risk_labels = db.Column(db.Text, default="[]")
    override_forbidden_risk_labels = db.Column(db.Text, default="[]")
    required_evidence_sides = db.Column(db.String(40), default="both")
    min_required_evidence_count = db.Column(db.Integer, default=2)
    status = db.Column(db.String(40), default="active")
    created_by = db.Column(db.Integer, nullable=True)
    updated_by = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


def validate_course_review_policy(policy):
    if not str(policy.course or "").strip():
        raise ValueError("course is required.")
    if policy.status not in course_review_policy_service.POLICY_STATUSES:
        policy.status = "disabled"
    if policy.required_evidence_sides not in course_review_policy_service.REQUIRED_EVIDENCE_SIDES:
        policy.required_evidence_sides = "both"
    policy.blocking_risk_labels = concept_card_serialize_list(
        concept_card_json_loads(policy.blocking_risk_labels, []),
        "blocking_risk_labels",
    )
    policy.override_allowed_risk_labels = concept_card_serialize_list(
        concept_card_json_loads(policy.override_allowed_risk_labels, []),
        "override_allowed_risk_labels",
    )
    policy.override_forbidden_risk_labels = concept_card_serialize_list(
        concept_card_json_loads(policy.override_forbidden_risk_labels, []),
        "override_forbidden_risk_labels",
    )
    policy.min_required_evidence_count = max(int(policy.min_required_evidence_count or 0), 0)
    return True


@event.listens_for(CourseReviewPolicy, "before_insert")
def before_insert_course_review_policy(mapper, connection, target):
    now = current_time_text()
    target.policy_uid = target.policy_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    validate_course_review_policy(target)


@event.listens_for(CourseReviewPolicy, "before_update")
def before_update_course_review_policy(mapper, connection, target):
    target.updated_at = current_time_text()
    validate_course_review_policy(target)


class CourseReviewPermission(db.Model):
    __tablename__ = "course_review_permission"

    id = db.Column(db.Integer, primary_key=True)
    permission_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    course = db.Column(db.String(160), nullable=False, default="")
    chapter = db.Column(db.String(160), default="")
    reviewer_id = db.Column(db.Integer, nullable=True)
    reviewer_role = db.Column(db.String(40), default="teacher")
    permission_level = db.Column(db.String(40), default="read")
    can_review = db.Column(db.Boolean, default=False)
    can_approve = db.Column(db.Boolean, default=False)
    can_override_risk = db.Column(db.Boolean, default=False)
    can_assign_reviewer = db.Column(db.Boolean, default=False)
    status = db.Column(db.String(40), default="active")
    granted_by = db.Column(db.Integer, nullable=True)
    granted_at = db.Column(db.String(40), default="")
    revoked_by = db.Column(db.Integer, nullable=True)
    revoked_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


def validate_course_review_permission(permission):
    if not str(permission.course or "").strip():
        raise ValueError("course is required.")
    if permission.reviewer_role not in course_review_policy_service.REVIEWER_ROLES:
        permission.reviewer_role = str(permission.reviewer_role or "")
    if permission.permission_level not in course_review_policy_service.PERMISSION_LEVELS:
        permission.permission_level = "read"
    if permission.status not in course_review_policy_service.PERMISSION_STATUSES:
        permission.status = "disabled"
    return True


@event.listens_for(CourseReviewPermission, "before_insert")
def before_insert_course_review_permission(mapper, connection, target):
    now = current_time_text()
    target.permission_uid = target.permission_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now
    validate_course_review_permission(target)


@event.listens_for(CourseReviewPermission, "before_update")
def before_update_course_review_permission(mapper, connection, target):
    target.updated_at = current_time_text()
    validate_course_review_permission(target)


class StudentTermRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False)
    term_id = db.Column(db.Integer, nullable=False)
    is_favorite = db.Column(db.Boolean, default=False)
    is_mastered = db.Column(db.Boolean, default=False)
    last_viewed_at = db.Column(db.String(40), default="")


class StudentCourseMembership(db.Model):
    __tablename__ = "student_course_membership"

    id = db.Column(db.Integer, primary_key=True)
    membership_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    user_id = db.Column(db.Integer, nullable=False)
    course = db.Column(db.String(160), nullable=False, default="")
    role_in_course = db.Column(db.String(40), default="student")
    status = db.Column(db.String(40), default="active")
    enrolled_by = db.Column(db.Integer, nullable=True)
    enrolled_at = db.Column(db.String(40), default="")
    revoked_by = db.Column(db.Integer, nullable=True)
    revoked_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")

    __table_args__ = (
        db.UniqueConstraint("user_id", "course", name="uq_student_course_membership_user_course"),
    )


@event.listens_for(StudentCourseMembership, "before_insert")
def before_insert_student_course_membership(mapper, connection, target):
    now = current_time_text()
    target.membership_uid = target.membership_uid or str(uuid.uuid4())
    target.course = str(target.course or "").strip()
    target.role_in_course = student_course_access_service.normalize_role(target.role_in_course)
    target.status = student_course_access_service.normalize_membership_status(target.status)
    target.enrolled_at = target.enrolled_at or now
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now


@event.listens_for(StudentCourseMembership, "before_update")
def before_update_student_course_membership(mapper, connection, target):
    target.course = str(target.course or "").strip()
    target.role_in_course = student_course_access_service.normalize_role(target.role_in_course)
    target.status = student_course_access_service.normalize_membership_status(target.status)
    target.updated_at = current_time_text()


class CourseStudentVisibilityPolicy(db.Model):
    __tablename__ = "course_student_visibility_policy"

    id = db.Column(db.Integer, primary_key=True)
    policy_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    course = db.Column(db.String(160), unique=True, nullable=False, default="")
    visibility = db.Column(db.String(40), default="enrolled_only")
    allow_auditor_view = db.Column(db.Boolean, default=False)
    allow_teacher_preview = db.Column(db.Boolean, default=True)
    allow_cross_course_search = db.Column(db.Boolean, default=False)
    status = db.Column(db.String(40), default="active")
    created_by = db.Column(db.Integer, nullable=True)
    updated_by = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


@event.listens_for(CourseStudentVisibilityPolicy, "before_insert")
def before_insert_course_student_visibility_policy(mapper, connection, target):
    now = current_time_text()
    target.policy_uid = target.policy_uid or str(uuid.uuid4())
    target.course = str(target.course or "").strip()
    target.visibility = student_course_access_service.normalize_visibility(target.visibility)
    if target.status not in student_course_access_service.POLICY_STATUSES:
        target.status = "disabled"
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now


@event.listens_for(CourseStudentVisibilityPolicy, "before_update")
def before_update_course_student_visibility_policy(mapper, connection, target):
    target.course = str(target.course or "").strip()
    target.visibility = student_course_access_service.normalize_visibility(target.visibility)
    if target.status not in student_course_access_service.POLICY_STATUSES:
        target.status = "disabled"
    target.updated_at = current_time_text()


class StudentConceptCardState(db.Model):
    __tablename__ = "student_concept_card_state"

    id = db.Column(db.Integer, primary_key=True)
    state_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    card_uid = db.Column(db.String(64), nullable=False, default="")
    user_id = db.Column(db.Integer, nullable=False)
    course = db.Column(db.String(160), default="")
    favorited = db.Column(db.Boolean, default=False)
    mastered = db.Column(db.Boolean, default=False)
    mastered_at = db.Column(db.String(40), default="")
    last_viewed_at = db.Column(db.String(40), default="")
    view_count = db.Column(db.Integer, default=0)
    personal_note = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")

    __table_args__ = (
        db.UniqueConstraint("user_id", "card_uid", name="uq_student_concept_card_state_user_card"),
    )


@event.listens_for(StudentConceptCardState, "before_insert")
def before_insert_student_concept_card_state(mapper, connection, target):
    now = current_time_text()
    target.state_uid = target.state_uid or str(uuid.uuid4())
    target.created_at = target.created_at or now
    target.updated_at = target.updated_at or now


@event.listens_for(StudentConceptCardState, "before_update")
def before_update_student_concept_card_state(mapper, connection, target):
    target.updated_at = current_time_text()


class AlignmentRun(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    document_id = db.Column(db.Integer, nullable=True)
    course_id = db.Column(db.Integer, nullable=True)
    triggered_by = db.Column(db.Integer, default=0)
    provider = db.Column(db.String(80), default="")
    model_name = db.Column(db.String(120), default="")
    ai_provider = db.Column(db.String(80), default="")
    ai_provider_mode = db.Column(db.String(40), default="")
    ai_model = db.Column(db.String(120), default="")
    prompt_key = db.Column(db.String(80), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    terms_extracted = db.Column(db.Integer, default=0)
    cards_created = db.Column(db.Integer, default=0)
    term_count = db.Column(db.Integer, default=0)
    card_created_count = db.Column(db.Integer, default=0)
    auto_approved_count = db.Column(db.Integer, default=0)
    qc_count = db.Column(db.Integer, default=0)
    needs_evidence_count = db.Column(db.Integer, default=0)
    conflict_count = db.Column(db.Integer, default=0)
    failed_count = db.Column(db.Integer, default=0)
    status = db.Column(db.String(40), default="running")
    metrics_json = db.Column(db.Text, default="{}")
    error_message = db.Column(db.Text, default="")
    started_at = db.Column(db.String(40), default="")
    finished_at = db.Column(db.String(40), default="")


class EvaluationSet(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(160), nullable=False, default="")
    course_id = db.Column(db.Integer, nullable=True)
    discipline = db.Column(db.String(120), default="")
    description = db.Column(db.Text, default="")
    split = db.Column(db.String(40), default="smoke")
    locked = db.Column(db.Boolean, default=False)
    is_locked = db.Column(db.Boolean, default=False)
    created_by = db.Column(db.Integer, default=0)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class EvaluationItem(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    set_id = db.Column(db.Integer, nullable=False)
    evaluation_set_id = db.Column(db.Integer, nullable=True)
    item_id = db.Column(db.String(80), default="")
    split = db.Column(db.String(40), default="test")
    discipline = db.Column(db.String(120), default="")
    course_id = db.Column(db.Integer, nullable=True)
    english_term = db.Column(db.String(220), nullable=False, default="")
    expected_chinese_term = db.Column(db.String(220), default="")
    expected_alignment_status = db.Column(db.String(64), default="")
    english_context = db.Column(db.Text, default="")
    english_evidence = db.Column(db.Text, default="")
    chinese_evidence = db.Column(db.Text, default="")
    expected_english_evidence = db.Column(db.Text, default="")
    expected_chinese_evidence = db.Column(db.Text, default="")
    negative_english_evidence = db.Column(db.Text, default="")
    negative_chinese_evidence = db.Column(db.Text, default="")
    difficulty = db.Column(db.String(40), default="medium")
    tags_json = db.Column(db.Text, default="[]")
    annotator = db.Column(db.String(120), default="")
    reviewed_by = db.Column(db.String(120), default="")
    disagreement_note = db.Column(db.Text, default="")
    version = db.Column(db.String(40), default="v1")
    created_at = db.Column(db.String(40), default="")


class EvaluationRun(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    evaluation_set_id = db.Column(db.Integer, nullable=True)
    triggered_by = db.Column(db.Integer, default=0)
    provider = db.Column(db.String(80), default="")
    provider_name = db.Column(db.String(80), default="")
    provider_mode = db.Column(db.String(40), default="")
    model_name = db.Column(db.String(120), default="")
    model_version = db.Column(db.String(120), default="")
    prompt_key = db.Column(db.String(80), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    alignment_version = db.Column(db.String(80), default="")
    commit_hash = db.Column(db.String(80), default="")
    split = db.Column(db.String(40), default="test")
    input_count = db.Column(db.Integer, default=0)
    skipped_count = db.Column(db.Integer, default=0)
    extraction_precision = db.Column(db.Float, default=0.0)
    extraction_recall = db.Column(db.Float, default=0.0)
    evidence_accuracy = db.Column(db.Float, default=0.0)
    english_evidence_accuracy = db.Column(db.Float, default=0.0)
    chinese_evidence_accuracy = db.Column(db.Float, default=0.0)
    alignment_accuracy = db.Column(db.Float, default=0.0)
    false_positive_rate = db.Column(db.Float, default=0.0)
    auto_approval_error_rate = db.Column(db.Float, default=0.0)
    ocr_noise_term_rate = db.Column(db.Float, nullable=True)
    no_evidence_forced_alignment_rate = db.Column(db.Float, default=0.0)
    created_by = db.Column(db.Integer, default=0)
    metrics_json = db.Column(db.Text, default="{}")
    report_json = db.Column(db.Text, default="{}")
    report_markdown = db.Column(db.Text, default="")
    status = db.Column(db.String(40), default="completed")
    error_message = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")
    finished_at = db.Column(db.String(40), default="")


class ModelPromptRegistry(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    provider = db.Column(db.String(80), default="")
    model_name = db.Column(db.String(120), default="")
    model_version = db.Column(db.String(120), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    enabled = db.Column(db.Boolean, default=True)
    allowed_workflows = db.Column(db.Text, default="")
    last_evaluation_run_id = db.Column(db.Integer, nullable=True)
    known_risks = db.Column(db.Text, default="")
    rollback_target = db.Column(db.String(120), default="")
    owner = db.Column(db.String(120), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class AIProviderConfig(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    provider_name = db.Column(db.String(80), default="")
    provider_mode = db.Column(db.String(40), default="none")
    base_url = db.Column(db.String(260), default="")
    default_model = db.Column(db.String(120), default="")
    is_enabled = db.Column(db.Boolean, default=True)
    is_default = db.Column(db.Boolean, default=False)
    supports_json_schema = db.Column(db.Boolean, default=True)
    supports_streaming = db.Column(db.Boolean, default=False)
    supports_vision = db.Column(db.Boolean, default=False)
    supports_formula_reasoning = db.Column(db.Boolean, default=False)
    max_input_tokens = db.Column(db.Integer, default=0)
    max_output_tokens = db.Column(db.Integer, default=0)
    timeout_seconds = db.Column(db.Integer, default=45)
    max_retries = db.Column(db.Integer, default=2)
    cost_per_1k_input_tokens = db.Column(db.Float, default=0.0)
    cost_per_1k_output_tokens = db.Column(db.Float, default=0.0)
    health_status = db.Column(db.String(40), default="unknown")
    last_healthcheck_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class AIModelRegistry(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    provider_name = db.Column(db.String(80), default="")
    model_name = db.Column(db.String(120), default="")
    model_version = db.Column(db.String(120), default="")
    model_display_name = db.Column(db.String(160), default="")
    provider_mode = db.Column(db.String(40), default="")
    supports_json_output = db.Column(db.Boolean, default=True)
    supports_tool_calling = db.Column(db.Boolean, default=False)
    supports_vision = db.Column(db.Boolean, default=False)
    max_input_tokens = db.Column(db.Integer, default=0)
    max_output_tokens = db.Column(db.Integer, default=0)
    cost_per_1k_input_tokens = db.Column(db.Float, default=0.0)
    cost_per_1k_output_tokens = db.Column(db.Float, default=0.0)
    is_enabled = db.Column(db.Boolean, default=True)
    is_default_for_provider = db.Column(db.Boolean, default=False)
    last_evaluation_run_id = db.Column(db.Integer, nullable=True)
    last_evaluation_score = db.Column(db.Float, default=0.0)
    known_risks_json = db.Column(db.Text, default="[]")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


class PromptTemplate(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    prompt_key = db.Column(db.String(80), default="")
    prompt_version = db.Column(db.String(40), default="")
    task_type = db.Column(db.String(80), default="")
    language = db.Column(db.String(40), default="")
    template_text = db.Column(db.Text, default="")
    json_schema = db.Column(db.Text, default="{}")
    is_active = db.Column(db.Boolean, default=True)
    is_default = db.Column(db.Boolean, default=False)
    created_by = db.Column(db.Integer, default=0)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")
    notes = db.Column(db.Text, default="")


class AICallLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    task_type = db.Column(db.String(80), default="")
    provider_name = db.Column(db.String(80), default="")
    provider_mode = db.Column(db.String(40), default="")
    model_name = db.Column(db.String(120), default="")
    prompt_key = db.Column(db.String(80), default="")
    prompt_version = db.Column(db.String(40), default="")
    user_id = db.Column(db.Integer, nullable=True)
    course_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    job_id = db.Column(db.Integer, nullable=True)
    alignment_run_id = db.Column(db.Integer, nullable=True)
    evaluation_run_id = db.Column(db.Integer, nullable=True)
    request_hash = db.Column(db.String(64), default="")
    response_hash = db.Column(db.String(64), default="")
    input_token_count = db.Column(db.Integer, default=0)
    output_token_count = db.Column(db.Integer, default=0)
    estimated_cost = db.Column(db.Float, default=0.0)
    latency_ms = db.Column(db.Integer, default=0)
    status = db.Column(db.String(40), default="")
    error_code = db.Column(db.String(80), default="")
    error_message = db.Column(db.Text, default="")
    redacted_prompt_preview = db.Column(db.Text, default="")
    redacted_response_preview = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


class AuditRecord(db.Model):
    __tablename__ = "audit_record"

    id = db.Column(db.Integer, primary_key=True)
    audit_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    event_type = db.Column(db.String(120), nullable=False, default="")
    target_type = db.Column(db.String(120), nullable=False, default="")
    target_uid = db.Column(db.String(120), default="")
    actor_id = db.Column(db.Integer, nullable=True)
    actor_role = db.Column(db.String(40), default="")
    actor_name = db.Column(db.String(160), default="")
    request_id = db.Column(db.String(120), default="")
    source = db.Column(db.String(40), default="service")
    before_snapshot = db.Column(db.Text, default="{}")
    after_snapshot = db.Column(db.Text, default="{}")
    input_payload = db.Column(db.Text, default="{}")
    output_payload = db.Column(db.Text, default="{}")
    changed_fields = db.Column(db.Text, default="[]")
    result = db.Column(db.String(40), default="success")
    error_code = db.Column(db.String(120), default="")
    error_message = db.Column(db.Text, default="")
    model_name = db.Column(db.String(120), default="")
    prompt_version = db.Column(db.String(80), default="")
    retrieval_version = db.Column(db.String(80), default="")
    latency_ms = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")


@event.listens_for(AuditRecord, "before_insert")
def before_insert_audit_record(mapper, connection, target):
    target.audit_uid = target.audit_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()


class AlignmentVerificationRun(db.Model):
    __tablename__ = "alignment_verification_run"

    id = db.Column(db.Integer, primary_key=True)
    run_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    card_uid = db.Column(db.String(64), default="")
    english_term = db.Column(db.String(220), nullable=False, default="")
    chinese_term = db.Column(db.String(220), default="")
    course = db.Column(db.String(160), default="")
    chapter = db.Column(db.String(160), default="")
    provider_name = db.Column(db.String(120), default="mock-rule-v1")
    provider_type = db.Column(db.String(40), default="mock")
    provider_version = db.Column(db.String(80), default="v1")
    input_payload = db.Column(db.Text, default="{}")
    output_payload = db.Column(db.Text, default="{}")
    english_evidence_count = db.Column(db.Integer, default=0)
    chinese_evidence_count = db.Column(db.Integer, default=0)
    top_english_chunk_uids = db.Column(db.Text, default="[]")
    top_chinese_chunk_uids = db.Column(db.Text, default="[]")
    retrieval_score_summary = db.Column(db.Text, default="{}")
    candidate_score_summary = db.Column(db.Text, default="{}")
    alignment_confidence = db.Column(db.Float, nullable=True)
    verification_status = db.Column(db.String(40), default="mock_only")
    recommendation = db.Column(db.String(80), default="needs_review")
    risk_labels = db.Column(db.Text, default="[]")
    prompt_version = db.Column(db.String(80), default="")
    prompt_summary = db.Column(db.Text, default="{}")
    raw_output_summary = db.Column(db.Text, default="{}")
    parser_version = db.Column(db.String(80), default="")
    output_schema_version = db.Column(db.String(80), default="")
    provider_response_status = db.Column(db.String(80), default="")
    error_code = db.Column(db.String(120), default="")
    error_message = db.Column(db.Text, default="")
    latency_ms = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")

    @validates("english_term")
    def validate_alignment_english_term(self, key, value):
        value = str(value or "").strip()
        if not value:
            raise ValueError("english_term is required.")
        return value

    @validates("alignment_confidence")
    def validate_alignment_confidence(self, key, value):
        if value is None:
            return None
        score = float(value)
        if score < 0 or score > 1:
            raise ValueError("alignment_confidence must be between 0 and 1.")
        return score


def validate_alignment_verification_run(run):
    if not str(run.english_term or "").strip():
        raise ValueError("english_term is required.")
    if run.alignment_confidence is not None:
        score = float(run.alignment_confidence)
        if score < 0 or score > 1:
            raise ValueError("alignment_confidence must be between 0 and 1.")
    if run.provider_type == "mock" and run.verification_status not in {"mock_only", "needs_review", "failed"}:
        raise ValueError("mock alignment verification run must use mock_only, needs_review, or failed status.")
    run.input_payload = concept_card_json_dumps(concept_card_json_loads(run.input_payload, {}))
    run.output_payload = concept_card_json_dumps(concept_card_json_loads(run.output_payload, {}))
    run.top_english_chunk_uids = concept_card_serialize_list(
        concept_card_json_loads(run.top_english_chunk_uids, []),
        "top_english_chunk_uids",
    )
    run.top_chinese_chunk_uids = concept_card_serialize_list(
        concept_card_json_loads(run.top_chinese_chunk_uids, []),
        "top_chinese_chunk_uids",
    )
    run.retrieval_score_summary = concept_card_json_dumps(concept_card_json_loads(run.retrieval_score_summary, {}))
    run.candidate_score_summary = concept_card_json_dumps(concept_card_json_loads(run.candidate_score_summary, {}))
    run.risk_labels = concept_card_serialize_list(concept_card_json_loads(run.risk_labels, []), "risk_labels")
    run.prompt_summary = concept_card_json_dumps(concept_card_json_loads(run.prompt_summary, {}))
    run.raw_output_summary = concept_card_json_dumps(concept_card_json_loads(run.raw_output_summary, {}))
    return True


@event.listens_for(AlignmentVerificationRun, "before_insert")
def before_insert_alignment_verification_run(mapper, connection, target):
    target.run_uid = target.run_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    validate_alignment_verification_run(target)


@event.listens_for(AlignmentVerificationRun, "before_update")
def before_update_alignment_verification_run(mapper, connection, target):
    validate_alignment_verification_run(target)


class AlignmentProviderPolicy(db.Model):
    __tablename__ = "alignment_provider_policy"

    id = db.Column(db.Integer, primary_key=True)
    policy_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    provider_name = db.Column(db.String(120), unique=True, nullable=False, default="")
    provider_type = db.Column(db.String(40), default="")
    enabled = db.Column(db.Boolean, default=False)
    replay_only = db.Column(db.Boolean, default=True)
    allow_external_calls = db.Column(db.Boolean, default=False)
    allow_attach_to_card = db.Column(db.Boolean, default=False)
    allow_production_result = db.Column(db.Boolean, default=False)
    allow_auto_approve = db.Column(db.Boolean, default=False)
    require_human_review = db.Column(db.Boolean, default=True)
    allowed_courses = db.Column(db.Text, default="[]")
    blocked_courses = db.Column(db.Text, default="[]")
    allowed_roles = db.Column(db.Text, default="[]")
    max_calls_per_day = db.Column(db.Integer, default=0)
    max_calls_per_month = db.Column(db.Integer, default=0)
    max_estimated_cost_per_call = db.Column(db.Float, nullable=True)
    max_estimated_cost_per_day = db.Column(db.Float, nullable=True)
    max_prompt_chars = db.Column(db.Integer, default=8000)
    max_output_chars = db.Column(db.Integer, default=4000)
    timeout_seconds = db.Column(db.Integer, default=30)
    max_retries = db.Column(db.Integer, default=0)
    status = db.Column(db.String(40), default="disabled")
    created_by = db.Column(db.Integer, nullable=True)
    updated_by = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")


def validate_alignment_provider_policy(policy):
    if not str(policy.provider_name or "").strip():
        raise ValueError("provider_name is required.")
    policy.allow_auto_approve = False
    policy.require_human_review = True
    policy.allowed_courses = concept_card_serialize_list(
        concept_card_json_loads(policy.allowed_courses, []),
        "allowed_courses",
    )
    policy.blocked_courses = concept_card_serialize_list(
        concept_card_json_loads(policy.blocked_courses, []),
        "blocked_courses",
    )
    policy.allowed_roles = concept_card_serialize_list(
        concept_card_json_loads(policy.allowed_roles, []),
        "allowed_roles",
    )
    if policy.status not in provider_governance_service.POLICY_STATUSES:
        policy.status = "disabled"
    policy.max_calls_per_day = max(int(policy.max_calls_per_day or 0), 0)
    policy.max_calls_per_month = max(int(policy.max_calls_per_month or 0), 0)
    policy.max_prompt_chars = max(int(policy.max_prompt_chars or 8000), 500)
    policy.max_output_chars = max(int(policy.max_output_chars or 4000), 500)
    policy.timeout_seconds = max(min(int(policy.timeout_seconds or 30), 120), 1)
    policy.max_retries = max(min(int(policy.max_retries or 0), 3), 0)
    return True


@event.listens_for(AlignmentProviderPolicy, "before_insert")
def before_insert_alignment_provider_policy(mapper, connection, target):
    target.policy_uid = target.policy_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    target.updated_at = target.updated_at or target.created_at
    validate_alignment_provider_policy(target)


@event.listens_for(AlignmentProviderPolicy, "before_update")
def before_update_alignment_provider_policy(mapper, connection, target):
    target.updated_at = target.updated_at or current_time_text()
    validate_alignment_provider_policy(target)


class AlignmentProviderUsageRecord(db.Model):
    __tablename__ = "alignment_provider_usage_record"

    id = db.Column(db.Integer, primary_key=True)
    usage_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    provider_name = db.Column(db.String(120), default="")
    provider_type = db.Column(db.String(40), default="")
    run_uid = db.Column(db.String(64), default="")
    card_uid = db.Column(db.String(64), default="")
    course = db.Column(db.String(160), default="")
    chapter = db.Column(db.String(160), default="")
    request_id = db.Column(db.String(120), default="")
    estimated_input_tokens = db.Column(db.Integer, default=0)
    estimated_output_tokens = db.Column(db.Integer, default=0)
    estimated_cost = db.Column(db.Float, default=0.0)
    actual_cost = db.Column(db.Float, nullable=True)
    provider_response_status = db.Column(db.String(80), default="")
    error_code = db.Column(db.String(120), default="")
    error_message = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


@event.listens_for(AlignmentProviderUsageRecord, "before_insert")
def before_insert_alignment_provider_usage_record(mapper, connection, target):
    target.usage_uid = target.usage_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    target.estimated_input_tokens = max(int(target.estimated_input_tokens or 0), 0)
    target.estimated_output_tokens = max(int(target.estimated_output_tokens or 0), 0)
    target.estimated_cost = float(target.estimated_cost or 0.0)


class AlignmentProviderPreflightRun(db.Model):
    __tablename__ = "alignment_provider_preflight_run"

    id = db.Column(db.Integer, primary_key=True)
    preflight_uid = db.Column(db.String(64), unique=True, nullable=False, default=lambda: str(uuid.uuid4()))
    provider_name = db.Column(db.String(120), default="")
    provider_type = db.Column(db.String(40), default="")
    policy_uid = db.Column(db.String(64), default="")
    course = db.Column(db.String(160), default="")
    requested_by = db.Column(db.String(120), default="")
    check_status = db.Column(db.String(40), default="failed")
    overall_ready = db.Column(db.Boolean, default=False)
    external_calls_enabled = db.Column(db.Boolean, default=False)
    replay_only = db.Column(db.Boolean, default=True)
    api_key_present = db.Column(db.Boolean, default=False)
    api_key_env_name = db.Column(db.String(120), default="")
    policy_summary = db.Column(db.Text, default="{}")
    check_results = db.Column(db.Text, default="{}")
    blocking_reasons = db.Column(db.Text, default="[]")
    warnings = db.Column(db.Text, default="[]")
    replay_dry_run_status = db.Column(db.String(40), default="not_run")
    estimated_cost_per_call = db.Column(db.Float, nullable=True)
    max_estimated_cost_per_call = db.Column(db.Float, nullable=True)
    max_calls_per_day = db.Column(db.Integer, default=0)
    max_calls_per_month = db.Column(db.Integer, default=0)
    require_human_review = db.Column(db.Boolean, default=True)
    allow_auto_approve = db.Column(db.Boolean, default=False)
    allow_production_result = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.String(40), default="")


def validate_alignment_provider_preflight_run(run):
    run.allow_auto_approve = False
    if run.check_status not in provider_preflight_service.PREFLIGHT_STATUSES:
        run.check_status = "failed"
    if run.replay_dry_run_status not in provider_preflight_service.REPLAY_DRY_RUN_STATUSES:
        run.replay_dry_run_status = "not_run"
    run.policy_summary = concept_card_json_dumps(concept_card_json_loads(run.policy_summary, {}))
    run.check_results = concept_card_json_dumps(concept_card_json_loads(run.check_results, {}))
    run.blocking_reasons = concept_card_serialize_list(
        concept_card_json_loads(run.blocking_reasons, []),
        "blocking_reasons",
    )
    run.warnings = concept_card_serialize_list(concept_card_json_loads(run.warnings, []), "warnings")
    run.max_calls_per_day = max(int(run.max_calls_per_day or 0), 0)
    run.max_calls_per_month = max(int(run.max_calls_per_month or 0), 0)
    return True


@event.listens_for(AlignmentProviderPreflightRun, "before_insert")
def before_insert_alignment_provider_preflight_run(mapper, connection, target):
    target.preflight_uid = target.preflight_uid or str(uuid.uuid4())
    target.created_at = target.created_at or current_time_text()
    validate_alignment_provider_preflight_run(target)


@event.listens_for(AlignmentProviderPreflightRun, "before_update")
def before_update_alignment_provider_preflight_run(mapper, connection, target):
    validate_alignment_provider_preflight_run(target)


class PersonalAccessAudit(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    actor_user_id = db.Column(db.Integer, nullable=False)
    target_user_id = db.Column(db.Integer, nullable=False)
    resource_type = db.Column(db.String(80), default="")
    resource_id = db.Column(db.Integer, nullable=True)
    reason = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


class SubscriptionPlan(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(80), nullable=False, default="")
    price_monthly = db.Column(db.Integer, default=0)
    monthly_pages = db.Column(db.Integer, default=0)
    monthly_ai_calls = db.Column(db.Integer, default=0)
    export_enabled = db.Column(db.Boolean, default=False)
    description = db.Column(db.Text, default="")
    is_active = db.Column(db.Boolean, default=True)


class UserSubscription(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False)
    plan_id = db.Column(db.Integer, nullable=False)
    start_date = db.Column(db.String(40), default="")
    end_date = db.Column(db.String(40), default="")
    status = db.Column(db.String(40), default="active")
    auto_renew = db.Column(db.Boolean, default=False)


class UsageRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False)
    action_type = db.Column(db.String(80), default="")
    units_used = db.Column(db.Integer, default=1)
    related_document_id = db.Column(db.Integer, nullable=True)
    related_term_id = db.Column(db.Integer, nullable=True)
    created_at = db.Column(db.String(40), default="")


class BillingRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, nullable=False)
    plan_id = db.Column(db.Integer, nullable=False)
    amount = db.Column(db.Integer, default=0)
    payment_method = db.Column(db.String(80), default="mock_payment")
    payment_status = db.Column(db.String(40), default="paid")
    created_at = db.Column(db.String(40), default="")


class IngestionJob(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    source_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    status = db.Column(db.String(40), default="queued")
    started_at = db.Column(db.String(40), default="")
    finished_at = db.Column(db.String(40), default="")
    error_message = db.Column(db.Text, default="")
    processed_pages = db.Column(db.Integer, default=0)
    created_by = db.Column(db.Integer, default=0)


class BackgroundJob(db.Model):
    """
    Local async job record for long-running MVP workflows.
    The worker is intentionally SQLite-friendly and single-process for this stage.
    """
    id = db.Column(db.Integer, primary_key=True)
    job_type = db.Column(db.String(80), nullable=False)
    status = db.Column(db.String(40), default="queued")
    priority = db.Column(db.Integer, default=100)
    created_by = db.Column(db.Integer, default=0)
    course_id = db.Column(db.Integer, nullable=True)
    document_id = db.Column(db.Integer, nullable=True)
    alignment_run_id = db.Column(db.Integer, nullable=True)
    evaluation_run_id = db.Column(db.Integer, nullable=True)
    scope_type = db.Column(db.String(30), default="")
    owner_user_id = db.Column(db.Integer, nullable=True)
    input_json = db.Column(db.Text, default="{}")
    result_json = db.Column(db.Text, default="{}")
    progress_current = db.Column(db.Integer, default=0)
    progress_total = db.Column(db.Integer, default=0)
    progress_message = db.Column(db.Text, default="")
    error_code = db.Column(db.String(80), default="")
    error_message = db.Column(db.Text, default="")
    attempt_count = db.Column(db.Integer, default=0)
    max_attempts = db.Column(db.Integer, default=JOB_MAX_ATTEMPTS)
    started_at = db.Column(db.String(40), default="")
    finished_at = db.Column(db.String(40), default="")
    canceled_at = db.Column(db.String(40), default="")
    created_at = db.Column(db.String(40), default="")
    updated_at = db.Column(db.String(40), default="")
    locked_by = db.Column(db.String(120), default="")
    locked_at = db.Column(db.String(40), default="")


class BackgroundJobEvent(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    job_id = db.Column(db.Integer, nullable=False)
    event_type = db.Column(db.String(80), default="")
    message = db.Column(db.Text, default="")
    progress_current = db.Column(db.Integer, default=0)
    progress_total = db.Column(db.Integer, default=0)
    metadata_json = db.Column(db.Text, default="{}")
    created_at = db.Column(db.String(40), default="")


class SystemLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    level = db.Column(db.String(30), default="info")
    module = db.Column(db.String(80), default="")
    message = db.Column(db.Text, default="")
    created_at = db.Column(db.String(40), default="")


# ============================================================
# 工具函数：文件类型、文本清洗、文本解析
# ============================================================

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def api_error(error_code, message, http_status=None, details=None):
    code = str(error_code or "INTERNAL_ERROR")
    status_code = int(http_status or ERROR_CODES.get(code, 500))
    details = details or {}
    payload = {
        "status": "error",
        "error_code": code,
        "message": str(message or ""),
        "details": details,
    }
    if isinstance(details, dict):
        payload.update({key: value for key, value in details.items() if key not in payload})
    return jsonify(payload), status_code


def api_success(data=None, message="Operation completed.", **extra):
    payload = {
        "status": "success",
        "message": message,
        "data": data or {},
    }
    payload.update(extra)
    return jsonify(payload)


@app.errorhandler(RequestEntityTooLarge)
def handle_file_too_large(exc):
    return api_error(
        "FILE_TOO_LARGE",
        f"File exceeds MAX_UPLOAD_SIZE_MB={MAX_UPLOAD_SIZE_MB}.",
        413,
        {"max_upload_size_mb": MAX_UPLOAD_SIZE_MB},
    )


@app.errorhandler(404)
def handle_not_found(exc):
    if request.path.startswith("/api/"):
        return api_error("RESOURCE_NOT_FOUND", "API resource not found.", 404)
    return exc


@app.errorhandler(500)
def handle_internal_error(exc):
    if request.path.startswith("/api/"):
        return api_error("INTERNAL_ERROR", "Internal server error.", 500)
    return exc


def normalize_knowledge_base_type(value, language="zh"):
    raw_value = str(value or "").strip()

    if raw_value in KNOWLEDGE_BASE_TYPES:
        return raw_value

    if str(language or "").strip().lower() == "en":
        return "en_course_kb"

    return "zh_course_kb"


def visibility_for_kb_type(knowledge_base_type):
    if knowledge_base_type == "student_personal_kb":
        return "private"

    return "course"


def visibility_for_scope(scope_type):
    if scope_type == "personal":
        return "private"
    if scope_type == "global":
        return "admin"
    return "course"


def create_storage_object_from_metadata(metadata, visibility="private", status="active"):
    storage = StorageObject(
        storage_backend=metadata.get("storage_backend", STORAGE_BACKEND),
        bucket=metadata.get("bucket", ""),
        storage_key=metadata.get("storage_key", ""),
        original_filename=metadata.get("original_filename", ""),
        content_type=metadata.get("content_type", ""),
        size_bytes=int(metadata.get("size_bytes", 0) or 0),
        sha256=metadata.get("sha256", ""),
        owner_user_id=metadata.get("owner_user_id"),
        course_id=metadata.get("course_id"),
        document_id=metadata.get("document_id"),
        visibility=visibility,
        purpose=metadata.get("purpose", "uploaded_document"),
        created_at=current_time_text(),
        updated_at=current_time_text(),
        status=status,
    )
    db.session.add(storage)
    db.session.flush()
    return storage


def storage_service():
    return StorageService()


def resolve_document_local_path(document):
    if getattr(document, "storage_key", ""):
        return storage_service().absolute_path(document.storage_key)
    return os.path.join(UPLOAD_FOLDER, document.saved_filename)


def ensure_schema_columns():
    """
    Minimal SQLite-compatible schema upgrade for the local demo database.
    """
    if not app.config["SQLALCHEMY_DATABASE_URI"].startswith("sqlite"):
        return

    table_columns = {
        "user": {
            "email": "VARCHAR(160) DEFAULT ''",
            "display_name": "VARCHAR(120) DEFAULT ''",
            "is_verified": "BOOLEAN DEFAULT 0",
            "verification_token": "VARCHAR(160) DEFAULT ''",
            "verification_token_expires_at": "VARCHAR(40) DEFAULT ''",
            "reset_token": "VARCHAR(160) DEFAULT ''",
            "reset_token_expires_at": "VARCHAR(40) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "last_login_at": "VARCHAR(40) DEFAULT ''"
        },
        "auth_token": {
            "expires_at": "VARCHAR(40) DEFAULT ''",
            "token_hash": "VARCHAR(128) DEFAULT ''",
            "last_used_at": "VARCHAR(40) DEFAULT ''"
        },
        "course": {
            "course_code": "VARCHAR(80) DEFAULT ''",
            "language_mode": "VARCHAR(60) DEFAULT 'bilingual'",
            "teacher_id": "INTEGER DEFAULT 0",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "deleted_at": "VARCHAR(40) DEFAULT ''"
        },
        "course_member": {
            "role_in_course": "VARCHAR(30) DEFAULT 'student'",
            "joined_at": "VARCHAR(40) DEFAULT ''",
            "status": "VARCHAR(40) DEFAULT 'active'"
        },
        "term": {
            "ai_status": "VARCHAR(40) DEFAULT 'pending'",
            "english_evidence": "TEXT DEFAULT ''",
            "chinese_evidence": "TEXT DEFAULT ''",
            "risk_note": "TEXT DEFAULT ''",
            "ai_model": "VARCHAR(80) DEFAULT ''",
            "courseware_sentence": "TEXT DEFAULT ''",
            "english_kb_evidence": "TEXT DEFAULT ''",
            "ai_translation_candidate": "VARCHAR(200) DEFAULT ''",
            "chinese_kb_evidence": "TEXT DEFAULT ''",
            "final_chinese_term": "VARCHAR(200) DEFAULT ''",
            "alignment_reason": "TEXT DEFAULT ''",
            "review_status": "VARCHAR(40) DEFAULT 'pending'",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "parse_block_uid": "VARCHAR(64) DEFAULT ''",
            "parse_quality_status": "VARCHAR(80) DEFAULT ''",
            "parse_quality_flags": "TEXT DEFAULT '[]'",
            "input_risk_labels": "TEXT DEFAULT '[]'",
            "source_uid": "VARCHAR(64) DEFAULT ''",
            "chunk_uid": "VARCHAR(64) DEFAULT ''"
        },
        "knowledge_document": {
            "knowledge_base_type": "VARCHAR(40) DEFAULT 'zh_course_kb'",
            "owner_user_id": "VARCHAR(80) DEFAULT ''",
            "visibility": "VARCHAR(30) DEFAULT 'course'",
            "parse_uid": "VARCHAR(64) DEFAULT ''"
        },
        "knowledge_chunk": {
            "chunk_uid": "VARCHAR(64) DEFAULT ''",
            "source_uid": "VARCHAR(64) DEFAULT ''",
            "language": "VARCHAR(30) DEFAULT ''",
            "knowledge_base_type": "VARCHAR(40) DEFAULT 'zh_course_kb'",
            "owner_user_id": "VARCHAR(80) DEFAULT ''",
            "visibility": "VARCHAR(30) DEFAULT 'course'",
            "source_id": "INTEGER",
            "knowledge_source_id": "INTEGER",
            "knowledge_base_version_id": "INTEGER",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "parse_block_uid": "VARCHAR(64) DEFAULT ''",
            "course_id": "INTEGER",
            "scope_type": "VARCHAR(30) DEFAULT 'course'",
            "discipline": "VARCHAR(120) DEFAULT ''",
            "chapter": "VARCHAR(120) DEFAULT ''",
            "page_number": "INTEGER",
            "source_slide": "VARCHAR(80) DEFAULT ''",
            "source_section": "VARCHAR(160) DEFAULT ''",
            "source_locator": "VARCHAR(160) DEFAULT ''",
            "slide_number": "INTEGER",
            "block_type": "VARCHAR(40) DEFAULT 'text'",
            "token_count": "INTEGER",
            "char_count": "INTEGER",
            "keywords": "TEXT DEFAULT ''",
            "source_citation": "TEXT DEFAULT ''",
            "embedding_id": "VARCHAR(120) DEFAULT ''",
            "normalized_text": "TEXT DEFAULT ''",
            "content_hash": "VARCHAR(64) DEFAULT ''",
            "formula_block_ids_json": "TEXT DEFAULT '[]'",
            "index_status": "VARCHAR(40) DEFAULT 'indexed'",
            "quality_status": "VARCHAR(80) DEFAULT ''",
            "quality_flags": "TEXT DEFAULT '[]'",
            "trust_level": "VARCHAR(60) DEFAULT 'unknown'",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "embedding_status": "VARCHAR(40) DEFAULT 'not_started'",
            "is_duplicate": "BOOLEAN DEFAULT 0",
            "duplicate_of_chunk_id": "INTEGER",
            "is_active": "BOOLEAN DEFAULT 1",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "document": {
            "ocr_provider": "VARCHAR(40) DEFAULT ''",
            "ocr_status": "VARCHAR(60) DEFAULT ''",
            "ocr_error": "TEXT DEFAULT ''",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "file_sha256": "VARCHAR(64) DEFAULT ''",
            "storage_object_id": "INTEGER",
            "storage_backend": "VARCHAR(40) DEFAULT ''",
            "storage_key": "VARCHAR(600) DEFAULT ''",
            "original_filename": "VARCHAR(260) DEFAULT ''",
            "content_type": "VARCHAR(160) DEFAULT ''",
            "size_bytes": "INTEGER DEFAULT 0",
            "sha256": "VARCHAR(64) DEFAULT ''",
            "quality_flags_json": "TEXT DEFAULT '[]'",
            "deleted_at": "VARCHAR(40) DEFAULT ''"
        },
        "document_chunk": {
            "ocr_provider": "VARCHAR(40) DEFAULT ''",
            "ocr_status": "VARCHAR(60) DEFAULT ''",
            "ocr_error": "TEXT DEFAULT ''",
            "owner_user_id": "INTEGER DEFAULT 0",
            "chunk_index": "INTEGER DEFAULT 0",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "parse_block_uid": "VARCHAR(64) DEFAULT ''",
            "quality_flags_json": "TEXT DEFAULT '[]'"
        },
        "courseware_upload": {
            "parse_uid": "VARCHAR(64) DEFAULT ''"
        },
        "document_parse_record": {
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "source_filename": "VARCHAR(260) DEFAULT ''",
            "stored_path": "VARCHAR(600) DEFAULT ''",
            "file_type": "VARCHAR(40) DEFAULT 'unknown'",
            "mime_type": "VARCHAR(160) DEFAULT ''",
            "file_size_bytes": "INTEGER",
            "parser_name": "VARCHAR(120) DEFAULT ''",
            "parser_version": "VARCHAR(80) DEFAULT ''",
            "parse_status": "VARCHAR(40) DEFAULT 'failed'",
            "quality_status": "VARCHAR(80) DEFAULT 'parse_failed'",
            "quality_flags": "TEXT DEFAULT '[]'",
            "page_count": "INTEGER",
            "block_count": "INTEGER DEFAULT 0",
            "extracted_text_chars": "INTEGER DEFAULT 0",
            "ocr_required": "BOOLEAN DEFAULT 0",
            "ocr_available": "BOOLEAN DEFAULT 0",
            "formula_detected": "BOOLEAN DEFAULT 0",
            "table_detected": "BOOLEAN DEFAULT 0",
            "image_only_suspected": "BOOLEAN DEFAULT 0",
            "error_code": "VARCHAR(120) DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "warnings": "TEXT DEFAULT '[]'",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "document_parse_block": {
            "block_uid": "VARCHAR(64) DEFAULT ''",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "page_number": "INTEGER",
            "slide_number": "INTEGER",
            "block_index": "INTEGER DEFAULT 0",
            "block_type": "VARCHAR(40) DEFAULT 'text'",
            "text": "TEXT DEFAULT ''",
            "confidence": "FLOAT",
            "parser_type": "VARCHAR(40) DEFAULT 'native'",
            "source_locator": "VARCHAR(160) DEFAULT ''",
            "quality_flags": "TEXT DEFAULT '[]'",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "formula_block": {
            "document_id": "INTEGER DEFAULT 0",
            "course_id": "INTEGER",
            "owner_user_id": "INTEGER DEFAULT 0",
            "scope_type": "VARCHAR(30) DEFAULT 'course'",
            "page_number": "INTEGER",
            "slide_number": "INTEGER",
            "bbox_json": "TEXT DEFAULT '{}'",
            "image_path": "VARCHAR(500) DEFAULT ''",
            "image_storage_object_id": "INTEGER",
            "image_storage_key": "VARCHAR(600) DEFAULT ''",
            "image_content_type": "VARCHAR(160) DEFAULT ''",
            "image_sha256": "VARCHAR(64) DEFAULT ''",
            "latex": "TEXT DEFAULT ''",
            "plain_text": "TEXT DEFAULT ''",
            "provider": "VARCHAR(80) DEFAULT ''",
            "confidence": "FLOAT DEFAULT 0",
            "status": "VARCHAR(80) DEFAULT ''",
            "error": "TEXT DEFAULT ''",
            "quality_flags_json": "TEXT DEFAULT '[]'",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "terminology_card": {
            "ai_model": "VARCHAR(80) DEFAULT ''",
            "risk_note": "TEXT DEFAULT ''",
            "normalized_english_term": "VARCHAR(255) DEFAULT ''",
            "normalized_chinese_term": "VARCHAR(255) DEFAULT ''",
            "english_evidence_snapshot": "TEXT DEFAULT ''",
            "chinese_evidence_snapshot": "TEXT DEFAULT ''",
            "english_evidence_score": "FLOAT DEFAULT 0",
            "chinese_evidence_score": "FLOAT DEFAULT 0",
            "alignment_status": "VARCHAR(64) DEFAULT 'unverified_translation'",
            "score_breakdown_json": "TEXT DEFAULT '{}'",
            "quality_flags_json": "TEXT DEFAULT '[]'",
            "ai_provider": "VARCHAR(80) DEFAULT ''",
            "ai_provider_mode": "VARCHAR(40) DEFAULT ''",
            "prompt_key": "VARCHAR(80) DEFAULT ''",
            "prompt_version": "VARCHAR(80) DEFAULT ''",
            "retrieval_version": "VARCHAR(80) DEFAULT ''",
            "knowledge_base_version_id": "INTEGER",
            "english_kb_version_id": "INTEGER",
            "chinese_kb_version_id": "INTEGER",
            "retrieval_run_id": "INTEGER",
            "index_version": "VARCHAR(80) DEFAULT ''",
            "evidence_content_hashes_json": "TEXT DEFAULT '[]'",
            "evidence_status": "VARCHAR(40) DEFAULT 'evidence_current'",
            "ai_call_log_id": "INTEGER",
            "alignment_run_id": "INTEGER",
            "source_alignment_run_id": "INTEGER",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "parse_block_uid": "VARCHAR(64) DEFAULT ''",
            "parse_quality_status": "VARCHAR(80) DEFAULT ''",
            "parse_quality_flags": "TEXT DEFAULT '[]'",
            "input_risk_labels": "TEXT DEFAULT '[]'",
            "approved_by": "INTEGER",
            "approved_at": "VARCHAR(40) DEFAULT ''",
            "rejected_reason": "TEXT DEFAULT ''",
            "reviewer_note": "TEXT DEFAULT ''"
        },
        "concept_alignment_card": {
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "english_term": "VARCHAR(220) DEFAULT ''",
            "chinese_term": "VARCHAR(220) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "concept_scope": "TEXT DEFAULT ''",
            "english_explanation": "TEXT DEFAULT ''",
            "chinese_explanation": "TEXT DEFAULT ''",
            "english_evidence": "TEXT DEFAULT '[]'",
            "chinese_evidence": "TEXT DEFAULT '[]'",
            "alignment_reason": "TEXT DEFAULT ''",
            "confidence_score": "FLOAT",
            "risk_labels": "TEXT DEFAULT '[]'",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "parse_block_uid": "VARCHAR(64) DEFAULT ''",
            "parse_quality_status": "VARCHAR(80) DEFAULT ''",
            "parse_quality_flags": "TEXT DEFAULT '[]'",
            "input_risk_labels": "TEXT DEFAULT '[]'",
            "status": "VARCHAR(40) DEFAULT 'draft'",
            "source_document_id": "INTEGER",
            "source_chunk_id": "INTEGER",
            "created_by": "INTEGER",
            "reviewed_by": "INTEGER",
            "reviewed_at": "VARCHAR(40) DEFAULT ''",
            "model_name": "VARCHAR(120) DEFAULT ''",
            "prompt_version": "VARCHAR(80) DEFAULT ''",
            "retrieval_version": "VARCHAR(80) DEFAULT ''",
            "version": "INTEGER DEFAULT 1",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "student_course_membership": {
            "membership_uid": "VARCHAR(64) DEFAULT ''",
            "user_id": "INTEGER DEFAULT 0",
            "course": "VARCHAR(160) DEFAULT ''",
            "role_in_course": "VARCHAR(40) DEFAULT 'student'",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "enrolled_by": "INTEGER",
            "enrolled_at": "VARCHAR(40) DEFAULT ''",
            "revoked_by": "INTEGER",
            "revoked_at": "VARCHAR(40) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "course_student_visibility_policy": {
            "policy_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "visibility": "VARCHAR(40) DEFAULT 'enrolled_only'",
            "allow_auditor_view": "BOOLEAN DEFAULT 0",
            "allow_teacher_preview": "BOOLEAN DEFAULT 1",
            "allow_cross_course_search": "BOOLEAN DEFAULT 0",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "created_by": "INTEGER",
            "updated_by": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "student_concept_card_state": {
            "state_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "user_id": "INTEGER DEFAULT 0",
            "course": "VARCHAR(160) DEFAULT ''",
            "favorited": "BOOLEAN DEFAULT 0",
            "mastered": "BOOLEAN DEFAULT 0",
            "mastered_at": "VARCHAR(40) DEFAULT ''",
            "last_viewed_at": "VARCHAR(40) DEFAULT ''",
            "view_count": "INTEGER DEFAULT 0",
            "personal_note": "TEXT DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "concept_card_review_record": {
            "review_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "reviewer_id": "INTEGER",
            "reviewer_role": "VARCHAR(40) DEFAULT ''",
            "reviewer_name": "VARCHAR(160) DEFAULT ''",
            "action": "VARCHAR(80) DEFAULT ''",
            "previous_status": "VARCHAR(40) DEFAULT ''",
            "new_status": "VARCHAR(40) DEFAULT ''",
            "decision": "VARCHAR(80) DEFAULT ''",
            "reason_code": "VARCHAR(120) DEFAULT ''",
            "review_comment": "TEXT DEFAULT ''",
            "evidence_assessment": "TEXT DEFAULT '{}'",
            "term_assessment": "TEXT DEFAULT '{}'",
            "risk_assessment": "TEXT DEFAULT '{}'",
            "required_changes": "TEXT DEFAULT '[]'",
            "resolved_risk_labels": "TEXT DEFAULT '[]'",
            "remaining_risk_labels": "TEXT DEFAULT '[]'",
            "verification_run_uid": "VARCHAR(64) DEFAULT ''",
            "request_id": "VARCHAR(120) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "concept_card_review_assignment": {
            "assignment_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "assigned_to": "VARCHAR(120) DEFAULT ''",
            "assigned_by": "INTEGER",
            "assignment_status": "VARCHAR(40) DEFAULT 'active'",
            "due_at": "VARCHAR(40) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "course_review_policy": {
            "policy_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "require_human_review": "BOOLEAN DEFAULT 1",
            "require_two_step_review": "BOOLEAN DEFAULT 0",
            "require_admin_for_override": "BOOLEAN DEFAULT 1",
            "allow_teacher_override": "BOOLEAN DEFAULT 0",
            "allow_approve_with_unverified_alignment": "BOOLEAN DEFAULT 0",
            "allow_approve_with_partial_text": "BOOLEAN DEFAULT 0",
            "allow_approve_with_missing_chinese_evidence": "BOOLEAN DEFAULT 0",
            "allow_approve_with_missing_english_evidence": "BOOLEAN DEFAULT 0",
            "blocking_risk_labels": "TEXT DEFAULT '[]'",
            "override_allowed_risk_labels": "TEXT DEFAULT '[]'",
            "override_forbidden_risk_labels": "TEXT DEFAULT '[]'",
            "required_evidence_sides": "VARCHAR(40) DEFAULT 'both'",
            "min_required_evidence_count": "INTEGER DEFAULT 2",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "created_by": "INTEGER",
            "updated_by": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "course_review_permission": {
            "permission_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "reviewer_id": "INTEGER",
            "reviewer_role": "VARCHAR(40) DEFAULT 'teacher'",
            "permission_level": "VARCHAR(40) DEFAULT 'read'",
            "can_review": "BOOLEAN DEFAULT 0",
            "can_approve": "BOOLEAN DEFAULT 0",
            "can_override_risk": "BOOLEAN DEFAULT 0",
            "can_assign_reviewer": "BOOLEAN DEFAULT 0",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "granted_by": "INTEGER",
            "granted_at": "VARCHAR(40) DEFAULT ''",
            "revoked_by": "INTEGER",
            "revoked_at": "VARCHAR(40) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "audit_record": {
            "audit_uid": "VARCHAR(64) DEFAULT ''",
            "event_type": "VARCHAR(120) DEFAULT ''",
            "target_type": "VARCHAR(120) DEFAULT ''",
            "target_uid": "VARCHAR(120) DEFAULT ''",
            "actor_id": "INTEGER",
            "actor_role": "VARCHAR(40) DEFAULT ''",
            "actor_name": "VARCHAR(160) DEFAULT ''",
            "request_id": "VARCHAR(120) DEFAULT ''",
            "source": "VARCHAR(40) DEFAULT 'service'",
            "before_snapshot": "TEXT DEFAULT '{}'",
            "after_snapshot": "TEXT DEFAULT '{}'",
            "input_payload": "TEXT DEFAULT '{}'",
            "output_payload": "TEXT DEFAULT '{}'",
            "changed_fields": "TEXT DEFAULT '[]'",
            "result": "VARCHAR(40) DEFAULT 'success'",
            "error_code": "VARCHAR(120) DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "model_name": "VARCHAR(120) DEFAULT ''",
            "prompt_version": "VARCHAR(80) DEFAULT ''",
            "retrieval_version": "VARCHAR(80) DEFAULT ''",
            "latency_ms": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "alignment_verification_run": {
            "run_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "english_term": "VARCHAR(220) DEFAULT ''",
            "chinese_term": "VARCHAR(220) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "provider_name": "VARCHAR(120) DEFAULT 'mock-rule-v1'",
            "provider_type": "VARCHAR(40) DEFAULT 'mock'",
            "provider_version": "VARCHAR(80) DEFAULT 'v1'",
            "input_payload": "TEXT DEFAULT '{}'",
            "output_payload": "TEXT DEFAULT '{}'",
            "english_evidence_count": "INTEGER DEFAULT 0",
            "chinese_evidence_count": "INTEGER DEFAULT 0",
            "top_english_chunk_uids": "TEXT DEFAULT '[]'",
            "top_chinese_chunk_uids": "TEXT DEFAULT '[]'",
            "retrieval_score_summary": "TEXT DEFAULT '{}'",
            "candidate_score_summary": "TEXT DEFAULT '{}'",
            "alignment_confidence": "FLOAT",
            "verification_status": "VARCHAR(40) DEFAULT 'mock_only'",
            "recommendation": "VARCHAR(80) DEFAULT 'needs_review'",
            "risk_labels": "TEXT DEFAULT '[]'",
            "prompt_version": "VARCHAR(80) DEFAULT ''",
            "prompt_summary": "TEXT DEFAULT '{}'",
            "raw_output_summary": "TEXT DEFAULT '{}'",
            "parser_version": "VARCHAR(80) DEFAULT ''",
            "output_schema_version": "VARCHAR(80) DEFAULT ''",
            "provider_response_status": "VARCHAR(80) DEFAULT ''",
            "error_code": "VARCHAR(120) DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "latency_ms": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "alignment_provider_policy": {
            "policy_uid": "VARCHAR(64) DEFAULT ''",
            "provider_name": "VARCHAR(120) DEFAULT ''",
            "provider_type": "VARCHAR(40) DEFAULT ''",
            "enabled": "BOOLEAN DEFAULT 0",
            "replay_only": "BOOLEAN DEFAULT 1",
            "allow_external_calls": "BOOLEAN DEFAULT 0",
            "allow_attach_to_card": "BOOLEAN DEFAULT 0",
            "allow_production_result": "BOOLEAN DEFAULT 0",
            "allow_auto_approve": "BOOLEAN DEFAULT 0",
            "require_human_review": "BOOLEAN DEFAULT 1",
            "allowed_courses": "TEXT DEFAULT '[]'",
            "blocked_courses": "TEXT DEFAULT '[]'",
            "allowed_roles": "TEXT DEFAULT '[]'",
            "max_calls_per_day": "INTEGER DEFAULT 0",
            "max_calls_per_month": "INTEGER DEFAULT 0",
            "max_estimated_cost_per_call": "FLOAT",
            "max_estimated_cost_per_day": "FLOAT",
            "max_prompt_chars": "INTEGER DEFAULT 8000",
            "max_output_chars": "INTEGER DEFAULT 4000",
            "timeout_seconds": "INTEGER DEFAULT 30",
            "max_retries": "INTEGER DEFAULT 0",
            "status": "VARCHAR(40) DEFAULT 'disabled'",
            "created_by": "INTEGER",
            "updated_by": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "alignment_provider_usage_record": {
            "usage_uid": "VARCHAR(64) DEFAULT ''",
            "provider_name": "VARCHAR(120) DEFAULT ''",
            "provider_type": "VARCHAR(40) DEFAULT ''",
            "run_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "request_id": "VARCHAR(120) DEFAULT ''",
            "estimated_input_tokens": "INTEGER DEFAULT 0",
            "estimated_output_tokens": "INTEGER DEFAULT 0",
            "estimated_cost": "FLOAT DEFAULT 0",
            "actual_cost": "FLOAT",
            "provider_response_status": "VARCHAR(80) DEFAULT ''",
            "error_code": "VARCHAR(120) DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "alignment_provider_preflight_run": {
            "preflight_uid": "VARCHAR(64) DEFAULT ''",
            "provider_name": "VARCHAR(120) DEFAULT ''",
            "provider_type": "VARCHAR(40) DEFAULT ''",
            "policy_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "requested_by": "VARCHAR(120) DEFAULT ''",
            "check_status": "VARCHAR(40) DEFAULT 'failed'",
            "overall_ready": "BOOLEAN DEFAULT 0",
            "external_calls_enabled": "BOOLEAN DEFAULT 0",
            "replay_only": "BOOLEAN DEFAULT 1",
            "api_key_present": "BOOLEAN DEFAULT 0",
            "api_key_env_name": "VARCHAR(120) DEFAULT ''",
            "policy_summary": "TEXT DEFAULT '{}'",
            "check_results": "TEXT DEFAULT '{}'",
            "blocking_reasons": "TEXT DEFAULT '[]'",
            "warnings": "TEXT DEFAULT '[]'",
            "replay_dry_run_status": "VARCHAR(40) DEFAULT 'not_run'",
            "estimated_cost_per_call": "FLOAT",
            "max_estimated_cost_per_call": "FLOAT",
            "max_calls_per_day": "INTEGER DEFAULT 0",
            "max_calls_per_month": "INTEGER DEFAULT 0",
            "require_human_review": "BOOLEAN DEFAULT 1",
            "allow_auto_approve": "BOOLEAN DEFAULT 0",
            "allow_production_result": "BOOLEAN DEFAULT 0",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "alignment_run": {
            "ai_provider": "VARCHAR(80) DEFAULT ''",
            "ai_provider_mode": "VARCHAR(40) DEFAULT ''",
            "ai_model": "VARCHAR(120) DEFAULT ''",
            "prompt_key": "VARCHAR(80) DEFAULT ''",
            "term_count": "INTEGER DEFAULT 0",
            "card_created_count": "INTEGER DEFAULT 0",
            "auto_approved_count": "INTEGER DEFAULT 0",
            "qc_count": "INTEGER DEFAULT 0",
            "needs_evidence_count": "INTEGER DEFAULT 0",
            "conflict_count": "INTEGER DEFAULT 0",
            "failed_count": "INTEGER DEFAULT 0"
        },
        "feedback": {
            "feedback_uid": "VARCHAR(64) DEFAULT ''",
            "user_id": "INTEGER DEFAULT 0",
            "course_id": "INTEGER",
            "document_id": "INTEGER",
            "terminology_card_id": "INTEGER",
            "formula_block_id": "INTEGER",
            "job_id": "INTEGER",
            "alignment_run_id": "INTEGER",
            "evaluation_run_id": "INTEGER",
            "user_role": "VARCHAR(40) DEFAULT ''",
            "course": "VARCHAR(120) DEFAULT ''",
            "chapter": "VARCHAR(120) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "english_term": "VARCHAR(200) DEFAULT ''",
            "chinese_term": "VARCHAR(200) DEFAULT ''",
            "status": "VARCHAR(30) DEFAULT 'open'",
            "resolved_at": "VARCHAR(40) DEFAULT ''",
            "severity": "VARCHAR(40) DEFAULT 'normal'",
            "priority": "VARCHAR(20) DEFAULT 'P2'",
            "message": "TEXT DEFAULT ''",
            "suggested_chinese_term": "VARCHAR(220) DEFAULT ''",
            "resolved_by": "INTEGER",
            "feedback_source": "VARCHAR(80) DEFAULT 'student_card_detail'",
            "reported_issue": "TEXT DEFAULT ''",
            "expected_result": "TEXT DEFAULT ''",
            "actual_result": "TEXT DEFAULT ''",
            "evidence_comment": "TEXT DEFAULT ''",
            "screenshot_path": "VARCHAR(500) DEFAULT ''",
            "classification": "VARCHAR(80) DEFAULT ''",
            "root_cause": "VARCHAR(80) DEFAULT ''",
            "resolution_action": "VARCHAR(80) DEFAULT ''",
            "resolution_note": "TEXT DEFAULT ''",
            "handled_by": "INTEGER",
            "handled_at": "VARCHAR(40) DEFAULT ''",
            "handler_role": "VARCHAR(40) DEFAULT ''",
            "teacher_note": "TEXT DEFAULT ''",
            "linked_review_uid": "VARCHAR(64) DEFAULT ''",
            "linked_card_uid": "VARCHAR(64) DEFAULT ''",
            "converted_to_evaluation_item_id": "INTEGER",
            "linked_backlog_item_id": "INTEGER",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''"
        },
        "concept_card_feedback_triage_record": {
            "triage_uid": "VARCHAR(64) DEFAULT ''",
            "feedback_uid": "VARCHAR(64) DEFAULT ''",
            "card_uid": "VARCHAR(64) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "action": "VARCHAR(80) DEFAULT ''",
            "previous_status": "VARCHAR(40) DEFAULT ''",
            "new_status": "VARCHAR(40) DEFAULT ''",
            "handled_by": "INTEGER",
            "handler_role": "VARCHAR(40) DEFAULT ''",
            "reason_code": "VARCHAR(120) DEFAULT ''",
            "teacher_note": "TEXT DEFAULT ''",
            "linked_review_uid": "VARCHAR(64) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''"
        },
        "iteration_backlog_item": {
            "title": "VARCHAR(220) DEFAULT ''",
            "description": "TEXT DEFAULT ''",
            "source_type": "VARCHAR(80) DEFAULT 'feedback'",
            "source_feedback_id": "INTEGER",
            "course_id": "INTEGER",
            "severity": "VARCHAR(40) DEFAULT 'medium'",
            "priority": "VARCHAR(20) DEFAULT 'P2'",
            "category": "VARCHAR(80) DEFAULT 'documentation'",
            "status": "VARCHAR(40) DEFAULT 'open'",
            "owner": "VARCHAR(120) DEFAULT ''",
            "target_pr": "VARCHAR(80) DEFAULT ''",
            "acceptance_criteria": "TEXT DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''",
            "closed_at": "VARCHAR(40) DEFAULT ''"
        },
        "storage_object": {
            "storage_backend": "VARCHAR(40) DEFAULT 'local'",
            "bucket": "VARCHAR(160) DEFAULT ''",
            "storage_key": "VARCHAR(600) DEFAULT ''",
            "original_filename": "VARCHAR(260) DEFAULT ''",
            "content_type": "VARCHAR(160) DEFAULT ''",
            "size_bytes": "INTEGER DEFAULT 0",
            "sha256": "VARCHAR(64) DEFAULT ''",
            "owner_user_id": "INTEGER",
            "course_id": "INTEGER",
            "document_id": "INTEGER",
            "visibility": "VARCHAR(40) DEFAULT 'private'",
            "purpose": "VARCHAR(80) DEFAULT 'uploaded_document'",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''",
            "status": "VARCHAR(40) DEFAULT 'active'"
        },
        "knowledge_base_version": {
            "owner_user_id": "INTEGER",
            "scope_type": "VARCHAR(30) DEFAULT 'course'",
            "version_number": "INTEGER DEFAULT 1",
            "status": "VARCHAR(40) DEFAULT 'draft'",
            "published_at": "VARCHAR(40) DEFAULT ''",
            "archived_at": "VARCHAR(40) DEFAULT ''",
            "parent_version_id": "INTEGER",
            "formula_block_count": "INTEGER DEFAULT 0",
            "deduped_chunk_count": "INTEGER DEFAULT 0",
            "index_backend": "VARCHAR(80) DEFAULT 'local_lexical'",
            "index_version": "VARCHAR(80) DEFAULT 'local_lexical_v1'",
            "retrieval_version": "VARCHAR(80) DEFAULT 'local_lexical_v1'",
            "embedding_provider": "VARCHAR(80) DEFAULT ''",
            "embedding_model": "VARCHAR(160) DEFAULT ''",
            "embedding_dimension": "INTEGER DEFAULT 0",
            "vector_index_status": "VARCHAR(40) DEFAULT ''",
            "vector_index_updated_at": "VARCHAR(40) DEFAULT ''",
            "evaluation_run_id": "INTEGER",
            "quality_gate_status": "VARCHAR(40) DEFAULT ''",
            "manifest_json": "TEXT DEFAULT '{}'"
        },
        "knowledge_source": {
            "source_uid": "VARCHAR(64) DEFAULT ''",
            "title": "VARCHAR(220) DEFAULT ''",
            "course": "VARCHAR(160) DEFAULT ''",
            "chapter": "VARCHAR(160) DEFAULT ''",
            "course_id": "INTEGER",
            "scope_type": "VARCHAR(30) DEFAULT 'course'",
            "owner_user_id": "INTEGER",
            "document_id": "INTEGER",
            "source_title": "VARCHAR(220) DEFAULT ''",
            "source_role": "VARCHAR(80) DEFAULT 'unknown'",
            "owner_type": "VARCHAR(40) DEFAULT 'unknown'",
            "owner_id": "VARCHAR(80) DEFAULT ''",
            "visibility": "VARCHAR(40) DEFAULT 'course'",
            "trust_level": "VARCHAR(60) DEFAULT 'unknown'",
            "parse_uid": "VARCHAR(64) DEFAULT ''",
            "source_filename": "VARCHAR(260) DEFAULT ''",
            "file_type": "VARCHAR(40) DEFAULT 'unknown'",
            "content_hash": "VARCHAR(64) DEFAULT ''",
            "version": "INTEGER DEFAULT 1",
            "license_note": "TEXT DEFAULT ''",
            "quality_status": "VARCHAR(80) DEFAULT ''",
            "quality_flags": "TEXT DEFAULT '[]'",
            "knowledge_base_type": "VARCHAR(40) DEFAULT ''",
            "license_type": "VARCHAR(80) DEFAULT 'unknown'",
            "authorization_status": "VARCHAR(80) DEFAULT 'unknown'",
            "source_quality": "FLOAT DEFAULT 0.4",
            "version_introduced_id": "INTEGER",
            "version_removed_id": "INTEGER",
            "status": "VARCHAR(40) DEFAULT 'active'",
            "effective_from": "VARCHAR(40) DEFAULT ''",
            "effective_to": "VARCHAR(40) DEFAULT ''"
        },
        "evaluation_set": {
            "course_id": "INTEGER",
            "discipline": "VARCHAR(120) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''",
            "is_locked": "BOOLEAN DEFAULT 0"
        },
        "evaluation_item": {
            "evaluation_set_id": "INTEGER",
            "item_id": "VARCHAR(80) DEFAULT ''",
            "split": "VARCHAR(40) DEFAULT 'test'",
            "discipline": "VARCHAR(120) DEFAULT ''",
            "expected_english_evidence": "TEXT DEFAULT ''",
            "expected_chinese_evidence": "TEXT DEFAULT ''",
            "negative_english_evidence": "TEXT DEFAULT ''",
            "negative_chinese_evidence": "TEXT DEFAULT ''",
            "tags_json": "TEXT DEFAULT '[]'",
            "annotator": "VARCHAR(120) DEFAULT ''",
            "reviewed_by": "VARCHAR(120) DEFAULT ''",
            "disagreement_note": "TEXT DEFAULT ''",
            "version": "VARCHAR(40) DEFAULT 'v1'"
        },
        "evaluation_run": {
            "provider_name": "VARCHAR(80) DEFAULT ''",
            "provider_mode": "VARCHAR(40) DEFAULT ''",
            "prompt_key": "VARCHAR(80) DEFAULT ''",
            "model_version": "VARCHAR(120) DEFAULT ''",
            "alignment_version": "VARCHAR(80) DEFAULT ''",
            "commit_hash": "VARCHAR(80) DEFAULT ''",
            "split": "VARCHAR(40) DEFAULT 'test'",
            "skipped_count": "INTEGER DEFAULT 0",
            "extraction_precision": "FLOAT DEFAULT 0",
            "extraction_recall": "FLOAT DEFAULT 0",
            "evidence_accuracy": "FLOAT DEFAULT 0",
            "english_evidence_accuracy": "FLOAT DEFAULT 0",
            "chinese_evidence_accuracy": "FLOAT DEFAULT 0",
            "alignment_accuracy": "FLOAT DEFAULT 0",
            "false_positive_rate": "FLOAT DEFAULT 0",
            "auto_approval_error_rate": "FLOAT DEFAULT 0",
            "ocr_noise_term_rate": "FLOAT",
            "no_evidence_forced_alignment_rate": "FLOAT DEFAULT 0",
            "created_by": "INTEGER DEFAULT 0",
            "report_markdown": "TEXT DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "finished_at": "VARCHAR(40) DEFAULT ''"
        },
        "background_job": {
            "job_type": "VARCHAR(80) DEFAULT ''",
            "status": "VARCHAR(40) DEFAULT 'queued'",
            "priority": "INTEGER DEFAULT 100",
            "created_by": "INTEGER DEFAULT 0",
            "course_id": "INTEGER",
            "document_id": "INTEGER",
            "alignment_run_id": "INTEGER",
            "evaluation_run_id": "INTEGER",
            "scope_type": "VARCHAR(30) DEFAULT ''",
            "owner_user_id": "INTEGER",
            "input_json": "TEXT DEFAULT '{}'",
            "result_json": "TEXT DEFAULT '{}'",
            "progress_current": "INTEGER DEFAULT 0",
            "progress_total": "INTEGER DEFAULT 0",
            "progress_message": "TEXT DEFAULT ''",
            "error_code": "VARCHAR(80) DEFAULT ''",
            "error_message": "TEXT DEFAULT ''",
            "attempt_count": "INTEGER DEFAULT 0",
            "max_attempts": "INTEGER DEFAULT 3",
            "started_at": "VARCHAR(40) DEFAULT ''",
            "finished_at": "VARCHAR(40) DEFAULT ''",
            "canceled_at": "VARCHAR(40) DEFAULT ''",
            "created_at": "VARCHAR(40) DEFAULT ''",
            "updated_at": "VARCHAR(40) DEFAULT ''",
            "locked_by": "VARCHAR(120) DEFAULT ''",
            "locked_at": "VARCHAR(40) DEFAULT ''"
        },
        "background_job_event": {
            "job_id": "INTEGER DEFAULT 0",
            "event_type": "VARCHAR(80) DEFAULT ''",
            "message": "TEXT DEFAULT ''",
            "progress_current": "INTEGER DEFAULT 0",
            "progress_total": "INTEGER DEFAULT 0",
            "metadata_json": "TEXT DEFAULT '{}'",
            "created_at": "VARCHAR(40) DEFAULT ''"
        }
    }

    for table_name, expected_columns in table_columns.items():
        existing = {
            row[1]
            for row in db.session.execute(db.text(f"PRAGMA table_info({table_name})")).fetchall()
        }

        for column_name, column_definition in expected_columns.items():
            if column_name not in existing:
                db.session.execute(
                    db.text(f"ALTER TABLE {table_name} ADD COLUMN {column_name} {column_definition}")
                )

    db.session.commit()


MATH_SYMBOL_MAP = {
    "∑": " sum ",
    "∫": " integral ",
    "√": " sqrt ",
    "∞": " infinity ",
    "π": " pi ",
    "θ": " theta ",
    "λ": " lambda ",
    "μ": " mu ",
    "σ": " sigma ",
    "ω": " omega ",
    "Ω": " ohm ",
    "Δ": " delta ",
    "∂": " partial ",
    "≤": " <= ",
    "≥": " >= ",
    "≠": " != ",
    "≈": " approx ",
    "→": " -> ",
    "⇒": " => ",
    "×": " x ",
    "÷": " / ",
}

MATH_PHYSICS_SYNONYMS = {
    "angular frequency": ["omega", "ω"],
    "wavelength": ["lambda", "λ"],
    "summation": ["sum", "sigma", "∑"],
    "integral": ["int", "∫"],
    "infinity": ["∞"],
    "pi": ["π"],
    "force": ["newton", "f=ma"],
    "energy": ["e=mc", "joule"],
    "velocity": ["speed", "v"],
    "acceleration": ["a", "m/s^2"],
    "frequency": ["hz", "hertz", "f"],
    "resistance": ["ohm", "Ω"],
    "voltage": ["potential difference"],
    "current": ["ampere"],
}


def expand_domain_query(text):
    expanded = normalize_math_symbols(str(text or ""))
    lowered = expanded.lower()
    additions = []

    for phrase, synonyms in MATH_PHYSICS_SYNONYMS.items():
        if phrase in lowered:
            additions.extend(synonyms)

    if additions:
        expanded = f"{expanded} {' '.join(additions)}"

    return expanded


def local_name(tag):
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def normalize_math_symbols(text):
    if not text:
        return ""

    for symbol, replacement in MATH_SYMBOL_MAP.items():
        text = text.replace(symbol, replacement)

    text = text.replace("−", "-")
    text = text.replace("–", "-")
    text = text.replace("—", "-")
    return text


def math_node_text(node):
    """
    Convert common Office Math (OMML) nodes into readable linear text.
    This is not a full TeX renderer; it preserves symbols enough for search/RAG.
    """
    tag = local_name(node.tag)

    if tag in {"t", "delim", "chr"}:
        return normalize_math_symbols(node.text or "")

    children = list(node)

    if tag == "f":
        numerator = ""
        denominator = ""
        for child in children:
            child_tag = local_name(child.tag)
            if child_tag == "num":
                numerator = math_node_text(child)
            elif child_tag == "den":
                denominator = math_node_text(child)
        return f"({numerator})/({denominator})"

    if tag == "sSup":
        base = ""
        sup = ""
        for child in children:
            child_tag = local_name(child.tag)
            if child_tag == "e":
                base = math_node_text(child)
            elif child_tag == "sup":
                sup = math_node_text(child)
        return f"{base}^{sup}"

    if tag == "sSub":
        base = ""
        sub = ""
        for child in children:
            child_tag = local_name(child.tag)
            if child_tag == "e":
                base = math_node_text(child)
            elif child_tag == "sub":
                sub = math_node_text(child)
        return f"{base}_{sub}"

    if tag == "sSubSup":
        base = ""
        sub = ""
        sup = ""
        for child in children:
            child_tag = local_name(child.tag)
            if child_tag == "e":
                base = math_node_text(child)
            elif child_tag == "sub":
                sub = math_node_text(child)
            elif child_tag == "sup":
                sup = math_node_text(child)
        return f"{base}_{sub}^{sup}"

    if tag == "rad":
        expr = ""
        for child in children:
            if local_name(child.tag) == "e":
                expr = math_node_text(child)
        return f"sqrt({expr})"

    if tag == "nary":
        operator = ""
        expr = ""
        sub = ""
        sup = ""
        for child in children:
            child_tag = local_name(child.tag)
            if child_tag == "naryPr":
                operator = math_node_text(child).strip()
            elif child_tag == "e":
                expr = math_node_text(child)
            elif child_tag == "sub":
                sub = math_node_text(child)
            elif child_tag == "sup":
                sup = math_node_text(child)
        operator = operator or "operator"
        bounds = ""
        if sub or sup:
            bounds = f"_{sub}^{sup}"
        return f"{operator}{bounds}({expr})"

    parts = []
    if node.text:
        parts.append(node.text)

    for child in children:
        value = math_node_text(child)
        if value:
            parts.append(value)

    if node.tail:
        parts.append(node.tail)

    return normalize_math_symbols(" ".join(part for part in parts if part).strip())


def extract_math_from_ooxml(path, member_prefix):
    formulas = []

    try:
        with zipfile.ZipFile(path) as archive:
            members = [
                name for name in archive.namelist()
                if name.startswith(member_prefix) and name.endswith(".xml")
            ]

            for member in members:
                try:
                    root = ET.fromstring(archive.read(member))
                except ET.ParseError:
                    continue

                for node in root.iter():
                    if local_name(node.tag) in {"oMath", "oMathPara"}:
                        formula = re.sub(r"\s+", " ", math_node_text(node)).strip()
                        if formula:
                            formulas.append(formula)
    except zipfile.BadZipFile:
        return []

    deduped = []
    seen = set()
    for formula in formulas:
        if formula not in seen:
            seen.add(formula)
            deduped.append(formula)

    return deduped


def clean_text(text):
    """
    对 PDF / Word / PPT 抽取出的文本做基础清洗。
    这里不做语义判断，只做格式规整。
    """
    if not text:
        return ""

    text = normalize_math_symbols(text.replace("\x00", " "))
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


OCR_PLACEHOLDER_RE = re.compile(r"\bOCR[_\s-]*(REQUIRED|FALLBACK)\b|\[(OCR_REQUIRED|OCR_FALLBACK)\]", re.I)
FORMULA_PLACEHOLDER_RE = re.compile(r"\[FormulaBlock\s+#?\d+\]|\bFormulaBlock\b", re.I)


def contains_ocr_placeholder(text):
    return bool(OCR_PLACEHOLDER_RE.search(str(text or "")))


def contains_formula_placeholder(text):
    return bool(FORMULA_PLACEHOLDER_RE.search(str(text or "")))


def safe_chunks_from_text(text, max_chars=700, overlap=80):
    if contains_ocr_placeholder(text) or contains_formula_placeholder(text):
        return []
    return split_text_into_chunks(text, max_chars=max_chars, overlap=overlap)


def extract_pdf_page_texts(path):
    """
    PDF 文本解析。
    依赖 PyMuPDF：pip install pymupdf
    """
    try:
        import fitz
    except ImportError as exc:
        raise ImportError("缺少 PyMuPDF，请运行：pip install pymupdf") from exc

    pages = []
    with fitz.open(path) as doc:
        for index, page in enumerate(doc, start=1):
            page_text = page.get_text("text") or ""
            raw_parts = []
            try:
                raw_dict = page.get_text("rawdict") or {}
                for block in raw_dict.get("blocks", []):
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            chars = span.get("chars", [])
                            if chars:
                                raw_parts.append("".join(ch.get("c", "") for ch in chars))
            except Exception:
                raw_parts = []

            raw_text = "\n".join(part for part in raw_parts if part.strip())
            combined = page_text
            if raw_text and raw_text.strip() not in page_text:
                combined = f"{page_text}\n{raw_text}"

            pages.append({
                "page_number": index,
                "text": clean_text(combined),
                "needs_ocr": len(clean_text(combined)) < 30
            })

    return pages


def extract_text_from_pdf(path):
    parts = []
    for page in extract_pdf_page_texts(path):
        if page["text"]:
            parts.append(f"[Page {page['page_number']}]\n{page['text']}")
    return clean_text("\n\n".join(parts))


def extract_text_from_docx(path):
    """
    DOCX 文本解析。
    依赖 python-docx：pip install python-docx
    """
    try:
        from docx import Document
    except ImportError as exc:
        raise ImportError("缺少 python-docx，请运行：pip install python-docx") from exc

    document = Document(path)
    parts = []

    for paragraph in document.paragraphs:
        content = paragraph.text.strip()
        if content:
            parts.append(content)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    formulas = extract_math_from_ooxml(path, "word/")
    if formulas:
        parts.append("[Office Math Formulas]")
        parts.extend(formulas)

    return clean_text("\n".join(parts))


def extract_text_from_pptx(path):
    """
    PPTX 文本解析。
    依赖 python-pptx：pip install python-pptx
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("缺少 python-pptx，请运行：pip install python-pptx") from exc

    presentation = Presentation(path)
    parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"[Slide {slide_index}]"]

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    slide_parts.append(content)

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    formulas = extract_math_from_ooxml(path, "ppt/slides/")
    if formulas:
        parts.append("[Office Math Formulas]")
        parts.extend(formulas)

    return clean_text("\n\n".join(parts))


def extract_text_from_txt(path):
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            with open(path, "r", encoding=encoding) as file:
                return clean_text(file.read())
        except UnicodeDecodeError:
            continue
    raise ValueError("无法识别文本文件编码。")


def extract_text_from_image(path):
    ocr = get_ocr_provider(OCR_PROVIDER).recognize_image(path)
    if not ocr.ok:
        raise ValueError(f"{ocr.status}: {ocr.error}")
    return clean_text(ocr.text)


def file_requires_ocr(path):
    ext = path.rsplit(".", 1)[1].lower()
    if ext in {"jpg", "jpeg", "png"}:
        return True
    if ext != "pdf":
        return False
    try:
        return any(page["needs_ocr"] for page in extract_pdf_page_texts(path))
    except Exception:
        return True

def extract_text(path):
    ext = path.rsplit(".", 1)[1].lower()

    if ext == "pdf":
        return extract_text_from_pdf(path)

    if ext == "docx":
        return extract_text_from_docx(path)

    if ext == "pptx":
        return extract_text_from_pptx(path)

    if ext in {"txt", "md", "markdown"}:
        return extract_text_from_txt(path)

    if ext in {"jpg", "jpeg", "png"}:
        return extract_text_from_image(path)

    raise ValueError("不支持的文件格式")


def derived_image_folder(document_id=None):
    folder_name = f"document_{document_id}" if document_id else f"pending_{uuid.uuid4().hex[:12]}"
    folder = os.path.join(UPLOAD_FOLDER, "derived", folder_name)
    os.makedirs(folder, exist_ok=True)
    return folder


def render_pdf_page_to_image(pdf_path, page_index, output_dir):
    try:
        import fitz
    except ImportError as exc:
        raise ImportError("缺少 PyMuPDF，请运行：pip install pymupdf") from exc
    with fitz.open(pdf_path) as doc:
        page = doc[page_index]
        zoom = max(PDF_OCR_DPI, 72) / 72
        matrix = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        image_path = os.path.join(output_dir, f"page_{page_index + 1}.png")
        pixmap.save(image_path)
        return image_path


def render_pdf_region_to_image(page, output_dir, page_number, region_index, bbox=None):
    import fitz
    zoom = max(PDF_OCR_DPI, 72) / 72
    matrix = fitz.Matrix(zoom, zoom)
    clip = None
    if bbox:
        clip = fitz.Rect(bbox)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False, clip=clip)
    image_path = os.path.join(output_dir, f"page_{page_number}_region_{region_index}.png")
    pixmap.save(image_path)
    return image_path


def extract_pdf_image_regions(page):
    if not OCR_ENABLE_REGION_EXTRACTION:
        return []

    regions = []
    seen = set()

    def add_region(bbox):
        if len(bbox) != 4:
            return
        width = abs(float(bbox[2]) - float(bbox[0]))
        height = abs(float(bbox[3]) - float(bbox[1]))
        if width < PDF_IMAGE_MIN_WIDTH or height < PDF_IMAGE_MIN_HEIGHT:
            return
        key = tuple(round(float(value), 2) for value in bbox)
        if key in seen:
            return
        seen.add(key)
        regions.append({
            "bbox": [float(value) for value in bbox],
            "width": width,
            "height": height
        })

    try:
        raw_dict = page.get_text("rawdict") or {}
    except Exception:
        raw_dict = {}

    for block in raw_dict.get("blocks", []):
        if block.get("type") != 1:
            continue
        add_region(block.get("bbox") or [])

    try:
        for image_info in page.get_images(full=True):
            xref = image_info[0]
            for rect in page.get_image_rects(xref):
                add_region([rect.x0, rect.y0, rect.x1, rect.y1])
    except Exception:
        pass

    return regions


def ocr_meta_from_result(result):
    return {
        "ocr_provider": getattr(result, "provider", ""),
        "ocr_status": getattr(result, "status", ""),
        "ocr_error": getattr(result, "error", ""),
        "ocr_confidence": getattr(result, "confidence", 0) or 0
    }


def build_ocr_chunks_from_text(text, language, source_type, source_location, page_number=None, slide_number=None, ocr_result=None):
    text = clean_text(text)
    if not text or contains_ocr_placeholder(text):
        return []

    chunks = []
    confidence = getattr(ocr_result, "confidence", 100) if ocr_result else 100
    provider = getattr(ocr_result, "provider", "") if ocr_result else ""
    status = getattr(ocr_result, "status", "not_required") if ocr_result else "not_required"
    error = getattr(ocr_result, "error", "") if ocr_result else ""
    chunk_texts = safe_chunks_from_text(text)
    if not chunk_texts and ocr_result is not None and len(text.strip()) >= 4:
        chunk_texts = [text.strip()]

    for index, content in enumerate(chunk_texts, start=1):
        chunks.append({
            "language": language,
            "page_number": page_number,
            "slide_number": slide_number,
            "section_title": "",
            "content": content,
            "source_type": source_type,
            "source_location": f"{source_location}, chunk {index}" if source_location else f"chunk {index}",
            "ocr_confidence": confidence,
            "ocr_provider": provider,
            "ocr_status": status,
            "ocr_error": error
        })
    return chunks


def process_visual_region(image_path, language, source_type, source_location, page_number=None, slide_number=None, bbox=None):
    ocr_provider = get_ocr_provider(OCR_PROVIDER, language=language)
    formula_provider = get_formula_ocr_provider(FORMULA_OCR_PROVIDER)
    warnings = []
    ocr_required = True
    chunks = []
    formula_blocks = []

    ocr_result = ocr_provider.recognize_image(image_path, language=language)
    if ocr_result.text and not contains_ocr_placeholder(ocr_result.text):
        chunks.extend(build_ocr_chunks_from_text(
            ocr_result.text,
            language,
            source_type,
            source_location,
            page_number=page_number,
            slide_number=slide_number,
            ocr_result=ocr_result
        ))
    elif ocr_result.status in {"ocr_unavailable", "ocr_failed"}:
        warnings.append(ocr_result.error or "OCR engine is unavailable.")

    formula_like = looks_like_formula_image(image_path, ocr_result.text)
    if formula_like:
        formula_result = formula_provider.recognize_formula(image_path, bbox={"bbox": bbox} if bbox else None)
        flags = list(getattr(formula_result, "quality_flags", []) or [])
        if formula_result.status in {"needs_formula_ocr_engine", "formula_ocr_failed", "low_confidence"}:
            flags.append(formula_result.status)
            warnings.append(formula_result.error or formula_result.status)
        formula_blocks.append({
            "page_number": page_number,
            "slide_number": slide_number,
            "bbox": bbox or {},
            "image_path": image_path,
            "latex": formula_result.latex,
            "plain_text": formula_result.plain_text,
            "provider": formula_result.provider,
            "confidence": formula_result.confidence,
            "status": formula_result.status,
            "error": formula_result.error,
            "quality_flags": sorted(set(flags))
        })

    meta = ocr_meta_from_result(ocr_result)
    return chunks, formula_blocks, ocr_required, meta, warnings


def extract_document_chunks(path, language="", source_type="upload", document_id=None):
    ext = path.rsplit(".", 1)[1].lower()
    structured = []
    formula_blocks = []
    ocr_required = False
    ocr_provider_name = ""
    ocr_status = "not_required"
    ocr_error = ""
    ocr_confidence = 100
    warnings = []
    text_parts = []
    output_dir = derived_image_folder(document_id)

    if ext in {"jpg", "jpeg", "png"}:
        chunks, formulas, needs_ocr, meta, region_warnings = process_visual_region(
            path,
            language,
            source_type,
            "image 1",
            page_number=1,
            bbox={}
        )
        structured.extend(chunks)
        formula_blocks.extend(formulas)
        ocr_required = needs_ocr
        ocr_provider_name = meta.get("ocr_provider", "")
        ocr_status = meta.get("ocr_status", "ocr_unavailable")
        ocr_error = meta.get("ocr_error", "")
        ocr_confidence = meta.get("ocr_confidence", 0)
        warnings.extend(region_warnings)
        if structured:
            text_parts.extend(chunk["content"] for chunk in structured)
        return structured, "\n\n".join(text_parts), True, {
            "ocr_provider": ocr_provider_name,
            "ocr_status": ocr_status,
            "ocr_error": ocr_error,
            "ocr_confidence": ocr_confidence,
            "warnings": warnings
        }, formula_blocks

    if ext == "pdf":
        try:
            import fitz
        except ImportError as exc:
            raise ImportError("缺少 PyMuPDF，请运行：pip install pymupdf") from exc

        with fitz.open(path) as doc:
            for page_index, page in enumerate(doc, start=1):
                page_text = clean_text(page.get_text("text") or "")
                if page_text:
                    text_parts.append(f"[Page {page_index}]\n{page_text}")
                    structured.extend(build_ocr_chunks_from_text(
                        page_text,
                        language,
                        source_type,
                        f"page {page_index}",
                        page_number=page_index
                    ))

                regions = extract_pdf_image_regions(page) if PDF_MIXED_PAGE_IMAGE_OCR else []
                for region_index, region in enumerate(regions, start=1):
                    image_path = render_pdf_region_to_image(
                        page,
                        output_dir,
                        page_index,
                        region_index,
                        bbox=region.get("bbox")
                    )
                    chunks, formulas, needs_ocr, meta, region_warnings = process_visual_region(
                        image_path,
                        language,
                        source_type,
                        f"page {page_index}, image region {region_index}",
                        page_number=page_index,
                        bbox=region.get("bbox")
                    )
                    structured.extend(chunks)
                    formula_blocks.extend(formulas)
                    ocr_required = ocr_required or needs_ocr
                    if needs_ocr:
                        ocr_provider_name = meta.get("ocr_provider", ocr_provider_name)
                        ocr_status = meta.get("ocr_status", ocr_status)
                        ocr_error = meta.get("ocr_error", ocr_error)
                        ocr_confidence = min(ocr_confidence, meta.get("ocr_confidence", ocr_confidence) or 0)
                    warnings.extend(region_warnings)

                if not page_text and not regions:
                    image_path = render_pdf_region_to_image(page, output_dir, page_index, 0, bbox=None)
                    chunks, formulas, needs_ocr, meta, region_warnings = process_visual_region(
                        image_path,
                        language,
                        source_type,
                        f"page {page_index}, full-page OCR",
                        page_number=page_index,
                        bbox={}
                    )
                    structured.extend(chunks)
                    formula_blocks.extend(formulas)
                    ocr_required = ocr_required or needs_ocr
                    ocr_provider_name = meta.get("ocr_provider", ocr_provider_name)
                    ocr_status = meta.get("ocr_status", ocr_status)
                    ocr_error = meta.get("ocr_error", ocr_error)
                    ocr_confidence = min(ocr_confidence, meta.get("ocr_confidence", ocr_confidence) or 0)
                    warnings.extend(region_warnings)

        if ocr_required and not structured:
            provider = get_ocr_provider(OCR_PROVIDER, language=language)
            return [], "\n\n".join(text_parts), True, {
                "ocr_provider": ocr_provider_name or provider.provider_name,
                "ocr_status": ocr_status or "ocr_unavailable",
                "ocr_error": ocr_error or "OCR was required but no text could be extracted.",
                "ocr_confidence": ocr_confidence if ocr_confidence != 100 else 0,
                "warnings": warnings
            }, formula_blocks
        return structured, "\n\n".join(text_parts), ocr_required, {
            "ocr_provider": ocr_provider_name,
            "ocr_status": ocr_status if ocr_required else "not_required",
            "ocr_error": ocr_error,
            "ocr_confidence": ocr_confidence if ocr_required else 100,
            "warnings": warnings
        }, formula_blocks

    text = extract_text(path)
    if contains_ocr_placeholder(text):
        return [], "", True, {
            "ocr_provider": "",
            "ocr_status": "parsing_failed",
            "ocr_error": "Parser produced OCR placeholder text; blocked from ingestion.",
            "warnings": ["Parser produced OCR placeholder text; blocked from ingestion."]
        }, formula_blocks
    confidence = 100
    structured = []

    if ext == "pptx":
        slide_splits = re.split(r"\[Slide\s+(\d+)\]", text)
        if len(slide_splits) > 1:
            pairs = list(zip(slide_splits[1::2], slide_splits[2::2]))
            for slide_no, slide_text in pairs:
                for index, content in enumerate(safe_chunks_from_text(slide_text), start=1):
                    structured.append({
                        "language": language,
                        "page_number": None,
                        "slide_number": int(slide_no),
                        "section_title": "",
                        "content": content,
                        "source_type": source_type,
                        "source_location": f"slide {slide_no}, chunk {index}",
                        "ocr_confidence": confidence,
                        "ocr_provider": "",
                        "ocr_status": "not_required",
                        "ocr_error": ""
                    })
            if structured:
                return structured, text, False, {
                    "ocr_provider": "",
                    "ocr_status": "not_required",
                    "ocr_error": "",
                    "ocr_confidence": 100,
                    "warnings": []
                }, formula_blocks

    for index, content in enumerate(safe_chunks_from_text(text), start=1):
        structured.append({
            "language": language,
            "page_number": None,
            "slide_number": None,
            "section_title": "",
            "content": content,
            "source_type": source_type,
            "source_location": f"chunk {index}",
            "ocr_confidence": confidence,
            "ocr_provider": "",
            "ocr_status": "not_required",
            "ocr_error": ""
        })

    return structured, text, False, {
        "ocr_provider": "",
        "ocr_status": "not_required",
        "ocr_error": "",
        "ocr_confidence": 100,
        "warnings": []
    }, formula_blocks

def split_text_into_chunks(text, max_chars=700, overlap=80):
    """
    将知识库文本切分为较小片段。
    后续 RAG 检索时，不应该直接检索整本文档，而应该检索 chunk。
    """
    text = clean_text(text)

    if not text:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    chunks = []
    current = ""

    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= max_chars:
            current = (current + "\n\n" + paragraph).strip()
        else:
            if current:
                chunks.append(current)

            if len(paragraph) <= max_chars:
                current = paragraph
            else:
                start = 0
                while start < len(paragraph):
                    end = min(start + max_chars, len(paragraph))
                    part = paragraph[start:end].strip()

                    if part:
                        chunks.append(part)

                    if end >= len(paragraph):
                        break

                    start = max(0, end - overlap)

                current = ""

    if current:
        chunks.append(current)

    cleaned_chunks = []

    for chunk in chunks:
        chunk = re.sub(r"\s+", " ", chunk).strip()

        if len(chunk) >= 30:
            cleaned_chunks.append(chunk)

    return cleaned_chunks


def embedding_tokens(text):
    text = expand_domain_query(str(text or "").lower())
    tokens = re.findall(r"[a-z0-9_+\-*/^=<>]+|[\u4e00-\u9fff]{1,}", text)
    expanded = []

    for token in tokens:
        expanded.append(token)
        if re.fullmatch(r"[\u4e00-\u9fff]+", token) and len(token) > 1:
            expanded.extend(token[i:i + 2] for i in range(len(token) - 1))

    return expanded


def local_embedding_vector(text, dim=None):
    dim = dim or LOCAL_EMBEDDING_DIM
    vector = [0.0] * dim

    for token in embedding_tokens(text):
        digest = hashlib.sha256(token.encode("utf-8")).digest()
        index = int.from_bytes(digest[:4], "big") % dim
        sign = 1.0 if digest[4] % 2 == 0 else -1.0
        vector[index] += sign

    norm = math.sqrt(sum(value * value for value in vector))
    if norm > 0:
        vector = [round(value / norm, 6) for value in vector]

    return vector


def cosine_similarity(left, right):
    if not left or not right:
        return 0.0

    size = min(len(left), len(right))
    return sum(left[i] * right[i] for i in range(size))


def ensure_chunk_embedding(chunk):
    record = EmbeddingRecord.query.filter_by(chunk_id=chunk.id, provider="local_hashing").first()
    if record is not None:
        try:
            return json.loads(record.vector_json)
        except json.JSONDecodeError:
            pass

    vector = local_embedding_vector(chunk.content)
    record = EmbeddingRecord(
        chunk_id=chunk.id,
        provider="local_hashing",
        dim=len(vector),
        vector_json=json.dumps(vector),
        created_at=current_time_text()
    )
    db.session.add(record)
    db.session.flush()
    return vector


def rebuild_embeddings_for_chunks(chunks):
    count = 0
    for chunk in chunks:
        ensure_chunk_embedding(chunk)
        count += 1
    db.session.commit()
    return count


DEMO_KB_ENTRIES = [
    {
        "language": "en",
        "course": "Data Structures and Algorithms",
        "chapter": "Hashing",
        "term": "Hash Table",
        "definition": "A hash table is a data structure that maps keys to values using a hash function. It supports efficient average-case search, insertion, and deletion.",
        "source": "English Textbook Demo: Data Structures, Chapter Hashing",
        "keywords": "hash table, hash function, key value, bucket, hashing"
    },
    {
        "language": "zh",
        "course": "Data Structures and Algorithms",
        "chapter": "Hashing",
        "term": "哈希表",
        "definition": "哈希表是一种通过哈希函数把关键字映射到存储位置的数据结构，通常用于实现快速查找、插入和删除。",
        "source": "中文教材演示：数据结构，散列表章节",
        "keywords": "哈希表, 散列表, 哈希函数, 关键字, 查找"
    },
    {
        "language": "en",
        "course": "Data Structures and Algorithms",
        "chapter": "Hashing",
        "term": "Collision Resolution",
        "definition": "Collision resolution handles the situation where two keys map to the same hash table index. Common strategies include chaining and open addressing.",
        "source": "English Textbook Demo: Hashing and Collision Handling",
        "keywords": "collision resolution, chaining, open addressing, hash table"
    },
    {
        "language": "zh",
        "course": "Data Structures and Algorithms",
        "chapter": "Hashing",
        "term": "冲突解决",
        "definition": "冲突解决用于处理多个关键字通过哈希函数映射到同一地址的情况，常见方法包括链地址法和开放定址法。",
        "source": "中文教材演示：散列表冲突处理",
        "keywords": "冲突解决, 哈希冲突, 链地址法, 开放定址法"
    },
    {
        "language": "en",
        "course": "Basic Physics I",
        "chapter": "Oscillation and Waves",
        "term": "Angular Frequency",
        "definition": "Angular frequency describes how fast an oscillation rotates in phase and is commonly denoted by omega, with omega equal to 2 pi times frequency.",
        "source": "English Textbook Demo: Oscillation and Waves",
        "keywords": "angular frequency, omega, oscillation, frequency, 2 pi f"
    },
    {
        "language": "zh",
        "course": "Basic Physics I",
        "chapter": "振动与波",
        "term": "角频率",
        "definition": "角频率表示振动相位变化的快慢，通常用 ω 表示，并满足 ω = 2πf。",
        "source": "中文教材演示：振动与波",
        "keywords": "角频率, 欧米伽, 振动, 频率"
    },
    {
        "language": "en",
        "course": "Basic Physics I",
        "chapter": "Oscillation and Waves",
        "term": "Wavelength",
        "definition": "Wavelength is the spatial period of a wave and is commonly denoted by lambda.",
        "source": "English Textbook Demo: Wave Concepts",
        "keywords": "wavelength, lambda, wave speed, frequency"
    },
    {
        "language": "zh",
        "course": "Basic Physics I",
        "chapter": "振动与波",
        "term": "波长",
        "definition": "波长表示波在空间中完成一个周期所对应的距离，通常用 λ 表示。",
        "source": "中文教材演示：波动基础",
        "keywords": "波长, lambda, 波速, 频率"
    }
]


LOCAL_TRANSLATION_FALLBACK = {
    "hash table": "哈希表",
    "collision resolution": "冲突解决",
    "hash function": "哈希函数",
    "angular frequency": "角频率",
    "wavelength": "波长",
    "frequency": "频率",
}


def seed_demo_knowledge_base():
    created = 0
    for item in DEMO_KB_ENTRIES:
        existing = KnowledgeBaseEntry.query.filter_by(
            language=item["language"],
            course=item["course"],
            term=item["term"]
        ).first()

        if existing is not None:
            continue

        entry = KnowledgeBaseEntry(
            language=item["language"],
            course=item["course"],
            chapter=item["chapter"],
            term=item["term"],
            definition=item["definition"],
            source=item["source"],
            keywords=item["keywords"],
            created_at=current_time_text()
        )
        db.session.add(entry)
        created += 1

    db.session.commit()
    return created


def score_kb_entry(entry, query):
    haystack = " ".join([
        entry.term or "",
        entry.definition or "",
        entry.keywords or "",
        entry.chapter or "",
    ])
    return score_knowledge_chunk(haystack, query)


def search_structured_kb(query, course="", language="", limit=3):
    if not query:
        return []

    db_query = KnowledgeBaseEntry.query
    if course:
        db_query = db_query.filter_by(course=course)
    if language:
        db_query = db_query.filter_by(language=language)

    scored = []
    expanded_query = expand_domain_query(query)

    for entry in db_query.all():
        score = score_kb_entry(entry, expanded_query)
        if score > 0:
            scored.append((score, entry))

    scored.sort(key=lambda item: item[0], reverse=True)
    return [
        {
            "score": score,
            "term": entry.term,
            "definition": entry.definition,
            "source": entry.source,
            "chapter": entry.chapter,
            "language": entry.language,
            "keywords": entry.keywords
        }
        for score, entry in scored[:limit]
    ]


def score_knowledge_chunk(content, query):
    """
    简单关键词检索评分。
    当前是 v0.1 检索版本：
    1. 精确包含 query，加高分
    2. query 中的英文词、数字、中文词片段命中，加分
    3. 后续可替换为 embedding / 向量检索
    """
    if not content or not query:
        return 0

    content_lower = normalize_math_symbols(content.lower())
    query_lower = expand_domain_query(query.lower().strip())

    score = 0

    if query_lower in content_lower:
        score += 100

    tokens = re.findall(r"[A-Za-z0-9]+|[\u4e00-\u9fff]{2,}", query_lower)

    for token in tokens:
        if token and token in content_lower:
            score += 20

    # 对中文短词做额外滑窗匹配，例如“收敛级数”
    chinese_chars = re.findall(r"[\u4e00-\u9fff]+", query_lower)

    for phrase in chinese_chars:
        if len(phrase) >= 2:
            for i in range(len(phrase) - 1):
                sub = phrase[i:i + 2]
                if sub in content_lower:
                    score += 6

    return score

# ============================================================
# 候选术语抽取
# ============================================================

def extract_context(text, term, window=140):
    """
    提取术语在原文中的上下文，方便质量控制。
    """
    lower_text = text.lower()
    lower_term = term.lower()
    index = lower_text.find(lower_term)

    if index == -1:
        return ""

    start = max(0, index - window)
    end = min(len(text), index + len(term) + window)

    context = text[start:end]
    context = re.sub(r"\s+", " ", context)
    return context.strip()


def extract_context_sentence(text, term):
    if not text or not term:
        return ""
    pattern = re.compile(r"[^.!?\n]*\b" + re.escape(term) + r"\b[^.!?\n]*[.!?]?", re.I)
    match = pattern.search(text)
    if match:
        return re.sub(r"\s+", " ", match.group(0)).strip()
    return extract_context(text, term)


def is_probably_noise(term):
    """
    通用噪声过滤。
    注意：这里不写死数学、物理、通信、计算机等具体学科关键词。
    只过滤明显不是术语的句子碎片、页面信息、学校信息、普通功能词。
    """
    lower = term.lower().strip()
    words = lower.split()

    if not lower:
        return True

    if contains_ocr_placeholder(term):
        return True

    if contains_formula_placeholder(term):
        return True

    if contains_formula_text(term):
        return True

    if re.search(r"\\[A-Za-z]+|[\^_]\{?[-A-Za-z0-9]+|[A-Za-z]\^\d+|e\^\{?-?[A-Za-z0-9]", term):
        return True

    if len(words) > 4:
        return True

    blacklist_exact = {
        "page", "slide", "chapter", "section", "example", "problem", "solution",
        "ocr required", "ocr fallback", "ocr", "required", "fallback",
        "formulablock", "formula block", "formula", "latex",
        "int", "sqrt", "frac", "infinity", "theta", "lambda", "x^2", "e^{-x^2}",
        "find", "suppose", "therefore", "then", "where", "when", "what",
        "beijing university", "posts and telecommunications",
        "international school", "jianhua yuan"
    }

    if lower in blacklist_exact:
        return True

    bad_starts = {
        "the", "this", "that", "these", "those",
        "a", "an", "and", "or", "but",
        "if", "when", "where", "why", "how", "what", "which",
        "we", "you", "they", "he", "she", "it",
        "suppose", "choose", "find", "show", "prove", "let",
        "in", "on", "at", "for", "with", "by", "from", "to", "of",
        "as", "is", "are", "was", "were"
    }

    bad_ends = {
        "the", "a", "an", "and", "or", "but",
        "is", "are", "was", "were", "be", "been",
        "in", "on", "at", "for", "with", "by", "from", "to", "of",
        "this", "that", "these", "those", "each"
    }

    if words and words[0] in bad_starts:
        return True

    if words and words[-1] in bad_ends:
        return True

    sentence_verbs = {
        "can", "will", "would", "could", "should", "may", "might",
        "must", "have", "has", "had", "do", "does", "did"
    }

    if any(word in sentence_verbs or word in TERM_ACTION_WORDS for word in words):
        return True

    # 全小写且过长，通常是句子碎片
    if len(words) >= 3 and term.islower():
        return True

    # 数字比例过高通常不是术语
    letters = sum(ch.isalpha() for ch in term)
    digits = sum(ch.isdigit() for ch in term)
    if digits > letters:
        return True

    return False


TERM_NGRAM_STOPWORDS = {
    "the", "this", "that", "these", "those", "a", "an", "and", "or", "but",
    "if", "when", "where", "why", "how", "what", "which", "who", "whom",
    "we", "you", "they", "he", "she", "it", "its", "their", "our", "your",
    "in", "on", "at", "for", "with", "by", "from", "to", "of", "as", "into",
    "than", "then", "there", "here", "same", "each", "every", "some", "any",
    "one", "two", "three", "four", "five", "six", "seven", "eight", "nine", "ten",
    "physics",
    "first", "second", "third", "fourth", "fifth", "sixth", "seventh", "eighth",
    "ninth", "tenth", "another", "other", "new", "old"
}

TERM_ACTION_WORDS = {
    "is", "are", "was", "were", "be", "been", "being",
    "use", "uses", "used", "using",
    "combine", "combines", "combined", "combining",
    "convert", "converts", "converted", "converting",
    "produce", "produces", "produced", "producing",
    "map", "maps", "mapped", "mapping",
    "handle", "handles", "handled", "handling",
    "equal", "equals", "equivalent",
    "where", "cases", "case",
    "means", "refers", "called", "defined",
    "has", "have", "had", "do", "does", "did",
    "can", "will", "would", "could", "should", "may", "might", "must"
}

TERM_SYMBOL_WORDS = {
    "omega", "lambda", "theta", "sigma", "delta", "mu", "pi", "sqrt", "integral",
    "sum", "infinity", "frac", "int", "lim", "sin", "cos", "tan", "log", "ln",
    "exp"
}


def is_ngram_noise(words):
    lower_words = [word.lower() for word in words]

    if not lower_words:
        return True

    if lower_words[0] in TERM_NGRAM_STOPWORDS or lower_words[-1] in TERM_NGRAM_STOPWORDS:
        return True

    if any(word in TERM_ACTION_WORDS for word in lower_words):
        return True

    if any(word in TERM_SYMBOL_WORDS for word in lower_words):
        return True

    if all(len(word) <= 2 for word in lower_words):
        return True

    # Single-letter symbols are useful in formulas but too noisy as extracted terms
    # unless they are part of a known KB term handled earlier.
    if any(len(word) == 1 for word in lower_words):
        return True

    return False


def add_candidate(raw_counts, display_form, contexts, boosts, text, term, score_boost=0):
    cleaned = " ".join(str(term or "").split()).strip()
    cleaned = re.sub(r"\s*-\s*", "-", cleaned)

    if len(cleaned) < 4 or len(cleaned) > 64:
        return

    if is_probably_noise(cleaned):
        return

    key = cleaned.lower()
    raw_counts[key] = raw_counts.get(key, 0) + 1

    if score_boost:
        boosts[key] = max(boosts.get(key, 0), score_boost)

    if key not in display_form:
        display_form[key] = cleaned

    if key not in contexts:
        contexts[key] = extract_context_sentence(text, cleaned)


def add_known_kb_terms(text, raw_counts, display_form, contexts, boosts):
    normalized_text = expand_domain_query(normalize_math_symbols(text.lower()))

    try:
        entries = KnowledgeBaseEntry.query.filter_by(language="en").all()
    except Exception:
        entries = []

    if not entries:
        entries = [
            type("DemoEntry", (), item)
            for item in DEMO_KB_ENTRIES
            if item.get("language") == "en"
        ]

    for entry in entries:
        term = getattr(entry, "term", "")
        triggers = [term]

        for trigger in triggers:
            normalized_trigger = expand_domain_query(normalize_math_symbols(trigger.lower())).strip()

            if normalized_trigger and normalized_trigger in normalized_text:
                add_candidate(raw_counts, display_form, contexts, boosts, text, term, score_boost=28)
                break

    for term in LOCAL_TRANSLATION_FALLBACK:
        normalized_term = expand_domain_query(normalize_math_symbols(term.lower())).strip()

        if normalized_term and normalized_term in normalized_text:
            display_term = " ".join(word.capitalize() for word in term.split())
            add_candidate(raw_counts, display_form, contexts, boosts, text, display_term, score_boost=24)


def score_candidate(term, count, context):
    """
    给候选术语打分。
    这是启发式评分，不代表真正的 AI 判断。
    """
    words = term.split()
    score = 45

    if count >= 2:
        score += min(25, 8 + count * 3)

    if len(words) >= 2:
        score += 12

    if len(words) >= 3:
        score += 4

    if term.isupper() and len(term) >= 3:
        score += 10

    if all(word[:1].isupper() for word in words if word):
        score += 8

    if "-" in term:
        score += 5

    # 领域无关的学术词形特征：不是学科关键词，只是术语形态线索
    morphology_suffixes = (
        "tion", "sion", "ment", "ness", "ity", "ics",
        "ism", "ance", "ence", "ing", "al", "ive", "ous"
    )
    if any(word.lower().endswith(morphology_suffixes) for word in words):
        score += 6

    definition_clues = (
        "is called", "is defined as", "means", "refers to",
        "denoted by", "known as", "consists of"
    )
    lower_context = context.lower()
    if any(clue in lower_context for clue in definition_clues):
        score += 5

    return max(0, min(score, 95))


def extract_terms_from_text(text):
    """
    通用候选术语抽取版本：
    1. 不写死具体专业术语；
    2. 不预设数学 / 物理 / 通信 / 计算机等学科关键词；
    3. 不伪造中文翻译；
    4. 只根据词组形态、出现频率、大小写、上下文等通用特征筛选候选术语；
    5. 最终术语是否成立，由质量控制或后续 AI API / RAG 决定。
    """
    if not text or not text.strip():
        return []

    if contains_ocr_placeholder(text):
        return []

    if contains_formula_placeholder(text):
        return []

    raw_counts = {}
    display_form = {}
    contexts = {}
    boosts = {}

    add_known_kb_terms(text, raw_counts, display_form, contexts, boosts)

    normalized_for_terms = normalize_math_symbols(text)
    sentences = re.split(r"(?<=[.!?])\s+|\n+", normalized_for_terms)

    for sentence in sentences:
        if contains_ocr_placeholder(sentence):
            continue
        tokens = re.findall(r"[A-Za-z][A-Za-z0-9]*(?:[-/][A-Za-z0-9]+)*", sentence)

        if not tokens:
            continue

        for ngram_size in (4, 3, 2, 1):
            if len(tokens) < ngram_size:
                continue

            for start in range(0, len(tokens) - ngram_size + 1):
                words = tokens[start:start + ngram_size]

                if is_ngram_noise(words):
                    continue

                term = " ".join(words)
                add_candidate(raw_counts, display_form, contexts, boosts, text, term)

    scored = []

    for key, count in raw_counts.items():
        term = display_form[key]
        words = term.split()
        context = contexts.get(key, "")
        boost = boosts.get(key, 0)
        is_known_or_seeded = boost > 0 or key in LOCAL_TRANSLATION_FALLBACK

        # 单词候选：保留较明确的单词术语或复数名词，但避免普通短词泛滥。
        if len(words) == 1 and not is_known_or_seeded:
            lower_word = words[0].lower()
            is_plural_noun = lower_word.endswith("s") and len(lower_word) >= 4
            is_long_academic = len(term) >= 7
            if not is_plural_noun and not is_long_academic:
                continue

        score = score_candidate(term, count, context) + boost

        if not is_known_or_seeded and score < 58:
            continue

        if is_known_or_seeded:
            score = max(score, 72)

        scored.append((score, count, term, context))

    # 高分优先；同分时出现次数多者优先；再按长度排序
    scored.sort(key=lambda item: (item[0], item[1], len(item[2])), reverse=True)

    candidates = []
    seen = set()

    for score, count, term, context in scored:
        key = term.lower()

        if key in seen:
            continue

        seen.add(key)

        candidates.append({
            "english_term": term,
            "chinese_term": "待质量控制",
            "explanation": "待质量控制：系统已抽取候选术语，等待证据对齐结果确认。",
            "context": context,
            "confidence": score,
            "status": "pending"
        })

    return candidates

def ai_status_from_confidence(confidence, has_conflict=False):
    if has_conflict:
        return "conflict"

    if confidence >= 85:
        return "auto_published"

    return "uncertain"


MIN_RETRIEVAL_SCORE = int(os.environ.get("MIN_RETRIEVAL_SCORE", "65"))


def retrieve_knowledge_evidence(query_text, course, knowledge_base_type, limit=4, min_score=MIN_RETRIEVAL_SCORE):
    if not query_text:
        return []

    course_obj = Course.query.filter_by(name=course).first() if course else None
    language = "en" if knowledge_base_type == "en_course_kb" else "zh" if knowledge_base_type == "zh_course_kb" else ""
    return retrieve_evidence_results(
        query_text,
        course_id=course_obj.id if course_obj else None,
        course_name=course,
        language=language,
        scope_type="course",
        owner_user_id=None,
        limit=limit,
        knowledge_base_type=knowledge_base_type
    )


def extract_json_object(text):
    if not text:
        raise ValueError("AI 返回内容为空。")

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        start = text.find("{")
        end = text.rfind("}")

        if start == -1 or end == -1 or end <= start:
            raise

        return json.loads(text[start:end + 1])


def call_deepseek_json(system_prompt, user_payload, max_tokens=1600):
    if AI_PROVIDER != "deepseek":
        raise RuntimeError("当前 AI_PROVIDER 不是 deepseek。")

    if not DEEPSEEK_API_KEY:
        raise RuntimeError("缺少 DEEPSEEK_API_KEY。")

    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [
            {
                "role": "system",
                "content": system_prompt
            },
            {
                "role": "user",
                "content": json.dumps(user_payload, ensure_ascii=False)
            }
        ],
        "response_format": {
            "type": "json_object"
        },
        "stream": False,
        "temperature": 0.2,
        "max_tokens": max_tokens
    }

    request_body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request_obj = urllib.request.Request(
        f"{DEEPSEEK_BASE_URL}/chat/completions",
        data=request_body,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}"
        },
        method="POST"
    )

    try:
        with urllib.request.urlopen(request_obj, timeout=60) as response:
            response_payload = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        error_body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"DeepSeek API HTTP {exc.code}: {error_body}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"DeepSeek API 网络错误：{exc.reason}") from exc

    content = (
        response_payload
        .get("choices", [{}])[0]
        .get("message", {})
        .get("content", "")
    )

    return extract_json_object(content)

def fallback_translation_candidate(english_term):
    key = str(english_term or "").strip().lower()
    return LOCAL_TRANSLATION_FALLBACK.get(key, english_term)


def first_evidence_text(items):
    if not items:
        return ""

    item = items[0]
    if "definition" in item:
        return f"{item.get('term', '')}: {item.get('definition', '')} ({item.get('source', '')})"

    if "content" in item:
        title = item.get("title", "") or item.get("source", "")
        content = item.get("content", "")
        return f"{title}: {content}".strip(": ")

    chunk = item.get("chunk") if isinstance(item, dict) else None
    if chunk:
        return f"{chunk.get('title', '')}: {chunk.get('content', '')}"

    return json.dumps(item, ensure_ascii=False)


def finalize_alignment_result(result, min_ocr_confidence=100):
    result["english_term"] = str(result.get("english_term", "")).strip()
    return finalize_alignment_decision(result, min_ocr_confidence=min_ocr_confidence)


def generate_alignment_result(english_term, courseware_sentence, course, chapter="", scope_type="course", owner_user_id=None, min_ocr_confidence=100):
    """
    Evidence-alignment workflow:
    1. retrieve English KB evidence
    2. produce a Chinese candidate
    3. retrieve Chinese KB evidence
    4. compare whether both sides refer to the same academic concept
    """
    english_term = str(english_term or "").strip()
    course_obj = Course.query.filter_by(name=course).first() if course else None
    course_id = course_obj.id if course_obj else None
    english_chunks = retrieve_evidence_results(
        english_term,
        course_id=course_id,
        course_name=course,
        language="en",
        scope_type=scope_type,
        owner_user_id=owner_user_id,
        limit=3,
        knowledge_base_type="student_personal_kb" if scope_type == "personal" else "en_course_kb"
    )
    english_evidence = english_chunks

    translation_candidate = fallback_translation_candidate(english_term)
    chinese_query = " ".join([translation_candidate, english_term, courseware_sentence or ""]).strip()
    chinese_chunks = retrieve_evidence_results(
        chinese_query,
        course_id=course_id,
        course_name=course,
        language="zh",
        scope_type=scope_type,
        owner_user_id=owner_user_id,
        limit=3,
        knowledge_base_type="student_personal_kb" if scope_type == "personal" else "zh_course_kb"
    )
    chinese_evidence = chinese_chunks

    try:
        ai_call = call_ai_task(
            task_type="term_alignment",
            prompt_key="term_alignment",
            prompt_version="v1",
            input_payload={
                "english_term": english_term,
                "course": course,
                "chapter": chapter,
                "courseware_sentence": courseware_sentence,
                "english_evidence": english_evidence,
                "translation_candidate_hint": translation_candidate,
                "chinese_evidence": chinese_evidence,
            },
            user_id=owner_user_id,
            course_id=course_id,
        )
        if ai_call.get("status") != "success":
            raise RuntimeError(f"{ai_call.get('error_code')}: {ai_call.get('message')}")
        ai_result = ai_call.get("result", {})
        raw_confidence = float(ai_result.get("ai_confidence", ai_result.get("confidence_score", 0)) or 0)
        confidence = int(raw_confidence * 100) if raw_confidence <= 1 else int(raw_confidence)
        confidence = max(0, min(confidence, 100))
        final_term = str(ai_result.get("candidate_chinese_term") or ai_result.get("final_chinese_term") or "").strip() or translation_candidate
        provider_mode = ai_call.get("provider_mode", "")
        provider_status = "real_provider" if provider_mode == "live" else provider_mode
        result = {
            "english_term": english_term,
            "course": course,
            "chapter": chapter,
            "ai_translation_candidate": final_term,
            "final_chinese_term": final_term,
            "chinese_term": final_term,
            "explanation": str(ai_result.get("concept_explanation") or ai_result.get("explanation") or "").strip(),
            "confidence_score": confidence,
            "alignment_status": str(ai_result.get("alignment_status", "")).strip(),
            "alignment_reason": str(ai_result.get("alignment_reason", "")).strip(),
            "review_status": "pending_quality_control" if ai_result.get("requires_human_review") else str(ai_result.get("review_status", "")).strip(),
            "risk_note": str(ai_result.get("risk_note", "")).strip(),
            "english_kb_evidence": first_evidence_text(english_evidence),
            "chinese_kb_evidence": first_evidence_text(chinese_evidence),
            "english_evidence_items": english_evidence,
            "chinese_evidence_items": chinese_evidence,
            "ai_provider": ai_call.get("provider_name", AI_PROVIDER),
            "ai_provider_mode": provider_mode,
            "ai_model": ai_call.get("model_name", ""),
            "provider_status": provider_status,
            "is_real_provider": provider_mode == "live",
            "prompt_key": ai_call.get("prompt_key", "term_alignment"),
            "prompt_version": ai_call.get("prompt_version", "v1"),
            "retrieval_version": RETRIEVAL_VERSION,
            "ai_call_log_id": ai_call.get("ai_call_log_id"),
        }
        if provider_mode != "live":
            result["risk_note"] = result["risk_note"] or "No live AI provider configured; local heuristic/mock result requires Quality Control."
        else:
            model_ok, model_reasons = can_use_model_for_auto_approval(
                result["ai_provider"],
                result["ai_model"],
                result["prompt_version"],
            )
            if not model_ok:
                result["quality_flags"] = sorted(set((result.get("quality_flags") or []) + ["model_not_evaluated"]))
                result["risk_note"] = result["risk_note"] or "; ".join(model_reasons[:3])
        return finalize_alignment_result(result, min_ocr_confidence=min_ocr_confidence)
    except Exception as exc:
        fallback_error = str(exc)
        add_system_log("warning", "ai_provider", f"{AI_PROVIDER} failed or unavailable; using local heuristic fallback: {fallback_error}")

    confidence = 58 if english_evidence and chinese_evidence else 35
    final_term = translation_candidate
    reason = (
        "Local demo fallback: English evidence and Chinese evidence were retrieved by keyword/semantic matching. "
        "The result should be reviewed before classroom publication."
    )
    if not english_evidence or not chinese_evidence:
        reason += f" Evidence is incomplete: {fallback_error}."

    result = {
        "english_term": english_term,
        "course": course,
        "chapter": chapter,
        "ai_translation_candidate": translation_candidate,
        "final_chinese_term": final_term,
        "chinese_term": final_term,
        "explanation": first_evidence_text(chinese_evidence) or first_evidence_text(english_evidence) or "本地 fallback 未找到足够教材证据。",
        "confidence_score": confidence,
        "alignment_reason": reason,
        "review_status": "pending_quality_control" if english_evidence and chinese_evidence else "needs_more_evidence",
        "risk_note": f"No live AI provider configured or provider failed: {fallback_error}",
        "english_kb_evidence": first_evidence_text(english_evidence),
        "chinese_kb_evidence": first_evidence_text(chinese_evidence),
        "english_evidence_items": english_evidence,
        "chinese_evidence_items": chinese_evidence,
        "ai_provider": AI_PROVIDER,
        "ai_provider_mode": "provider_failed" if AI_PROVIDER not in {"none", "mock"} else "local_heuristic",
        "ai_model": "local_heuristic",
        "provider_status": "provider_failed" if AI_PROVIDER not in {"none", "mock"} else "local_heuristic",
        "prompt_key": "term_alignment",
        "prompt_version": "v1",
        "retrieval_version": RETRIEVAL_VERSION
    }
    return finalize_alignment_result(result, min_ocr_confidence=min_ocr_confidence)


def generate_student_learning_answer(question, course, owner_user_id="local_student"):
    english_evidence = retrieve_knowledge_evidence(question, course, "en_course_kb", limit=4)
    chinese_evidence = retrieve_knowledge_evidence(question, course, "zh_course_kb", limit=4)

    personal_query = KnowledgeChunk.query.filter_by(knowledge_base_type="student_personal_kb")

    if course:
        personal_query = personal_query.filter_by(course=course)

    if owner_user_id:
        personal_query = personal_query.filter_by(owner_user_id=owner_user_id)

    personal_chunks = personal_query.all()
    scored_personal = []

    for chunk in personal_chunks:
        score = score_knowledge_chunk(chunk.content, question)

        if score >= MIN_RETRIEVAL_SCORE:
            scored_personal.append((score, chunk))

    scored_personal.sort(key=lambda item: item[0], reverse=True)
    personal_evidence = [
        {
            "score": score,
            "chunk_id": chunk.id,
            "title": chunk.title,
            "course": chunk.course,
            "chunk_index": chunk.chunk_index,
            "content": chunk.content[:900]
        }
        for score, chunk in scored_personal[:4]
    ]

    meta = current_provider_metadata()
    answer = {
        "answer": "Retrieved bilingual evidence is available below. This Local MVP does not treat local/mock AI as a verified tutor answer.",
        "key_terms": [],
        "personalized_note": "Personal workspace evidence is included only when it belongs to the current user.",
        "study_suggestion": "Compare the English and Chinese evidence cards, then mark confusing terms for teacher feedback.",
        "evidence_summary": f"English evidence: {len(english_evidence)}; Chinese evidence: {len(chinese_evidence)}; personal evidence: {len(personal_evidence)}.",
        "confidence": 45 if meta.get("provider_mode") != "live" else 60,
        "risk_note": "Local heuristic answer; not a verified live AI response." if meta.get("provider_mode") != "live" else "",
    }

    confidence = int(answer.get("confidence", 0) or 0)
    confidence = max(0, min(confidence, 100))

    return {
        "answer": answer.get("answer", "").strip(),
        "key_terms": answer.get("key_terms", []),
        "personalized_note": answer.get("personalized_note", "").strip(),
        "study_suggestion": answer.get("study_suggestion", "").strip(),
        "evidence_summary": answer.get("evidence_summary", "").strip(),
        "confidence": confidence,
        "risk_note": answer.get("risk_note", "").strip(),
        "english_evidence": english_evidence,
        "chinese_evidence": chinese_evidence,
        "personal_evidence": personal_evidence,
        "ai_provider": meta.get("provider"),
        "ai_provider_mode": meta.get("provider_mode"),
        "ai_model": meta.get("model_name")
    }


def apply_ai_suggestion_to_term(term, suggestion, auto_publish=False):
    explanation_parts = []

    if suggestion.get("chinese_explanation"):
        explanation_parts.append(suggestion["chinese_explanation"])

    if suggestion.get("english_definition"):
        explanation_parts.append(f"英文定义：{suggestion['english_definition']}")

    if suggestion.get("evidence_summary"):
        explanation_parts.append(f"依据摘要：{suggestion['evidence_summary']}")

    if suggestion.get("explanation"):
        explanation_parts.append(suggestion["explanation"])

    final_term = suggestion.get("final_chinese_term") or suggestion.get("chinese_term")
    term.chinese_term = final_term
    term.explanation = "\n".join(explanation_parts).strip() or "AI 未能生成有效解释。"
    term.confidence = suggestion.get("confidence", suggestion.get("confidence_score", 60))
    term.ai_status = suggestion.get("ai_status", ai_status_from_confidence(term.confidence))
    term.english_evidence = json.dumps(suggestion.get("english_evidence", suggestion.get("english_evidence_items", [])), ensure_ascii=False)
    term.chinese_evidence = json.dumps(suggestion.get("chinese_evidence", suggestion.get("chinese_evidence_items", [])), ensure_ascii=False)
    term.risk_note = suggestion.get("risk_note", "")
    term.ai_model = suggestion.get("ai_model", "")
    term.ai_translation_candidate = suggestion.get("ai_translation_candidate", "")
    term.final_chinese_term = final_term
    term.english_kb_evidence = suggestion.get("english_kb_evidence", "")
    term.chinese_kb_evidence = suggestion.get("chinese_kb_evidence", "")
    term.alignment_reason = suggestion.get("alignment_reason", "")
    term.review_status = suggestion.get("review_status", "pending")

    if auto_publish and term.ai_status == "auto_published":
        term.status = "approved"


def serialize_term(term):
    return {
        "id": term.id,
        "course": term.course,
        "chapter": term.chapter,
        "english_term": term.english_term,
        "chinese_term": term.chinese_term,
        "explanation": term.explanation,
        "context": term.context,
        "courseware_sentence": getattr(term, "courseware_sentence", ""),
        "english_kb_evidence": getattr(term, "english_kb_evidence", ""),
        "ai_translation_candidate": getattr(term, "ai_translation_candidate", ""),
        "chinese_kb_evidence": getattr(term, "chinese_kb_evidence", ""),
        "final_chinese_term": getattr(term, "final_chinese_term", ""),
        "alignment_reason": getattr(term, "alignment_reason", ""),
        "review_status": getattr(term, "review_status", "pending"),
        "confidence": term.confidence,
        "status": term.status,
        "ai_status": getattr(term, "ai_status", "pending"),
        "english_evidence": getattr(term, "english_evidence", ""),
        "chinese_evidence": getattr(term, "chinese_evidence", ""),
        "risk_note": getattr(term, "risk_note", ""),
        "parse_uid": getattr(term, "parse_uid", ""),
        "parse_block_uid": getattr(term, "parse_block_uid", ""),
        "parse_quality_status": getattr(term, "parse_quality_status", ""),
        "parse_quality_flags": safe_json_loads(getattr(term, "parse_quality_flags", "[]"), []),
        "input_risk_labels": safe_json_loads(getattr(term, "input_risk_labels", "[]"), []),
        "source_uid": getattr(term, "source_uid", ""),
        "chunk_uid": getattr(term, "chunk_uid", ""),
        "ai_model": getattr(term, "ai_model", "")
    }
def serialize_feedback(feedback):
    return {
        "id": feedback.id,
        "feedback_uid": getattr(feedback, "feedback_uid", "") or f"feedback-{feedback.id}",
        "term_id": feedback.term_id,
        "terminology_card_id": getattr(feedback, "terminology_card_id", None) or feedback.term_id,
        "user_id": getattr(feedback, "user_id", 0),
        "student_display": anonymize_user(getattr(feedback, "user_id", 0)),
        "user_role": getattr(feedback, "user_role", ""),
        "course_id": getattr(feedback, "course_id", None),
        "document_id": getattr(feedback, "document_id", None),
        "formula_block_id": getattr(feedback, "formula_block_id", None),
        "job_id": getattr(feedback, "job_id", None),
        "alignment_run_id": getattr(feedback, "alignment_run_id", None),
        "evaluation_run_id": getattr(feedback, "evaluation_run_id", None),
        "card_uid": getattr(feedback, "card_uid", "") or (getattr(feedback, "actual_result", "") if getattr(feedback, "feedback_source", "") == "student_concept_card" else ""),
        "course": feedback.course,
        "chapter": feedback.chapter,
        "english_term": feedback.english_term,
        "chinese_term": feedback.chinese_term,
        "feedback_type": feedback.feedback_type,
        "feedback_source": getattr(feedback, "feedback_source", "student_card_detail"),
        "severity": getattr(feedback, "severity", "medium"),
        "priority": getattr(feedback, "priority", "P2"),
        "message": getattr(feedback, "message", "") or getattr(feedback, "reported_issue", "") or feedback.feedback_content,
        "suggested_chinese_term": getattr(feedback, "suggested_chinese_term", "") or getattr(feedback, "expected_result", ""),
        "feedback_content": feedback.feedback_content,
        "reported_issue": getattr(feedback, "reported_issue", "") or feedback.feedback_content,
        "expected_result": getattr(feedback, "expected_result", ""),
        "actual_result": getattr(feedback, "actual_result", ""),
        "evidence_comment": getattr(feedback, "evidence_comment", ""),
        "classification": getattr(feedback, "classification", ""),
        "root_cause": getattr(feedback, "root_cause", ""),
        "resolution_action": getattr(feedback, "resolution_action", ""),
        "resolution_note": getattr(feedback, "resolution_note", ""),
        "handled_by": getattr(feedback, "handled_by", None),
        "handled_at": getattr(feedback, "handled_at", ""),
        "handler_role": getattr(feedback, "handler_role", ""),
        "teacher_note": getattr(feedback, "teacher_note", ""),
        "linked_review_uid": getattr(feedback, "linked_review_uid", ""),
        "linked_card_uid": getattr(feedback, "linked_card_uid", ""),
        "converted_to_evaluation_item_id": getattr(feedback, "converted_to_evaluation_item_id", None),
        "linked_backlog_item_id": getattr(feedback, "linked_backlog_item_id", None),
        "status": feedback.status,
        "resolved_by": getattr(feedback, "resolved_by", None),
        "created_at": feedback.created_at,
        "updated_at": getattr(feedback, "updated_at", ""),
        "resolved_at": getattr(feedback, "resolved_at", "")
    }


def serialize_backlog_item(item):
    return {
        "id": item.id,
        "title": item.title,
        "description": item.description,
        "source_type": item.source_type,
        "source_feedback_id": item.source_feedback_id,
        "course_id": item.course_id,
        "severity": item.severity,
        "priority": item.priority,
        "category": item.category,
        "status": item.status,
        "owner": item.owner,
        "target_pr": item.target_pr,
        "acceptance_criteria": item.acceptance_criteria,
        "created_at": item.created_at,
        "updated_at": item.updated_at,
        "closed_at": item.closed_at,
    }
def serialize_knowledge_document(doc):
    return {
        "id": doc.id,
        "course": doc.course,
        "title": doc.title,
        "filename": doc.filename,
        "saved_filename": doc.saved_filename,
        "file_type": doc.file_type,
        "parse_uid": getattr(doc, "parse_uid", ""),
        "language": doc.language,
        "source_type": doc.source_type,
        "knowledge_base_type": getattr(doc, "knowledge_base_type", "zh_course_kb"),
        "owner_user_id": getattr(doc, "owner_user_id", ""),
        "visibility": getattr(doc, "visibility", "course"),
        "text_length": doc.text_length,
        "chunk_count": doc.chunk_count,
        "created_at": doc.created_at
    }


def serialize_knowledge_chunk(chunk):
    return {
        "id": chunk.id,
        "chunk_uid": getattr(chunk, "chunk_uid", ""),
        "source_uid": getattr(chunk, "source_uid", ""),
        "document_id": chunk.document_id,
        "source_id": getattr(chunk, "source_id", None),
        "course_id": getattr(chunk, "course_id", None),
        "parse_uid": getattr(chunk, "parse_uid", ""),
        "parse_block_uid": getattr(chunk, "parse_block_uid", ""),
        "course": chunk.course,
        "title": chunk.title,
        "discipline": getattr(chunk, "discipline", ""),
        "chapter": getattr(chunk, "chapter", ""),
        "chunk_index": chunk.chunk_index,
        "content": chunk.content,
        "source_page": chunk.source_page,
        "page_number": getattr(chunk, "page_number", None),
        "slide_number": getattr(chunk, "slide_number", None),
        "source_locator": getattr(chunk, "source_locator", "") or getattr(chunk, "source_section", "") or getattr(chunk, "source_page", ""),
        "block_type": getattr(chunk, "block_type", "text"),
        "token_count": getattr(chunk, "token_count", None),
        "char_count": getattr(chunk, "char_count", None),
        "keywords": getattr(chunk, "keywords", ""),
        "source_citation": getattr(chunk, "source_citation", ""),
        "embedding_id": getattr(chunk, "embedding_id", ""),
        "language": getattr(chunk, "language", ""),
        "knowledge_base_type": getattr(chunk, "knowledge_base_type", "zh_course_kb"),
        "owner_user_id": getattr(chunk, "owner_user_id", ""),
        "visibility": getattr(chunk, "visibility", "course"),
        "knowledge_base_version_id": getattr(chunk, "knowledge_base_version_id", None),
        "knowledge_source_id": getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None),
        "scope_type": getattr(chunk, "scope_type", "course"),
        "normalized_text": getattr(chunk, "normalized_text", ""),
        "content_hash": getattr(chunk, "content_hash", ""),
        "source_slide": getattr(chunk, "source_slide", ""),
        "source_section": getattr(chunk, "source_section", ""),
        "formula_block_ids": safe_json_loads(getattr(chunk, "formula_block_ids_json", "[]"), []),
        "index_status": getattr(chunk, "index_status", "indexed"),
        "quality_status": getattr(chunk, "quality_status", ""),
        "quality_flags": safe_json_loads(getattr(chunk, "quality_flags", "[]"), []),
        "trust_level": getattr(chunk, "trust_level", "unknown"),
        "status": getattr(chunk, "status", "active"),
        "embedding_status": getattr(chunk, "embedding_status", "not_started"),
        "is_duplicate": bool(getattr(chunk, "is_duplicate", False)),
        "duplicate_of_chunk_id": getattr(chunk, "duplicate_of_chunk_id", None),
        "is_active": bool(getattr(chunk, "is_active", True)),
        "created_at": chunk.created_at
    }


def serialize_user(user):
    if user is None:
        return None

    return {
        "id": user.id,
        "username": user.username,
        "email": getattr(user, "email", ""),
        "role": user.role,
        "display_name": user.display_name,
        "is_verified": bool(getattr(user, "is_verified", False)),
        "created_at": user.created_at
    }


def serialize_course(course):
    course_document_count = 0
    course_has_courseware = False
    try:
        course_document_count = Document.query.filter(
            Document.course_id == course.id,
            Document.scope_type == "course",
            db.or_(Document.deleted_at == "", Document.deleted_at.is_(None)),
        ).count()
        course_has_courseware = course_document_count > 0
    except Exception:
        course_document_count = 0
        course_has_courseware = False
    return {
        "id": course.id,
        "name": course.name,
        "course_name": course.name,
        "course_code": getattr(course, "course_code", ""),
        "semester": course.semester,
        "description": course.description,
        "language_mode": getattr(course, "language_mode", "bilingual"),
        "teacher_id": getattr(course, "teacher_id", 0),
        "status": getattr(course, "status", "active"),
        "document_count": course_document_count,
        "has_courseware": course_has_courseware,
        "deleted_at": getattr(course, "deleted_at", ""),
        "created_at": course.created_at
    }


def serialize_document(doc):
    formula_count = 0
    formula_statuses = []
    formula_providers = []
    try:
        formula_blocks = FormulaBlock.query.filter_by(document_id=doc.id).all()
        formula_count = len(formula_blocks)
        formula_statuses = sorted({block.status for block in formula_blocks if block.status})
        formula_providers = sorted({block.provider for block in formula_blocks if block.provider})
    except Exception:
        formula_blocks = []
    return {
        "id": doc.id,
        "owner_user_id": doc.owner_user_id,
        "course_id": doc.course_id,
        "scope_type": doc.scope_type,
        "filename": doc.filename,
        "saved_filename": doc.saved_filename,
        "file_sha256": getattr(doc, "file_sha256", ""),
        "storage_object_id": getattr(doc, "storage_object_id", None),
        "storage_backend": getattr(doc, "storage_backend", ""),
        "storage_key": getattr(doc, "storage_key", ""),
        "original_filename": getattr(doc, "original_filename", "") or doc.filename,
        "content_type": getattr(doc, "content_type", ""),
        "size_bytes": int(getattr(doc, "size_bytes", 0) or 0),
        "sha256": getattr(doc, "sha256", "") or getattr(doc, "file_sha256", ""),
        "file_type": doc.file_type,
        "language": doc.language,
        "upload_time": doc.upload_time,
        "parsing_status": doc.parsing_status,
        "parse_uid": getattr(doc, "parse_uid", ""),
        "ocr_required": bool(doc.ocr_required),
        "ocr_provider": getattr(doc, "ocr_provider", ""),
        "ocr_status": getattr(doc, "ocr_status", ""),
        "ocr_error": getattr(doc, "ocr_error", ""),
        "formula_blocks_count": formula_count,
        "formula_statuses": formula_statuses,
        "formula_providers": formula_providers,
        "formula_status": next((status for status in formula_statuses if status != "ok"), formula_statuses[0] if formula_statuses else "not_required"),
        "formula_provider": formula_providers[0] if formula_providers else "not_required",
        "quality_flags": safe_json_loads(getattr(doc, "quality_flags_json", "[]"), []),
        "source_type": doc.source_type,
        "error_message": doc.error_message,
        "deleted_at": getattr(doc, "deleted_at", "")
    }


def serialize_document_chunk(chunk):
    return {
        "id": chunk.id,
        "document_id": chunk.document_id,
        "course_id": chunk.course_id,
        "user_id": chunk.user_id,
        "owner_user_id": getattr(chunk, "owner_user_id", 0),
        "chunk_index": getattr(chunk, "chunk_index", 0),
        "parse_uid": getattr(chunk, "parse_uid", ""),
        "parse_block_uid": getattr(chunk, "parse_block_uid", ""),
        "language": chunk.language,
        "page_number": chunk.page_number,
        "slide_number": chunk.slide_number,
        "section_title": chunk.section_title,
        "content": chunk.content,
        "source_type": chunk.source_type,
        "source_location": chunk.source_location,
        "ocr_confidence": chunk.ocr_confidence,
        "ocr_provider": getattr(chunk, "ocr_provider", ""),
        "ocr_status": getattr(chunk, "ocr_status", ""),
        "ocr_error": getattr(chunk, "ocr_error", ""),
        "quality_flags": safe_json_loads(getattr(chunk, "quality_flags_json", "[]"), []),
        "created_at": chunk.created_at
    }


def serialize_formula_block(block):
    return {
        "id": block.id,
        "document_id": block.document_id,
        "course_id": block.course_id,
        "owner_user_id": block.owner_user_id,
        "scope_type": block.scope_type,
        "page_number": block.page_number,
        "slide_number": block.slide_number,
        "bbox": safe_json_loads(getattr(block, "bbox_json", "{}"), {}),
        "image_path": block.image_path,
        "latex": block.latex,
        "plain_text": block.plain_text,
        "provider": block.provider,
        "confidence": block.confidence,
        "status": block.status,
        "error": block.error,
        "quality_flags": safe_json_loads(getattr(block, "quality_flags_json", "[]"), []),
        "created_at": block.created_at
    }


def serialize_knowledge_source(source):
    return {
        "id": source.id,
        "source_uid": getattr(source, "source_uid", ""),
        "title": getattr(source, "title", "") or getattr(source, "source_title", "") or source.name,
        "name": source.name,
        "source_title": getattr(source, "source_title", "") or source.name,
        "course": getattr(source, "course", ""),
        "chapter": getattr(source, "chapter", ""),
        "course_id": getattr(source, "course_id", None),
        "scope_type": getattr(source, "scope_type", "course"),
        "owner_user_id": getattr(source, "owner_user_id", None),
        "owner_type": getattr(source, "owner_type", "unknown"),
        "owner_id": getattr(source, "owner_id", ""),
        "document_id": getattr(source, "document_id", None),
        "language": source.language,
        "discipline": source.discipline,
        "source_type": source.source_type,
        "source_role": getattr(source, "source_role", "unknown"),
        "visibility": getattr(source, "visibility", "course"),
        "trust_level": getattr(source, "trust_level", "unknown"),
        "parse_uid": getattr(source, "parse_uid", ""),
        "source_filename": getattr(source, "source_filename", ""),
        "file_type": getattr(source, "file_type", "unknown"),
        "content_hash": getattr(source, "content_hash", ""),
        "version": getattr(source, "version", 1),
        "license_note": getattr(source, "license_note", ""),
        "quality_status": getattr(source, "quality_status", ""),
        "quality_flags": safe_json_loads(getattr(source, "quality_flags", "[]"), []),
        "knowledge_base_type": getattr(source, "knowledge_base_type", ""),
        "access_method": source.access_method,
        "license_status": source.license_status,
        "license_type": getattr(source, "license_type", "unknown"),
        "authorization_status": getattr(source, "authorization_status", "unknown"),
        "source_quality": float(getattr(source, "source_quality", 0.4) or 0.4),
        "version_introduced_id": getattr(source, "version_introduced_id", None),
        "version_removed_id": getattr(source, "version_removed_id", None),
        "status": getattr(source, "status", "active"),
        "effective_from": getattr(source, "effective_from", ""),
        "effective_to": getattr(source, "effective_to", ""),
        "update_frequency": source.update_frequency,
        "allow_full_text_indexing": bool(source.allow_full_text_indexing),
        "allow_student_search": bool(source.allow_student_search),
        "allow_derivative_cards": bool(source.allow_derivative_cards),
        "created_by": source.created_by,
        "created_at": source.created_at,
        "updated_at": source.updated_at
    }


def serialize_kb_version(version):
    return {
        "id": version.id,
        "kb_scope": version.kb_scope,
        "scope_type": getattr(version, "scope_type", "") or getattr(version, "kb_scope", ""),
        "course_id": version.course_id,
        "owner_user_id": getattr(version, "owner_user_id", None),
        "version_name": version.version_name,
        "version_number": getattr(version, "version_number", 1),
        "status": getattr(version, "status", "draft"),
        "description": version.description,
        "source_count": version.source_count,
        "chunk_count": version.chunk_count,
        "formula_block_count": getattr(version, "formula_block_count", 0),
        "deduped_chunk_count": getattr(version, "deduped_chunk_count", 0),
        "index_backend": getattr(version, "index_backend", "local_lexical"),
        "index_version": getattr(version, "index_version", "local_lexical_v1"),
        "retrieval_version": getattr(version, "retrieval_version", "local_lexical_v1"),
        "embedding_provider": getattr(version, "embedding_provider", ""),
        "embedding_model": getattr(version, "embedding_model", ""),
        "embedding_dimension": getattr(version, "embedding_dimension", 0),
        "vector_index_status": getattr(version, "vector_index_status", ""),
        "vector_index_updated_at": getattr(version, "vector_index_updated_at", ""),
        "evaluation_run_id": getattr(version, "evaluation_run_id", None),
        "quality_gate_status": getattr(version, "quality_gate_status", ""),
        "manifest": safe_json_loads(getattr(version, "manifest_json", "{}"), {}),
        "created_at": version.created_at,
        "created_by": version.created_by,
        "published_at": getattr(version, "published_at", ""),
        "archived_at": getattr(version, "archived_at", ""),
        "parent_version_id": getattr(version, "parent_version_id", None),
        "is_active": bool(version.is_active)
    }


def serialize_terminology_card(card):
    return {
        "id": card.id,
        "scope_type": card.scope_type,
        "course_id": card.course_id,
        "owner_user_id": card.owner_user_id,
        "english_term": card.english_term,
        "normalized_english_term": getattr(card, "normalized_english_term", ""),
        "final_chinese_term": card.final_chinese_term,
        "normalized_chinese_term": getattr(card, "normalized_chinese_term", ""),
        "ai_translation_candidate": card.ai_translation_candidate,
        "courseware_sentence": card.courseware_sentence,
        "english_kb_evidence": card.english_kb_evidence,
        "chinese_kb_evidence": card.chinese_kb_evidence,
        "english_evidence_snapshot": safe_json_loads(getattr(card, "english_evidence_snapshot", "[]"), []),
        "chinese_evidence_snapshot": safe_json_loads(getattr(card, "chinese_evidence_snapshot", "[]"), []),
        "english_evidence_score": float(getattr(card, "english_evidence_score", 0.0) or 0.0),
        "chinese_evidence_score": float(getattr(card, "chinese_evidence_score", 0.0) or 0.0),
        "concept_explanation": card.concept_explanation,
        "alignment_reason": card.alignment_reason,
        "alignment_status": getattr(card, "alignment_status", "unverified_translation"),
        "score_breakdown": safe_json_loads(getattr(card, "score_breakdown_json", "{}"), {}),
        "quality_flags": safe_json_loads(getattr(card, "quality_flags_json", "[]"), []),
        "confidence_score": card.confidence_score,
        "status": card.status,
        "ai_provider": getattr(card, "ai_provider", ""),
        "ai_provider_mode": getattr(card, "ai_provider_mode", ""),
        "ai_model": getattr(card, "ai_model", ""),
        "prompt_key": getattr(card, "prompt_key", ""),
        "prompt_version": getattr(card, "prompt_version", ""),
        "retrieval_version": getattr(card, "retrieval_version", ""),
        "knowledge_base_version_id": getattr(card, "knowledge_base_version_id", None),
        "english_kb_version_id": getattr(card, "english_kb_version_id", None),
        "chinese_kb_version_id": getattr(card, "chinese_kb_version_id", None),
        "retrieval_run_id": getattr(card, "retrieval_run_id", None),
        "index_version": getattr(card, "index_version", ""),
        "evidence_content_hashes": safe_json_loads(getattr(card, "evidence_content_hashes_json", "[]"), []),
        "evidence_status": getattr(card, "evidence_status", "evidence_current"),
        "ai_call_log_id": getattr(card, "ai_call_log_id", None),
        "alignment_run_id": getattr(card, "alignment_run_id", None),
        "source_alignment_run_id": getattr(card, "source_alignment_run_id", None),
        "risk_note": getattr(card, "risk_note", ""),
        "parse_uid": getattr(card, "parse_uid", ""),
        "parse_block_uid": getattr(card, "parse_block_uid", ""),
        "parse_quality_status": getattr(card, "parse_quality_status", ""),
        "parse_quality_flags": safe_json_loads(getattr(card, "parse_quality_flags", "[]"), []),
        "input_risk_labels": safe_json_loads(getattr(card, "input_risk_labels", "[]"), []),
        "source_document_id": card.source_document_id,
        "english_evidence_chunk_id": card.english_evidence_chunk_id,
        "chinese_evidence_chunk_id": card.chinese_evidence_chunk_id,
        "approved_by": getattr(card, "approved_by", None),
        "approved_at": getattr(card, "approved_at", ""),
        "rejected_reason": getattr(card, "rejected_reason", ""),
        "reviewer_note": getattr(card, "reviewer_note", ""),
        "feedback_count": getattr(card, "feedback_count", 0),
        "created_at": card.created_at,
        "updated_at": card.updated_at
    }


def serialize_alignment_run(run):
    term_count = getattr(run, "term_count", None)
    if term_count is None:
        term_count = getattr(run, "terms_extracted", 0)
    card_count = getattr(run, "card_created_count", None)
    if card_count is None:
        card_count = getattr(run, "cards_created", 0)
    return {
        "id": run.id,
        "document_id": run.document_id,
        "course_id": run.course_id,
        "triggered_by": run.triggered_by,
        "provider": run.provider,
        "model_name": run.model_name,
        "ai_provider": getattr(run, "ai_provider", "") or getattr(run, "provider", ""),
        "ai_provider_mode": getattr(run, "ai_provider_mode", ""),
        "ai_model": getattr(run, "ai_model", "") or getattr(run, "model_name", ""),
        "prompt_key": getattr(run, "prompt_key", ""),
        "prompt_version": run.prompt_version,
        "retrieval_version": run.retrieval_version,
        "terms_extracted": run.terms_extracted,
        "cards_created": run.cards_created,
        "term_count": term_count,
        "card_created_count": card_count,
        "auto_approved_count": getattr(run, "auto_approved_count", 0),
        "qc_count": getattr(run, "qc_count", 0),
        "needs_evidence_count": getattr(run, "needs_evidence_count", 0),
        "conflict_count": getattr(run, "conflict_count", 0),
        "failed_count": getattr(run, "failed_count", 0),
        "status": run.status,
        "metrics": safe_json_loads(run.metrics_json, {}),
        "error_message": run.error_message,
        "started_at": run.started_at,
        "finished_at": run.finished_at,
    }


def update_alignment_run_stats(run, cards=None, term_count=None, failed_count=0):
    cards = cards or []
    status_counts = {}
    alignment_status_counts = {}
    for card in cards:
        status_counts[card.status] = status_counts.get(card.status, 0) + 1
        alignment_status_counts[card.alignment_status] = alignment_status_counts.get(card.alignment_status, 0) + 1

    total_terms = len(cards) if term_count is None else int(term_count or 0)
    run.terms_extracted = total_terms
    run.term_count = total_terms
    run.cards_created = len(cards)
    run.card_created_count = len(cards)
    run.auto_approved_count = status_counts.get("auto_approved", 0)
    run.qc_count = (
        status_counts.get("pending_quality_control", 0)
        + status_counts.get("conflict_detected", 0)
    )
    run.needs_evidence_count = status_counts.get("needs_more_evidence", 0)
    run.conflict_count = status_counts.get("conflict_detected", 0)
    run.failed_count = int(failed_count or 0)
    run.metrics_json = json.dumps({
        "status_counts": status_counts,
        "alignment_status_counts": alignment_status_counts,
    }, ensure_ascii=False)
    return run


def serialize_evaluation_set(evaluation_set):
    return {
        "id": evaluation_set.id,
        "name": evaluation_set.name,
        "course_id": getattr(evaluation_set, "course_id", None),
        "discipline": getattr(evaluation_set, "discipline", ""),
        "description": evaluation_set.description,
        "split": evaluation_set.split,
        "locked": bool(evaluation_set.locked or getattr(evaluation_set, "is_locked", False)),
        "is_locked": bool(getattr(evaluation_set, "is_locked", False) or evaluation_set.locked),
        "created_by": evaluation_set.created_by,
        "created_at": evaluation_set.created_at,
        "updated_at": getattr(evaluation_set, "updated_at", ""),
    }


def serialize_evaluation_item(item):
    return {
        "id": item.id,
        "set_id": item.set_id,
        "evaluation_set_id": getattr(item, "evaluation_set_id", None) or item.set_id,
        "item_id": getattr(item, "item_id", "") or str(item.id),
        "split": getattr(item, "split", "test"),
        "discipline": getattr(item, "discipline", ""),
        "course_id": item.course_id,
        "english_term": item.english_term,
        "expected_chinese_term": item.expected_chinese_term,
        "expected_alignment_status": item.expected_alignment_status,
        "english_context": item.english_context,
        "english_evidence": item.english_evidence,
        "chinese_evidence": item.chinese_evidence,
        "expected_english_evidence": getattr(item, "expected_english_evidence", "") or item.english_evidence,
        "expected_chinese_evidence": getattr(item, "expected_chinese_evidence", "") or item.chinese_evidence,
        "negative_english_evidence": getattr(item, "negative_english_evidence", ""),
        "negative_chinese_evidence": getattr(item, "negative_chinese_evidence", ""),
        "difficulty": item.difficulty,
        "tags": safe_json_loads(getattr(item, "tags_json", "[]"), []),
        "annotator": getattr(item, "annotator", ""),
        "reviewed_by": getattr(item, "reviewed_by", ""),
        "disagreement_note": getattr(item, "disagreement_note", ""),
        "version": getattr(item, "version", "v1"),
        "created_at": item.created_at,
    }


def serialize_evaluation_run(run):
    return {
        "id": run.id,
        "evaluation_set_id": run.evaluation_set_id,
        "triggered_by": run.triggered_by,
        "provider": run.provider,
        "provider_name": getattr(run, "provider_name", "") or getattr(run, "provider", ""),
        "provider_mode": getattr(run, "provider_mode", ""),
        "model_name": run.model_name,
        "model_version": getattr(run, "model_version", ""),
        "prompt_key": getattr(run, "prompt_key", ""),
        "prompt_version": run.prompt_version,
        "retrieval_version": run.retrieval_version,
        "alignment_version": getattr(run, "alignment_version", ""),
        "commit_hash": getattr(run, "commit_hash", ""),
        "split": getattr(run, "split", "test"),
        "input_count": run.input_count,
        "skipped_count": getattr(run, "skipped_count", 0),
        "extraction_precision": getattr(run, "extraction_precision", 0.0),
        "extraction_recall": getattr(run, "extraction_recall", 0.0),
        "evidence_accuracy": getattr(run, "evidence_accuracy", 0.0),
        "english_evidence_accuracy": getattr(run, "english_evidence_accuracy", 0.0),
        "chinese_evidence_accuracy": getattr(run, "chinese_evidence_accuracy", 0.0),
        "alignment_accuracy": getattr(run, "alignment_accuracy", 0.0),
        "false_positive_rate": getattr(run, "false_positive_rate", 0.0),
        "auto_approval_error_rate": getattr(run, "auto_approval_error_rate", 0.0),
        "ocr_noise_term_rate": getattr(run, "ocr_noise_term_rate", None),
        "no_evidence_forced_alignment_rate": getattr(run, "no_evidence_forced_alignment_rate", 0.0),
        "metrics": safe_json_loads(run.metrics_json, {}),
        "report": safe_json_loads(run.report_json, {}),
        "report_json": safe_json_loads(run.report_json, {}),
        "report_markdown": getattr(run, "report_markdown", ""),
        "status": run.status,
        "error_message": getattr(run, "error_message", ""),
        "created_by": getattr(run, "created_by", 0),
        "created_at": run.created_at,
        "finished_at": getattr(run, "finished_at", ""),
    }


def serialize_background_job(job):
    progress_total = int(getattr(job, "progress_total", 0) or 0)
    progress_current = int(getattr(job, "progress_current", 0) or 0)
    progress_percent = 0
    if progress_total > 0:
        progress_percent = int(round(max(0, min(progress_current, progress_total)) * 100 / progress_total))
    return {
        "id": job.id,
        "job_type": job.job_type,
        "status": job.status,
        "priority": job.priority,
        "created_by": job.created_by,
        "course_id": job.course_id,
        "document_id": job.document_id,
        "alignment_run_id": job.alignment_run_id,
        "evaluation_run_id": job.evaluation_run_id,
        "scope_type": job.scope_type,
        "owner_user_id": job.owner_user_id,
        "input": safe_json_loads(job.input_json, {}),
        "result": safe_json_loads(job.result_json, {}),
        "progress_current": progress_current,
        "progress_total": progress_total,
        "progress_percent": progress_percent,
        "progress_message": job.progress_message,
        "error_code": job.error_code,
        "error_message": job.error_message,
        "attempt_count": job.attempt_count,
        "max_attempts": job.max_attempts,
        "started_at": job.started_at,
        "finished_at": job.finished_at,
        "canceled_at": job.canceled_at,
        "created_at": job.created_at,
        "updated_at": job.updated_at,
        "locked_by": job.locked_by,
        "locked_at": job.locked_at,
    }


def serialize_background_job_event(event):
    return {
        "id": event.id,
        "job_id": event.job_id,
        "event_type": event.event_type,
        "message": event.message,
        "progress_current": event.progress_current,
        "progress_total": event.progress_total,
        "metadata": safe_json_loads(event.metadata_json, {}),
        "created_at": event.created_at,
    }


def can_manage_evaluation_set(user, evaluation_set):
    if user.role == "admin":
        return True
    if evaluation_set is None:
        return False
    if evaluation_set.created_by == user.id:
        return True
    course_id = getattr(evaluation_set, "course_id", None)
    if course_id:
        course = db.session.get(Course, course_id)
        return bool(course and can_manage_course(user, course))
    return False


class JobExecutionError(Exception):
    def __init__(self, message, error_code="INTERNAL_ERROR", retryable=True):
        super().__init__(message)
        self.error_code = error_code
        self.retryable = retryable


def job_input(job):
    return safe_json_loads(getattr(job, "input_json", "{}"), {})


def add_job_event(job, event_type, message, progress_current=None, progress_total=None, metadata=None):
    if job is None:
        return None
    event = BackgroundJobEvent(
        job_id=job.id,
        event_type=str(event_type or "")[:80],
        message=redact_for_log(message or ""),
        progress_current=int(progress_current if progress_current is not None else (job.progress_current or 0)),
        progress_total=int(progress_total if progress_total is not None else (job.progress_total or 0)),
        metadata_json=json.dumps(metadata or {}, ensure_ascii=False),
        created_at=current_time_text()
    )
    db.session.add(event)
    db.session.flush()
    return event


def create_background_job(
    job_type,
    created_by,
    course_id=None,
    document_id=None,
    alignment_run_id=None,
    evaluation_run_id=None,
    scope_type="",
    owner_user_id=None,
    input_data=None,
    priority=100,
    max_attempts=None,
):
    job_type = str(job_type or "").strip()
    if job_type not in JOB_TYPES:
        raise ValueError(f"Unsupported job_type: {job_type}")
    now = current_time_text()
    creator_id = created_by.id if hasattr(created_by, "id") else int(created_by or 0)
    job = BackgroundJob(
        job_type=job_type,
        status="queued",
        priority=int(priority or 100),
        created_by=creator_id,
        course_id=course_id,
        document_id=document_id,
        alignment_run_id=alignment_run_id,
        evaluation_run_id=evaluation_run_id,
        scope_type=str(scope_type or ""),
        owner_user_id=owner_user_id,
        input_json=json.dumps(input_data or {}, ensure_ascii=False),
        result_json="{}",
        progress_current=0,
        progress_total=100,
        progress_message="Queued",
        attempt_count=0,
        max_attempts=int(max_attempts or JOB_MAX_ATTEMPTS),
        created_at=now,
        updated_at=now,
    )
    db.session.add(job)
    db.session.flush()
    add_job_event(job, "created", f"{job_type} queued.", metadata={"job_type": job_type})
    return job


def can_view_job(user, job):
    if not user or job is None:
        return False
    if user.role == "admin":
        return True
    if job.created_by == user.id or job.owner_user_id == user.id:
        return True
    if user.role == "teacher" and job.course_id:
        course = db.session.get(Course, job.course_id)
        return can_manage_course(user, course)
    return False


def can_mutate_job(user, job):
    return can_view_job(user, job)


def visible_jobs_query(user):
    query = BackgroundJob.query
    if user.role == "admin":
        return query
    if user.role == "student":
        return query.filter(db.or_(
            BackgroundJob.created_by == user.id,
            BackgroundJob.owner_user_id == user.id
        ))
    manageable_ids = [
        course.id for course in Course.query.all()
        if can_manage_course(user, course)
    ]
    return query.filter(db.or_(
        BackgroundJob.created_by == user.id,
        BackgroundJob.course_id.in_(manageable_ids or [-1])
    ))


def update_job_progress(job, current=None, total=None, message=""):
    if current is not None:
        job.progress_current = max(0, int(current))
    if total is not None:
        job.progress_total = max(0, int(total))
    if message:
        job.progress_message = str(message)
    job.updated_at = current_time_text()
    add_job_event(job, "progress", job.progress_message, job.progress_current, job.progress_total)
    db.session.flush()


def copy_evaluation_run_result(source_run, target_run):
    fields = [
        "provider", "model_name", "model_version", "prompt_version", "retrieval_version",
        "alignment_version", "commit_hash", "split", "input_count", "skipped_count",
        "extraction_precision", "extraction_recall", "evidence_accuracy",
        "english_evidence_accuracy", "chinese_evidence_accuracy", "alignment_accuracy",
        "false_positive_rate", "auto_approval_error_rate", "ocr_noise_term_rate",
        "no_evidence_forced_alignment_rate", "metrics_json", "report_json",
        "report_markdown", "status", "error_message", "finished_at"
    ]
    for field in fields:
        setattr(target_run, field, getattr(source_run, field))
    if not target_run.created_at:
        target_run.created_at = getattr(source_run, "created_at", current_time_text())
    db.session.flush()
    return target_run


def resolve_project_file_path(file_path):
    raw_path = str(file_path or "").strip()
    if not raw_path:
        raise ValueError("file_path 不能为空。")
    candidate = os.path.abspath(raw_path if os.path.isabs(raw_path) else os.path.join(PROJECT_ROOT, raw_path))
    allowed_roots = [
        os.path.abspath(PROJECT_ROOT),
        os.path.abspath(os.path.join(PROJECT_ROOT, "docs")),
        os.path.abspath(os.path.join(PROJECT_ROOT, "tests")),
    ]
    if not any(candidate == root or candidate.startswith(root + os.sep) for root in allowed_roots):
        raise ValueError("file_path 必须位于项目目录内。")
    if not os.path.exists(candidate):
        raise FileNotFoundError(f"评估文件不存在：{file_path}")
    return candidate


def import_evaluation_items(file_path, evaluation_set_id, user):
    evaluation_set = db.session.get(EvaluationSet, int(evaluation_set_id))
    if evaluation_set is None:
        raise ValueError("evaluation_set_id 无效。")
    if not can_manage_evaluation_set(user, evaluation_set):
        raise PermissionError("无权导入该 EvaluationSet。")
    if bool(getattr(evaluation_set, "is_locked", False) or getattr(evaluation_set, "locked", False)):
        raise PermissionError("EvaluationSet 已锁定，不能导入新项目。")

    records, errors = read_evaluation_jsonl(resolve_project_file_path(file_path))
    imported = 0
    skipped = 0
    for raw in records:
        item = normalize_evaluation_record(raw)
        if not item["english_term"]:
            skipped += 1
            errors.append({"line": None, "error": "english_term is required", "item_id": item.get("item_id", "")})
            continue
        course_id = item.get("course_id") or evaluation_set.course_id
        course_id = int(course_id) if str(course_id or "").isdigit() else None
        if user.role != "admin" and course_id:
            course = db.session.get(Course, course_id)
            if course and not can_manage_course(user, course):
                skipped += 1
                errors.append({"line": None, "error": "no course permission", "item_id": item.get("item_id", "")})
                continue
        db.session.add(EvaluationItem(
            set_id=evaluation_set.id,
            evaluation_set_id=evaluation_set.id,
            item_id=item["item_id"],
            split=item["split"],
            discipline=item["discipline"] or evaluation_set.discipline,
            course_id=course_id,
            english_term=item["english_term"],
            expected_chinese_term=item["expected_chinese_term"],
            expected_alignment_status=item["expected_alignment_status"],
            english_context=item["english_context"],
            english_evidence=item["expected_english_evidence"],
            chinese_evidence=item["expected_chinese_evidence"],
            expected_english_evidence=item["expected_english_evidence"],
            expected_chinese_evidence=item["expected_chinese_evidence"],
            negative_english_evidence=item["negative_english_evidence"],
            negative_chinese_evidence=item["negative_chinese_evidence"],
            difficulty=item["difficulty"],
            tags_json=json.dumps(item["tags"], ensure_ascii=False),
            annotator=item["annotator"],
            reviewed_by=item["reviewed_by"],
            disagreement_note=item["disagreement_note"],
            version=item["version"],
            created_at=current_time_text()
        ))
        imported += 1
    evaluation_set.updated_at = current_time_text()
    db.session.flush()
    return {"imported_count": imported, "skipped_count": skipped, "errors": errors}


def evaluation_item_to_dict(item):
    return {
        "id": item.id,
        "item_id": getattr(item, "item_id", "") or str(item.id),
        "split": getattr(item, "split", "test"),
        "discipline": getattr(item, "discipline", ""),
        "course_id": item.course_id,
        "english_term": item.english_term,
        "expected_chinese_term": item.expected_chinese_term,
        "english_context": item.english_context,
        "expected_english_evidence": getattr(item, "expected_english_evidence", "") or item.english_evidence,
        "expected_chinese_evidence": getattr(item, "expected_chinese_evidence", "") or item.chinese_evidence,
        "expected_alignment_status": item.expected_alignment_status,
        "negative_english_evidence": getattr(item, "negative_english_evidence", ""),
        "negative_chinese_evidence": getattr(item, "negative_chinese_evidence", ""),
        "difficulty": item.difficulty,
        "tags": safe_json_loads(getattr(item, "tags_json", "[]"), []),
    }


def evaluate_single_evaluation_item(item):
    payload = evaluation_item_to_dict(item)
    context = payload["english_context"] or payload["english_term"]
    extracted_terms = extract_terms_from_text(context)
    course = db.session.get(Course, item.course_id) if item.course_id else None
    alignment = generate_alignment_result(
        english_term=item.english_term,
        courseware_sentence=context,
        course=course.name if course else "",
        chapter="",
        scope_type="course"
    )
    return evaluate_single_item_result(payload, extracted_terms, alignment)


def run_evaluation_set(evaluation_set, user, split="test", model_version="", prompt_version="", retrieval_version=""):
    if not can_manage_evaluation_set(user, evaluation_set):
        raise PermissionError("无权运行该 EvaluationSet。")
    split = str(split or "test").strip().lower() or "test"
    items = EvaluationItem.query.filter_by(set_id=evaluation_set.id, split=split).order_by(EvaluationItem.id.asc()).all()
    if not items:
        items = EvaluationItem.query.filter_by(set_id=evaluation_set.id).order_by(EvaluationItem.id.asc()).all()
    if not items:
        raise ValueError("Evaluation set 没有测试项。")

    meta = current_provider_metadata()
    started = current_time_text()
    run = EvaluationRun(
        evaluation_set_id=evaluation_set.id,
        triggered_by=user.id,
        created_by=user.id,
        provider=meta["provider"],
        provider_name=meta["provider"],
        provider_mode=meta.get("provider_mode", ""),
        model_name=meta["model_name"],
        model_version=model_version or meta["model_name"],
        prompt_key="term_alignment",
        prompt_version=prompt_version or "v1",
        retrieval_version=retrieval_version or meta["retrieval_version"],
        alignment_version=ALIGNMENT_PROMPT_VERSION,
        commit_hash=os.environ.get("COMMIT_HASH", "local"),
        split=split,
        status="running",
        created_at=started
    )
    db.session.add(run)
    db.session.flush()
    try:
        results = [evaluate_single_evaluation_item(item) for item in items]
        metrics = compute_evaluation_metrics(results)
        run_info = {
            "evaluation_set_id": evaluation_set.id,
            "evaluation_set_name": evaluation_set.name,
            "split": split,
            "model_version": run.model_version,
            "prompt_version": run.prompt_version,
            "retrieval_version": run.retrieval_version,
            "created_at": started,
        }
        report_markdown, report_json = generate_evaluation_report(run_info, results, metrics)
        run.input_count = metrics["input_count"]
        run.skipped_count = metrics["skipped_count"]
        run.extraction_precision = metrics["extraction_precision"] or 0
        run.extraction_recall = metrics["extraction_recall"] or 0
        run.evidence_accuracy = metrics["evidence_accuracy"] or 0
        run.english_evidence_accuracy = metrics["english_evidence_accuracy"] or 0
        run.chinese_evidence_accuracy = metrics["chinese_evidence_accuracy"] or 0
        run.alignment_accuracy = metrics["alignment_accuracy"] or 0
        run.false_positive_rate = metrics["false_positive_rate"] or 0
        run.auto_approval_error_rate = metrics["auto_approval_error_rate"] or 0
        run.ocr_noise_term_rate = metrics["ocr_noise_term_rate"]
        run.no_evidence_forced_alignment_rate = metrics["no_evidence_forced_alignment_rate"] or 0
        run.metrics_json = json.dumps(metrics, ensure_ascii=False)
        run.report_json = json.dumps(report_json, ensure_ascii=False)
        run.report_markdown = report_markdown
        run.status = "completed"
        run.finished_at = current_time_text()
    except Exception as exc:
        run.status = "failed"
        run.error_message = redact_for_log(exc)
        run.finished_at = current_time_text()
        add_system_log("error", "evaluation", f"EvaluationRun {run.id} failed: {exc}")
        db.session.flush()
        return run
    db.session.flush()
    return run


def process_document_ingestion_job(job):
    data = job_input(job)
    document = db.session.get(Document, job.document_id)
    if document is None:
        raise JobExecutionError("Document not found.", "RESOURCE_NOT_FOUND", retryable=False)
    user = db.session.get(User, job.created_by)
    course = db.session.get(Course, document.course_id) if document.course_id else None
    audit_context = audit_context_service.normalize_audit_context({
        "request_id": f"background-job-{job.id}",
        "actor_id": getattr(user, "id", None),
        "actor_role": getattr(user, "role", ""),
        "actor_name": getattr(user, "display_name", "") or getattr(user, "username", "") or getattr(user, "email", ""),
        "source": "system",
    })
    try:
        save_path = storage_service().absolute_path(data.get("storage_key") or document.storage_key) if (data.get("storage_key") or getattr(document, "storage_key", "")) else os.path.join(UPLOAD_FOLDER, document.saved_filename)
    except Exception:
        save_path = data.get("save_path") or os.path.join(UPLOAD_FOLDER, document.saved_filename)
    if not os.path.exists(save_path):
        raise JobExecutionError("Saved upload file is missing.", "PARSING_FAILED", retryable=False)

    ingestion_job = IngestionJob.query.filter_by(document_id=document.id).order_by(IngestionJob.id.desc()).first()
    now = current_time_text()
    document.parsing_status = "processing"
    if ingestion_job:
        ingestion_job.status = "running"
        ingestion_job.started_at = ingestion_job.started_at or now
    update_job_progress(job, 5, 100, "Parsing document")

    DocumentChunk.query.filter_by(document_id=document.id).delete()
    FormulaBlock.query.filter_by(document_id=document.id).delete()
    KnowledgeChunk.query.filter_by(document_id=document.id).delete()
    db.session.flush()

    parse_record = DocumentParseRecord.query.filter_by(parse_uid=document.parse_uid).first() if document.parse_uid else None
    if parse_record is None:
        parse_result, parse_record, parse_blocks = create_parse_record_for_saved_file(
            save_path,
            document,
            document.original_filename or document.filename,
            document.content_type,
            audit_context=audit_context,
        )
    else:
        parse_blocks = DocumentParseBlock.query.filter_by(parse_uid=parse_record.parse_uid).order_by(DocumentParseBlock.block_index.asc()).all()
        document.parse_uid = parse_record.parse_uid
        document.parsed_text = "\n\n".join(clean_text(block.text) for block in parse_blocks if clean_text(block.text))
        document.ocr_required = bool(parse_record.ocr_required)
        document.ocr_status = parse_record.quality_status if parse_record.ocr_required else "not_required"
        document.ocr_error = parse_record.error_message or ""
        document.quality_flags_json = parse_record.quality_flags
    if not document_parse_quality_service.should_allow_term_extraction(parse_record):
        blocked = block_document_by_quality_gate(document, ingestion_job, parse_record, audit_context=audit_context)
        raise JobExecutionError(blocked["blocked_reason"], quality_gate_error_code(parse_record), retryable=False)
    parsed_chunks = parse_block_records_to_chunk_items(document, parse_record, parse_blocks)
    parsed_text = document.parsed_text
    ocr_required = bool(parse_record.ocr_required)
    ocr_meta = {
        "ocr_provider": "",
        "ocr_status": document.ocr_status or "not_required",
        "ocr_error": document.ocr_error or "",
        "ocr_confidence": 100,
        "warnings": safe_json_loads(getattr(parse_record, "warnings", "[]"), []),
    }
    parsed_formula_blocks = []
    parser_warnings = list(ocr_meta.get("warnings", []) or [])
    update_job_progress(job, 25, 100, "Saving OCR and formula blocks")

    formula_records = []
    formula_statuses = []
    formula_providers = []
    formula_flags = set()
    for formula_item in parsed_formula_blocks or []:
        item_flags = list(formula_item.get("quality_flags", []) or [])
        formula_status = formula_item.get("status", "")
        if formula_status:
            formula_statuses.append(formula_status)
            if formula_status != "ok":
                formula_flags.add(formula_status)
        provider = formula_item.get("provider", "")
        if provider:
            formula_providers.append(provider)
        formula_record = FormulaBlock(
            document_id=document.id,
            course_id=document.course_id,
            owner_user_id=document.owner_user_id,
            scope_type=document.scope_type,
            page_number=formula_item.get("page_number"),
            slide_number=formula_item.get("slide_number"),
            bbox_json=json.dumps(formula_item.get("bbox", {}) or {}, ensure_ascii=False),
            image_path=formula_item.get("image_path", ""),
            latex=formula_item.get("latex", ""),
            plain_text=formula_item.get("plain_text", ""),
            provider=provider,
            confidence=float(formula_item.get("confidence", 0) or 0),
            status=formula_status,
            error=formula_item.get("error", ""),
            quality_flags_json=json.dumps(item_flags, ensure_ascii=False),
            created_at=current_time_text()
        )
        db.session.add(formula_record)
        formula_records.append(formula_record)
    db.session.flush()

    if not parsed_chunks:
        if document.ocr_required:
            document.parsing_status = "needs_ocr_engine"
            if formula_flags and document.ocr_status in {"", "ocr_unavailable", "empty_result", "ocr_failed"}:
                document.parsing_status = "needs_formula_ocr_engine"
            message = document.ocr_error or "OCR required but unavailable or no text detected."
            document.error_message = message
            if ingestion_job:
                ingestion_job.status = "failed"
                ingestion_job.error_message = message
                ingestion_job.finished_at = current_time_text()
            error_code = "FORMULA_OCR_UNAVAILABLE" if document.parsing_status == "needs_formula_ocr_engine" else "OCR_UNAVAILABLE"
            add_system_log("warning", "ocr", f"Document {document.id} OCR unavailable in background job: {message}")
            raise JobExecutionError(message, error_code, retryable=False)
        raise JobExecutionError("No valid text chunks were parsed from document.", "PARSING_FAILED", retryable=False)

    page_units = len({
        chunk.get("page_number") or chunk.get("slide_number") or index
        for index, chunk in enumerate(parsed_chunks, start=1)
    })
    if document.scope_type == "personal" and user is not None:
        allowed, reasons, plan, subscription, totals = check_usage_quota(
            user,
            page_units=page_units,
            ai_units=min(12, len(parsed_chunks))
        )
        if not allowed:
            document.parsing_status = "quota_blocked"
            message = "；".join(reasons)
            if ingestion_job:
                ingestion_job.status = "failed"
                ingestion_job.error_message = message
                ingestion_job.finished_at = current_time_text()
            raise JobExecutionError(message, "QUOTA_EXCEEDED", retryable=False)

    update_job_progress(job, 50, 100, "Creating structured text chunks")
    chunk_records = []
    document_flags = set()
    if ocr_required:
        document_flags.add("ocr_triggered")
    if formula_records:
        document_flags.add("formula_region_detected")
    for flag in formula_flags:
        document_flags.add(flag)
    for warning in parser_warnings:
        if "formula" in str(warning).lower():
            document_flags.add("formula_warning")

    for index, item in enumerate(parsed_chunks, start=1):
        item_content = item.get("content", "")
        item_ocr_confidence = int(item.get("ocr_confidence", 100) or 100)
        item_flags = list(item.get("quality_flags", []) or [])
        if contains_ocr_placeholder(item_content) or contains_formula_placeholder(item_content):
            item_flags.append("ocr_placeholder_blocked")
        if item_ocr_confidence < 60:
            item_flags.append("ocr_low_confidence")
            document_flags.add("ocr_low_confidence")
        for flag in item_flags:
            if flag and flag != "native_text_ok":
                document_flags.add(flag)
        chunk = DocumentChunk(
            document_id=document.id,
            course_id=document.course_id,
            user_id=document.owner_user_id,
            owner_user_id=document.owner_user_id,
            chunk_index=index,
            parse_uid=item.get("parse_uid", getattr(document, "parse_uid", "")),
            parse_block_uid=item.get("parse_block_uid", ""),
            language=item.get("language", document.language),
            page_number=item.get("page_number"),
            slide_number=item.get("slide_number"),
            section_title=item.get("section_title", ""),
            content=item_content,
            source_type=item.get("source_type", document.source_type),
            source_location=item.get("source_location", ""),
            ocr_confidence=item_ocr_confidence,
            ocr_provider=item.get("ocr_provider", ""),
            ocr_status=item.get("ocr_status", "not_required"),
            ocr_error=item.get("ocr_error", ""),
            quality_flags_json=json.dumps(item_flags, ensure_ascii=False),
            created_at=current_time_text()
        )
        db.session.add(chunk)
        chunk_records.append(chunk)
    db.session.flush()
    document.quality_flags_json = json.dumps(sorted(document_flags), ensure_ascii=False)

    update_job_progress(job, 70, 100, "Indexing knowledge chunks")
    knowledge_metadata = build_governed_ingestion_metadata(
        parse_record=parse_record,
        title=str(data.get("source_name", "")).strip() or document.filename,
        course=course,
        chapter=str(data.get("chapter", "")).strip(),
        language=document.language,
        source_type=document.source_type,
        scope_type=document.scope_type,
        owner_user=user,
        owner_user_id=document.owner_user_id,
        document_id=document.id,
        knowledge_base_type="student_personal_kb" if document.scope_type == "personal" else ("en_course_kb" if document.language == "en" else "zh_course_kb"),
        visibility=visibility_for_scope(document.scope_type),
        content_hash=getattr(document, "sha256", "") or getattr(document, "file_sha256", ""),
        extra_quality_flags=document_flags,
    )
    governed_ingestion = knowledge_ingestion_service.ingest_parse_record_to_governed_knowledge(
        db.session,
        knowledge_ingestion_models(),
        parse_record,
        parse_blocks,
        knowledge_metadata,
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=False,
    )
    source = governed_ingestion.source
    knowledge_chunks = governed_ingestion.chunks
    version = create_kb_version(
        document.scope_type,
        document.course_id,
        document.owner_user_id,
        source_count=1,
        chunk_count=len(knowledge_chunks),
        description=f"Document upload: {document.filename}"
    )

    document.parsing_status = "parsed_with_warnings" if document_flags else "parsed"
    if ingestion_job:
        ingestion_job.status = "completed"
        ingestion_job.finished_at = current_time_text()
        ingestion_job.processed_pages = page_units
        ingestion_job.source_id = source.id
    if document.scope_type == "personal" and user is not None:
        record_usage(user.id, "document_parse_page", page_units, related_document_id=document.id)
        if document.ocr_required:
            record_usage(user.id, "ocr_page", page_units, related_document_id=document.id)
    update_job_progress(job, 100, 100, "Document parsed and indexed")
    ingestion_status = "partial" if parse_record.quality_status == "partial_text" else "ingested"
    record_document_ingestion_audit(
        document,
        parse_record,
        ingestion_status,
        audit_context=audit_context,
        commit=False,
    )
    return {
        "document_id": document.id,
        "parsing_status": document.parsing_status,
        **parse_quality_summary(parse_record),
        "ingestion_status": ingestion_status,
        "chunks_created": len(chunk_records),
        "formula_blocks_created": len(formula_records),
        "knowledge_chunks_created": len(knowledge_chunks),
        "knowledge_source_id": source.id,
        "source_uid": getattr(source, "source_uid", ""),
        "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in knowledge_chunks[:20] if getattr(chunk, "chunk_uid", "")],
        "knowledge_base_version_id": version.id,
        "ocr_status": document.ocr_status or "not_required",
        "ocr_provider": document.ocr_provider or "not_required",
        "formula_status": next((status for status in formula_statuses if status != "ok"), formula_statuses[0] if formula_statuses else "not_required"),
        "formula_provider": next(iter(formula_providers), "not_required") if formula_providers else "not_required",
        "warnings": parser_warnings,
    }


def process_alignment_job(job):
    data = job_input(job)
    user = db.session.get(User, job.created_by)
    alignment_run = db.session.get(AlignmentRun, job.alignment_run_id) if job.alignment_run_id else None
    if alignment_run is None:
        raise JobExecutionError("AlignmentRun not found.", "RESOURCE_NOT_FOUND", retryable=False)
    alignment_run.status = "running"
    alignment_run.started_at = alignment_run.started_at or current_time_text()
    update_job_progress(job, 10, 100, "Starting terminology alignment")

    document_id = data.get("document_id") or job.document_id
    scope_type = str(data.get("scope_type") or job.scope_type or "course").strip()
    course = db.session.get(Course, job.course_id) if job.course_id else None
    if document_id:
        document = db.session.get(Document, int(document_id))
        if document is None:
            raise JobExecutionError("Document not found.", "RESOURCE_NOT_FOUND", retryable=False)
        chunks = DocumentChunk.query.filter_by(document_id=document.id).order_by(DocumentChunk.chunk_index.asc()).all()
        if not chunks:
            raise JobExecutionError("Document has no parsed chunks. Run document ingestion first.", "PARSING_FAILED", retryable=False)
        update_job_progress(job, 35, 100, "Extracting terms and retrieving evidence")
        cards = run_alignment_for_chunks(
            chunks,
            course=course or (db.session.get(Course, document.course_id) if document.course_id else None),
            scope_type=document.scope_type if document.scope_type == "personal" else "course",
            owner_user_id=document.owner_user_id if document.scope_type == "personal" else None,
            source_document_id=document.id,
            triggered_by_user_id=job.created_by
        )
        for card in cards:
            card.alignment_run_id = alignment_run.id
            card.source_alignment_run_id = alignment_run.id
        update_alignment_run_stats(alignment_run, cards=cards, term_count=len(cards))
        alignment_run.status = "completed"
        alignment_run.finished_at = current_time_text()
        update_job_progress(job, 100, 100, "Alignment completed")
        return {
            "alignment_run_id": alignment_run.id,
            "cards_created": len(cards),
            "cards": [serialize_terminology_card(card) for card in cards],
        }

    english_term = str(data.get("english_term", "")).strip()
    if not english_term:
        raise JobExecutionError("english_term or document_id is required.", "VALIDATION_ERROR", retryable=False)
    update_job_progress(job, 35, 100, "Retrieving bilingual evidence")
    alignment = generate_alignment_result(
        english_term=english_term,
        courseware_sentence=str(data.get("courseware_sentence", "")).strip(),
        course=course.name if course else "",
        chapter=str(data.get("chapter", "")).strip(),
        scope_type=scope_type,
        owner_user_id=job.owner_user_id if scope_type == "personal" else None
    )
    alignment["alignment_run_id"] = alignment_run.id
    card = create_or_update_card_from_alignment(
        english_term=english_term,
        alignment=alignment,
        scope_type=scope_type,
        course_id=course.id if course else None,
        owner_user_id=job.owner_user_id if scope_type == "personal" else None,
        courseware_sentence=str(data.get("courseware_sentence", "")).strip()
    )
    if scope_type == "personal" and card and user:
        record_usage(user.id, "ai_alignment", 1, related_term_id=card.id)
    update_alignment_run_stats(alignment_run, cards=[card] if card else [], term_count=1)
    alignment_run.status = "completed"
    alignment_run.finished_at = current_time_text()
    update_job_progress(job, 100, 100, "Alignment completed")
    return {
        "alignment_run_id": alignment_run.id,
        "card": serialize_terminology_card(card) if card else None,
        "alignment": alignment,
    }


def process_evaluation_job(job):
    data = job_input(job)
    user = db.session.get(User, job.created_by)
    if user is None:
        raise JobExecutionError("Evaluation user not found.", "RESOURCE_NOT_FOUND", retryable=False)
    evaluation_run = db.session.get(EvaluationRun, job.evaluation_run_id) if job.evaluation_run_id else None
    if evaluation_run is None:
        raise JobExecutionError("EvaluationRun not found.", "RESOURCE_NOT_FOUND", retryable=False)
    evaluation_set = db.session.get(EvaluationSet, evaluation_run.evaluation_set_id)
    if evaluation_set is None:
        raise JobExecutionError("EvaluationSet not found.", "RESOURCE_NOT_FOUND", retryable=False)
    update_job_progress(job, 10, 100, "Running evaluation harness")
    evaluation_run.status = "running"
    evaluation_run.created_at = evaluation_run.created_at or current_time_text()
    produced_run = run_evaluation_set(
        evaluation_set,
        user,
        split=data.get("split", evaluation_run.split or "test"),
        model_version=str(data.get("model_version", evaluation_run.model_version or "")).strip(),
        prompt_version=str(data.get("prompt_version", evaluation_run.prompt_version or "")).strip(),
        retrieval_version=str(data.get("retrieval_version", evaluation_run.retrieval_version or "")).strip(),
    )
    update_job_progress(job, 85, 100, "Saving evaluation report")
    if produced_run.id != evaluation_run.id:
        copy_evaluation_run_result(produced_run, evaluation_run)
        db.session.delete(produced_run)
    registry = ensure_model_registry_seed(owner_user_id=user.id)
    registry.last_evaluation_run_id = evaluation_run.id
    registry.updated_at = current_time_text()
    update_job_progress(job, 100, 100, "Evaluation completed")
    return {
        "evaluation_run_id": evaluation_run.id,
        "status": evaluation_run.status,
        "metrics": safe_json_loads(evaluation_run.metrics_json, {}),
    }


def run_background_job(job_id, worker_id=JOB_WORKER_ID):
    job = db.session.get(BackgroundJob, int(job_id))
    if job is None:
        raise JobExecutionError("BackgroundJob not found.", "RESOURCE_NOT_FOUND", retryable=False)
    if job.status == "canceled":
        add_job_event(job, "canceled", "Job was canceled before execution.")
        return job
    if job.status in TERMINAL_JOB_STATUSES:
        return job

    now = current_time_text()
    if job.status != "running":
        job.attempt_count = int(job.attempt_count or 0) + 1
    elif int(job.attempt_count or 0) == 0:
        job.attempt_count = 1
    job.status = "running"
    job.locked_by = worker_id
    job.locked_at = now
    job.started_at = job.started_at or now
    job.updated_at = now
    add_job_event(job, "started", f"{job.job_type} started by {worker_id}.")
    db.session.flush()

    try:
        if job.job_type == "document_ingestion":
            result = process_document_ingestion_job(job)
        elif job.job_type == "alignment_run":
            result = process_alignment_job(job)
        elif job.job_type == "evaluation_run":
            result = process_evaluation_job(job)
        else:
            raise JobExecutionError(f"Unsupported job type: {job.job_type}", "VALIDATION_ERROR", retryable=False)

        if job.status != "canceled":
            job.status = "completed"
            job.error_code = ""
            job.error_message = ""
            job.result_json = json.dumps(result or {}, ensure_ascii=False)
            job.progress_current = job.progress_total or 100
            job.progress_message = "Completed"
            job.finished_at = current_time_text()
            job.updated_at = job.finished_at
            add_job_event(job, "completed", "Job completed.", job.progress_current, job.progress_total, result)
        db.session.commit()
        return job
    except JobExecutionError as exc:
        error_code = exc.error_code
        error_message = redact_for_log(str(exc))
        job.error_code = error_code
        job.error_message = error_message
        job.updated_at = current_time_text()
        should_retry = bool(exc.retryable) and int(job.attempt_count or 0) < int(job.max_attempts or JOB_MAX_ATTEMPTS)
        if should_retry:
            job.status = "retrying"
            job.progress_message = "Retry scheduled"
            add_job_event(job, "retrying", error_message, metadata={"error_code": error_code, "attempt": job.attempt_count})
        else:
            job.status = "failed"
            job.finished_at = current_time_text()
            job.progress_message = "Failed"
            add_job_event(job, "failed", error_message, metadata={"error_code": error_code, "attempt": job.attempt_count})
        add_system_log("error", "background_job", f"BackgroundJob {job.id} failed: {error_message}")
        db.session.commit()
        return job
    except Exception as exc:
        error_message = redact_for_log(exc)
        job.error_code = "INTERNAL_ERROR"
        job.error_message = error_message
        job.updated_at = current_time_text()
        if int(job.attempt_count or 0) < int(job.max_attempts or JOB_MAX_ATTEMPTS):
            job.status = "retrying"
            job.progress_message = "Retry scheduled"
            add_job_event(job, "retrying", error_message, metadata={"error_code": "INTERNAL_ERROR", "attempt": job.attempt_count})
        else:
            job.status = "failed"
            job.finished_at = current_time_text()
            job.progress_message = "Failed"
            add_job_event(job, "failed", error_message, metadata={"error_code": "INTERNAL_ERROR", "attempt": job.attempt_count})
        add_system_log("error", "background_job", f"BackgroundJob {job.id} failed: {error_message}")
        db.session.commit()
        return job


def claim_next_background_job(worker_id=JOB_WORKER_ID):
    job = BackgroundJob.query.filter(
        BackgroundJob.status.in_(["queued", "retrying"])
    ).order_by(BackgroundJob.priority.asc(), BackgroundJob.id.asc()).first()
    if job is None:
        return None
    now = current_time_text()
    job.status = "running"
    job.attempt_count = int(job.attempt_count or 0) + 1
    job.locked_by = worker_id
    job.locked_at = now
    job.started_at = job.started_at or now
    job.updated_at = now
    add_job_event(job, "claimed", f"Job claimed by {worker_id}.")
    db.session.commit()
    return job


def run_worker_once(worker_id=JOB_WORKER_ID):
    job = claim_next_background_job(worker_id)
    if job is None:
        return None
    return run_background_job(job.id, worker_id=worker_id)


def serialize_model_registry(registry):
    return {
        "id": registry.id,
        "provider": registry.provider,
        "model_name": registry.model_name,
        "model_version": registry.model_version,
        "prompt_version": registry.prompt_version,
        "retrieval_version": registry.retrieval_version,
        "enabled": bool(registry.enabled),
        "allowed_workflows": safe_json_loads(registry.allowed_workflows, registry.allowed_workflows.split(",") if registry.allowed_workflows else []),
        "last_evaluation_run_id": registry.last_evaluation_run_id,
        "known_risks": registry.known_risks,
        "rollback_target": registry.rollback_target,
        "owner": registry.owner,
        "created_at": registry.created_at,
        "updated_at": registry.updated_at,
    }


def serialize_ai_provider_config(config):
    return {
        "id": config.id,
        "provider_name": config.provider_name,
        "provider_mode": config.provider_mode,
        "base_url": config.base_url,
        "default_model": config.default_model,
        "is_enabled": bool(config.is_enabled),
        "is_default": bool(config.is_default),
        "supports_json_schema": bool(config.supports_json_schema),
        "supports_streaming": bool(config.supports_streaming),
        "supports_vision": bool(config.supports_vision),
        "supports_formula_reasoning": bool(config.supports_formula_reasoning),
        "max_input_tokens": config.max_input_tokens,
        "max_output_tokens": config.max_output_tokens,
        "timeout_seconds": config.timeout_seconds,
        "max_retries": config.max_retries,
        "cost_per_1k_input_tokens": float(config.cost_per_1k_input_tokens or 0),
        "cost_per_1k_output_tokens": float(config.cost_per_1k_output_tokens or 0),
        "health_status": config.health_status,
        "last_healthcheck_at": config.last_healthcheck_at,
        "created_at": config.created_at,
        "updated_at": config.updated_at,
    }


def serialize_ai_model_registry(model):
    return {
        "id": model.id,
        "provider_name": model.provider_name,
        "model_name": model.model_name,
        "model_version": model.model_version,
        "model_display_name": model.model_display_name,
        "provider_mode": model.provider_mode,
        "supports_json_output": bool(model.supports_json_output),
        "supports_tool_calling": bool(model.supports_tool_calling),
        "supports_vision": bool(model.supports_vision),
        "max_input_tokens": model.max_input_tokens,
        "max_output_tokens": model.max_output_tokens,
        "cost_per_1k_input_tokens": float(model.cost_per_1k_input_tokens or 0),
        "cost_per_1k_output_tokens": float(model.cost_per_1k_output_tokens or 0),
        "is_enabled": bool(model.is_enabled),
        "is_default_for_provider": bool(model.is_default_for_provider),
        "last_evaluation_run_id": model.last_evaluation_run_id,
        "last_evaluation_score": float(model.last_evaluation_score or 0),
        "known_risks": safe_json_loads(model.known_risks_json, []),
        "created_at": model.created_at,
        "updated_at": model.updated_at,
    }


def serialize_prompt_template(prompt):
    return {
        "id": prompt.id,
        "prompt_key": prompt.prompt_key,
        "prompt_version": prompt.prompt_version,
        "task_type": prompt.task_type,
        "language": prompt.language,
        "json_schema": safe_json_loads(prompt.json_schema, {}),
        "is_active": bool(prompt.is_active),
        "is_default": bool(prompt.is_default),
        "created_by": prompt.created_by,
        "created_at": prompt.created_at,
        "updated_at": prompt.updated_at,
        "notes": prompt.notes,
    }


def serialize_ai_call_log(log):
    return {
        "id": log.id,
        "task_type": log.task_type,
        "provider_name": log.provider_name,
        "provider_mode": log.provider_mode,
        "model_name": log.model_name,
        "prompt_key": log.prompt_key,
        "prompt_version": log.prompt_version,
        "user_id": log.user_id,
        "course_id": log.course_id,
        "document_id": log.document_id,
        "job_id": log.job_id,
        "alignment_run_id": log.alignment_run_id,
        "evaluation_run_id": log.evaluation_run_id,
        "request_hash": log.request_hash,
        "response_hash": log.response_hash,
        "input_token_count": log.input_token_count,
        "output_token_count": log.output_token_count,
        "estimated_cost": float(log.estimated_cost or 0),
        "latency_ms": log.latency_ms,
        "status": log.status,
        "error_code": log.error_code,
        "error_message": log.error_message,
        "redacted_prompt_preview": log.redacted_prompt_preview,
        "redacted_response_preview": log.redacted_response_preview,
        "created_at": log.created_at,
    }


def serialize_personal_access_audit(audit):
    return {
        "id": audit.id,
        "actor_user_id": audit.actor_user_id,
        "target_user_id": audit.target_user_id,
        "resource_type": audit.resource_type,
        "resource_id": audit.resource_id,
        "reason": audit.reason,
        "created_at": audit.created_at,
    }


def serialize_subscription_plan(plan):
    return {
        "id": plan.id,
        "name": plan.name,
        "price_monthly": plan.price_monthly,
        "monthly_pages": plan.monthly_pages,
        "monthly_ai_calls": plan.monthly_ai_calls,
        "export_enabled": bool(plan.export_enabled),
        "description": plan.description,
        "is_active": bool(plan.is_active)
    }


def serialize_subscription(subscription):
    plan = db.session.get(SubscriptionPlan, subscription.plan_id) if subscription else None
    return {
        "id": subscription.id,
        "user_id": subscription.user_id,
        "plan_id": subscription.plan_id,
        "plan": serialize_subscription_plan(plan) if plan else None,
        "start_date": subscription.start_date,
        "end_date": subscription.end_date,
        "status": subscription.status,
        "auto_renew": bool(subscription.auto_renew)
    }


def serialize_usage(record):
    return {
        "id": record.id,
        "user_id": record.user_id,
        "action_type": record.action_type,
        "units_used": record.units_used,
        "related_document_id": record.related_document_id,
        "related_term_id": record.related_term_id,
        "created_at": record.created_at
    }


def current_time_text():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def future_time_text(minutes):
    return (datetime.now() + timedelta(minutes=minutes)).strftime("%Y-%m-%d %H:%M:%S")


def parse_time_text(value):
    if not value:
        return None
    try:
        return datetime.strptime(value, "%Y-%m-%d %H:%M:%S")
    except ValueError:
        return None


def normalize_role(role):
    value = str(role or "student").strip().lower()
    if value not in {"student", "teacher", "admin"}:
        return "student"
    return value


def validate_password_strength(password):
    errors = []
    if len(password) < 8:
        errors.append("至少 8 位")
    if not re.search(r"[A-Za-z]", password):
        errors.append("至少包含一个英文字母")
    if not re.search(r"\d", password):
        errors.append("至少包含一个数字")
    return errors


def make_mock_email_token():
    return f"{secrets.randbelow(1000000):06d}"


def log_mock_email(kind, email, token):
    message = f"{kind} mock email for {email}: verification_code={token}"
    print(f"[LexiBridge MockEmail] {message}")
    try:
        db.session.add(SystemLog(
            level="info",
            module="mock_email",
            message=message,
            created_at=current_time_text()
        ))
        db.session.commit()
    except Exception:
        db.session.rollback()


def find_user_by_email(email):
    return User.query.filter(db.func.lower(User.email) == str(email or "").strip().lower()).first()


def get_current_user():
    setattr(get_current_user, "_last_error_code", "")
    auth_header = request.headers.get("Authorization", "").strip()

    if not auth_header.lower().startswith("bearer "):
        return None

    token_value = auth_header.split(" ", 1)[1].strip()
    if not token_value:
        return None

    hashed = token_hash(token_value)
    token = AuthToken.query.filter_by(token_hash=hashed, revoked=False).first()
    if token is None:
        token = AuthToken.query.filter_by(token=token_value, revoked=False).first()
    if token is None:
        return None

    expires_at = parse_time_text(getattr(token, "expires_at", ""))
    if expires_at and expires_at < datetime.now():
        token.revoked = True
        db.session.commit()
        setattr(get_current_user, "_last_error_code", "TOKEN_EXPIRED")
        return None

    token.last_used_at = current_time_text()
    db.session.flush()
    return db.session.get(User, token.user_id)


def require_current_user(roles=None):
    user = get_current_user()

    if user is None:
        error_code = getattr(get_current_user, "_last_error_code", "") or "AUTH_REQUIRED"
        return None, api_error(error_code, "请先登录。", 401)

    if roles and user.role not in roles:
        return None, api_error("PERMISSION_DENIED", "当前账号没有权限执行该操作。", 403)

    return user, None


def create_auth_token(user):
    token_value = secrets.token_urlsafe(32)
    token = AuthToken(
        user_id=user.id,
        token=token_value,
        token_hash=token_hash(token_value),
        created_at=current_time_text(),
        expires_at=future_time_text(60 * 24 * 14),
        revoked=False
    )
    db.session.add(token)
    db.session.commit()
    return token_value


def add_system_log(level, module, message):
    try:
        db.session.add(SystemLog(
            level=str(level or "info"),
            module=str(module or "system")[:80],
            message=redact_for_log(message),
            created_at=current_time_text()
        ))
        db.session.flush()
    except Exception:
        db.session.rollback()


def record_personal_access(actor, target_user_id, resource_type, resource_id=None, reason=""):
    if actor is None or not target_user_id:
        return None
    if int(actor.id) == int(target_user_id):
        return None
    audit = PersonalAccessAudit(
        actor_user_id=actor.id,
        target_user_id=int(target_user_id),
        resource_type=str(resource_type or "")[:80],
        resource_id=resource_id,
        reason=redact_for_log(reason or "admin_personal_access"),
        created_at=current_time_text()
    )
    db.session.add(audit)
    db.session.flush()
    return audit


def current_provider_metadata():
    try:
        ensure_ai_registry_seed()
        selection = ai_selection_from_config()
    except Exception:
        selection = env_provider_selection(os.environ)
    is_real = selection.provider_mode == "live"
    return {
        "provider": selection.provider_name,
        "provider_mode": selection.provider_mode,
        "model_name": selection.model_name if is_real else (selection.model_name or selection.provider_mode),
        "is_real_provider": is_real,
        "prompt_version": ALIGNMENT_PROMPT_VERSION,
        "prompt_key": "term_alignment",
        "term_extraction_prompt_version": TERM_EXTRACTION_PROMPT_VERSION,
        "retrieval_version": RETRIEVAL_VERSION,
    }


def current_system_status(user=None):
    meta = current_provider_metadata()
    provider_name = str(meta.get("provider") or AI_PROVIDER or "none").lower()
    is_real = bool(meta.get("is_real_provider"))
    if is_real and provider_name == "deepseek":
        ai_status_text = "Live AI: DeepSeek active"
    elif provider_name == "mock":
        ai_status_text = "Mock AI: demonstration only, cannot auto-approve"
    elif provider_name in {"none", "local_heuristic"} or not DEEPSEEK_API_KEY:
        ai_status_text = "AI unavailable: local heuristic only"
    else:
        ai_status_text = "No AI provider configured"

    ocr_name = str(OCR_PROVIDER or "none").lower()
    if ocr_name in {"none", "mock"}:
        ocr_status_text = "OCR disabled" if ocr_name == "none" else "OCR unavailable: mock provider does not fabricate text"
    elif ocr_name == "tesseract":
        ocr_status_text = "Tesseract configured"
    elif ocr_name == "paddle":
        ocr_status_text = "PaddleOCR configured"
    elif ocr_name == "auto":
        ocr_status_text = "OCR auto-detection enabled"
    else:
        ocr_status_text = f"OCR provider configured: {OCR_PROVIDER}"

    formula_name = str(FORMULA_OCR_PROVIDER or "none").lower()
    if formula_name in {"none", "mock"}:
        formula_status_text = "Formula OCR disabled"
    elif formula_name == "mathpix":
        formula_status_text = "Formula OCR provider configured"
    else:
        formula_status_text = f"Formula OCR provider configured: {FORMULA_OCR_PROVIDER}"

    job_summary = None
    if user is not None:
        try:
            jobs = visible_jobs_query(user).all()
            job_summary = {
                "total": len(jobs),
                "queued": len([job for job in jobs if job.status == "queued"]),
                "running": len([job for job in jobs if job.status == "running"]),
                "failed": len([job for job in jobs if job.status == "failed"]),
                "completed": len([job for job in jobs if job.status == "completed"]),
            }
        except Exception:
            job_summary = {"total": 0, "queued": 0, "running": 0, "failed": 0, "completed": 0}

    return {
        "ai_provider": {
            **meta,
            "status_text": ai_status_text,
            "allow_mock_ai": bool(ALLOW_MOCK_AI),
        },
        "ocr": {
            "provider": OCR_PROVIDER,
            "langs": OCR_LANGS,
            "min_confidence": OCR_MIN_CONFIDENCE,
            "status_text": ocr_status_text,
        },
        "formula_ocr": {
            "provider": FORMULA_OCR_PROVIDER,
            "min_confidence": FORMULA_OCR_MIN_CONFIDENCE,
            "status_text": formula_status_text,
        },
        "job_summary": job_summary,
    }


def ensure_model_registry_seed(owner_user_id=0):
    meta = current_provider_metadata()
    registry = ModelPromptRegistry.query.filter_by(
        provider=meta["provider"],
        model_name=meta["model_name"],
        prompt_version=meta["prompt_version"],
        retrieval_version=meta["retrieval_version"]
    ).first()
    if registry is None:
        registry = ModelPromptRegistry(
            provider=meta["provider"],
            model_name=meta["model_name"],
            model_version=os.environ.get("MODEL_VERSION", "local-mvp-v1"),
            prompt_version=meta["prompt_version"],
            retrieval_version=meta["retrieval_version"],
            enabled=True,
            allowed_workflows=json.dumps(["term_extraction", "bilingual_alignment", "student_answer"], ensure_ascii=False),
            known_risks="Mock/local heuristic providers are demonstration-only and cannot auto approve terminology cards.",
            rollback_target="local_heuristic",
            owner=str(owner_user_id or "system"),
            created_at=current_time_text(),
            updated_at=current_time_text()
        )
        db.session.add(registry)
        db.session.flush()
    return registry


def ai_selection_from_config(config=None, provider_name=None, model_name=None):
    if config is None:
        if provider_name:
            config = AIProviderConfig.query.filter_by(provider_name=provider_name, is_enabled=True).order_by(AIProviderConfig.id.desc()).first()
        if config is None:
            config = AIProviderConfig.query.filter_by(is_default=True, is_enabled=True).order_by(AIProviderConfig.id.desc()).first()
        if config is None:
            config = ensure_ai_registry_seed().get("provider_config")
    selection = env_provider_selection(os.environ)
    if config is not None:
        selection.provider_name = config.provider_name
        selection.provider_mode = config.provider_mode
        selection.base_url = config.base_url or selection.base_url
        selection.model_name = model_name or config.default_model or selection.model_name
        selection.timeout_seconds = config.timeout_seconds or selection.timeout_seconds
        selection.max_retries = config.max_retries or selection.max_retries
        selection.cost_per_1k_input_tokens = float(config.cost_per_1k_input_tokens or selection.cost_per_1k_input_tokens or 0)
        selection.cost_per_1k_output_tokens = float(config.cost_per_1k_output_tokens or selection.cost_per_1k_output_tokens or 0)
    if provider_name:
        selection.provider_name = provider_name
    if model_name:
        selection.model_name = model_name
    if selection.provider_name == "deepseek":
        selection.api_key = DEEPSEEK_API_KEY
        selection.base_url = DEEPSEEK_BASE_URL
        selection.model_name = model_name or selection.model_name or DEEPSEEK_MODEL
    elif selection.provider_name in {"openai", "custom_openai_compatible"}:
        selection.api_key = OPENAI_API_KEY
        selection.base_url = OPENAI_BASE_URL
        selection.model_name = model_name or selection.model_name or OPENAI_MODEL
    return selection


def legacy_provider_credential_present(provider_name):
    provider = str(provider_name or "").strip().lower()
    if provider == "deepseek":
        value = DEEPSEEK_API_KEY
    elif provider in {"openai", "custom_openai_compatible"}:
        value = OPENAI_API_KEY
    else:
        value = ""
    return bool(value and not is_placeholder_secret(value))


def legacy_provider_config_credential_present(config):
    return legacy_provider_credential_present(getattr(config, "provider_name", ""))


def legacy_provider_local_readiness_service(*, request, provider):
    return evaluate_legacy_provider_local_readiness(request=request, provider=provider)


def ensure_ai_registry_seed(owner_user_id=0):
    result = ensure_legacy_provider_registry_seed(
        db=db,
        models=LegacyProviderRegistrySeedModels(
            AIProviderConfig=AIProviderConfig,
            AIModelRegistry=AIModelRegistry,
            PromptTemplate=PromptTemplate,
        ),
        selection=env_provider_selection(os.environ),
        default_prompts=DEFAULT_PROMPTS,
        current_time_text=current_time_text,
        model_version=os.environ.get("MODEL_VERSION", "local-mvp-v1"),
        owner_user_id=owner_user_id,
    )
    return {
        "provider_config": result.provider_config,
        "model": result.model,
        "prompts": list(result.prompts),
    }


def legacy_prompt_mutation_seed_registry(owner_user_id=0):
    return ensure_legacy_provider_registry_seed(
        db=db,
        models=LegacyProviderRegistrySeedModels(
            AIProviderConfig=AIProviderConfig,
            AIModelRegistry=AIModelRegistry,
            PromptTemplate=PromptTemplate,
        ),
        selection=env_provider_selection(os.environ),
        default_prompts=DEFAULT_PROMPTS,
        current_time_text=current_time_text,
        model_version=os.environ.get("MODEL_VERSION", "local-mvp-v1"),
        owner_user_id=owner_user_id,
    )


def legacy_prompt_mutation_dependencies():
    return LegacyPromptMutationDependencies(
        db=db,
        PromptTemplate=PromptTemplate,
        current_time_text=current_time_text,
        safe_json_loads=safe_json_loads,
        seed_registry=legacy_prompt_mutation_seed_registry,
    )


def get_prompt_template(prompt_key, prompt_version=None, task_type=None):
    query = PromptTemplate.query.filter_by(prompt_key=prompt_key, is_active=True)
    if prompt_version:
        query = query.filter_by(prompt_version=prompt_version)
    if task_type:
        query = query.filter_by(task_type=task_type)
    prompt = query.order_by(PromptTemplate.is_default.desc(), PromptTemplate.id.desc()).first()
    if prompt is not None:
        return prompt
    default = default_prompt_lookup(prompt_key, prompt_version=prompt_version, task_type=task_type)
    if default:
        return PromptTemplate(
            prompt_key=default["prompt_key"],
            prompt_version=default["prompt_version"],
            task_type=default["task_type"],
            language=default["language"],
            template_text=default["template_text"],
            json_schema=json.dumps(default["json_schema"], ensure_ascii=False),
            is_active=True,
            is_default=True,
            notes=default.get("notes", ""),
        )
    return None


def can_use_model_for_auto_approval(provider_name, model_name, prompt_version):
    provider = AIProviderConfig.query.filter_by(provider_name=provider_name, is_enabled=True).order_by(AIProviderConfig.id.desc()).first()
    if provider is None:
        return False, ["provider is not registered"]
    if provider.provider_mode != "live":
        return False, ["provider mode is not live"]
    model = AIModelRegistry.query.filter_by(provider_name=provider_name, model_name=model_name, is_enabled=True).order_by(AIModelRegistry.id.desc()).first()
    if model is None:
        return False, ["model is not registered or not enabled"]
    prompt = get_prompt_template("term_alignment", prompt_version=prompt_version, task_type="term_alignment")
    if prompt is None:
        return False, ["prompt version is not active"]
    if not model.last_evaluation_run_id:
        return False, ["model has no evaluation run"]
    run = db.session.get(EvaluationRun, model.last_evaluation_run_id)
    if run is None:
        return False, ["last evaluation run not found"]
    if float(getattr(run, "no_evidence_forced_alignment_rate", 1.0) or 0) != 0:
        return False, ["last evaluation has no-evidence forced alignment"]
    if float(getattr(run, "alignment_accuracy", 0.0) or 0) < float(os.environ.get("AI_AUTO_APPROVAL_MIN_ALIGNMENT_ACCURACY", "0.70")):
        return False, ["alignment accuracy below threshold"]
    if float(getattr(run, "auto_approval_error_rate", 1.0) or 0) > float(os.environ.get("AI_AUTO_APPROVAL_MAX_ERROR_RATE", "0.05")):
        return False, ["auto approval error rate above threshold"]
    return True, []


def call_ai_task(
    task_type,
    input_payload,
    prompt_key,
    prompt_version=None,
    provider_name=None,
    model_name=None,
    user_id=None,
    course_id=None,
    document_id=None,
    job_id=None,
    alignment_run_id=None,
    evaluation_run_id=None,
):
    ensure_ai_registry_seed(owner_user_id=user_id or 0)
    prompt = get_prompt_template(prompt_key, prompt_version=prompt_version, task_type=task_type)
    if prompt is None:
        return {"status": "error", "error_code": "VALIDATION_ERROR", "message": f"Prompt not found: {prompt_key}"}

    selection = ai_selection_from_config(provider_name=provider_name, model_name=model_name)
    prompt_text = prompt.template_text
    input_tokens = estimate_ai_tokens({"prompt": prompt_text, "input": input_payload})
    estimated_cost = estimate_ai_cost(
        input_tokens,
        0,
        selection.cost_per_1k_input_tokens,
        selection.cost_per_1k_output_tokens,
    )

    if user_id:
        allowed, quota_code, quota_message = check_ai_quota(
            UsageRecord.query.filter_by(user_id=user_id).all(),
            daily_call_limit=AI_DAILY_CALL_LIMIT_PER_USER,
            monthly_call_limit=AI_MONTHLY_CALL_LIMIT_PER_USER,
            daily_cost_limit=AI_DAILY_COST_LIMIT_PER_USER,
            new_cost=estimated_cost,
        )
        if not allowed:
            log = AICallLog(
                task_type=task_type,
                provider_name=selection.provider_name,
                provider_mode=selection.provider_mode,
                model_name=selection.model_name,
                prompt_key=prompt.prompt_key,
                prompt_version=prompt.prompt_version,
                user_id=user_id,
                course_id=course_id,
                document_id=document_id,
                job_id=job_id,
                alignment_run_id=alignment_run_id,
                evaluation_run_id=evaluation_run_id,
                request_hash=hash_ai_payload(input_payload or {}),
                response_hash="",
                input_token_count=input_tokens,
                output_token_count=0,
                estimated_cost=0,
                latency_ms=0,
                status="error",
                error_code=quota_code,
                error_message=quota_message,
                redacted_prompt_preview=preview_ai_payload(prompt_text),
                redacted_response_preview="",
                created_at=current_time_text(),
            )
            db.session.add(log)
            db.session.flush()
            return {
                "status": "error",
                "error_code": quota_code,
                "message": quota_message,
                "provider_name": selection.provider_name,
                "provider_mode": selection.provider_mode,
                "model_name": selection.model_name,
                "ai_call_log_id": log.id,
            }

    provider = provider_from_selection(selection)
    started = datetime.utcnow()
    response = provider.call(task_type, prompt_text, input_payload or {}, json_schema=safe_json_loads(prompt.json_schema, {}))
    latency_ms = int((datetime.utcnow() - started).total_seconds() * 1000)
    result = response.get("result") if isinstance(response, dict) else None
    status = response.get("status", "error") if isinstance(response, dict) else "error"
    error_code = response.get("error_code", "") if isinstance(response, dict) else "AI_PROVIDER_FAILED"
    error_message = response.get("message", "") if isinstance(response, dict) else "AI provider returned invalid response."
    if status == "success":
        ok, reason = validate_ai_json(task_type, result)
        if not ok:
            status = "error"
            error_code = "AI_INVALID_RESPONSE"
            error_message = reason

    output_tokens = estimate_ai_tokens(result or response)
    estimated_cost = estimate_ai_cost(
        input_tokens,
        output_tokens,
        selection.cost_per_1k_input_tokens,
        selection.cost_per_1k_output_tokens,
    )
    log = AICallLog(
        task_type=task_type,
        provider_name=selection.provider_name,
        provider_mode=selection.provider_mode,
        model_name=selection.model_name,
        prompt_key=prompt.prompt_key,
        prompt_version=prompt.prompt_version,
        user_id=user_id,
        course_id=course_id,
        document_id=document_id,
        job_id=job_id,
        alignment_run_id=alignment_run_id,
        evaluation_run_id=evaluation_run_id,
        request_hash=hash_ai_payload(input_payload or {}),
        response_hash=hash_ai_payload(result or response),
        input_token_count=input_tokens,
        output_token_count=output_tokens,
        estimated_cost=estimated_cost,
        latency_ms=latency_ms,
        status=status,
        error_code=error_code,
        error_message=redact_for_log(error_message),
        redacted_prompt_preview=preview_ai_payload(prompt_text),
        redacted_response_preview=preview_ai_payload(result or response),
        created_at=current_time_text(),
    )
    db.session.add(log)
    if user_id:
        db.session.add(UsageRecord(
            user_id=user_id,
            action_type=AI_EVENT_BY_TASK.get(task_type, "ai_alignment_call"),
            units_used=1,
            related_document_id=document_id,
            related_term_id=None,
            created_at=current_time_text(),
        ))
    db.session.flush()

    if status != "success":
        add_system_log("warning", "ai_provider", f"{selection.provider_name} {task_type} failed: {error_code} {error_message}")
        return {
            "status": "error",
            "error_code": error_code or "AI_PROVIDER_FAILED",
            "provider_name": selection.provider_name,
            "provider_mode": selection.provider_mode,
            "model_name": selection.model_name,
            "message": error_message or "AI provider request failed.",
            "fallback_used": False,
            "ai_call_log_id": log.id,
        }

    return {
        "status": "success",
        "provider_name": selection.provider_name,
        "provider_mode": selection.provider_mode,
        "model_name": selection.model_name,
        "prompt_key": prompt.prompt_key,
        "prompt_version": prompt.prompt_version,
        "result": result,
        "usage": {
            "input_tokens": input_tokens,
            "output_tokens": output_tokens,
            "estimated_cost": estimated_cost,
        },
        "latency_ms": latency_ms,
        "ai_call_log_id": log.id,
    }


def get_course_by_id_or_name(course_id=None, course_name=""):
    course = None
    if course_id:
        try:
            course = db.session.get(Course, int(course_id))
        except (TypeError, ValueError):
            course = None
    if course is None and course_name:
        course = Course.query.filter_by(name=str(course_name).strip()).first()
    return course


def is_course_member(user, course_id):
    if not user or not course_id:
        return False
    if user.role == "admin":
        return True
    return CourseMember.query.filter_by(user_id=user.id, course_id=course_id).first() is not None


def can_manage_course(user, course):
    if not user or course is None:
        return False
    if user.role == "admin":
        return True
    if user.role != "teacher":
        return False
    if getattr(course, "teacher_id", 0) == user.id:
        return True
    member = CourseMember.query.filter_by(user_id=user.id, course_id=course.id).first()
    return member is not None and getattr(member, "role_in_course", "") in {"teacher", "owner", "admin"}


def can_manage_course_name(user, course_name):
    if not user:
        return False
    if user.role == "admin":
        return True
    course = Course.query.filter_by(name=str(course_name or "").strip()).first()
    if course is None:
        return user.role == "teacher"
    return can_manage_course(user, course)


def can_view_course_name(user, course_name):
    if not user:
        return False
    if user.role in {"admin", "teacher"} and can_manage_course_name(user, course_name):
        return True
    course = Course.query.filter_by(name=str(course_name or "").strip()).first()
    if course is None:
        return False
    return is_course_member(user, course.id)


def get_free_plan():
    plan = SubscriptionPlan.query.filter_by(name="Free").first()
    if plan is None:
        plan = SubscriptionPlan(
            name="Free",
            price_monthly=0,
            monthly_pages=5,
            monthly_ai_calls=20,
            export_enabled=False,
            description="Free local demo plan: 5 pages and 20 AI/search units per month.",
            is_active=True
        )
        db.session.add(plan)
        db.session.flush()
    return plan


def get_active_subscription(user_id):
    now = datetime.now()
    subscriptions = (
        UserSubscription.query
        .filter_by(user_id=user_id, status="active")
        .order_by(UserSubscription.id.desc())
        .all()
    )
    for subscription in subscriptions:
        end_date = parse_time_text(subscription.end_date)
        if end_date is None or end_date >= now:
            return subscription
    return None


def get_effective_plan(user_id):
    subscription = get_active_subscription(user_id)
    if subscription:
        plan = db.session.get(SubscriptionPlan, subscription.plan_id)
        if plan and plan.is_active:
            return plan, subscription
    return get_free_plan(), None


def month_start_text():
    now = datetime.now()
    return datetime(now.year, now.month, 1).strftime("%Y-%m-%d %H:%M:%S")


def get_usage_totals(user_id):
    start = month_start_text()
    records = UsageRecord.query.filter(
        UsageRecord.user_id == user_id,
        UsageRecord.created_at >= start
    ).all()
    page_actions = {"document_parse_page", "ocr_page", "personal_upload_page"}
    ai_actions = {"ai_alignment", "ai_term_extract", "term_search", "knowledge_search"}
    pages_used = sum(record.units_used for record in records if record.action_type in page_actions)
    ai_used = sum(record.units_used for record in records if record.action_type in ai_actions)
    return {
        "period_start": start,
        "pages_used": pages_used,
        "ai_calls_used": ai_used,
        "records": records
    }


def check_usage_quota(user, page_units=0, ai_units=0):
    plan, subscription = get_effective_plan(user.id)
    totals = get_usage_totals(user.id)
    allowed = True
    reasons = []
    if page_units and totals["pages_used"] + page_units > plan.monthly_pages:
        allowed = False
        reasons.append(f"文档解析页数额度不足：{totals['pages_used']}/{plan.monthly_pages}")
    if ai_units and totals["ai_calls_used"] + ai_units > plan.monthly_ai_calls:
        allowed = False
        reasons.append(f"AI/search 额度不足：{totals['ai_calls_used']}/{plan.monthly_ai_calls}")
    return allowed, reasons, plan, subscription, totals


def record_usage(user_id, action_type, units_used=1, related_document_id=None, related_term_id=None):
    record = UsageRecord(
        user_id=user_id,
        action_type=action_type,
        units_used=max(1, int(units_used or 1)),
        related_document_id=related_document_id,
        related_term_id=related_term_id,
        created_at=current_time_text()
    )
    db.session.add(record)
    db.session.flush()
    return record


def normalize_english_term(term):
    value = normalize_math_symbols(str(term or "").strip().lower())
    value = re.sub(r"[_\s]+", " ", value)
    value = re.sub(r"\s*-\s*", "-", value)
    return re.sub(r"\s+", " ", value).strip()


def safe_json_loads(value, fallback):
    try:
        return json.loads(value or "")
    except Exception:
        return fallback


def redact_for_log(value):
    text = str(value or "")
    text = re.sub(r"sk-[A-Za-z0-9_-]{8,}", "sk-***REDACTED***", text)
    text = re.sub(r"Bearer\s+[A-Za-z0-9._~+/=-]{12,}", "Bearer ***REDACTED***", text, flags=re.I)
    text = re.sub(r"token=([A-Za-z0-9._~+/=-]{12,})", "token=***REDACTED***", text, flags=re.I)
    return text[:1200]


def token_hash(value):
    return hashlib.sha256((TOKEN_HASH_SECRET + "::" + str(value or "")).encode("utf-8")).hexdigest()


def file_sha256_for_path(path):
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def detect_magic_type(path):
    with open(path, "rb") as handle:
        header = handle.read(12)
    if header.startswith(b"%PDF"):
        return "pdf"
    if header.startswith(b"PK\x03\x04") or header.startswith(b"PK\x05\x06") or header.startswith(b"PK\x07\x08"):
        return "zip_office"
    if header.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if header.startswith(b"\xff\xd8\xff"):
        return "jpg"
    try:
        with open(path, "rb") as handle:
            handle.read(4096).decode("utf-8")
        return "text"
    except UnicodeDecodeError:
        return "binary"


def validate_upload_magic(filename, path):
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    magic_type = detect_magic_type(path)
    expected = {
        "pdf": {"pdf"},
        "docx": {"zip_office"},
        "pptx": {"zip_office"},
        "png": {"png"},
        "jpg": {"jpg"},
        "jpeg": {"jpg"},
        "txt": {"text"},
        "md": {"text"},
        "markdown": {"text"},
    }
    if ext not in expected:
        return False, f"文件扩展名不在允许列表中：{ext}"
    if magic_type not in expected[ext]:
        return False, f"文件内容类型与扩展名不匹配：extension={ext}, magic={magic_type}"
    return True, ""


def evidence_score_from_items(items):
    scores = []
    for item in items or []:
        if not isinstance(item, dict):
            continue
        for key in ("similarity_score", "score", "evidence_score"):
            if key in item:
                try:
                    score = float(item.get(key) or 0)
                    scores.append(score if score <= 1 else score / 100.0)
                    break
                except (TypeError, ValueError):
                    pass
    if not scores:
        return 0.0
    return round(max(0.0, min(max(scores), 1.0)), 4)


def evidence_snapshot(items):
    snapshot = []
    for item in items or []:
        if not isinstance(item, dict):
            continue
        snapshot.append({
            "chunk_id": item.get("chunk_id"),
            "source": item.get("source", ""),
            "source_citation": item.get("source_citation", ""),
            "language": item.get("language", ""),
            "page_number": item.get("page_number"),
            "chapter": item.get("chapter", ""),
            "similarity_score": item.get("similarity_score", item.get("score", 0)),
            "retrieval_reason": item.get("retrieval_reason", ""),
            "content": item.get("content", item.get("definition", ""))[:900],
        })
    return snapshot


def quality_flags_for_alignment(alignment, min_ocr_confidence=100):
    flags = []
    if not str(alignment.get("english_kb_evidence", "")).strip():
        flags.append("no_en_evidence")
    if not str(alignment.get("chinese_kb_evidence", "")).strip():
        flags.append("no_zh_evidence")
    evidence_items = list(alignment.get("english_evidence_items", []) or []) + list(alignment.get("chinese_evidence_items", []) or [])
    for item in evidence_items:
        if not isinstance(item, dict):
            continue
        try:
            evidence_score = float(item.get("evidence_score", item.get("score", 0)) or 0)
        except (TypeError, ValueError):
            evidence_score = 0
        if item.get("evidence_strength") == "weak" or (evidence_score and evidence_score < 0.80):
            flags.append("weak_evidence")
        if "domain_mismatch" in set(item.get("risk_flags", []) or []):
            flags.append("domain_mismatch")
    if min_ocr_confidence < 60:
        flags.append("ocr_low_confidence")
    ai_model = str(alignment.get("ai_model", "")).strip().lower()
    provider_status = str(alignment.get("provider_status", "")).strip().lower()
    if ai_model in {"mock", "local_mock_fallback", "local_heuristic"} or provider_status in {"mock", "local_heuristic", "provider_failed", "provider_unavailable"}:
        flags.append("mock_or_local_ai")
    review_status = str(alignment.get("review_status", "")).strip()
    if review_status in {"conflict_detected", "multi_translation_conflict"}:
        flags.append("multi_translation_conflict")
    if review_status in {"ambiguous_candidate"}:
        flags.append("ambiguous_candidate")
    if "domain mismatch" in str(alignment.get("risk_note", "")).lower():
        flags.append("domain_mismatch")
    if not str(alignment.get("english_term", "")).strip() or is_probably_noise(str(alignment.get("english_term", ""))):
        flags.append("invalid_term_candidate")
    return sorted(set(flags))


def alignment_status_from_flags(flags, confidence=0):
    flags = set(flags or [])
    if "no_en_evidence" in flags:
        return "no_en_evidence"
    if "no_zh_evidence" in flags:
        return "no_zh_evidence"
    if "domain_mismatch" in flags:
        return "domain_mismatch"
    if "multi_translation_conflict" in flags:
        return "multi_translation_conflict"
    if "ambiguous_candidate" in flags:
        return "ambiguous_candidate"
    if "ocr_low_confidence" in flags:
        return "ocr_low_confidence"
    if "weak_evidence" in flags:
        return "unverified_translation"
    if confidence >= 85:
        return "exact_match"
    if confidence >= 70:
        return "accepted_translation"
    return "unverified_translation"


def score_breakdown_for_alignment(alignment, min_ocr_confidence=100):
    english_score = float(alignment.get("english_evidence_score", 0) or 0)
    chinese_score = float(alignment.get("chinese_evidence_score", 0) or 0)
    ai_score = max(0.0, min(float(alignment.get("confidence_score", 0) or 0) / 100.0, 1.0))
    term_quality = max(0.0, min(float(alignment.get("term_quality_score", 0.75) or 0.75), 1.0))
    course_scope_score = 1.0 if alignment.get("course") else 0.6
    source_quality_score = 0.8
    flags = quality_flags_for_alignment(alignment, min_ocr_confidence=min_ocr_confidence)
    penalties = {
        "no_zh_evidence": 0.40,
        "no_en_evidence": 0.40,
        "domain_mismatch": 0.50,
        "ocr_low_confidence": 0.25,
        "mock_or_local_ai": 0.30,
        "ambiguous_candidate": 0.20,
        "multi_translation_conflict": 0.30,
        "invalid_term_candidate": 0.60,
    }
    risk_penalty = sum(penalties.get(flag, 0) for flag in flags)
    weighted = (
        0.25 * term_quality
        + 0.25 * english_score
        + 0.25 * chinese_score
        + 0.15 * ai_score
        + 0.05 * course_scope_score
        + 0.05 * source_quality_score
        - risk_penalty
    )
    confidence = int(round(max(0.0, min(weighted, 1.0)) * 100))
    if "no_en_evidence" in flags or "no_zh_evidence" in flags:
        confidence = min(confidence, 45)
    if "weak_evidence" in flags:
        confidence = min(confidence, 79)
    return {
        "term_quality_score": round(term_quality, 4),
        "english_evidence_score": round(english_score, 4),
        "chinese_evidence_score": round(chinese_score, 4),
        "ai_alignment_score": round(ai_score, 4),
        "course_scope_score": round(course_scope_score, 4),
        "source_quality_score": round(source_quality_score, 4),
        "risk_penalty": round(risk_penalty, 4),
        "weighted_confidence": confidence,
        "quality_flags": flags,
    }


def card_status_from_alignment(alignment, min_ocr_confidence=100):
    confidence = int(alignment.get("confidence_score", 0) or 0)
    has_english = bool(str(alignment.get("english_kb_evidence", "")).strip())
    has_chinese = bool(str(alignment.get("chinese_kb_evidence", "")).strip())
    ai_model = str(alignment.get("ai_model", "")).strip().lower()
    provider_status = str(alignment.get("provider_status", "")).strip().lower()
    flags = set(alignment.get("quality_flags") or quality_flags_for_alignment(alignment, min_ocr_confidence=min_ocr_confidence))
    input_risks = set(alignment.get("input_risk_labels") or [])
    english_score = float(alignment.get("english_evidence_score", 0) or 0)
    chinese_score = float(alignment.get("chinese_evidence_score", 0) or 0)

    if flags & parse_quality_risk_service.FORCE_REVIEW_RISK_LABELS or input_risks & parse_quality_risk_service.FORCE_REVIEW_RISK_LABELS:
        return "pending_quality_control"
    if "invalid_term_candidate" in flags:
        return "rejected"
    if min_ocr_confidence < 60 or "ocr_low_confidence" in flags:
        return "pending_quality_control"
    if not has_english or not has_chinese or "no_en_evidence" in flags or "no_zh_evidence" in flags:
        return "needs_more_evidence"
    if "domain_mismatch" in flags:
        return "pending_quality_control"
    if "weak_evidence" in flags:
        return "pending_quality_control"
    if "multi_translation_conflict" in flags:
        return "conflict_detected"
    if ai_model in {"mock", "local_mock_fallback", "local_heuristic"} or provider_status in {"mock", "local_heuristic", "provider_failed", "provider_unavailable"}:
        return "pending_quality_control"
    if str(alignment.get("review_status", "")).strip() in {"conflict_detected", "rejected"}:
        return alignment.get("review_status")
    alignment_status = str(alignment.get("alignment_status", "")).strip()
    if (
        confidence >= 85
        and english_score >= 0.80
        and chinese_score >= 0.80
        and alignment_status in {"exact_match", "accepted_translation", ""}
        and not flags
    ):
        return "auto_approved"
    return "pending_quality_control"


def parse_quality_metadata_from_chunks(chunks):
    parse_uid = ""
    parse_block_uid = ""
    flags = set()
    for chunk in chunks or []:
        parse_uid = parse_uid or getattr(chunk, "parse_uid", "")
        parse_block_uid = parse_block_uid or getattr(chunk, "parse_block_uid", "")
        for flag in safe_json_loads(getattr(chunk, "quality_flags_json", "[]"), []):
            if flag:
                flags.add(flag)
    status_priority = [
        "mixed_quality",
        "partial_text",
        "ocr_low_confidence",
        "formula_ocr_unavailable",
        "formula_ocr_required",
        "formula_detected",
        "native_text_ok",
    ]
    quality_status = next((status for status in status_priority if status in flags), "")
    return parse_quality_risk_service.build_parse_quality_metadata({
        "parse_uid": parse_uid,
        "parse_block_uid": parse_block_uid,
        "parse_quality_status": quality_status,
        "parse_quality_flags": sorted(flags),
    })


def apply_parse_quality_to_alignment(alignment, parse_quality_metadata):
    if not parse_quality_metadata:
        return alignment
    if parse_quality_risk_service.should_block_downstream_creation(parse_quality_metadata):
        raise ValueError(f"blocked parse quality status: {parse_quality_metadata.get('parse_quality_status')}")
    risk_labels = parse_quality_risk_service.parse_quality_to_risk_labels(parse_quality_metadata)
    if not risk_labels:
        return alignment
    alignment["input_risk_labels"] = parse_quality_risk_service.merge_risk_labels(
        alignment.get("input_risk_labels", []),
        risk_labels,
    )
    alignment["quality_flags"] = parse_quality_risk_service.merge_risk_labels(
        alignment.get("quality_flags", []),
        [*risk_labels, *parse_quality_metadata.get("parse_quality_flags", [])],
    )
    alignment["parse_uid"] = parse_quality_metadata.get("parse_uid", "")
    alignment["parse_block_uid"] = parse_quality_metadata.get("parse_block_uid", "")
    alignment["parse_quality_status"] = parse_quality_metadata.get("parse_quality_status", "")
    alignment["parse_quality_flags"] = parse_quality_metadata.get("parse_quality_flags", [])
    if parse_quality_risk_service.should_force_needs_review(parse_quality_metadata):
        if int(alignment.get("confidence_score", 0) or 0) > 79:
            alignment["confidence_score"] = 79
        if alignment.get("review_status") in {"approved", "auto_approved"}:
            alignment["review_status"] = "pending_quality_control"
    risk_note = str(alignment.get("risk_note", "") or "").strip()
    quality_note = f"parse_quality: {parse_quality_metadata.get('parse_quality_status')}"
    alignment["risk_note"] = "; ".join(part for part in [risk_note, quality_note] if part)
    return alignment


def create_or_update_card_from_alignment(
    english_term,
    alignment,
    scope_type,
    course_id=None,
    owner_user_id=None,
    source_document_id=None,
    courseware_sentence="",
    min_ocr_confidence=100
):
    english_term = str(english_term or "").strip()
    if not english_term:
        return None
    normalized_term = normalize_english_term(english_term)

    query = TerminologyCard.query.filter_by(
        normalized_english_term=normalized_term,
        scope_type=scope_type
    )
    if scope_type == "course":
        query = query.filter_by(course_id=course_id)
    elif scope_type == "personal":
        query = query.filter_by(owner_user_id=owner_user_id)
        if source_document_id:
            query = query.filter_by(source_document_id=source_document_id)

    card = query.first()
    if card is None:
        legacy_query = TerminologyCard.query.filter_by(english_term=english_term, scope_type=scope_type)
        if scope_type == "course":
            legacy_query = legacy_query.filter_by(course_id=course_id)
        elif scope_type == "personal":
            legacy_query = legacy_query.filter_by(owner_user_id=owner_user_id)
        card = legacy_query.first()
    now = current_time_text()
    requested_status = alignment.get("review_status") or card_status_from_alignment(alignment, min_ocr_confidence=min_ocr_confidence)
    input_risks = parse_quality_risk_service.merge_risk_labels(alignment.get("input_risk_labels", []), [])
    if set(input_risks) & parse_quality_risk_service.FORCE_REVIEW_RISK_LABELS and requested_status in {"auto_approved", "approved"}:
        requested_status = "pending_quality_control"
    old_status = card.status if card else "draft"
    status = requested_status
    if not validate_card_status_transition(old_status, requested_status, "system", system_action=True):
        if old_status == "rejected":
            status = old_status
            alignment["risk_note"] = (
                alignment.get("risk_note", "")
                or "This card was previously rejected and cannot be system auto-approved."
            )
        elif requested_status == "auto_approved":
            status = "pending_quality_control"
            reasons = alignment.get("auto_approve_reasons") or []
            if "system transition from previous status is not allowed" not in reasons:
                reasons.append("system transition from previous status is not allowed")
            alignment["auto_approve_reasons"] = reasons
        else:
            status = old_status if old_status in {"approved", "rejected", "archived"} else "pending_quality_control"

    if card is None:
        card = TerminologyCard(
            english_term=english_term,
            scope_type=scope_type,
            course_id=course_id,
            owner_user_id=owner_user_id,
            source_document_id=source_document_id,
            created_at=now
        )
        db.session.add(card)

    card.normalized_english_term = normalized_term
    card.final_chinese_term = alignment.get("final_chinese_term") or alignment.get("chinese_term") or english_term
    card.normalized_chinese_term = normalize_english_term(card.final_chinese_term)
    card.ai_translation_candidate = alignment.get("ai_translation_candidate", "")
    card.courseware_sentence = courseware_sentence or alignment.get("courseware_sentence", "")
    card.english_kb_evidence = alignment.get("english_kb_evidence", "")
    card.chinese_kb_evidence = alignment.get("chinese_kb_evidence", "")
    card.english_evidence_snapshot = json.dumps(alignment.get("english_evidence_snapshot", evidence_snapshot(alignment.get("english_evidence_items", []))), ensure_ascii=False)
    card.chinese_evidence_snapshot = json.dumps(alignment.get("chinese_evidence_snapshot", evidence_snapshot(alignment.get("chinese_evidence_items", []))), ensure_ascii=False)
    card.english_evidence_score = float(alignment.get("english_evidence_score", 0) or 0)
    card.chinese_evidence_score = float(alignment.get("chinese_evidence_score", 0) or 0)
    card.concept_explanation = alignment.get("explanation", "")
    card.alignment_reason = alignment.get("alignment_reason", "")
    card.alignment_status = alignment.get("alignment_status", "unverified_translation")
    card.score_breakdown_json = json.dumps(alignment.get("score_breakdown", {}), ensure_ascii=False)
    card.quality_flags_json = json.dumps(parse_quality_risk_service.merge_risk_labels(alignment.get("quality_flags", []), input_risks), ensure_ascii=False)
    confidence_score = int(alignment.get("confidence_score", 0) or 0)
    if set(input_risks) & parse_quality_risk_service.FORCE_REVIEW_RISK_LABELS:
        confidence_score = min(confidence_score, 79)
    card.confidence_score = confidence_score
    card.status = status
    card.ai_provider = alignment.get("ai_provider", alignment.get("provider_status", ""))
    card.ai_provider_mode = alignment.get("ai_provider_mode", alignment.get("provider_status", ""))
    card.ai_model = alignment.get("ai_model", "")
    card.prompt_key = alignment.get("prompt_key", "term_alignment")
    card.prompt_version = alignment.get("prompt_version", ALIGNMENT_PROMPT_VERSION)
    card.retrieval_version = alignment.get("retrieval_version", RETRIEVAL_VERSION)
    card.knowledge_base_version_id = alignment.get("knowledge_base_version_id")
    card.english_kb_version_id = alignment.get("english_kb_version_id")
    card.chinese_kb_version_id = alignment.get("chinese_kb_version_id")
    card.retrieval_run_id = alignment.get("retrieval_run_id")
    card.index_version = alignment.get("index_version", "")
    card.evidence_content_hashes_json = json.dumps(alignment.get("evidence_content_hashes", []), ensure_ascii=False)
    card.evidence_status = alignment.get("evidence_status", "evidence_current")
    card.ai_call_log_id = alignment.get("ai_call_log_id")
    card.alignment_run_id = alignment.get("alignment_run_id")
    card.source_alignment_run_id = alignment.get("source_alignment_run_id") or alignment.get("alignment_run_id")
    card.risk_note = alignment.get("risk_note", "")
    card.parse_uid = alignment.get("parse_uid", "")
    card.parse_block_uid = alignment.get("parse_block_uid", "")
    card.parse_quality_status = alignment.get("parse_quality_status", "")
    card.parse_quality_flags = json.dumps(alignment.get("parse_quality_flags", []), ensure_ascii=False)
    card.input_risk_labels = json.dumps(input_risks, ensure_ascii=False)
    if status == "rejected":
        card.rejected_reason = alignment.get("rejected_reason", card.rejected_reason or "Rejected by alignment status machine.")
    if alignment.get("reviewer_note"):
        card.reviewer_note = alignment.get("reviewer_note")
    en_items = alignment.get("english_evidence_items") or []
    zh_items = alignment.get("chinese_evidence_items") or []
    if en_items and isinstance(en_items[0], dict):
        card.english_evidence_chunk_id = en_items[0].get("chunk_id")
    if zh_items and isinstance(zh_items[0], dict):
        card.chinese_evidence_chunk_id = zh_items[0].get("chunk_id")
    if status in {"auto_approved", "approved"} and not card.approved_at:
        card.approved_at = now
    card.updated_at = now
    db.session.flush()
    if input_risks:
        record_parse_quality_risk_audit(
            "parse_quality_risk_propagated",
            "terminology_card",
            card.id,
            parse_uid=card.parse_uid,
            quality_status=card.parse_quality_status,
            risk_labels=input_risks,
            forced_status=card.status,
            commit=False,
        )
    return card


def sync_term_to_card(term, course_id=None, owner_user_id=None, source_document_id=None, min_ocr_confidence=100):
    if term is None:
        return None
    alignment = {
        "english_term": term.english_term,
        "final_chinese_term": term.final_chinese_term or term.chinese_term,
        "chinese_term": term.chinese_term,
        "ai_translation_candidate": term.ai_translation_candidate,
        "english_kb_evidence": term.english_kb_evidence,
        "chinese_kb_evidence": term.chinese_kb_evidence,
        "explanation": term.explanation,
        "alignment_reason": term.alignment_reason,
        "confidence_score": term.confidence,
        "review_status": term.review_status,
        "parse_uid": getattr(term, "parse_uid", ""),
        "parse_block_uid": getattr(term, "parse_block_uid", ""),
        "parse_quality_status": getattr(term, "parse_quality_status", ""),
        "parse_quality_flags": safe_json_loads(getattr(term, "parse_quality_flags", "[]"), []),
        "input_risk_labels": safe_json_loads(getattr(term, "input_risk_labels", "[]"), []),
        "quality_flags": safe_json_loads(getattr(term, "input_risk_labels", "[]"), []),
        "risk_note": getattr(term, "risk_note", ""),
    }
    return create_or_update_card_from_alignment(
        english_term=term.english_term,
        alignment=alignment,
        scope_type="course",
        course_id=course_id,
        owner_user_id=owner_user_id,
        source_document_id=source_document_id,
        courseware_sentence=term.courseware_sentence or term.context,
        min_ocr_confidence=min_ocr_confidence
    )


def knowledge_result_from_structured(item):
    return {
        "content": item.get("definition", ""),
        "source": item.get("source", "Seed Knowledge Base"),
        "language": item.get("language", ""),
        "page_number": None,
        "chapter": item.get("chapter", ""),
        "similarity_score": item.get("score", 0),
        "retrieval_reason": "Seed structured knowledge base keyword match",
        "source_citation": item.get("source", ""),
        "chunk_id": None
    }


def knowledge_result_from_chunk(score, chunk, reason="Keyword/simple similarity match"):
    return {
        "content": chunk.content[:900],
        "source": chunk.source_citation or chunk.title or chunk.course or "KnowledgeChunk",
        "language": getattr(chunk, "language", ""),
        "page_number": getattr(chunk, "page_number", None),
        "chapter": getattr(chunk, "chapter", "") or chunk.source_page,
        "similarity_score": round(score, 4),
        "retrieval_reason": reason,
        "source_citation": getattr(chunk, "source_citation", "") or chunk.title,
        "chunk_id": chunk.id
    }


def annotate_chunks_with_source_metadata(chunks):
    source_ids = sorted({
        getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None)
        for chunk in chunks
        if getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None)
    })
    sources = {}
    if source_ids:
        for source in KnowledgeSource.query.filter(KnowledgeSource.id.in_(source_ids)).all():
            sources[source.id] = source
    for chunk in chunks:
        source = sources.get(getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None))
        if source is None:
            continue
        setattr(chunk, "_source_type", source.source_type)
        setattr(chunk, "_license_status", source.license_status)
        setattr(chunk, "_allow_derivative_cards", bool(source.allow_derivative_cards))
        setattr(chunk, "_source_status", getattr(source, "status", "active"))
        setattr(chunk, "_authorization_status", getattr(source, "authorization_status", "unknown"))
        setattr(chunk, "_source_quality", source_quality_from_governance(source))
        setattr(chunk, "_source_governance_flags", source_status_flags(source))
    return chunks


def get_published_kb_version(course_id=None, scope_type="course", owner_user_id=None):
    query = KnowledgeBaseVersion.query.filter_by(status="published")
    if scope_type == "personal":
        query = query.filter_by(scope_type="personal", owner_user_id=owner_user_id)
    elif scope_type == "global":
        query = query.filter_by(scope_type="global")
    else:
        query = query.filter_by(scope_type="course", course_id=course_id)
    return query.order_by(KnowledgeBaseVersion.version_number.desc(), KnowledgeBaseVersion.id.desc()).first()


def create_knowledge_base_version(course_id, scope_type, owner_user_id, description, parent_version_id=None, created_by=0):
    scope_type = str(scope_type or "course").strip().lower()
    if scope_type not in {"course", "personal", "global"}:
        raise ValueError("scope_type must be course, personal, or global")
    if scope_type == "course" and not course_id:
        raise ValueError("course KB version requires course_id")
    if scope_type == "personal" and not owner_user_id:
        raise ValueError("personal KB version requires owner_user_id")
    query = KnowledgeBaseVersion.query.filter_by(scope_type=scope_type)
    if scope_type == "course":
        query = query.filter_by(course_id=course_id)
    elif scope_type == "personal":
        query = query.filter_by(owner_user_id=owner_user_id)
    version_number = next_version_number(query.all())
    version = KnowledgeBaseVersion(
        kb_scope=scope_type,
        scope_type=scope_type,
        course_id=course_id if scope_type == "course" else None,
        owner_user_id=owner_user_id if scope_type == "personal" else None,
        version_number=version_number,
        version_name=default_version_name(scope_type, version_number, course_id=course_id, owner_user_id=owner_user_id),
        status="draft",
        description=description or "",
        parent_version_id=parent_version_id,
        index_backend=INDEX_BACKEND,
        index_version=INDEX_VERSION,
        retrieval_version=KB_RETRIEVAL_VERSION,
        created_by=created_by or 0,
        created_at=current_time_text(),
        is_active=False,
    )
    db.session.add(version)
    db.session.flush()
    return version


def create_or_update_knowledge_source_for_version(document, kb_version, creator=None, discipline="", source_name=""):
    source = KnowledgeSource.query.filter_by(document_id=document.id, version_introduced_id=kb_version.id).first()
    scope_type = kb_version.scope_type or kb_version.kb_scope or document.scope_type
    kb_type = "student_personal_kb" if scope_type == "personal" else ("en_course_kb" if document.language == "en" else "zh_course_kb")
    parse_record = DocumentParseRecord.query.filter_by(parse_uid=document.parse_uid).first() if getattr(document, "parse_uid", "") else None
    quality_status = getattr(parse_record, "quality_status", "") if parse_record else ""
    quality_flags = safe_json_loads(getattr(parse_record, "quality_flags", "[]") if parse_record else getattr(document, "quality_flags_json", "[]"), [])
    course_obj = db.session.get(Course, document.course_id) if document.course_id else None
    source_role = "student_private_material" if scope_type == "personal" else ("english_course_material" if document.language == "en" else "chinese_reference_material")
    trust_level = "low_quality" if quality_status in (knowledge_governance_service.REVIEW_PARSE_STATUSES | knowledge_governance_service.BLOCKED_PARSE_STATUSES) else ("student_uploaded" if scope_type == "personal" else "teacher_verified")
    governance_status = "blocked" if quality_status in knowledge_governance_service.BLOCKED_PARSE_STATUSES else ("needs_review" if quality_status in knowledge_governance_service.REVIEW_PARSE_STATUSES else "active")
    if source is None:
        source = KnowledgeSource(
            title=source_name or document.filename,
            name=source_name or document.filename,
            source_title=source_name or document.filename,
            course=course_obj.name if course_obj else "",
            chapter="",
            course_id=document.course_id if scope_type == "course" else None,
            scope_type=scope_type,
            owner_user_id=document.owner_user_id if scope_type == "personal" else None,
            document_id=document.id,
            language=document.language,
            discipline=discipline,
            source_type=document.source_type or ("student_personal_upload" if scope_type == "personal" else "teacher_upload"),
            source_role=source_role,
            owner_type="student" if scope_type == "personal" else "teacher",
            owner_id=str(document.owner_user_id if scope_type == "personal" else (creator.id if creator else "")),
            visibility="private" if scope_type == "personal" else ("public" if scope_type == "global" else "course"),
            trust_level=trust_level,
            parse_uid=getattr(document, "parse_uid", ""),
            source_filename=getattr(document, "original_filename", "") or getattr(document, "filename", ""),
            file_type=getattr(document, "file_type", "") or "unknown",
            content_hash=getattr(document, "sha256", "") or getattr(document, "file_sha256", ""),
            version=1,
            license_note="",
            quality_status=quality_status,
            quality_flags=json.dumps(quality_flags, ensure_ascii=False),
            knowledge_base_type=kb_type,
            access_method="manual_upload",
            license_status="restricted" if scope_type == "personal" else "authorized",
            license_type="teacher_provided" if scope_type == "course" else "unknown",
            authorization_status="allowed_for_private_use" if scope_type == "personal" else "allowed_for_course_use",
            source_quality=0.6 if scope_type == "personal" else 0.9,
            allow_full_text_indexing=scope_type != "personal",
            allow_student_search=scope_type != "personal",
            allow_derivative_cards=scope_type != "personal",
            version_introduced_id=kb_version.id,
            status=governance_status,
            created_by=creator.id if creator else 0,
            created_at=current_time_text(),
            updated_at=current_time_text(),
        )
        db.session.add(source)
    else:
        source.quality_status = quality_status or getattr(source, "quality_status", "")
        source.quality_flags = json.dumps(quality_flags, ensure_ascii=False)
        source.trust_level = trust_level
        source.status = governance_status
        source.updated_at = current_time_text()
    db.session.flush()
    return source


def index_document_into_kb_version(document_id, kb_version_id):
    document = db.session.get(Document, int(document_id))
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id))
    if document is None or version is None:
        raise ValueError("Document or KB version not found")
    scope_type = version.scope_type or version.kb_scope
    if scope_type == "course" and document.scope_type == "personal":
        raise PermissionError("personal document cannot be indexed into course KB")
    if scope_type == "course" and document.course_id != version.course_id:
        raise PermissionError("document course does not match KB version")
    if scope_type == "personal" and document.owner_user_id != version.owner_user_id:
        raise PermissionError("document owner does not match personal KB version")
    source = create_or_update_knowledge_source_for_version(document, version)
    doc_chunks = DocumentChunk.query.filter_by(document_id=document.id).order_by(DocumentChunk.chunk_index.asc(), DocumentChunk.id.asc()).all()
    existing = KnowledgeChunk.query.filter_by(knowledge_base_version_id=version.id).all()
    created = []
    duplicate_count = 0
    skipped_count = 0
    for index, doc_chunk in enumerate(doc_chunks, start=1):
        fields = build_knowledge_chunk_fields(document, doc_chunk, source, version.id, index)
        course_obj = db.session.get(Course, document.course_id) if document.course_id else None
        fields["course"] = course_obj.name if course_obj else ""
        fields["title"] = document.filename
        fields["discipline"] = getattr(source, "discipline", "")
        duplicate = find_duplicate_chunk(fields["content_hash"], version.id, existing + created)
        record = KnowledgeChunk(**{key: value for key, value in fields.items() if hasattr(KnowledgeChunk, key)})
        record.created_at = current_time_text()
        record.updated_at = current_time_text()
        if duplicate is not None:
            mark_duplicate_chunk(record, duplicate.id)
            duplicate_count += 1
        elif getattr(record, "status", "") == "blocked":
            record.index_status = "blocked"
            record.is_active = False
            skipped_count += 1
        elif record.index_status == "skipped":
            skipped_count += 1
        else:
            record.index_status = "indexed"
            record.is_active = True
        db.session.add(record)
        db.session.flush()
        created.append(record)
    formula_count = FormulaBlock.query.filter_by(document_id=document.id).count()
    version.source_count = KnowledgeSource.query.filter_by(version_introduced_id=version.id).count()
    version.chunk_count = KnowledgeChunk.query.filter_by(knowledge_base_version_id=version.id).count()
    version.deduped_chunk_count = duplicate_count
    version.formula_block_count = formula_count
    version.status = "ready" if version.chunk_count > 0 else "failed"
    version.quality_gate_status = "pass" if version.status == "ready" else "fail"
    version.manifest_json = json.dumps(build_kb_version_manifest(version), ensure_ascii=False)
    db.session.flush()
    return {
        "document_id": document.id,
        "knowledge_base_version_id": version.id,
        "chunks_created": len(created),
        "duplicate_count": duplicate_count,
        "skipped_count": skipped_count,
        "source_id": source.id,
        "status": version.status,
    }


def publish_kb_version(kb_version_id, actor_user_id=0):
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id))
    if version is None:
        return {"status": "error", "message": "KB version not found"}
    ok, reasons = can_publish_version(version)
    health = run_knowledge_health_check(version.course_id, version.id, scope_type=version.scope_type, owner_user_id=version.owner_user_id)
    if health["status"] == "FAIL":
        ok = False
        reasons.extend(health["issues"])
    if not ok:
        version.quality_gate_status = "fail"
        db.session.flush()
        return {"status": "error", "message": "KB version cannot be published", "reasons": reasons, "health": health}
    query = KnowledgeBaseVersion.query.filter_by(scope_type=version.scope_type, status="published")
    if version.scope_type == "course":
        query = query.filter_by(course_id=version.course_id)
    elif version.scope_type == "personal":
        query = query.filter_by(owner_user_id=version.owner_user_id)
    for old in query.all():
        if old.id != version.id:
            old.status = "archived"
            old.is_active = False
            old.archived_at = current_time_text()
    version.status = "published"
    version.is_active = True
    version.published_at = current_time_text()
    version.quality_gate_status = "pass"
    version.manifest_json = json.dumps(build_kb_version_manifest(version), ensure_ascii=False)
    add_system_log("info", "knowledge_versioning", f"Published KB version {version.id} by user {actor_user_id}")
    db.session.flush()
    return {"status": "success", "version": serialize_kb_version(version), "health": health}


def rollback_kb_version(course_id, target_version_id, actor_user_id=0):
    target = db.session.get(KnowledgeBaseVersion, int(target_version_id))
    if target is None or target.course_id != course_id or target.status == "failed":
        return {"status": "error", "message": "Invalid rollback target"}
    current = get_published_kb_version(course_id=course_id, scope_type=target.scope_type, owner_user_id=target.owner_user_id)
    if current and current.id != target.id:
        current.status = "archived"
        current.is_active = False
        current.archived_at = current_time_text()
    target.status = "published"
    target.is_active = True
    target.published_at = current_time_text()
    add_system_log("warning", "knowledge_versioning", f"Rolled back KB to version {target.id} by user {actor_user_id}")
    db.session.flush()
    return {"status": "success", "version": serialize_kb_version(target)}


def build_kb_version_manifest(version):
    chunks = KnowledgeChunk.query.filter_by(knowledge_base_version_id=version.id).all()
    sources = KnowledgeSource.query.filter_by(version_introduced_id=version.id).all()
    health = summarize_kb_health(version, chunks, sources)
    return {
        "knowledge_base_version_id": version.id,
        "course_id": version.course_id,
        "scope_type": version.scope_type,
        "version_name": version.version_name,
        "status": version.status,
        "source_count": len(sources),
        "chunk_count": len(chunks),
        "deduped_chunk_count": len([chunk for chunk in chunks if chunk.is_duplicate]),
        "index_backend": version.index_backend,
        "index_version": version.index_version,
        "retrieval_version": version.retrieval_version,
        "sources": [
            {
                "source_title": source.source_title or source.name,
                "language": source.language,
                "source_type": source.source_type,
                "authorization_status": source.authorization_status,
                "status": source.status,
                "chunk_count": KnowledgeChunk.query.filter_by(knowledge_source_id=source.id).count(),
            }
            for source in sources
        ],
        "quality": {
            "duplicate_ratio": health["metrics"].get("duplicate_ratio", 0),
            "unknown_authorization_count": health["metrics"].get("unknown_authorization_count", 0),
            "health_status": health["status"],
        },
    }


def run_knowledge_health_check(course_id=None, kb_version_id=None, scope_type="course", owner_user_id=None):
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id)) if kb_version_id else get_published_kb_version(course_id=course_id, scope_type=scope_type, owner_user_id=owner_user_id)
    chunks = KnowledgeChunk.query.filter_by(knowledge_base_version_id=version.id).all() if version else []
    source_ids = sorted({getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None) for chunk in chunks if getattr(chunk, "knowledge_source_id", None) or getattr(chunk, "source_id", None)})
    sources = KnowledgeSource.query.filter(KnowledgeSource.id.in_(source_ids)).all() if source_ids else []
    return summarize_kb_health(version, chunks, sources)


def run_retrieval_regression_for_course(course_id=None, kb_version_id=None):
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id)) if kb_version_id else get_published_kb_version(course_id=course_id)
    if version is None:
        return {"status": "failed", "case_count": 0, "passed": 0, "failed": 0, "negative_match_errors": 0, "no_evidence_forced_match": 0, "results": [], "message": "No KB version found"}
    items = EvaluationItem.query
    if course_id:
        items = items.filter(EvaluationItem.course_id == course_id)
    cases = []
    for item in items.limit(80).all():
        cases.append({
            "query": item.english_term,
            "expected_chinese_evidence": item.expected_chinese_evidence,
            "negative_chinese_evidence": item.negative_chinese_evidence,
            "expect_no_evidence": item.expected_alignment_status in {"no_zh_evidence", "no_en_evidence"},
        })
    if not cases:
        cases = [{"query": "Fourier Transform", "expected_chinese_evidence": "傅里叶"}]

    def retrieve_case(case, query):
        return retrieve_evidence_results(
            query,
            course_id=version.course_id,
            language="zh",
            scope_type=version.scope_type,
            owner_user_id=version.owner_user_id,
            knowledge_base_type="zh_course_kb" if version.scope_type != "personal" else "student_personal_kb",
            knowledge_base_version_id=version.id,
            limit=5,
        )

    result = evaluate_retrieval_cases(cases, retrieve_case)
    if version.status != "published":
        version.quality_gate_status = "pass" if result.get("failed", 0) == 0 else "fail"
        version.manifest_json = json.dumps({"retrieval_regression": {key: result.get(key) for key in ["case_count", "passed", "failed", "negative_match_errors", "no_evidence_forced_match"]}}, ensure_ascii=False)
        db.session.commit()
    return result


def active_chunks_for_kb_version(kb_version_id):
    return KnowledgeChunk.query.filter_by(knowledge_base_version_id=int(kb_version_id)).filter(
        db.or_(KnowledgeChunk.is_active == True, KnowledgeChunk.is_active.is_(None)),  # noqa: E712
        db.or_(KnowledgeChunk.is_duplicate == False, KnowledgeChunk.is_duplicate.is_(None)),  # noqa: E712
        db.or_(KnowledgeChunk.index_status == "indexed", KnowledgeChunk.index_status == "", KnowledgeChunk.index_status.is_(None)),
    ).all()


def build_vector_index_for_kb_version(kb_version_id, apply=False, embedding_provider_name=None, vector_backend_name=None):
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id))
    if version is None:
        return {"status": "error", "error_code": "RESOURCE_NOT_FOUND", "message": "KB version not found."}
    chunks = active_chunks_for_kb_version(version.id)
    provider = get_embedding_provider(embedding_provider_name or EMBEDDING_PROVIDER)
    vector_backend = get_vector_index_backend(vector_backend_name or VECTOR_INDEX_BACKEND)
    result = {
        "status": "dry_run" if not apply else "ready",
        "kb_version_id": version.id,
        "chunks_scanned": len(chunks),
        "chunks_embedded": 0,
        "skipped": 0,
        "failed": 0,
        "backend": vector_backend.backend_name,
        "embedding_provider": provider.provider_name(),
        "embedding_model": EMBEDDING_MODEL,
        "embedding_dimension": provider.dimension(),
    }
    if not chunks:
        result["status"] = "empty"
        return result
    if not provider.is_available() or vector_backend.backend_name == "none":
        result["status"] = "unavailable"
        result["message"] = "Vector index or embedding provider is not configured."
        result["skipped"] = len(chunks)
        return result
    texts = [chunk.content for chunk in chunks]
    if not apply:
        result["chunks_embedded"] = len(chunks)
        return result
    try:
        embeddings = provider.embed_texts(texts, model=EMBEDDING_MODEL or None)
        items = []
        for chunk, embedding in zip(chunks, embeddings):
            if not embedding:
                result["skipped"] += 1
                continue
            items.append({
                "chunk_id": chunk.id,
                "kb_version_id": version.id,
                "embedding": embedding,
                "metadata": {
                    "course_id": chunk.course_id,
                    "scope_type": chunk.scope_type,
                    "owner_user_id": chunk.owner_user_id,
                    "language": chunk.language,
                    "knowledge_base_type": chunk.knowledge_base_type,
                    "visibility": chunk.visibility,
                    "source_status": getattr(chunk, "_source_status", ""),
                    "authorization_status": getattr(chunk, "_authorization_status", ""),
                },
            })
        write_result = vector_backend.upsert(version.id, items)
        result["chunks_embedded"] = len(items)
        result["index_result"] = write_result
        version.index_backend = vector_backend.backend_name
        version.index_version = f"{vector_backend.backend_name}_v1"
        version.retrieval_version = RETRIEVAL_VERSION if RETRIEVAL_BACKEND == "lexical" else f"{RETRIEVAL_BACKEND}_v1"
        version.embedding_provider = provider.provider_name()
        version.embedding_model = EMBEDDING_MODEL
        version.embedding_dimension = provider.dimension()
        version.vector_index_status = "ready"
        version.vector_index_updated_at = current_time_text()
        db.session.commit()
        return result
    except Exception as exc:  # pragma: no cover - defensive external-provider boundary
        db.session.rollback()
        result["status"] = "error"
        result["failed"] = len(chunks)
        result["message"] = str(exc)
        return result


def vector_index_health(kb_version_id=None, vector_backend_name=None):
    return get_vector_index_backend(vector_backend_name or VECTOR_INDEX_BACKEND).healthcheck(kb_version_id=kb_version_id)


def _evaluation_cases_for_course(course_id=None, evaluation_set_id=None, limit=80):
    query = EvaluationItem.query
    if evaluation_set_id:
        query = query.filter_by(set_id=int(evaluation_set_id))
    elif course_id:
        query = query.filter(EvaluationItem.course_id == int(course_id))
    items = query.limit(limit).all()
    return [
        {
            "query": item.english_term,
            "course_id": item.course_id or course_id,
            "expected_chinese_evidence": item.expected_chinese_evidence,
            "expected_english_evidence": item.expected_english_evidence,
            "negative_chinese_evidence": item.negative_chinese_evidence,
            "negative_english_evidence": item.negative_english_evidence,
            "expect_no_evidence": item.expected_alignment_status in {"no_zh_evidence", "no_en_evidence"},
        }
        for item in items
    ]


def run_retrieval_experiment(course_id=None, evaluation_set_id=None, kb_version_id=None, created_by=0):
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id)) if kb_version_id else get_published_kb_version(course_id=course_id)
    if version is None:
        return {"status": "failed", "message": "No KB version found.", "metrics": {}}
    cases = _evaluation_cases_for_course(course_id=course_id or version.course_id, evaluation_set_id=evaluation_set_id)
    if not cases:
        cases = [{"query": "Fourier Transform", "expected_chinese_evidence": "傅里叶", "course_id": course_id or version.course_id}]

    def search_for_backend(backend_name, case):
        return retrieve_evidence_results(
            case.get("query"),
            course_id=course_id or version.course_id,
            language="zh",
            scope_type=version.scope_type,
            owner_user_id=version.owner_user_id,
            limit=int(os.environ.get("FINAL_EVIDENCE_TOP_K", "5")),
            knowledge_base_type="zh_course_kb" if version.scope_type != "personal" else "student_personal_kb",
            knowledge_base_version_id=version.id,
            retrieval_backend=backend_name,
        )

    metrics = {}
    for backend_name in BACKENDS_TO_COMPARE:
        if backend_name in {"vector", "hybrid", "hybrid_rerank"}:
            health = vector_index_health(version.id)
            provider = get_embedding_provider(EMBEDDING_PROVIDER)
            if health.get("status") in {"unavailable", "empty"} or not provider.is_available():
                metrics[backend_name] = {"skipped": True, "reason": "Vector index or embedding provider unavailable."}
                continue
        metrics[backend_name] = evaluate_backend_cases(cases, backend_name, search_for_backend)
    recommendation = recommend_backend(metrics)
    experiment_payload = {
        "course_id": course_id or version.course_id,
        "evaluation_set_id": evaluation_set_id,
        "kb_version_id": version.id,
        "results": metrics,
        "recommendation": recommendation,
    }
    run = RetrievalExperimentRun(
        course_id=course_id or version.course_id,
        evaluation_set_id=int(evaluation_set_id) if evaluation_set_id else None,
        kb_version_id=version.id,
        experiment_name=f"retrieval_exp_{datetime.utcnow():%Y%m%d_%H%M%S}",
        backends_tested_json=json.dumps(BACKENDS_TO_COMPARE, ensure_ascii=False),
        best_backend="lexical",
        recommendation=recommendation,
        metrics_json=json.dumps(metrics, ensure_ascii=False),
        created_by=created_by or 0,
        created_at=current_time_text(),
        finished_at=current_time_text(),
        status="completed",
        report_markdown=retrieval_experiment_markdown(experiment_payload),
    )
    db.session.add(run)
    db.session.commit()
    return {"status": "completed", "experiment_id": run.id, "metrics": metrics, "recommendation": recommendation, "report_markdown": run.report_markdown}


def retrieve_evidence_results(
    query_text,
    course_id=None,
    course_name="",
    language="",
    scope_type="course",
    owner_user_id=None,
    limit=4,
    min_score=MIN_RETRIEVAL_SCORE,
    knowledge_base_type="",
    discipline=None,
    knowledge_base_version_id=None,
    retrieval_backend=None
):
    if not query_text:
        return []

    kb_type = knowledge_base_type or ("en_course_kb" if language == "en" else "zh_course_kb" if language == "zh" else "")
    query = KnowledgeChunk.query
    selected_version = db.session.get(KnowledgeBaseVersion, int(knowledge_base_version_id)) if knowledge_base_version_id else get_published_kb_version(
        course_id=course_id,
        scope_type=scope_type,
        owner_user_id=owner_user_id,
    )
    if selected_version is not None:
        query = query.filter_by(knowledge_base_version_id=selected_version.id)
    if scope_type == "personal":
        query = query.filter_by(
            knowledge_base_type="student_personal_kb",
            visibility="private",
            owner_user_id=str(owner_user_id or "")
        )
        if course_id:
            query = query.filter_by(course_id=course_id)
    elif scope_type == "global":
        query = query.filter_by(visibility="global")
    else:
        if course_id:
            query = query.filter_by(course_id=course_id, visibility="course")
        elif course_name:
            query = query.filter_by(course=course_name, visibility="course")
        else:
            return []
    if kb_type:
        query = query.filter_by(knowledge_base_type=kb_type)
    query = query.filter(db.or_(KnowledgeChunk.is_active == True, KnowledgeChunk.is_active.is_(None)))  # noqa: E712
    query = query.filter(db.or_(KnowledgeChunk.is_duplicate == False, KnowledgeChunk.is_duplicate.is_(None)))  # noqa: E712
    query = query.filter(db.or_(KnowledgeChunk.index_status == "indexed", KnowledgeChunk.index_status == "", KnowledgeChunk.index_status.is_(None)))

    chunks = []
    for chunk in annotate_chunks_with_source_metadata(query.all()):
        source_status = getattr(chunk, "_source_status", "active")
        auth_status = getattr(chunk, "_authorization_status", "unknown")
        if source_status in {"removed", "archived"}:
            continue
        if scope_type != "personal" and auth_status == "restricted_no_derivative":
            continue
        chunks.append(chunk)
    backend_name = (retrieval_backend or RETRIEVAL_BACKEND or "lexical").strip().lower()
    if backend_name not in VALID_RETRIEVAL_BACKENDS:
        backend_name = "lexical"
    backend = get_retrieval_backend(backend_name)
    results = backend.search(
        query_text,
        {
            "course_id": course_id,
            "language": language,
            "knowledge_base_type": "student_personal_kb" if scope_type == "personal" else kb_type,
            "scope_type": scope_type,
            "owner_user_id": owner_user_id,
            "discipline": discipline,
        },
        selected_version.id if selected_version else None,
        top_k=limit,
        context={"chunks": chunks},
    )
    retrieval_run = RetrievalRun(
        query=query_text,
        course_id=course_id,
        scope_type=scope_type,
        owner_user_id=int(owner_user_id) if str(owner_user_id or "").isdigit() else None,
        knowledge_base_version_id=selected_version.id if selected_version else None,
        retrieval_version=(results[0].get("retrieval_version") if results else (selected_version.retrieval_version if selected_version else KB_RETRIEVAL_VERSION)),
        index_version=(selected_version.index_version if selected_version else INDEX_VERSION),
        result_count=len(results),
        top_score=float(results[0].get("evidence_score", 0) if results else 0),
        status="completed",
        created_at=current_time_text(),
        metadata_json=json.dumps({"knowledge_base_type": kb_type, "language": language, "retrieval_backend": backend_name}, ensure_ascii=False),
    )
    db.session.add(retrieval_run)
    db.session.flush()
    for item in results:
        item.pop("_chunk", None)
        item["retrieval_run_id"] = retrieval_run.id
        item["knowledge_base_version_id"] = selected_version.id if selected_version else item.get("knowledge_base_version_id")
        item["retrieval_version"] = retrieval_run.retrieval_version
        item["index_version"] = retrieval_run.index_version

    db.session.flush()
    return results


def create_knowledge_source_for_document(document, creator, discipline="", source_name=""):
    if document.scope_type == "personal":
        license_status = "restricted"
        allow_search = False
        allow_cards = False
        source_type = "student_upload"
    elif document.scope_type == "global":
        license_status = "open_licensed"
        allow_search = True
        allow_cards = True
        source_type = document.source_type or "platform_seed"
    else:
        license_status = "authorized"
        allow_search = True
        allow_cards = True
        source_type = document.source_type or "teacher_upload"

    parse_record = DocumentParseRecord.query.filter_by(parse_uid=document.parse_uid).first() if getattr(document, "parse_uid", "") else None
    quality_status = getattr(parse_record, "quality_status", "") if parse_record else ""
    quality_flags = safe_json_loads(getattr(parse_record, "quality_flags", "[]") if parse_record else getattr(document, "quality_flags_json", "[]"), [])
    course_obj = db.session.get(Course, document.course_id) if document.course_id else None
    source_role = "student_private_material" if document.scope_type == "personal" else ("english_course_material" if document.language == "en" else "chinese_reference_material")
    trust_level = "low_quality" if quality_status in (knowledge_governance_service.REVIEW_PARSE_STATUSES | knowledge_governance_service.BLOCKED_PARSE_STATUSES) else ("student_uploaded" if document.scope_type == "personal" else "teacher_verified")
    governance_status = "blocked" if quality_status in knowledge_governance_service.BLOCKED_PARSE_STATUSES else ("needs_review" if quality_status in knowledge_governance_service.REVIEW_PARSE_STATUSES else "active")
    source = KnowledgeSource(
        title=source_name or document.filename,
        name=source_name or document.filename,
        source_title=source_name or document.filename,
        course=course_obj.name if course_obj else "",
        chapter="",
        course_id=document.course_id if document.scope_type == "course" else None,
        scope_type=document.scope_type,
        owner_user_id=document.owner_user_id if document.scope_type == "personal" else None,
        document_id=document.id,
        language=document.language,
        discipline=discipline,
        source_type=source_type,
        source_role=source_role,
        owner_type="student" if document.scope_type == "personal" else "teacher",
        owner_id=str(document.owner_user_id if document.scope_type == "personal" else (creator.id if creator else "")),
        visibility="private" if document.scope_type == "personal" else ("public" if document.scope_type == "global" else "course"),
        trust_level=trust_level,
        parse_uid=getattr(document, "parse_uid", ""),
        source_filename=getattr(document, "original_filename", "") or getattr(document, "filename", ""),
        file_type=getattr(document, "file_type", "") or "unknown",
        content_hash=getattr(document, "sha256", "") or getattr(document, "file_sha256", ""),
        version=1,
        license_note="",
        quality_status=quality_status,
        quality_flags=json.dumps(quality_flags, ensure_ascii=False),
        knowledge_base_type="student_personal_kb" if document.scope_type == "personal" else ("en_course_kb" if document.language == "en" else "zh_course_kb"),
        access_method="manual_upload",
        license_status=license_status,
        license_type="teacher_provided" if document.scope_type == "course" else "unknown",
        authorization_status="allowed_for_private_use" if document.scope_type == "personal" else "allowed_for_course_use",
        source_quality=0.6 if document.scope_type == "personal" else 0.9,
        status=governance_status,
        update_frequency="manual",
        allow_full_text_indexing=document.scope_type != "personal",
        allow_student_search=allow_search,
        allow_derivative_cards=allow_cards,
        created_by=creator.id if creator else 0,
        created_at=current_time_text(),
        updated_at=current_time_text()
    )
    db.session.add(source)
    db.session.flush()
    return source


def index_document_chunks_to_knowledge(document, chunks, source=None, course_name="", discipline="", chapter=""):
    created = []
    if document.scope_type == "personal":
        kb_type = "student_personal_kb"
        visibility = "private"
    else:
        kb_type = "en_course_kb" if document.language == "en" else "zh_course_kb"
        visibility = "global" if document.scope_type == "global" else "course"

    parse_record = DocumentParseRecord.query.filter_by(parse_uid=document.parse_uid).first() if getattr(document, "parse_uid", "") else None
    quality_status = getattr(parse_record, "quality_status", "") if parse_record else ""
    source_quality_flags = set(safe_json_loads(getattr(parse_record, "quality_flags", "[]") if parse_record else getattr(document, "quality_flags_json", "[]"), []))
    if quality_status:
        source_quality_flags.add(quality_status)
    trust_level = getattr(source, "trust_level", "") or ("low_quality" if quality_status in knowledge_governance_service.REVIEW_PARSE_STATUSES else "unknown")
    chunk_status = "blocked" if quality_status in knowledge_governance_service.BLOCKED_PARSE_STATUSES else ("needs_review" if quality_status in knowledge_governance_service.REVIEW_PARSE_STATUSES else "active")
    for index, chunk in enumerate(chunks, start=1):
        page_label = chunk.source_location or ""
        if chunk.page_number:
            page_label = f"page {chunk.page_number}"
        elif chunk.slide_number:
            page_label = f"slide {chunk.slide_number}"
        chunk_flags = set(source_quality_flags)
        chunk_flags.update(safe_json_loads(getattr(chunk, "quality_flags_json", "[]"), []))
        record = KnowledgeChunk(
            source_uid=getattr(source, "source_uid", "") if source else "",
            document_id=document.id,
            source_id=source.id if source else None,
            knowledge_source_id=source.id if source else None,
            parse_uid=getattr(chunk, "parse_uid", ""),
            parse_block_uid=getattr(chunk, "parse_block_uid", ""),
            course_id=document.course_id,
            scope_type=document.scope_type,
            course=course_name,
            title=document.filename,
            discipline=discipline,
            chapter=chapter or chunk.section_title or "",
            chunk_index=index,
            content=chunk.content,
            normalized_text=normalize_chunk_text(chunk.content),
            content_hash=compute_content_hash(chunk.content),
            source_page=page_label,
            source_slide=str(chunk.slide_number or ""),
            source_section=chapter or chunk.section_title or "",
            source_locator=chunk.source_location or page_label,
            page_number=chunk.page_number,
            slide_number=chunk.slide_number,
            block_type="text",
            token_count=len(normalize_chunk_text(chunk.content).split()),
            char_count=len(chunk.content or ""),
            keywords="",
            source_citation=f"{document.filename} {page_label}".strip(),
            language=document.language,
            knowledge_base_type=kb_type,
            owner_user_id=str(document.owner_user_id or ""),
            visibility=visibility,
            index_status="indexed",
            quality_status=quality_status,
            quality_flags=json.dumps(sorted(flag for flag in chunk_flags if flag), ensure_ascii=False),
            trust_level=trust_level,
            status=chunk_status,
            embedding_status="not_started",
            is_active=chunk_status != "blocked",
            created_at=current_time_text(),
            updated_at=current_time_text()
        )
        db.session.add(record)
        created.append(record)

    db.session.flush()
    rebuild_embeddings_for_chunks(created)
    return created


def create_kb_version(scope_type, course_id, created_by, source_count=1, chunk_count=0, description=""):
    version_count = KnowledgeBaseVersion.query.filter_by(kb_scope=scope_type, course_id=course_id).count()
    version = KnowledgeBaseVersion(
        kb_scope=scope_type,
        course_id=course_id,
        owner_user_id=created_by if scope_type == "personal" else None,
        version_name=f"{scope_type}-v{version_count + 1}",
        description=description,
        source_count=source_count,
        chunk_count=chunk_count,
        created_at=current_time_text(),
        created_by=created_by,
        is_active=True
    )
    db.session.add(version)
    db.session.flush()
    return version


def run_alignment_for_chunks(chunks, course=None, scope_type="course", owner_user_id=None, source_document_id=None, limit_terms=12, triggered_by_user_id=None):
    parse_quality_metadata = parse_quality_metadata_from_chunks(chunks)
    if parse_quality_risk_service.should_block_downstream_creation(parse_quality_metadata):
        return []
    text = "\n\n".join(
        chunk.content
        for chunk in chunks
        if chunk.content and not contains_formula_placeholder(chunk.content)
    )
    meta = current_provider_metadata()
    alignment_run = AlignmentRun(
        document_id=source_document_id,
        course_id=course.id if course else None,
        triggered_by=triggered_by_user_id or owner_user_id or 0,
        provider=meta["provider"],
        model_name=meta["model_name"],
        ai_provider=meta["provider"],
        ai_provider_mode=meta.get("provider_mode", ""),
        ai_model=meta["model_name"],
        prompt_key="term_alignment",
        prompt_version="v1",
        retrieval_version=RETRIEVAL_VERSION,
        status="running",
        started_at=current_time_text()
    )
    db.session.add(alignment_run)
    db.session.flush()
    if contains_ocr_placeholder(text) or contains_formula_placeholder(text):
        alignment_run.status = "failed"
        alignment_run.error_message = "OCR or formula placeholder content is not eligible for term extraction."
        alignment_run.finished_at = current_time_text()
        alignment_run.failed_count = 1
        alignment_run.metrics_json = json.dumps({"reason": "ocr_or_formula_placeholder_blocked"}, ensure_ascii=False)
        db.session.flush()
        return []
    terms = extract_terms_from_text(text)
    if bool(meta.get("is_real_provider")):
        try:
            chunk_payload = [
                {
                    "content": chunk.content,
                    "source_location": chunk.source_location,
                    "ocr_confidence": chunk.ocr_confidence
                }
                for chunk in chunks[:20]
                if chunk.content and not contains_ocr_placeholder(chunk.content) and not contains_formula_placeholder(chunk.content)
            ]
            ai_call = call_ai_task(
                task_type="term_extraction",
                prompt_key="term_extraction",
                prompt_version="v1",
                input_payload={
                    "chunks": chunk_payload,
                    "course": course.name if course else "",
                    "scope_type": scope_type,
                },
                user_id=triggered_by_user_id or owner_user_id,
                course_id=course.id if course else None,
                document_id=source_document_id,
                alignment_run_id=alignment_run.id,
            )
            if ai_call.get("status") != "success":
                raise RuntimeError(f"{ai_call.get('error_code')}: {ai_call.get('message')}")
            ai_terms = ai_call.get("result", {}).get("terms", [])
            existing = {item["english_term"].strip().lower() for item in terms}
            for item in ai_terms:
                term = str(item.get("english_term", "")).strip()
                if not term or term.lower() in existing or is_probably_noise(term):
                    continue
                terms.append({
                    "english_term": term,
                    "context": item.get("context_sentence") or extract_context_sentence(text, term),
                    "confidence": int(item.get("confidence", 70) or 70),
                    "status": "pending",
                    "extraction_method": "ai"
                })
                existing.add(term.lower())
        except Exception as exc:
            add_system_log("warning", "ai_term_extraction", f"{AI_PROVIDER} term extraction failed; local heuristic used: {exc}")
    cards = []
    min_ocr = min([chunk.ocr_confidence for chunk in chunks], default=100)
    course_name = course.name if course else ""
    course_id = course.id if course else None

    try:
        for item in terms[:limit_terms]:
            context = item.get("context", "") or extract_context(text, item["english_term"])
            alignment = generate_alignment_result(
                english_term=item["english_term"],
                courseware_sentence=context,
                course=course_name,
                chapter="",
                scope_type=scope_type,
                owner_user_id=owner_user_id,
                min_ocr_confidence=min_ocr
            )
            alignment = apply_parse_quality_to_alignment(alignment, parse_quality_metadata)
            alignment["alignment_run_id"] = alignment_run.id
            card = create_or_update_card_from_alignment(
                english_term=item["english_term"],
                alignment=alignment,
                scope_type=scope_type,
                course_id=course_id,
                owner_user_id=owner_user_id,
                source_document_id=source_document_id,
                courseware_sentence=context,
                min_ocr_confidence=min_ocr
            )
            if card:
                cards.append(card)
                if owner_user_id:
                    record_usage(owner_user_id, "ai_alignment", 1, related_document_id=source_document_id, related_term_id=card.id)
        update_alignment_run_stats(alignment_run, cards=cards, term_count=len(terms))
        alignment_run.status = "completed"
        metrics = safe_json_loads(alignment_run.metrics_json, {})
        metrics.update({
            "min_ocr_confidence": min_ocr,
            "provider_is_real": bool(meta.get("is_real_provider")),
        })
        alignment_run.metrics_json = json.dumps(metrics, ensure_ascii=False)
        alignment_run.finished_at = current_time_text()
    except Exception as exc:
        alignment_run.status = "failed"
        alignment_run.error_message = redact_for_log(exc)
        alignment_run.failed_count = 1
        alignment_run.finished_at = current_time_text()
        add_system_log("error", "alignment_run", f"Alignment run {alignment_run.id} failed: {exc}")
        raise

    db.session.flush()
    return cards


def find_pdf_font():
    candidates = [
        os.environ.get("PDF_FONT_PATH", ""),
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simsun.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc"
    ]
    for path in candidates:
        if path and os.path.exists(path):
            return path
    return ""


def build_cards_pdf(cards):
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    except ImportError as exc:
        raise RuntimeError("缺少 reportlab，无法导出 PDF。") from exc

    font_path = find_pdf_font()
    if not font_path:
        raise RuntimeError("未找到可用中文字体。请设置 PDF_FONT_PATH 或安装 Arial Unicode / Noto CJK 字体。")

    font_name = "LexiBridgeCJK"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(TTFont(font_name, font_path))

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=16 * mm,
        rightMargin=16 * mm,
        topMargin=16 * mm,
        bottomMargin=16 * mm
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "LexiTitle",
        parent=styles["Heading1"],
        fontName=font_name,
        fontSize=16,
        leading=22,
        textColor=colors.HexColor("#07172c")
    )
    heading_style = ParagraphStyle(
        "LexiHeading",
        parent=styles["Heading2"],
        fontName=font_name,
        fontSize=12,
        leading=16,
        textColor=colors.HexColor("#0f6f8f")
    )
    body_style = ParagraphStyle(
        "LexiBody",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=9,
        leading=13
    )
    story = [
        Paragraph("LexiBridge AI 双语课程术语知识卡片导出", title_style),
        Spacer(1, 8)
    ]
    for card in cards:
        story.append(Paragraph(f"{card.english_term} -> {card.final_chinese_term}", heading_style))
        fields = [
            ("Status", card.status),
            ("Confidence", str(card.confidence_score)),
            ("AI Model", getattr(card, "ai_model", "")),
            ("Risk Note", getattr(card, "risk_note", "")),
            ("Context", card.courseware_sentence),
            ("Concept Explanation", card.concept_explanation),
            ("English Evidence", card.english_kb_evidence),
            ("Chinese Evidence", card.chinese_kb_evidence),
            ("Alignment Reason", card.alignment_reason)
        ]
        for label, value in fields:
            safe_value = str(value or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(f"<b>{label}:</b> {safe_value}", body_style))
        story.append(Spacer(1, 8))
    doc.build(story)
    buffer.seek(0)
    return buffer


# ============================================================
# API 路由
# ============================================================

@app.route("/")
def home():
    return jsonify({
        "message": "LexiBridge AI backend is running.",
        "version": "0.1.0"
    })


@app.route("/api/test", methods=["GET"])
def test():
    return jsonify({
        "status": "success",
        "project": "LexiBridge AI",
        "message": "AI 双语课程知识对齐平台前后端连接成功。",
        "ai_provider": AI_PROVIDER,
        "ai_configured": bool(DEEPSEEK_API_KEY) if AI_PROVIDER == "deepseek" else False,
        "auth_required": AUTH_REQUIRED,
        "embedding_dim": LOCAL_EMBEDDING_DIM
    })


@app.route("/api/auth/register", methods=["POST"])
def register_user():
    data = request.get_json() or {}
    username = str(data.get("username", "")).strip()
    email = str(data.get("email", "")).strip().lower()
    password = str(data.get("password", "")).strip()
    role = normalize_role(data.get("role", "student"))

    # Public registration always creates students; admins can promote roles later.
    if role != "student":
        role = "student"

    if len(username) < 3:
        return api_error("VALIDATION_ERROR", "用户名至少需要 3 个字符。", 400)

    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", email):
        return api_error("VALIDATION_ERROR", "请输入有效邮箱。", 400)

    password_errors = validate_password_strength(password)
    if password_errors:
        return jsonify({
            "status": "error",
            "message": "密码强度不足：" + "，".join(password_errors)
        }), 400

    if User.query.filter_by(username=username).first() is not None:
        return api_error("VALIDATION_ERROR", "用户名已存在。", 400)

    if find_user_by_email(email) is not None:
        return api_error("VALIDATION_ERROR", "邮箱已注册。", 400)

    token = make_mock_email_token()
    user = User(
        username=username,
        email=email,
        password_hash=generate_password_hash(password, method="pbkdf2:sha256"),
        role=role,
        display_name=username,
        is_verified=False,
        verification_token=token,
        verification_token_expires_at=future_time_text(60 * 24),
        created_at=current_time_text()
    )
    db.session.add(user)
    db.session.commit()
    log_mock_email("verify", email, token)

    return jsonify({
        "status": "success",
        "message": "注册成功。开发阶段已生成 mock email 验证码。",
        "verification_code": token if MOCK_EMAIL_ENABLED else "",
        "verification_token": token,
        "verification_url": f"/api/auth/verify-email?token={token}",
        "user": serialize_user(user)
    })


@app.route("/api/auth/verify-email", methods=["GET", "POST"])
def verify_email():
    token = request.args.get("token", "").strip()
    if request.method == "POST":
        token = str((request.get_json() or {}).get("token", token)).strip()

    user = User.query.filter_by(verification_token=token).first()
    if user is None:
        return api_error("RESOURCE_NOT_FOUND", "邮箱验证码不存在。", 404)

    expires_at = parse_time_text(getattr(user, "verification_token_expires_at", ""))
    if expires_at and expires_at < datetime.now():
        return api_error("VALIDATION_ERROR", "邮箱验证码已过期。", 400)

    user.is_verified = True
    user.verification_token = ""
    user.verification_token_expires_at = ""
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": "邮箱验证成功。",
        "user": serialize_user(user)
    })


@app.route("/api/auth/login", methods=["POST"])
def login_user():
    data = request.get_json() or {}
    email = str(data.get("email", "")).strip().lower()
    password = str(data.get("password", "")).strip()

    user = find_user_by_email(email)
    if user is None or not check_password_hash(user.password_hash, password):
        return api_error("AUTH_REQUIRED", "邮箱或密码错误。", 401)

    if not bool(getattr(user, "is_verified", False)):
        return jsonify({
            "status": "error",
            "message": "邮箱尚未验证。请先使用邮箱验证码完成验证。",
            "verification_code": getattr(user, "verification_token", "") if MOCK_EMAIL_ENABLED else "",
            "verification_token": getattr(user, "verification_token", "")
        }), 403

    user.last_login_at = current_time_text()
    token = create_auth_token(user)
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": "登录成功。",
        "token": token,
        "user": serialize_user(user)
    })


@app.route("/api/auth/me", methods=["GET"])
def get_me():
    user, error_response = require_current_user()
    if error_response:
        return error_response

    return jsonify({
        "status": "success",
        "user": serialize_user(user),
        "system_status": current_system_status(user)
    })


@app.route("/api/auth/logout", methods=["POST"])
def logout_user():
    auth_header = request.headers.get("Authorization", "").strip()
    if auth_header.lower().startswith("bearer "):
        token_value = auth_header.split(" ", 1)[1].strip()
        token = AuthToken.query.filter_by(token_hash=token_hash(token_value), revoked=False).first()
        if token is None:
            token = AuthToken.query.filter_by(token=token_value, revoked=False).first()
        if token is not None:
            token.revoked = True
            db.session.commit()

    return jsonify({
        "status": "success",
        "message": "已退出登录。"
    })


@app.route("/api/auth/password-reset/request", methods=["POST"])
def request_password_reset():
    data = request.get_json() or {}
    email = str(data.get("email", "")).strip().lower()
    user = find_user_by_email(email)

    if user is not None:
        token = make_mock_email_token()
        user.reset_token = token
        user.reset_token_expires_at = future_time_text(60)
        db.session.commit()
        log_mock_email("password-reset", email, token)

    return jsonify({
        "status": "success",
        "message": "如果邮箱存在，系统已生成 mock password reset 验证码。",
        "reset_code": getattr(user, "reset_token", "") if user is not None and MOCK_EMAIL_ENABLED else "",
        "reset_token": getattr(user, "reset_token", "") if user is not None else ""
    })


@app.route("/api/auth/password-reset/confirm", methods=["POST"])
def confirm_password_reset():
    data = request.get_json() or {}
    token = str(data.get("token", "")).strip()
    password = str(data.get("password", "")).strip()

    user = User.query.filter_by(reset_token=token).first()
    if user is None:
        return api_error("RESOURCE_NOT_FOUND", "重置 token 不存在。", 404)

    expires_at = parse_time_text(getattr(user, "reset_token_expires_at", ""))
    if expires_at and expires_at < datetime.now():
        return api_error("VALIDATION_ERROR", "重置 token 已过期。", 400)

    password_errors = validate_password_strength(password)
    if password_errors:
        return jsonify({
            "status": "error",
            "message": "密码强度不足：" + "，".join(password_errors)
        }), 400

    user.password_hash = generate_password_hash(password, method="pbkdf2:sha256")
    user.reset_token = ""
    user.reset_token_expires_at = ""
    db.session.commit()

    return jsonify({"status": "success", "message": "密码已重置。"})


@app.route("/api/courses", methods=["GET", "POST"])
def courses():
    current_user = None
    if AUTH_REQUIRED:
        roles = {"student", "teacher", "admin"} if request.method == "GET" else {"teacher", "admin"}
        current_user, error_response = require_current_user(roles)
        if error_response:
            return error_response

    if request.method == "POST":
        data = request.get_json() or {}
        name = str(data.get("course_name") or data.get("name", "")).strip()
        course_code = str(data.get("course_code", "")).strip()
        semester = str(data.get("semester", "")).strip()
        description = str(data.get("description", "")).strip()
        language_mode = str(data.get("language_mode", "bilingual")).strip() or "bilingual"

        if not name:
            return api_error("VALIDATION_ERROR", "课程名称不能为空。", 400)

        course = Course.query.filter_by(name=name).first()
        if course is None:
            course = Course(
                name=name,
                course_code=course_code,
                semester=semester,
                description=description,
                language_mode=language_mode,
                teacher_id=current_user.id if current_user else 0,
                status="draft",
                created_at=current_time_text()
            )
            db.session.add(course)
            db.session.commit()

            if current_user:
                db.session.add(CourseMember(
                    course_id=course.id,
                    user_id=current_user.id,
                    role=current_user.role,
                    role_in_course="teacher" if current_user.role in {"teacher", "admin"} else "student",
                    created_at=current_time_text(),
                    joined_at=current_time_text()
                ))
                db.session.commit()

        return jsonify({
            "status": "success",
            "course": serialize_course(course)
        })

    search_text = str(request.args.get("q", "")).strip()
    query = Course.query.filter(db.or_(Course.deleted_at == "", Course.deleted_at.is_(None)))
    if search_text:
        pattern = f"%{search_text}%"
        query = query.filter(db.or_(
            Course.name.ilike(pattern),
            Course.course_code.ilike(pattern),
            Course.description.ilike(pattern),
        ))

    if current_user and current_user.role == "student":
        courseware_course_ids = db.session.query(Document.course_id).filter(
            Document.scope_type == "course",
            Document.course_id.isnot(None),
            db.or_(Document.deleted_at == "", Document.deleted_at.is_(None)),
        ).distinct()
        query = query.filter(Course.status == "active", Course.id.in_(courseware_course_ids))

    courses_list = query.order_by(Course.name.asc()).all()

    return jsonify({
        "status": "success",
        "count": len(courses_list),
        "courses": [serialize_course(course) for course in courses_list]
    })


@app.route("/api/courses/mine", methods=["GET"])
def my_courses():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    if user.role == "admin":
        courses_list = Course.query.order_by(Course.name.asc()).all()
    else:
        memberships = CourseMember.query.filter_by(user_id=user.id).all()
        course_ids = [item.course_id for item in memberships]
        courses_list = Course.query.filter(Course.id.in_(course_ids)).order_by(Course.name.asc()).all() if course_ids else []

    return jsonify({
        "status": "success",
        "count": len(courses_list),
        "courses": [serialize_course(course) for course in courses_list]
    })


@app.route("/api/courses/<int:course_id>/join", methods=["POST"])
def join_course(course_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    course = db.session.get(Course, course_id)
    if course is None:
        return api_error("RESOURCE_NOT_FOUND", "课程不存在。", 404)

    existing = CourseMember.query.filter_by(course_id=course_id, user_id=user.id).first()
    if existing is None:
        db.session.add(CourseMember(
            course_id=course_id,
            user_id=user.id,
            role=user.role,
            role_in_course="student" if user.role == "student" else user.role,
            created_at=current_time_text(),
            joined_at=current_time_text()
        ))
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": "已加入课程空间。",
        "course": serialize_course(course)
    })


@app.route("/api/documents/upload", methods=["POST"])
def upload_document():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)

    sync_requested = str(request.args.get("sync", "")).strip().lower() in {"1", "true", "yes"}
    file = request.files.get("file")
    if file is None or file.filename == "":
        return api_error_with_audit_context("VALIDATION_ERROR", "没有收到文件。", 400, audit_context)
    if not allowed_file(file.filename):
        ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
        details = {"extension": ext, "dangerous": ext in DANGEROUS_EXTENSIONS}
        return api_error_with_audit_context(
            "UNSUPPORTED_FILE_TYPE",
            "文件格式不支持。支持 PDF、DOCX、PPTX、JPG、PNG、TXT、Markdown。",
            415,
            audit_context,
            details,
        )

    scope_type = str(request.form.get("scope_type", "")).strip().lower()
    if not scope_type:
        scope_type = "personal" if user.role == "student" else "course"
    if scope_type not in KB_SCOPES:
        return api_error_with_audit_context("VALIDATION_ERROR", "scope_type 只能是 global、course 或 personal。", 400, audit_context)

    course = get_course_by_id_or_name(
        request.form.get("course_id", "").strip(),
        request.form.get("course_name", request.form.get("course", "")).strip()
    )
    if scope_type == "global" and user.role != "admin":
        return api_error_with_audit_context("PERMISSION_DENIED", "只有管理员可以上传平台级资料。", 403, audit_context)
    if scope_type == "course":
        if course is None:
            return api_error_with_audit_context("VALIDATION_ERROR", "课程资料上传必须选择课程。", 400, audit_context)
        if not can_manage_course(user, course):
            return api_error_with_audit_context("PERMISSION_DENIED", "当前账号不能管理该课程。", 403, audit_context)
    if scope_type == "personal" and course is not None and not is_course_member(user, course.id):
        return api_error_with_audit_context("PERMISSION_DENIED", "个人资料只能关联已加入课程。", 403, audit_context)

    if scope_type == "personal":
        allowed, reasons, plan, subscription, totals = check_usage_quota(user, page_units=1, ai_units=1)
        if not allowed:
            return api_error_with_audit_context(
                "QUOTA_EXCEEDED",
                "个人学习工作区额度不足，请开通或升级会员。",
                402,
                audit_context,
                {
                    "reasons": reasons,
                    "plan": serialize_subscription_plan(plan),
                    "usage": {
                    "pages_used": totals["pages_used"],
                    "ai_calls_used": totals["ai_calls_used"]
                    }
                },
            )

    original_filename = secure_filename(file.filename)
    if not original_filename or "." not in original_filename:
        return api_error_with_audit_context("UNSUPPORTED_FILE_TYPE", "文件名无有效扩展名。", 415, audit_context)
    ext = original_filename.rsplit(".", 1)[1].lower()
    saved_filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}_{original_filename}"
    fd, temp_upload_path = tempfile.mkstemp(prefix="lexibridge_upload_", suffix=f".{ext}")
    os.close(fd)
    file.save(temp_upload_path)
    magic_ok, magic_error = validate_upload_magic(original_filename, temp_upload_path)
    if not magic_ok:
        try:
            os.remove(temp_upload_path)
        except OSError:
            pass
        return api_error_with_audit_context(
            "UNSUPPORTED_FILE_TYPE",
            "上传文件校验失败，文件内容类型与扩展名不一致。",
            415,
            audit_context,
            {"error": magic_error},
        )
    try:
        storage_meta = storage_service().save_file(
            temp_upload_path,
            purpose="uploaded_document",
            owner_user_id=user.id,
            course_id=course.id if course else None,
            original_filename=saved_filename,
        )
    except Exception as exc:
        try:
            os.remove(temp_upload_path)
        except OSError:
            pass
        return api_error_with_audit_context("INTERNAL_ERROR", "文件存储失败。", 500, audit_context, {"error": str(exc)})
    finally:
        try:
            os.remove(temp_upload_path)
        except OSError:
            pass
    save_path = storage_meta["absolute_path"]
    file_hash = storage_meta["sha256"]

    now = current_time_text()
    document = Document(
        owner_user_id=user.id,
        course_id=course.id if course else None,
        scope_type=scope_type,
        filename=original_filename,
        saved_filename=saved_filename,
        file_sha256=file_hash,
        storage_backend=storage_meta["storage_backend"],
        storage_key=storage_meta["storage_key"],
        original_filename=original_filename,
        content_type=storage_meta["content_type"],
        size_bytes=storage_meta["size_bytes"],
        sha256=file_hash,
        file_type=ext,
        language=str(request.form.get("language", "en")).strip() or "en",
        upload_time=now,
        parsing_status="processing" if sync_requested else "queued",
        ocr_required=False,
        source_type=str(request.form.get("source_type", "")).strip() or ("student_upload" if scope_type == "personal" else "teacher_upload")
    )
    db.session.add(document)
    if scope_type == "course" and course is not None:
        course.status = "active"
    db.session.flush()
    storage_meta["document_id"] = document.id
    storage_record = create_storage_object_from_metadata(
        storage_meta,
        visibility=visibility_for_scope(scope_type),
    )
    document.storage_object_id = storage_record.id

    job = IngestionJob(
        document_id=document.id,
        status="running" if sync_requested else "queued",
        started_at=now if sync_requested else "",
        created_by=user.id
    )
    db.session.add(job)
    db.session.flush()

    parse_result, parse_record, parse_blocks = create_parse_record_for_saved_file(
        save_path,
        document,
        original_filename,
        storage_meta.get("content_type", ""),
        audit_context=audit_context,
    )
    parse_summary = parse_quality_summary(parse_record)
    if not document_parse_quality_service.should_allow_term_extraction(parse_record):
        formula_blocks = create_formula_block_for_quality_gate(document, save_path)
        blocked = block_document_by_quality_gate(document, job, parse_record, audit_context=audit_context)
        return api_error_with_audit_context(
            quality_gate_error_code(parse_record),
            "文档已解析，但被解析质量门禁阻止，未进入知识库入库或术语抽取。",
            422,
            audit_context,
            {
                **blocked,
                "blocked_by_quality_gate": True,
                "document": serialize_document(document),
                "formula_status": formula_blocks[0].status if formula_blocks else "not_required",
                "formula_blocks": [serialize_formula_block(block) for block in formula_blocks],
                "cards": [],
            },
        )

    if not sync_requested:
        background_job = create_background_job(
            "document_ingestion",
            user,
            course_id=document.course_id,
            document_id=document.id,
            scope_type=document.scope_type,
            owner_user_id=document.owner_user_id,
            input_data={
                "storage_key": document.storage_key,
                "discipline": str(request.form.get("discipline", "")).strip(),
                "source_name": str(request.form.get("source_name", "")).strip(),
                "chapter": str(request.form.get("chapter", "")).strip(),
                "ingestion_job_id": job.id,
            }
        )
        db.session.commit()
        return api_success_with_audit_context({
            "document_id": document.id,
            "job_id": background_job.id,
            "job_type": background_job.job_type,
            "job_status": background_job.status,
            **parse_summary,
            "ingestion_status": "queued",
            "document": serialize_document(document),
            "job": serialize_background_job(background_job),
        }, "文档已保存，解析任务已进入后台队列。", audit_context)

    try:
        parsed_chunks = parse_block_records_to_chunk_items(document, parse_record, parse_blocks)
        parsed_text = parse_result.raw_text
        ocr_required = bool(parse_record.ocr_required)
        ocr_meta = {
            "ocr_provider": "",
            "ocr_status": document.ocr_status or "not_required",
            "ocr_error": document.ocr_error or "",
            "ocr_confidence": 100,
            "warnings": safe_json_loads(getattr(parse_record, "warnings", "[]"), []),
        }
        parser_warnings = list(ocr_meta.get("warnings", []) or [])
        parsed_formula_blocks = []
        formula_records = []
        formula_statuses = []
        formula_providers = []
        formula_flags = set()
        for formula_item in parsed_formula_blocks or []:
            item_flags = list(formula_item.get("quality_flags", []) or [])
            formula_status = formula_item.get("status", "")
            if formula_status:
                formula_statuses.append(formula_status)
                if formula_status != "ok":
                    formula_flags.add(formula_status)
            provider = formula_item.get("provider", "")
            if provider:
                formula_providers.append(provider)
            formula_record = FormulaBlock(
                document_id=document.id,
                course_id=document.course_id,
                owner_user_id=user.id,
                scope_type=scope_type,
                page_number=formula_item.get("page_number"),
                slide_number=formula_item.get("slide_number"),
                bbox_json=json.dumps(formula_item.get("bbox", {}) or {}, ensure_ascii=False),
                image_path=formula_item.get("image_path", ""),
                latex=formula_item.get("latex", ""),
                plain_text=formula_item.get("plain_text", ""),
                provider=provider,
                confidence=float(formula_item.get("confidence", 0) or 0),
                status=formula_status,
                error=formula_item.get("error", ""),
                quality_flags_json=json.dumps(item_flags, ensure_ascii=False),
                created_at=current_time_text()
            )
            db.session.add(formula_record)
            formula_records.append(formula_record)
        if formula_records:
            db.session.flush()

        if not parsed_chunks:
            if document.ocr_required:
                document.parsing_status = "needs_ocr_engine"
                if formula_flags and document.ocr_status in {"", "ocr_unavailable", "empty_result", "ocr_failed"}:
                    document.parsing_status = "needs_formula_ocr_engine"
                job.status = "failed"
                job.error_message = document.ocr_error or "OCR required but unavailable or no text detected."
                job.finished_at = current_time_text()
                add_system_log("warning", "ocr", f"Document {document.id} OCR failed: {job.error_message}")
                db.session.commit()
                error_code = "FORMULA_OCR_UNAVAILABLE" if document.parsing_status == "needs_formula_ocr_engine" else "OCR_UNAVAILABLE"
                return api_error_with_audit_context(
                    error_code,
                    "该文件需要 OCR，但当前 OCR 引擎不可用或未识别到文字；未生成术语卡片。",
                    422,
                    audit_context,
                    {
                        **parse_summary,
                        "ingestion_status": "blocked",
                        "blocked_by_quality_gate": True,
                        "error": job.error_message,
                        "document": serialize_document(document),
                        "formula_blocks": [serialize_formula_block(block) for block in formula_records],
                        "formula_status": next(iter(formula_statuses), "not_required") if formula_statuses else "not_required",
                        "warnings": parser_warnings,
                        "cards": [],
                    },
                )
            raise ValueError("文件解析完成，但没有得到有效文本块。")

        page_units = len({
            chunk.get("page_number") or chunk.get("slide_number") or index
            for index, chunk in enumerate(parsed_chunks, start=1)
        })
        if scope_type == "personal":
            allowed, reasons, plan, subscription, totals = check_usage_quota(user, page_units=page_units, ai_units=min(12, len(parsed_chunks)))
            if not allowed:
                document.parsing_status = "quota_blocked"
                job.status = "failed"
                job.error_message = "；".join(reasons)
                job.finished_at = current_time_text()
                db.session.commit()
                return api_error_with_audit_context(
                    "QUOTA_EXCEEDED",
                    "个人资料解析需要更多会员额度。",
                    402,
                    audit_context,
                    {"reasons": reasons, "document": serialize_document(document)},
                )

        chunk_records = []
        document_flags = set()
        if ocr_required:
            document_flags.add("ocr_triggered")
        if formula_records:
            document_flags.add("formula_region_detected")
        for flag in formula_flags:
            document_flags.add(flag)
        for warning in parser_warnings:
            if "formula" in str(warning).lower():
                document_flags.add("formula_warning")
        for index, item in enumerate(parsed_chunks, start=1):
            item_content = item.get("content", "")
            item_ocr_confidence = int(item.get("ocr_confidence", 100) or 100)
            item_flags = list(item.get("quality_flags", []) or [])
            if contains_ocr_placeholder(item_content) or contains_formula_placeholder(item_content):
                item_flags.append("ocr_placeholder_blocked")
            if item_ocr_confidence < 60:
                item_flags.append("ocr_low_confidence")
                document_flags.add("ocr_low_confidence")
            for flag in item_flags:
                if flag and flag != "native_text_ok":
                    document_flags.add(flag)
            chunk = DocumentChunk(
                document_id=document.id,
                course_id=document.course_id,
                user_id=user.id,
                owner_user_id=user.id,
                chunk_index=index,
                parse_uid=item.get("parse_uid", getattr(document, "parse_uid", "")),
                parse_block_uid=item.get("parse_block_uid", ""),
                language=item.get("language", document.language),
                page_number=item.get("page_number"),
                slide_number=item.get("slide_number"),
                section_title=item.get("section_title", ""),
                content=item_content,
                source_type=item.get("source_type", document.source_type),
                source_location=item.get("source_location", ""),
                ocr_confidence=item_ocr_confidence,
                ocr_provider=item.get("ocr_provider", ""),
                ocr_status=item.get("ocr_status", "not_required"),
                ocr_error=item.get("ocr_error", ""),
                quality_flags_json=json.dumps(item_flags, ensure_ascii=False),
                created_at=current_time_text()
            )
            db.session.add(chunk)
            chunk_records.append(chunk)
        db.session.flush()
        document.quality_flags_json = json.dumps(sorted(document_flags), ensure_ascii=False)

        knowledge_metadata = build_governed_ingestion_metadata(
            parse_record=parse_record,
            title=str(request.form.get("source_name", "")).strip() or document.filename,
            course=course,
            chapter=str(request.form.get("chapter", "")).strip(),
            language=document.language,
            source_type=document.source_type,
            scope_type=scope_type,
            owner_user=user,
            owner_user_id=user.id,
            document_id=document.id,
            knowledge_base_type="student_personal_kb" if scope_type == "personal" else ("en_course_kb" if document.language == "en" else "zh_course_kb"),
            visibility=visibility_for_scope(scope_type),
            content_hash=getattr(document, "sha256", "") or getattr(document, "file_sha256", ""),
            extra_quality_flags=document_flags,
        )
        governed_ingestion = knowledge_ingestion_service.ingest_parse_record_to_governed_knowledge(
            db.session,
            knowledge_ingestion_models(),
            parse_record,
            parse_blocks,
            knowledge_metadata,
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=False,
        )
        source = governed_ingestion.source
        knowledge_chunks = governed_ingestion.chunks
        version = create_kb_version(
            scope_type,
            document.course_id,
            user.id,
            source_count=1,
            chunk_count=len(knowledge_chunks),
            description=f"Document upload: {document.filename}"
        )

        cards = run_alignment_for_chunks(
            chunk_records,
            course=course,
            scope_type="personal" if scope_type == "personal" else "course",
            owner_user_id=user.id if scope_type == "personal" else None,
            source_document_id=document.id,
            triggered_by_user_id=user.id
        )

        if document_flags:
            document.parsing_status = "parsed_with_warnings"
        else:
            document.parsing_status = "parsed"
        job.status = "completed"
        job.finished_at = current_time_text()
        job.processed_pages = page_units
        job.source_id = source.id

        if document.ocr_required:
            add_system_log(
                "warning",
                "ocr",
                f"OCR used for document {document.id}; provider={document.ocr_provider}; status={document.ocr_status}; low confidence chunks route to QC."
            )

        if scope_type == "personal":
            record_usage(user.id, "document_parse_page", page_units, related_document_id=document.id)
            if document.ocr_required:
                record_usage(user.id, "ocr_page", page_units, related_document_id=document.id)

        ingestion_status = "partial" if parse_record.quality_status == "partial_text" else "ingested"
        record_document_ingestion_audit(
            document,
            parse_record,
            ingestion_status,
            audit_context=audit_context,
            commit=False,
        )
        db.session.commit()

        return jsonify({
            "status": "success",
            "request_id": audit_context.get("request_id", ""),
            "message": "文档已解析、切分、索引，并生成术语知识卡片。",
            **parse_summary,
            "ingestion_status": ingestion_status,
            "source_uid": getattr(source, "source_uid", ""),
            "chunk_count": len(knowledge_chunks),
            "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in knowledge_chunks[:20] if getattr(chunk, "chunk_uid", "")],
            "document": serialize_document(document),
            "chunks": [serialize_document_chunk(chunk) for chunk in chunk_records[:8]],
            "ocr_status": document.ocr_status or "not_required",
            "ocr_provider": document.ocr_provider or "not_required",
            "ocr_confidence": ocr_meta.get("ocr_confidence", 100),
            "formula_status": next((status for status in formula_statuses if status != "ok"), formula_statuses[0] if formula_statuses else "not_required"),
            "formula_provider": next(iter(formula_providers), "not_required") if formula_providers else "not_required",
            "formula_blocks_created": len(formula_records),
            "formula_blocks": [serialize_formula_block(block) for block in formula_records],
            "warnings": parser_warnings,
            "knowledge_source": serialize_knowledge_source(source),
            "knowledge_chunks": [serialize_knowledge_chunk(chunk) for chunk in knowledge_chunks[:8]],
            "knowledge_base_version": serialize_kb_version(version),
            "cards": [serialize_terminology_card(card) for card in cards],
            "job": {
                "id": job.id,
                "status": job.status,
                "processed_pages": job.processed_pages
            }
        })
    except Exception as exc:
        document.parsing_status = "failed"
        document.error_message = str(exc)
        if document.ocr_required and not document.ocr_status:
            document.ocr_status = "parsing_failed"
            document.ocr_error = str(exc)
        job.status = "failed"
        job.error_message = str(exc)
        job.finished_at = current_time_text()
        add_system_log("error", "ingestion", f"Document {document.id} failed: {exc}")
        db.session.commit()
        return api_error_with_audit_context(
            "PARSING_FAILED",
            "文档解析失败。",
            422,
            audit_context,
            {"error": redact_for_log(exc), "document": serialize_document(document)},
        )


@app.route("/api/documents", methods=["GET"])
def list_documents():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    scope_type = request.args.get("scope_type", "").strip()
    query = Document.query
    if user.role == "student":
        course_ids = [item.course_id for item in CourseMember.query.filter_by(user_id=user.id).all()]
        query = query.filter(
            db.or_(
                db.and_(Document.scope_type == "personal", Document.owner_user_id == user.id),
                db.and_(Document.scope_type == "course", Document.course_id.in_(course_ids or [-1]))
            )
        )
    elif user.role == "teacher":
        manageable_ids = [
            course.id for course in Course.query.all()
            if can_manage_course(user, course)
        ]
        query = query.filter(
            db.or_(
                Document.owner_user_id == user.id,
                db.and_(Document.scope_type != "personal", Document.course_id.in_(manageable_ids or [-1]))
            )
        )
    if scope_type:
        query = query.filter_by(scope_type=scope_type)

    documents = query.order_by(Document.id.desc()).limit(100).all()
    return jsonify({
        "status": "success",
        "count": len(documents),
        "documents": [serialize_document(document) for document in documents]
    })


@app.route("/api/documents/<int:document_id>/chunks", methods=["GET"])
def list_document_chunks(document_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    document = db.session.get(Document, document_id)
    if document is None:
        return api_error("RESOURCE_NOT_FOUND", "文档不存在。", 404)
    course = db.session.get(Course, document.course_id) if document.course_id else None
    if document.scope_type == "personal":
        if user.role != "admin" and document.owner_user_id != user.id:
            return api_error("PERMISSION_DENIED", "无权查看该个人文档。", 403)
        if user.role == "admin" and document.owner_user_id != user.id:
            record_personal_access(user, document.owner_user_id, "document", document.id, "admin viewed personal document chunks")
    elif user.role == "student" and not is_course_member(user, document.course_id):
        return api_error("PERMISSION_DENIED", "无权查看该课程文档。", 403)
    elif user.role == "teacher" and not (document.owner_user_id == user.id or can_manage_course(user, course)):
        return api_error("PERMISSION_DENIED", "无权查看该课程文档。", 403)

    chunks = DocumentChunk.query.filter_by(document_id=document.id).order_by(DocumentChunk.id.asc()).all()
    formula_blocks = FormulaBlock.query.filter_by(document_id=document.id).order_by(FormulaBlock.id.asc()).all()
    db.session.commit()
    return jsonify({
        "status": "success",
        "document": serialize_document(document),
        "count": len(chunks),
        "chunks": [serialize_document_chunk(chunk) for chunk in chunks],
        "formula_blocks_count": len(formula_blocks),
        "formula_blocks": [serialize_formula_block(block) for block in formula_blocks]
    })


@app.route("/api/knowledge/sources", methods=["GET", "POST"])
def knowledge_sources():
    if request.method == "POST":
        user, error_response = require_current_user({"admin", "teacher"})
    else:
        user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    if request.method == "POST":
        data = request.get_json() or {}
        license_status = str(data.get("license_status", "unknown")).strip()
        allow_full_text = bool(data.get("allow_full_text_indexing", False))
        if license_status not in {"authorized", "open_licensed", "public_domain", "restricted", "unknown"}:
            return api_error("VALIDATION_ERROR", "license_status 不合法。", 400)
        if license_status not in {"authorized", "open_licensed", "public_domain"} and allow_full_text:
            return api_error("VALIDATION_ERROR", "未授权资料不能开启全文索引。", 400)

        source = KnowledgeSource(
            name=str(data.get("name", "")).strip(),
            language=str(data.get("language", "")).strip(),
            discipline=str(data.get("discipline", "")).strip(),
            source_type=str(data.get("source_type", "manual_upload")).strip(),
            access_method=str(data.get("access_method", "manual_upload")).strip(),
            license_status=license_status,
            update_frequency=str(data.get("update_frequency", "manual")).strip(),
            allow_full_text_indexing=allow_full_text,
            allow_student_search=bool(data.get("allow_student_search", False)),
            allow_derivative_cards=bool(data.get("allow_derivative_cards", False)),
            created_by=user.id,
            created_at=current_time_text(),
            updated_at=current_time_text()
        )
        if not source.name:
            return api_error("VALIDATION_ERROR", "知识源名称不能为空。", 400)
        db.session.add(source)
        db.session.commit()
        return jsonify({"status": "success", "source": serialize_knowledge_source(source)})

    query = KnowledgeSource.query
    if user.role == "student":
        query = query.filter_by(allow_student_search=True)
    sources = query.order_by(KnowledgeSource.id.desc()).all()
    return jsonify({
        "status": "success",
        "count": len(sources),
        "sources": [serialize_knowledge_source(source) for source in sources]
    })


@app.route("/api/knowledge/versions", methods=["GET", "POST"])
def knowledge_versions():
    if request.method == "POST":
        user, error_response = require_current_user({"admin", "teacher"})
    else:
        user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    if request.method == "POST":
        data = request.get_json() or {}
        course = get_course_by_id_or_name(data.get("course_id"), data.get("course_name", ""))
        kb_scope = str(data.get("scope_type", data.get("kb_scope", "course"))).strip()
        if kb_scope not in KB_SCOPES:
            return api_error("VALIDATION_ERROR", "kb_scope 不合法。", 400)
        if kb_scope == "course" and (course is None or not can_manage_course(user, course)):
            return api_error("PERMISSION_DENIED", "无权创建该课程知识库版本。", 403)
        if kb_scope == "global" and user.role != "admin":
            return api_error("PERMISSION_DENIED", "只有管理员可以创建平台级知识库版本。", 403)
        try:
            version = create_knowledge_base_version(
                course_id=course.id if course else None,
                scope_type=kb_scope,
                owner_user_id=user.id if kb_scope == "personal" else None,
                description=str(data.get("description", "")).strip(),
                parent_version_id=data.get("parent_version_id"),
                created_by=user.id,
            )
        except ValueError as exc:
            return api_error("VALIDATION_ERROR", str(exc), 400)
        requested_name = str(data.get("version_name", "")).strip()
        if requested_name:
            version.version_name = requested_name
        db.session.commit()
        return jsonify({"status": "success", "version": serialize_kb_version(version)})

    query = KnowledgeBaseVersion.query
    if user.role == "student":
        course_ids = [item.course_id for item in CourseMember.query.filter_by(user_id=user.id).all()]
        query = query.filter(db.or_(KnowledgeBaseVersion.kb_scope == "global", KnowledgeBaseVersion.course_id.in_(course_ids or [-1])))
    versions = query.order_by(KnowledgeBaseVersion.id.desc()).all()
    return jsonify({
        "status": "success",
        "count": len(versions),
        "versions": [serialize_kb_version(version) for version in versions]
    })


@app.route("/api/knowledge/versions/<int:version_id>", methods=["GET"])
def knowledge_version_detail(version_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    version = db.session.get(KnowledgeBaseVersion, version_id)
    if version is None:
        return api_error("RESOURCE_NOT_FOUND", "Knowledge base version not found.", 404)
    if version.scope_type == "course":
        course = db.session.get(Course, version.course_id)
        if user.role == "student" and not is_course_member(user, version.course_id):
            return api_error("PERMISSION_DENIED", "无权查看该课程知识库版本。", 403)
        if user.role == "teacher" and not can_manage_course(user, course):
            return api_error("PERMISSION_DENIED", "无权查看该课程知识库版本。", 403)
    if version.scope_type == "personal" and user.role != "admin" and version.owner_user_id != user.id:
        return api_error("PERMISSION_DENIED", "无权查看该个人知识库版本。", 403)
    return api_success({"version": serialize_kb_version(version), "manifest": build_kb_version_manifest(version)})


@app.route("/api/knowledge/versions/<int:version_id>/publish", methods=["POST"])
def knowledge_version_publish(version_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    version = db.session.get(KnowledgeBaseVersion, version_id)
    if version is None:
        return api_error("RESOURCE_NOT_FOUND", "Knowledge base version not found.", 404)
    if version.scope_type == "course":
        course = db.session.get(Course, version.course_id)
        if not can_manage_course(user, course):
            return api_error("PERMISSION_DENIED", "无权发布该课程知识库版本。", 403)
    elif version.scope_type == "personal" and not (user.role == "admin" or version.owner_user_id == user.id):
        return api_error("PERMISSION_DENIED", "无权发布该个人知识库版本。", 403)
    elif version.scope_type == "global" and user.role != "admin":
        return api_error("PERMISSION_DENIED", "只有管理员可以发布全局知识库。", 403)
    result = publish_kb_version(version_id, actor_user_id=user.id)
    db.session.commit()
    if result["status"] != "success":
        return api_error("VALIDATION_ERROR", result["message"], 400, result)
    return api_success(result, "Knowledge base version published.")


@app.route("/api/knowledge/versions/<int:version_id>/rollback", methods=["POST"])
def knowledge_version_rollback(version_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    version = db.session.get(KnowledgeBaseVersion, version_id)
    if version is None:
        return api_error("RESOURCE_NOT_FOUND", "Knowledge base version not found.", 404)
    course = db.session.get(Course, version.course_id)
    if version.scope_type != "course" or not can_manage_course(user, course):
        return api_error("PERMISSION_DENIED", "无权回滚该课程知识库。", 403)
    result = rollback_kb_version(version.course_id, version.id, actor_user_id=user.id)
    db.session.commit()
    if result["status"] != "success":
        return api_error("VALIDATION_ERROR", result["message"], 400, result)
    return api_success(result, "Knowledge base version rolled back.")


@app.route("/api/knowledge/versions/<int:version_id>/rebuild", methods=["POST"])
def knowledge_version_rebuild(version_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    version = db.session.get(KnowledgeBaseVersion, version_id)
    if version is None:
        return api_error("RESOURCE_NOT_FOUND", "Knowledge base version not found.", 404)
    course = db.session.get(Course, version.course_id) if version.course_id else None
    if version.scope_type == "course" and not can_manage_course(user, course):
        return api_error("PERMISSION_DENIED", "无权重建该课程知识库。", 403)
    data = request.get_json(silent=True) or {}
    dry_run = bool(data.get("dry_run", True))
    if dry_run:
        documents = Document.query.filter_by(course_id=version.course_id, scope_type=version.scope_type).count()
        return api_success({"dry_run": True, "document_count": documents, "target_version_id": version.id})
    documents = Document.query.filter_by(course_id=version.course_id, scope_type=version.scope_type).all()
    reports = []
    for document in documents:
        reports.append(index_document_into_kb_version(document.id, version.id))
    db.session.commit()
    return api_success({"dry_run": False, "reports": reports, "version": serialize_kb_version(version)})


@app.route("/api/knowledge/versions/<int:version_id>/manifest", methods=["GET"])
def knowledge_version_manifest(version_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    version = db.session.get(KnowledgeBaseVersion, version_id)
    if version is None:
        return api_error("RESOURCE_NOT_FOUND", "Knowledge base version not found.", 404)
    return api_success(build_kb_version_manifest(version))


@app.route("/api/knowledge/health", methods=["GET"])
def knowledge_health_api():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    course_id = request.args.get("course_id", type=int)
    version_id = request.args.get("kb_version_id", type=int)
    if course_id and user.role == "teacher":
        course = db.session.get(Course, course_id)
        if not can_manage_course(user, course):
            return api_error("PERMISSION_DENIED", "无权查看该课程知识库健康状态。", 403)
    health = run_knowledge_health_check(course_id=course_id, kb_version_id=version_id)
    return api_success({"health": health})


@app.route("/api/knowledge/retrieval-regression", methods=["POST"])
def knowledge_retrieval_regression_api():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    data = request.get_json() or {}
    course_id = data.get("course_id")
    if course_id and user.role == "teacher":
        course = db.session.get(Course, int(course_id))
        if not can_manage_course(user, course):
            return api_error("PERMISSION_DENIED", "无权运行该课程检索回归。", 403)
    result = run_retrieval_regression_for_course(course_id=int(course_id) if course_id else None, kb_version_id=data.get("kb_version_id"))
    return api_success({"regression": result})


def _can_access_kb_version_for_retrieval(user, version):
    if version is None:
        return False
    if user.role == "admin":
        return True
    if version.scope_type == "course":
        return can_manage_course(user, db.session.get(Course, version.course_id))
    if version.scope_type == "personal":
        return version.owner_user_id == user.id
    return False


def serialize_retrieval_experiment(run):
    return {
        "id": run.id,
        "course_id": run.course_id,
        "evaluation_set_id": run.evaluation_set_id,
        "kb_version_id": run.kb_version_id,
        "experiment_name": run.experiment_name,
        "backends_tested": safe_json_loads(run.backends_tested_json, []),
        "best_backend": run.best_backend,
        "recommendation": run.recommendation,
        "metrics": safe_json_loads(run.metrics_json, {}),
        "created_by": run.created_by,
        "created_at": run.created_at,
        "finished_at": run.finished_at,
        "status": run.status,
        "report_markdown": run.report_markdown,
    }


@app.route("/api/admin/retrieval/vector-index/build", methods=["POST"])
def admin_retrieval_vector_index_build():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    data = request.get_json() or {}
    kb_version_id = data.get("kb_version_id")
    if not kb_version_id:
        return api_error("VALIDATION_ERROR", "kb_version_id is required.", 400)
    version = db.session.get(KnowledgeBaseVersion, int(kb_version_id))
    if not _can_access_kb_version_for_retrieval(user, version):
        return api_error("PERMISSION_DENIED", "无权构建该知识库版本的向量索引。", 403)
    result = build_vector_index_for_kb_version(
        int(kb_version_id),
        apply=bool(data.get("apply", False)),
        embedding_provider_name=data.get("embedding_provider"),
        vector_backend_name=data.get("vector_index_backend"),
    )
    if result.get("status") == "error":
        return api_error(result.get("error_code", "VALIDATION_ERROR"), result.get("message", "Vector index build failed."), 400, result)
    return api_success({"vector_index": result})


@app.route("/api/admin/retrieval/experiments/run", methods=["POST"])
def admin_retrieval_experiment_run():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    data = request.get_json() or {}
    course_id = data.get("course_id")
    kb_version_id = data.get("kb_version_id")
    if course_id and user.role == "teacher" and not can_manage_course(user, db.session.get(Course, int(course_id))):
        return api_error("PERMISSION_DENIED", "无权运行该课程检索实验。", 403)
    if kb_version_id:
        version = db.session.get(KnowledgeBaseVersion, int(kb_version_id))
        if not _can_access_kb_version_for_retrieval(user, version):
            return api_error("PERMISSION_DENIED", "无权运行该知识库版本检索实验。", 403)
    result = run_retrieval_experiment(
        course_id=int(course_id) if course_id else None,
        evaluation_set_id=data.get("evaluation_set_id"),
        kb_version_id=kb_version_id,
        created_by=user.id,
    )
    if result.get("status") == "failed":
        return api_error("VALIDATION_ERROR", result.get("message", "Retrieval experiment failed."), 400, result)
    return api_success({"experiment": result})


@app.route("/api/admin/retrieval/experiments", methods=["GET"])
def admin_retrieval_experiments_list():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    query = RetrievalExperimentRun.query
    if user.role == "teacher":
        manageable_ids = [course.id for course in Course.query.all() if can_manage_course(user, course)]
        query = query.filter(RetrievalExperimentRun.course_id.in_(manageable_ids or [-1]))
    runs = query.order_by(RetrievalExperimentRun.id.desc()).limit(100).all()
    return api_success({"items": [serialize_retrieval_experiment(run) for run in runs]})


@app.route("/api/admin/retrieval/experiments/<int:experiment_id>", methods=["GET"])
def admin_retrieval_experiment_detail(experiment_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    run = db.session.get(RetrievalExperimentRun, experiment_id)
    if run is None:
        return api_error("RESOURCE_NOT_FOUND", "Retrieval experiment not found.", 404)
    if user.role == "teacher" and not can_manage_course(user, db.session.get(Course, run.course_id)):
        return api_error("PERMISSION_DENIED", "无权查看该检索实验。", 403)
    return api_success({"experiment": serialize_retrieval_experiment(run)})


@app.route("/api/admin/retrieval/health", methods=["GET"])
def admin_retrieval_health():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    kb_version_id = request.args.get("kb_version_id", type=int)
    if kb_version_id:
        version = db.session.get(KnowledgeBaseVersion, kb_version_id)
        if not _can_access_kb_version_for_retrieval(user, version):
            return api_error("PERMISSION_DENIED", "无权查看该索引健康状态。", 403)
    return api_success({
        "retrieval_backend": RETRIEVAL_BACKEND,
        "embedding_provider": EMBEDDING_PROVIDER,
        "vector_index_backend": VECTOR_INDEX_BACKEND,
        "vector_index_health": vector_index_health(kb_version_id),
    })


@app.route("/api/alignment/run", methods=["POST"])
def run_alignment():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    sync_requested = str(request.args.get("sync", "")).strip().lower() in {"1", "true", "yes"}
    data = request.get_json() or {}
    scope_type = str(data.get("scope_type", "course")).strip()
    document_id = data.get("document_id")
    english_term = str(data.get("english_term", "")).strip()
    course = get_course_by_id_or_name(data.get("course_id"), data.get("course_name", ""))

    if scope_type not in {"course", "personal"}:
        return api_error("VALIDATION_ERROR", "alignment scope_type 只能是 course 或 personal。", 400)
    if scope_type == "course" and (course is None or not can_manage_course(user, course)):
        return api_error("PERMISSION_DENIED", "无权运行该课程对齐。", 403)
    if scope_type == "personal" and user.role == "student":
        allowed, reasons, plan, subscription, totals = check_usage_quota(user, ai_units=1)
        if not allowed:
            return api_error("QUOTA_EXCEEDED", "AI 对齐额度不足。", 402, {"reasons": reasons})

    if not sync_requested:
        document = None
        if document_id:
            document = db.session.get(Document, int(document_id))
            if document is None:
                return api_error("RESOURCE_NOT_FOUND", "文档不存在。", 404)
            if document.scope_type == "personal" and document.owner_user_id != user.id and user.role != "admin":
                return api_error("PERMISSION_DENIED", "无权对齐该个人文档。", 403)
            if document.scope_type == "course" and document.course_id:
                doc_course = db.session.get(Course, document.course_id)
                if doc_course and not can_manage_course(user, doc_course):
                    return api_error("PERMISSION_DENIED", "无权对齐该课程文档。", 403)
            if course is None and document.course_id:
                course = db.session.get(Course, document.course_id)
        elif not english_term:
            return api_error("VALIDATION_ERROR", "english_term 或 document_id 必须提供一个。", 400)

        meta = current_provider_metadata()
        alignment_run = AlignmentRun(
            document_id=document.id if document else None,
            course_id=course.id if course else None,
            triggered_by=user.id,
            provider=meta["provider"],
            model_name=meta["model_name"],
            ai_provider=meta["provider"],
            ai_provider_mode=meta.get("provider_mode", ""),
            ai_model=meta["model_name"],
            prompt_key="term_alignment",
            prompt_version="v1",
            retrieval_version=meta["retrieval_version"],
            term_count=1 if english_term else 0,
            status="queued",
            started_at="",
        )
        db.session.add(alignment_run)
        db.session.flush()
        background_job = create_background_job(
            "alignment_run",
            user,
            course_id=course.id if course else None,
            document_id=document.id if document else None,
            alignment_run_id=alignment_run.id,
            scope_type=document.scope_type if document else scope_type,
            owner_user_id=(document.owner_user_id if document and document.scope_type == "personal" else (user.id if scope_type == "personal" else None)),
            input_data={
                "document_id": document.id if document else None,
                "english_term": english_term,
                "courseware_sentence": str(data.get("courseware_sentence", "")).strip(),
                "chapter": str(data.get("chapter", "")).strip(),
                "scope_type": document.scope_type if document else scope_type,
                "course_id": course.id if course else None,
            }
        )
        db.session.commit()
        return api_success({
            "alignment_run_id": alignment_run.id,
            "job_id": background_job.id,
            "job_type": background_job.job_type,
            "job_status": background_job.status,
            "run": serialize_alignment_run(alignment_run),
            "job": serialize_background_job(background_job),
        }, "术语对齐任务已进入后台队列。")

    if document_id:
        document = db.session.get(Document, int(document_id))
        if document is None:
            return api_error("RESOURCE_NOT_FOUND", "文档不存在。", 404)
        if document.scope_type == "personal" and document.owner_user_id != user.id and user.role != "admin":
            return api_error("PERMISSION_DENIED", "无权对齐该个人文档。", 403)
        chunks = DocumentChunk.query.filter_by(document_id=document.id).all()
        cards = run_alignment_for_chunks(
            chunks,
            course=course or (db.session.get(Course, document.course_id) if document.course_id else None),
            scope_type=document.scope_type if document.scope_type == "personal" else "course",
            owner_user_id=document.owner_user_id if document.scope_type == "personal" else None,
            source_document_id=document.id,
            triggered_by_user_id=user.id
        )
        db.session.commit()
        return jsonify({
            "status": "success",
            "message": "文档术语对齐已完成。",
            "cards": [serialize_terminology_card(card) for card in cards]
        })

    if not english_term:
        return api_error("VALIDATION_ERROR", "english_term 或 document_id 必须提供一个。", 400)

    course_name = course.name if course else ""
    meta = current_provider_metadata()
    alignment_run = AlignmentRun(
        document_id=None,
        course_id=course.id if course else None,
        triggered_by=user.id,
        provider=meta["provider"],
        model_name=meta["model_name"],
        ai_provider=meta["provider"],
        ai_provider_mode=meta.get("provider_mode", ""),
        ai_model=meta["model_name"],
        prompt_key="term_alignment",
        prompt_version="v1",
        retrieval_version=RETRIEVAL_VERSION,
        terms_extracted=1,
        term_count=1,
        status="running",
        started_at=current_time_text()
    )
    db.session.add(alignment_run)
    db.session.flush()
    alignment = generate_alignment_result(
        english_term=english_term,
        courseware_sentence=str(data.get("courseware_sentence", "")).strip(),
        course=course_name,
        chapter=str(data.get("chapter", "")).strip(),
        scope_type=scope_type,
        owner_user_id=user.id if scope_type == "personal" else None
    )
    alignment["alignment_run_id"] = alignment_run.id
    card = create_or_update_card_from_alignment(
        english_term=english_term,
        alignment=alignment,
        scope_type=scope_type,
        course_id=course.id if course else None,
        owner_user_id=user.id if scope_type == "personal" else None,
        courseware_sentence=str(data.get("courseware_sentence", "")).strip()
    )
    if scope_type == "personal" and card:
        record_usage(user.id, "ai_alignment", 1, related_term_id=card.id)
    update_alignment_run_stats(alignment_run, cards=[card] if card else [], term_count=1)
    alignment_run.status = "completed"
    alignment_run.finished_at = current_time_text()
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "术语证据对齐已完成。",
        "alignment": alignment,
        "card": serialize_terminology_card(card)
    })


@app.route("/api/alignment/runs/<int:run_id>", methods=["GET"])
def alignment_run_detail(run_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    run = db.session.get(AlignmentRun, run_id)
    if run is None:
        return api_error("RESOURCE_NOT_FOUND", "AlignmentRun 不存在。", 404)
    if user.role == "admin":
        return jsonify({"status": "success", "run": serialize_alignment_run(run)})
    if user.role == "teacher":
        course = db.session.get(Course, run.course_id) if run.course_id else None
        if course and can_manage_course(user, course):
            return jsonify({"status": "success", "run": serialize_alignment_run(run)})
    if run.triggered_by == user.id:
        return jsonify({"status": "success", "run": serialize_alignment_run(run)})
    return api_error("PERMISSION_DENIED", "无权查看该对齐运行记录。", 403)


@app.route("/api/alignment/runs", methods=["GET"])
def alignment_runs():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    page = max(1, int(request.args.get("page", "1") or 1))
    page_size = max(1, min(int(request.args.get("page_size", "20") or 20), 100))
    query = AlignmentRun.query
    if user.role == "student":
        query = query.filter_by(triggered_by=user.id)
    elif user.role == "teacher":
        manageable_ids = [
            course.id for course in Course.query.all()
            if can_manage_course(user, course)
        ]
        query = query.filter(db.or_(
            AlignmentRun.triggered_by == user.id,
            AlignmentRun.course_id.in_(manageable_ids or [-1])
        ))
    total = query.count()
    runs = query.order_by(AlignmentRun.id.desc()).offset((page - 1) * page_size).limit(page_size).all()
    return api_success({
        "items": [serialize_alignment_run(run) for run in runs],
        "pagination": {
            "page": page,
            "page_size": page_size,
            "total": total,
            "has_next": page * page_size < total,
        }
    })


def get_route_audit_context(user=None):
    if not hasattr(g, "lexibridge_request_id"):
        g.lexibridge_request_id = audit_context_service.get_request_id(request)
    return audit_context_service.build_audit_context_from_request(
        request,
        user,
        source="api",
        request_id=g.lexibridge_request_id,
    )


def api_success_with_audit_context(data=None, message="Operation completed.", audit_context=None):
    audit_context = audit_context_service.normalize_audit_context(audit_context)
    return api_success(data, message, request_id=audit_context.get("request_id", ""))


def api_error_with_audit_context(error_code, message, http_status=None, audit_context=None, details=None):
    audit_context = audit_context_service.normalize_audit_context(audit_context)
    details = dict(details or {})
    details["request_id"] = audit_context.get("request_id", "")
    return api_error(error_code, message, http_status, details)


def attach_request_id_to_response(response_or_tuple, audit_context):
    audit_context = audit_context_service.normalize_audit_context(audit_context)
    request_id = audit_context.get("request_id", "")
    if isinstance(response_or_tuple, tuple):
        response = response_or_tuple[0]
        status_code = response_or_tuple[1]
    else:
        response = response_or_tuple
        status_code = getattr(response, "status_code", 200)
    payload = response.get_json(silent=True) if hasattr(response, "get_json") else None
    if isinstance(payload, dict):
        payload.setdefault("request_id", request_id)
        return jsonify(payload), status_code
    return response_or_tuple


def record_concept_card_api_failure(target_uid, error, audit_context, input_payload=None, error_code=None):
    try:
        audit_record_service.record_concept_card_operation_failed(
            db.session,
            AuditRecord,
            target_uid=target_uid,
            error=error,
            error_code=error_code or concept_card_service.classify_concept_card_error(error),
            input_payload=input_payload or {},
            audit_context=audit_context,
            source="api",
            now_fn=current_time_text,
            commit=True,
        )
    except Exception:
        db.session.rollback()


def record_parse_quality_risk_audit(
    event_type,
    target_type,
    target_uid,
    *,
    parse_uid="",
    quality_status="",
    risk_labels=None,
    forced_status="",
    blocked_reason="",
    audit_context=None,
    commit=True,
):
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": target_type,
                "target_uid": str(target_uid or ""),
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "service",
                "input_payload": {
                    "parse_uid": parse_uid,
                    "quality_status": quality_status,
                },
                "output_payload": {
                    "risk_labels": risk_labels or [],
                    "forced_status": forced_status,
                    "blocked_reason": blocked_reason,
                },
                "changed_fields": ["risk_labels", "status"] if forced_status else ["risk_labels"],
                "result": "error" if blocked_reason else "success",
                "error_code": event_type if blocked_reason else "",
                "error_message": blocked_reason,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def concept_card_error_response(exc, audit_context=None):
    details = {
        "reason": str(exc),
        "audit_error_code": concept_card_service.classify_concept_card_error(exc),
    }
    if isinstance(exc, concept_card_service.ConceptCardNotFoundError):
        return api_error_with_audit_context("RESOURCE_NOT_FOUND", str(exc), 404, audit_context, details)
    return api_error_with_audit_context("VALIDATION_ERROR", str(exc), 400, audit_context, details)


def course_review_policy_error_response(exc, audit_context=None):
    reason = getattr(exc, "reason", "course_review_policy_error")
    if reason in {"policy_not_found", "permission_not_found"}:
        return api_error_with_audit_context("RESOURCE_NOT_FOUND", str(exc), 404, audit_context, {"audit_error_code": reason})
    return api_error_with_audit_context("VALIDATION_ERROR", str(exc), 400, audit_context, {"audit_error_code": reason})


def record_course_review_governance_audit(
    event_type,
    *,
    policy=None,
    permission=None,
    input_data=None,
    audit_context=None,
    error_code="",
    error_message="",
    commit=True,
):
    input_data = dict(input_data or {})
    policy_data = course_review_policy_service.serialize_course_review_policy(policy) if policy is not None else {}
    permission_data = course_review_policy_service.serialize_course_review_permission(permission) if permission is not None else {}
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "course_review_policy",
                "target_uid": policy_data.get("policy_uid") or permission_data.get("permission_uid", ""),
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "course": input_data.get("course") or policy_data.get("course") or permission_data.get("course", ""),
                    "chapter": input_data.get("chapter") or policy_data.get("chapter") or permission_data.get("chapter", ""),
                    "reviewer_id": input_data.get("reviewer_id") or permission_data.get("reviewer_id"),
                    "reviewer_role": input_data.get("reviewer_role") or permission_data.get("reviewer_role", ""),
                },
                "output_payload": {
                    "policy_uid": policy_data.get("policy_uid", ""),
                    "permission_uid": permission_data.get("permission_uid", ""),
                    "course": policy_data.get("course") or permission_data.get("course", ""),
                    "status": policy_data.get("status") or permission_data.get("status", ""),
                    "permission_level": permission_data.get("permission_level", ""),
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


@app.route("/api/review-policies", methods=["GET"])
def list_course_review_policies_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    result = course_review_policy_service.list_course_review_policies(
        db.session,
        CourseReviewPolicy,
        request.args.to_dict(),
    )
    return api_success_with_audit_context(
        {
            "items": [course_review_policy_service.serialize_course_review_policy(item) for item in result.items],
            "pagination": result.pagination,
        },
        audit_context=audit_context,
    )


@app.route("/api/review-policies/<policy_uid>", methods=["GET"])
def get_course_review_policy_api(policy_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    policy = course_review_policy_service.get_course_review_policy_by_uid(db.session, CourseReviewPolicy, policy_uid)
    if policy is None:
        return api_error_with_audit_context(
            "RESOURCE_NOT_FOUND",
            "CourseReviewPolicy not found.",
            404,
            audit_context,
            {"audit_error_code": "policy_not_found"},
        )
    return api_success_with_audit_context(
        {"policy": course_review_policy_service.serialize_course_review_policy(policy)},
        audit_context=audit_context,
    )


@app.route("/api/review-policies", methods=["POST"])
def upsert_course_review_policy_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json(silent=True) or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    try:
        policy, created = course_review_policy_service.create_or_update_course_review_policy(
            db.session,
            CourseReviewPolicy,
            data.get("course", ""),
            data,
            actor=user,
            now_fn=current_time_text,
            commit=True,
        )
    except course_review_policy_service.CourseReviewPolicyError as exc:
        db.session.rollback()
        return course_review_policy_error_response(exc, audit_context)
    event_type = "course_review_policy_created" if created else "course_review_policy_updated"
    record_course_review_governance_audit(
        event_type,
        policy=policy,
        input_data=data,
        audit_context=audit_context,
    )
    return api_success_with_audit_context(
        {"policy": course_review_policy_service.serialize_course_review_policy(policy), "created": created},
        audit_context=audit_context,
    )


@app.route("/api/review-permissions", methods=["GET"])
def list_course_review_permissions_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    result = course_review_policy_service.list_course_review_permissions(
        db.session,
        CourseReviewPermission,
        request.args.to_dict(),
    )
    return api_success_with_audit_context(
        {
            "items": [course_review_policy_service.serialize_course_review_permission(item) for item in result.items],
            "pagination": result.pagination,
        },
        audit_context=audit_context,
    )


@app.route("/api/review-permissions", methods=["POST"])
def grant_course_review_permission_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json(silent=True) or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    try:
        permission, created = course_review_policy_service.grant_course_review_permission(
            db.session,
            CourseReviewPermission,
            data.get("course", ""),
            data.get("reviewer_id"),
            data,
            actor=user,
            now_fn=current_time_text,
            commit=True,
        )
    except course_review_policy_service.CourseReviewPolicyError as exc:
        db.session.rollback()
        return course_review_policy_error_response(exc, audit_context)
    record_course_review_governance_audit(
        "course_review_permission_granted",
        permission=permission,
        input_data=data,
        audit_context=audit_context,
    )
    return api_success_with_audit_context(
        {"permission": course_review_policy_service.serialize_course_review_permission(permission), "created": created},
        audit_context=audit_context,
    )


@app.route("/api/review-permissions/<permission_uid>/revoke", methods=["POST"])
def revoke_course_review_permission_api(permission_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        permission = course_review_policy_service.revoke_course_review_permission(
            db.session,
            CourseReviewPermission,
            permission_uid,
            actor=user,
            now_fn=current_time_text,
            commit=True,
        )
    except course_review_policy_service.CourseReviewPolicyError as exc:
        db.session.rollback()
        return course_review_policy_error_response(exc, audit_context)
    record_course_review_governance_audit(
        "course_review_permission_revoked",
        permission=permission,
        input_data={"permission_uid": permission_uid},
        audit_context=audit_context,
    )
    return api_success_with_audit_context(
        {"permission": course_review_policy_service.serialize_course_review_permission(permission)},
        audit_context=audit_context,
    )


def record_student_course_access_audit(
    event_type,
    *,
    course="",
    card=None,
    membership=None,
    policy=None,
    user_id=None,
    denied_reason="",
    audit_context=None,
    result="success",
):
    audit_record_service.create_audit_record(
        db.session,
        AuditRecord,
        {
            "event_type": event_type,
            "target_type": "student_course_access",
            "target_uid": getattr(card, "card_uid", "") or getattr(membership, "membership_uid", "") or getattr(policy, "policy_uid", ""),
            "result": result,
            "input_payload": {
                "user_id": user_id,
                "course": course or getattr(card, "course", "") or getattr(membership, "course", "") or getattr(policy, "course", ""),
                "card_uid": getattr(card, "card_uid", ""),
                "policy_uid": getattr(policy, "policy_uid", ""),
                "membership_uid": getattr(membership, "membership_uid", ""),
                "denied_reason": denied_reason,
            },
            "output_payload": {},
            "changed_fields": [],
        },
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=False,
    )


def student_course_access_error_response(exc, audit_context):
    reason = getattr(exc, "reason", "student_course_access_error")
    status = 404 if reason in {"membership_not_found"} else 400
    error_code = "RESOURCE_NOT_FOUND" if status == 404 else "VALIDATION_ERROR"
    return api_error_with_audit_context(
        error_code,
        str(exc),
        status,
        audit_context,
        {"audit_error_code": reason},
    )


def student_visible_course_names(user):
    known_courses = [
        row[0]
        for row in db.session.query(ConceptAlignmentCard.course)
        .filter(ConceptAlignmentCard.status == "approved")
        .distinct()
        .all()
        if str(row[0] or "").strip()
    ]
    return student_course_access_service.visible_courses_for_user(
        db.session,
        StudentCourseMembership,
        CourseStudentVisibilityPolicy,
        user,
        known_courses=known_courses,
    )


def record_student_learning_audit(event_type, *, user_id=None, course="", audit_context=None, result="success"):
    audit_record_service.create_audit_record(
        db.session,
        AuditRecord,
        {
            "event_type": event_type,
            "target_type": "student_learning_progress",
            "target_uid": str(user_id or ""),
            "result": result,
            "input_payload": {
                "user_id": user_id,
                "course": course,
            },
            "output_payload": {},
            "changed_fields": [],
        },
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=False,
    )


@app.route("/api/student/courses", methods=["GET"])
def list_student_visible_courses_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    courses = student_visible_course_names(user)
    items = []
    for course in courses:
        decision = student_course_access_service.can_student_view_course(
            db.session,
            StudentCourseMembership,
            CourseStudentVisibilityPolicy,
            user,
            course,
        )
        policy = student_course_access_service.get_course_student_visibility_policy(
            db.session,
            CourseStudentVisibilityPolicy,
            course,
        )
        items.append({
            "course": course,
            "visibility": decision.visibility,
            "membership_status": "active" if decision.membership_uid else "",
            "membership_uid": decision.membership_uid,
            "role_in_course": decision.role_in_course,
            "access_reason": decision.reason,
            "policy": student_course_access_service.serialize_course_student_visibility_policy(policy, course),
        })
    return api_success_with_audit_context(
        {"items": items, "total": len(items)},
        audit_context=audit_context,
    )


@app.route("/api/student/progress", methods=["GET"])
def get_student_progress_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    filters = request.args.to_dict()
    result = student_learning_progress_service.get_student_progress(
        db.session,
        ConceptAlignmentCard,
        StudentConceptCardState,
        Feedback,
        StudentCourseMembership,
        CourseStudentVisibilityPolicy,
        user,
        filters=filters,
    )
    record_student_learning_audit(
        "student_learning_progress_viewed",
        user_id=user.id,
        course=str(filters.get("course") or "").strip(),
        audit_context=audit_context,
    )
    db.session.commit()
    return api_success_with_audit_context(
        student_learning_progress_service.serialize_student_progress_summary(result),
        audit_context=audit_context,
    )


@app.route("/api/student/course-memberships", methods=["GET"])
def list_student_course_memberships_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    filters = request.args.to_dict()
    requested_user_id = filters.get("user_id")
    try:
        user_id = int(requested_user_id) if user.role == "admin" and requested_user_id else user.id
    except (TypeError, ValueError):
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            "user_id must be an integer.",
            400,
            audit_context,
            {"audit_error_code": "invalid_user_id"},
        )
    result = student_course_access_service.get_student_course_memberships(
        db.session,
        StudentCourseMembership,
        user_id=user_id,
        status=filters.get("status", "active"),
        filters=filters,
    )
    return api_success_with_audit_context(
        {
            "items": [student_course_access_service.serialize_student_course_membership(item) for item in result.items],
            "pagination": result.pagination,
        },
        audit_context=audit_context,
    )


@app.route("/api/student/course-memberships", methods=["POST"])
def create_student_course_membership_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    data = request.get_json(silent=True) or {}
    if not isinstance(data, dict):
        data = {}
    try:
        membership = student_course_access_service.add_student_course_membership(
            db.session,
            StudentCourseMembership,
            data.get("user_id"),
            data.get("course", ""),
            data,
            actor=user,
            now_fn=current_time_text,
            commit=False,
        )
        record_student_course_access_audit(
            "student_course_membership_created",
            membership=membership,
            user_id=membership.user_id,
            audit_context=audit_context,
        )
        db.session.commit()
    except student_course_access_service.StudentCourseAccessError as exc:
        db.session.rollback()
        return student_course_access_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"membership": student_course_access_service.serialize_student_course_membership(membership)},
        audit_context=audit_context,
    )


@app.route("/api/student/course-memberships/<membership_uid>/revoke", methods=["POST"])
def revoke_student_course_membership_api(membership_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        membership = student_course_access_service.revoke_student_course_membership(
            db.session,
            StudentCourseMembership,
            membership_uid,
            actor=user,
            now_fn=current_time_text,
            commit=False,
        )
        record_student_course_access_audit(
            "student_course_membership_revoked",
            membership=membership,
            user_id=membership.user_id,
            audit_context=audit_context,
        )
        db.session.commit()
    except student_course_access_service.StudentCourseAccessError as exc:
        db.session.rollback()
        return student_course_access_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"membership": student_course_access_service.serialize_student_course_membership(membership)},
        audit_context=audit_context,
    )


@app.route("/api/course-student-visibility-policies", methods=["GET"])
def list_course_student_visibility_policies_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    filters = request.args.to_dict()
    try:
        page = max(1, int(filters.get("page") or 1))
        per_page = min(100, max(1, int(filters.get("per_page") or 50)))
    except (TypeError, ValueError):
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            "page and per_page must be integers.",
            400,
            audit_context,
            {"audit_error_code": "invalid_pagination"},
        )
    query = CourseStudentVisibilityPolicy.query
    if str(filters.get("course") or "").strip():
        query = query.filter_by(course=str(filters.get("course")).strip())
    if str(filters.get("status") or "").strip():
        query = query.filter_by(status=str(filters.get("status")).strip())
    total = query.count()
    items = (
        query.order_by(CourseStudentVisibilityPolicy.course.asc(), CourseStudentVisibilityPolicy.id.asc())
        .offset((page - 1) * per_page)
        .limit(per_page)
        .all()
    )
    return api_success_with_audit_context(
        {
            "items": [student_course_access_service.serialize_course_student_visibility_policy(item) for item in items],
            "pagination": {
                "page": page,
                "per_page": per_page,
                "total": total,
                "has_next": page * per_page < total,
            },
        },
        audit_context=audit_context,
    )


@app.route("/api/course-student-visibility-policies", methods=["POST"])
def upsert_course_student_visibility_policy_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    data = request.get_json(silent=True) or {}
    if not isinstance(data, dict):
        data = {}
    try:
        policy, created = student_course_access_service.create_or_update_course_student_visibility_policy(
            db.session,
            CourseStudentVisibilityPolicy,
            data.get("course", ""),
            data,
            actor=user,
            now_fn=current_time_text,
            commit=False,
        )
        record_student_course_access_audit(
            "course_student_visibility_policy_created" if created else "course_student_visibility_policy_updated",
            policy=policy,
            audit_context=audit_context,
        )
        db.session.commit()
    except student_course_access_service.StudentCourseAccessError as exc:
        db.session.rollback()
        return student_course_access_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {
            "policy": student_course_access_service.serialize_course_student_visibility_policy(policy),
            "created": created,
        },
        audit_context=audit_context,
    )


route_core = RouteCoreDependencies(
    db=db,
    audit_record_model=AuditRecord,
    audit_record_service=audit_record_service,
    current_time_text=current_time_text,
    require_current_user=require_current_user,
    get_route_audit_context=get_route_audit_context,
    attach_request_id_to_response=attach_request_id_to_response,
    api_success_with_audit_context=api_success_with_audit_context,
    api_error_with_audit_context=api_error_with_audit_context,
)


register_student_concept_card_routes(
    app,
    core=route_core,
    models=StudentConceptCardModels(
        ConceptAlignmentCard=ConceptAlignmentCard,
        StudentConceptCardState=StudentConceptCardState,
        Feedback=Feedback,
        StudentCourseMembership=StudentCourseMembership,
        CourseStudentVisibilityPolicy=CourseStudentVisibilityPolicy,
    ),
    student_visible_course_names=student_visible_course_names,
    student_course_access_service=student_course_access_service,
    record_student_course_access_audit=record_student_course_access_audit,
)


register_concept_card_feedback_routes(
    app,
    core=route_core,
    models=ConceptCardFeedbackModels(
        Feedback=Feedback,
        ConceptAlignmentCard=ConceptAlignmentCard,
        ConceptCardReviewRecord=ConceptCardReviewRecord,
        ConceptCardFeedbackTriageRecord=ConceptCardFeedbackTriageRecord,
        CourseReviewPermission=CourseReviewPermission,
        CourseReviewPolicy=CourseReviewPolicy,
    ),
)


register_teacher_learning_analytics_routes(
    app,
    core=route_core,
    models=TeacherLearningAnalyticsModels(
        ConceptAlignmentCard=ConceptAlignmentCard,
        StudentConceptCardState=StudentConceptCardState,
        Feedback=Feedback,
        StudentCourseMembership=StudentCourseMembership,
        CourseReviewPermission=CourseReviewPermission,
        CourseStudentVisibilityPolicy=CourseStudentVisibilityPolicy,
    ),
)


register_provider_governance_routes(
    app,
    core=route_core,
    models=ProviderGovernanceModels(
        AlignmentProviderPolicy=AlignmentProviderPolicy,
        AlignmentProviderUsageRecord=AlignmentProviderUsageRecord,
        AlignmentProviderPreflightRun=AlignmentProviderPreflightRun,
    ),
)


@app.route("/api/concept-cards", methods=["POST"])
def create_concept_card_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json() or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    data["created_by"] = user.id
    try:
        card = concept_card_service.create_concept_card(
            db.session,
            ConceptAlignmentCard,
            data,
            audit_model=AuditRecord,
            actor=user,
            audit_context=audit_context,
            source="api",
            now_fn=current_time_text,
        )
    except concept_card_service.ConceptCardError as exc:
        db.session.rollback()
        if isinstance(exc, concept_card_service.ConceptCardQualityGateError):
            record_parse_quality_risk_audit(
                "concept_card_quality_gate_blocked",
                "concept_alignment_card",
                "",
                parse_uid=data.get("parse_uid", ""),
                quality_status=data.get("parse_quality_status", ""),
                risk_labels=data.get("risk_labels", []) + data.get("input_risk_labels", []) if isinstance(data.get("risk_labels", []), list) else data.get("input_risk_labels", []),
                blocked_reason=str(exc),
                audit_context=audit_context,
            )
        return concept_card_error_response(exc, audit_context)
    except ValueError as exc:
        db.session.rollback()
        return concept_card_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"card": concept_card_service.serialize_concept_card(card)},
        "Concept card created.",
        audit_context,
    )


@app.route("/api/concept-cards/draft-from-evidence", methods=["POST"])
def create_concept_card_draft_from_evidence_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json(silent=True) or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    data["created_by"] = user.id
    create_card = data.get("create", True)
    if isinstance(create_card, str):
        create_card = create_card.strip().lower() not in {"0", "false", "no"}
    else:
        create_card = bool(create_card)
    force_create = data.get("force_create", False)
    if isinstance(force_create, str):
        force_create = force_create.strip().lower() in {"1", "true", "yes"}
    else:
        force_create = bool(force_create)
    try:
        if not create_card:
            draft_payload, bilingual_result = concept_card_draft_service.build_draft_payload_from_bilingual_evidence(
                db.session,
                KnowledgeChunk,
                KnowledgeSource,
                data,
                concept_card_model=ConceptAlignmentCard,
                term_model=Term,
                terminology_card_model=TerminologyCard,
                audit_context=audit_context,
            )
            concept_card_draft_service.record_draft_audit(
                db.session,
                AuditRecord,
                event_type="concept_card_draft_payload_created",
                draft_payload=draft_payload,
                bilingual_result=bilingual_result,
                created=False,
                audit_context=audit_context,
                now_fn=current_time_text,
                commit=True,
            )
            if data.get("auto_generate_chinese_candidates") and not str(raw_data.get("chinese_term") or "").strip():
                candidate_event = "chinese_term_candidates_generated" if bilingual_result.chinese_term_candidates else "chinese_term_candidates_not_found"
                record_chinese_candidate_audit(
                    candidate_event,
                    input_data=data,
                    candidates=bilingual_result.chinese_term_candidates,
                    selected_candidate=bilingual_result.selected_chinese_candidate,
                    audit_context=audit_context,
                )
                if bilingual_result.selected_chinese_candidate:
                    record_chinese_candidate_audit(
                        "chinese_candidate_selected_for_draft",
                        input_data=data,
                        candidates=bilingual_result.chinese_term_candidates,
                        selected_candidate=bilingual_result.selected_chinese_candidate,
                        audit_context=audit_context,
                    )
            payload = concept_card_draft_service.serialize_concept_card_draft_result(
                None,
                bilingual_result,
                draft_payload=draft_payload,
                card_serializer=concept_card_service.serialize_concept_card,
                created=False,
                reused=False,
            )
            return api_success_with_audit_context(payload, "Concept card draft payload created.", audit_context)

        result = concept_card_draft_service.create_concept_card_draft_from_evidence(
            db.session,
            card_model=ConceptAlignmentCard,
            chunk_model=KnowledgeChunk,
            source_model=KnowledgeSource,
            term_model=Term,
            terminology_card_model=TerminologyCard,
            input_data=data,
            audit_model=AuditRecord,
            actor=user,
            audit_context=audit_context,
            now_fn=current_time_text,
            force_create=force_create,
            commit=True,
        )
    except (concept_card_draft_service.ConceptCardDraftError, bilingual_evidence_service.BilingualEvidenceWorkflowError, concept_card_service.ConceptCardError) as exc:
        db.session.rollback()
        if data.get("auto_generate_chinese_candidates") and not str(raw_data.get("chinese_term") or "").strip():
            record_chinese_candidate_audit(
                "chinese_term_candidate_generation_failed",
                input_data=data,
                error_code="candidate_generation_failed",
                error_message=str(exc),
                audit_context=audit_context,
            )
        concept_card_draft_service.record_draft_audit(
            db.session,
            AuditRecord,
            event_type="concept_card_draft_creation_failed",
            draft_payload=data,
            result="error",
            error_code="concept_card_draft_creation_failed",
            error_message=str(exc),
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=True,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            str(exc),
            400,
            audit_context,
            {"audit_error_code": "concept_card_draft_creation_failed"},
        )
    except Exception as exc:
        db.session.rollback()
        if data.get("auto_generate_chinese_candidates") and not str(raw_data.get("chinese_term") or "").strip():
            record_chinese_candidate_audit(
                "chinese_term_candidate_generation_failed",
                input_data=data,
                error_code="candidate_generation_failed",
                error_message=str(exc),
                audit_context=audit_context,
            )
        concept_card_draft_service.record_draft_audit(
            db.session,
            AuditRecord,
            event_type="concept_card_draft_creation_failed",
            draft_payload=data,
            result="error",
            error_code="concept_card_draft_creation_failed",
            error_message=str(exc),
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=True,
        )
        return api_error_with_audit_context(
            "INTERNAL_ERROR",
            "Concept card draft creation failed.",
            500,
            audit_context,
            {"audit_error_code": "concept_card_draft_creation_failed"},
        )
    payload = concept_card_draft_service.serialize_concept_card_draft_result(
        result.card,
        result.bilingual_result,
        draft_payload=result.draft_payload,
        card_serializer=concept_card_service.serialize_concept_card,
        created=result.created,
        reused=result.reused,
    )
    if data.get("auto_generate_chinese_candidates") and not str(raw_data.get("chinese_term") or "").strip():
        candidate_event = "chinese_term_candidates_generated" if result.bilingual_result.chinese_term_candidates else "chinese_term_candidates_not_found"
        record_chinese_candidate_audit(
            candidate_event,
            input_data=data,
            candidates=result.bilingual_result.chinese_term_candidates,
            selected_candidate=result.bilingual_result.selected_chinese_candidate,
            audit_context=audit_context,
        )
        if result.bilingual_result.selected_chinese_candidate:
            record_chinese_candidate_audit(
                "chinese_candidate_selected_for_draft",
                input_data=data,
                candidates=result.bilingual_result.chinese_term_candidates,
                selected_candidate=result.bilingual_result.selected_chinese_candidate,
                audit_context=audit_context,
            )
    return api_success_with_audit_context(payload, "Concept card draft from evidence handled.", audit_context)


@app.route("/api/concept-cards", methods=["GET"])
def list_concept_cards_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        result = concept_card_service.list_concept_cards(
            db.session,
            ConceptAlignmentCard,
            request.args.to_dict(),
        )
    except concept_card_service.ConceptCardError as exc:
        return concept_card_error_response(exc, audit_context)
    cards = [concept_card_service.serialize_concept_card(card) for card in result.items]
    return api_success_with_audit_context({
            "items": cards,
            "pagination": result.pagination,
        },
        audit_context=audit_context,
    )


register_concept_card_review_routes(
    app,
    core=route_core,
    models=ConceptCardReviewModels(
        ConceptAlignmentCard=ConceptAlignmentCard,
        ConceptCardReviewRecord=ConceptCardReviewRecord,
        ConceptCardReviewAssignment=ConceptCardReviewAssignment,
        CourseReviewPolicy=CourseReviewPolicy,
        CourseReviewPermission=CourseReviewPermission,
        AlignmentVerificationRun=AlignmentVerificationRun,
    ),
)


@app.route("/api/concept-cards/<card_uid>", methods=["GET"])
def get_concept_card_api(card_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        card = concept_card_service.get_concept_card(db.session, ConceptAlignmentCard, card_uid)
    except concept_card_service.ConceptCardError as exc:
        record_concept_card_api_failure(
            card_uid,
            exc,
            audit_context,
            {"card_uid": card_uid},
            error_code="concept_card_not_found",
        )
        return concept_card_error_response(exc, audit_context)
    return api_success_with_audit_context({"card": concept_card_service.serialize_concept_card(card)}, audit_context=audit_context)


@app.route("/api/concept-cards/<card_uid>", methods=["PATCH"])
def update_concept_card_api(card_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json() or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    invalid_fields = sorted(set(data) - concept_card_service.UPDATE_FIELDS)
    if invalid_fields:
        message = f"PATCH contains unsupported fields: {', '.join(invalid_fields)}."
        record_concept_card_api_failure(
            card_uid,
            message,
            audit_context,
            data,
            error_code="invalid_patch_field",
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            message,
            400,
            audit_context,
            {"invalid_fields": invalid_fields, "audit_error_code": "invalid_patch_field"},
        )
    try:
        card = concept_card_service.update_concept_card(
            db.session,
            ConceptAlignmentCard,
            card_uid,
            data,
            audit_model=AuditRecord,
            actor=user,
            audit_context=audit_context,
            source="api",
            now_fn=current_time_text,
        )
    except concept_card_service.ConceptCardError as exc:
        db.session.rollback()
        if isinstance(exc, concept_card_service.ConceptCardQualityGateError):
            record_parse_quality_risk_audit(
                "concept_card_quality_gate_blocked",
                "concept_alignment_card",
                card_uid,
                parse_uid=data.get("parse_uid", ""),
                quality_status=data.get("parse_quality_status", ""),
                risk_labels=data.get("risk_labels", []) + data.get("input_risk_labels", []) if isinstance(data.get("risk_labels", []), list) else data.get("input_risk_labels", []),
                blocked_reason=str(exc),
                audit_context=audit_context,
            )
        if isinstance(exc, concept_card_service.ConceptCardNotFoundError):
            record_concept_card_api_failure(
                card_uid,
                exc,
                audit_context,
                data,
                error_code="concept_card_not_found",
            )
        return concept_card_error_response(exc, audit_context)
    except ValueError as exc:
        db.session.rollback()
        return concept_card_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"card": concept_card_service.serialize_concept_card(card)},
        "Concept card updated.",
        audit_context,
    )


@app.route("/api/concept-cards/<card_uid>/status", methods=["POST"])
def change_concept_card_status_api(card_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    raw_data = request.get_json() or {}
    data = dict(raw_data) if isinstance(raw_data, dict) else {}
    try:
        card = concept_card_service.change_concept_card_status(
            db.session,
            ConceptAlignmentCard,
            card_uid,
            str(data.get("status", "")).strip(),
            reviewer=user,
            audit_model=AuditRecord,
            actor=user,
            audit_context=audit_context,
            source="api",
            now_fn=current_time_text,
        )
    except concept_card_service.ConceptCardError as exc:
        db.session.rollback()
        if isinstance(exc, concept_card_service.ConceptCardQualityGateError):
            record_parse_quality_risk_audit(
                "concept_card_quality_gate_blocked",
                "concept_alignment_card",
                card_uid,
                parse_uid=data.get("parse_uid", ""),
                quality_status=data.get("parse_quality_status", ""),
                risk_labels=data.get("risk_labels", []) + data.get("input_risk_labels", []) if isinstance(data.get("risk_labels", []), list) else data.get("input_risk_labels", []),
                blocked_reason=str(exc),
                audit_context=audit_context,
            )
        if isinstance(exc, concept_card_service.ConceptCardNotFoundError):
            record_concept_card_api_failure(
                card_uid,
                exc,
                audit_context,
                data,
                error_code="concept_card_not_found",
            )
        return concept_card_error_response(exc, audit_context)
    except ValueError as exc:
        db.session.rollback()
        return concept_card_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"card": concept_card_service.serialize_concept_card(card)},
        "Concept card status updated.",
        audit_context,
    )


def audit_record_error_response(exc, audit_context=None):
    if isinstance(exc, audit_record_service.AuditRecordNotFoundError):
        return api_error_with_audit_context("RESOURCE_NOT_FOUND", str(exc), 404, audit_context)
    return api_error_with_audit_context("VALIDATION_ERROR", str(exc), 400, audit_context, {"reason": str(exc)})


@app.route("/api/audit-records", methods=["GET"])
def list_audit_records_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        result = audit_record_service.list_audit_records(
            db.session,
            AuditRecord,
            request.args.to_dict(),
        )
    except audit_record_service.AuditRecordError as exc:
        return audit_record_error_response(exc, audit_context)
    records = [audit_record_service.serialize_audit_record(record) for record in result.items]
    return api_success_with_audit_context({
            "items": records,
            "pagination": result.pagination,
        },
        audit_context=audit_context,
    )


@app.route("/api/audit-records/<audit_uid>", methods=["GET"])
def get_audit_record_api(audit_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        record = audit_record_service.get_audit_record(db.session, AuditRecord, audit_uid)
    except audit_record_service.AuditRecordError as exc:
        return audit_record_error_response(exc, audit_context)
    return api_success_with_audit_context(
        {"audit_record": audit_record_service.serialize_audit_record(record)},
        audit_context=audit_context,
    )


def persist_document_parse_result(parse_result, *, audit_context=None, stored_path=""):
    data = dict(parse_result.parse_record_data)
    if stored_path:
        data["stored_path"] = stored_path
    else:
        data["stored_path"] = ""
    data["quality_flags"] = document_parse_quality_service.json_dumps(data.get("quality_flags", []))
    data["warnings"] = document_parse_quality_service.json_dumps(data.get("warnings", []))
    record = DocumentParseRecord(**data)
    db.session.add(record)
    db.session.flush()

    block_records = []
    for index, block_data in enumerate(parse_result.blocks or [], start=1):
        block = DocumentParseBlock(
            block_uid=block_data.get("block_uid") or str(uuid.uuid4()),
            parse_uid=record.parse_uid,
            page_number=block_data.get("page_number"),
            slide_number=block_data.get("slide_number"),
            block_index=block_data.get("block_index") or index,
            block_type=block_data.get("block_type", "text"),
            text=block_data.get("text", ""),
            confidence=block_data.get("confidence"),
            parser_type=block_data.get("parser_type", "native"),
            source_locator=block_data.get("source_locator", ""),
            quality_flags=document_parse_quality_service.json_dumps(block_data.get("quality_flags", [])),
            created_at=current_time_text(),
        )
        db.session.add(block)
        block_records.append(block)
    db.session.flush()

    event_type = "document_parse_failed" if record.parse_status == "failed" else "document_parse_created"
    audit_record_service.create_audit_record(
        db.session,
        AuditRecord,
        {
            "event_type": event_type,
            "target_type": "document_parse_record",
            "target_uid": record.parse_uid,
            "source": "api",
            "input_payload": {
                "source_filename": record.source_filename,
                "file_type": record.file_type,
                "mime_type": record.mime_type,
                "file_size_bytes": record.file_size_bytes,
            },
            "output_payload": {
                "parse_uid": record.parse_uid,
                "parse_status": record.parse_status,
                "quality_status": record.quality_status,
                "error_code": record.error_code,
            },
            "changed_fields": [],
            "result": "error" if record.parse_status == "failed" else "success",
            "error_code": record.error_code if record.parse_status == "failed" else "",
            "error_message": record.error_message if record.parse_status == "failed" else "",
        },
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=False,
    )
    db.session.commit()
    return record, block_records


def parse_quality_summary(record):
    serialized = document_parse_quality_service.serialize_parse_record(record)
    return {
        "parse_uid": serialized["parse_uid"],
        "parse_status": serialized["parse_status"],
        "quality_status": serialized["quality_status"],
        "quality_flags": serialized["quality_flags"],
        "should_allow_term_extraction": serialized["allow_term_extraction"],
        "extracted_text_chars": serialized["extracted_text_chars"],
        "block_count": serialized["block_count"],
        "warnings": serialized["warnings"],
    }


def quality_gate_error_code(record):
    status = getattr(record, "quality_status", "")
    if status in {"ocr_required", "ocr_unavailable"}:
        return "OCR_UNAVAILABLE"
    if status == "unsupported_file_type":
        return "UNSUPPORTED_FILE_TYPE"
    return "PARSING_FAILED"


def record_document_ingestion_audit(
    document,
    parse_record,
    ingestion_status,
    blocked_reason="",
    audit_context=None,
    commit=False,
    target_type="document",
    target_uid=None,
):
    event_type = "document_ingestion_blocked" if ingestion_status == "blocked" else "document_ingestion_completed"
    normalized_context = audit_context_service.normalize_audit_context(audit_context)
    audit_record_service.create_audit_record(
        db.session,
        AuditRecord,
        {
            "event_type": event_type,
            "target_type": target_type,
            "target_uid": str(target_uid if target_uid is not None else (getattr(document, "id", "") or getattr(parse_record, "parse_uid", ""))),
            "source": normalized_context.get("source") or "service",
            "input_payload": {
                "document_id": getattr(document, "id", None),
                "parse_uid": getattr(parse_record, "parse_uid", ""),
                "source_filename": getattr(parse_record, "source_filename", ""),
                "file_type": getattr(parse_record, "file_type", ""),
            },
            "output_payload": {
                "parse_status": getattr(parse_record, "parse_status", ""),
                "quality_status": getattr(parse_record, "quality_status", ""),
                "ingestion_status": ingestion_status,
                "blocked_reason": blocked_reason,
            },
            "changed_fields": [],
            "result": "error" if ingestion_status == "blocked" else "success",
            "error_code": "blocked_by_quality_gate" if ingestion_status == "blocked" else "",
            "error_message": blocked_reason if ingestion_status == "blocked" else "",
        },
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=commit,
    )


def record_evidence_retrieval_audit(
    event_type,
    *,
    query_text="",
    filters=None,
    candidates=None,
    result_count=0,
    error_code="",
    error_message="",
    latency_ms=None,
    audit_context=None,
    commit=True,
):
    candidates = candidates or []
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "evidence_retrieval",
                "target_uid": "",
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "query": str(query_text or "")[:240],
                    "filters": filters or {},
                },
                "output_payload": {
                    "result_count": int(result_count or 0),
                    "top_chunk_uids": [item.get("chunk_uid", "") for item in candidates[:5]],
                    "top_source_uids": [item.get("source_uid", "") for item in candidates[:5]],
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
                "latency_ms": latency_ms,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def record_bilingual_evidence_audit(
    event_type,
    *,
    input_data=None,
    result=None,
    error_code="",
    error_message="",
    latency_ms=None,
    audit_context=None,
    commit=True,
):
    input_data = dict(input_data or {})
    english_candidates = getattr(result, "english_evidence_candidates", []) if result is not None else []
    chinese_candidates = getattr(result, "chinese_evidence_candidates", []) if result is not None else []
    risk_labels = getattr(result, "risk_labels", []) if result is not None else []
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "bilingual_evidence_workflow",
                "target_uid": "",
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "english_term": str(input_data.get("english_term") or "")[:240],
                    "chinese_term": str(input_data.get("chinese_term") or "")[:240],
                    "course": str(input_data.get("course") or "")[:160],
                    "chapter": str(input_data.get("chapter") or "")[:160],
                    "limit": input_data.get("limit"),
                },
                "output_payload": {
                    "english_result_count": len(english_candidates),
                    "chinese_result_count": len(chinese_candidates),
                    "top_english_chunk_uids": [item.get("chunk_uid", "") for item in english_candidates[:5]],
                    "top_chinese_chunk_uids": [item.get("chunk_uid", "") for item in chinese_candidates[:5]],
                    "risk_labels": risk_labels,
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
                "latency_ms": latency_ms,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def record_chinese_candidate_audit(
    event_type,
    *,
    input_data=None,
    result=None,
    candidates=None,
    selected_candidate=None,
    error_code="",
    error_message="",
    latency_ms=None,
    audit_context=None,
    commit=True,
):
    input_data = dict(input_data or {})
    if result is not None:
        candidates = getattr(result, "candidates", candidates or [])
        risk_labels = getattr(result, "risk_labels", [])
    else:
        candidates = candidates or []
        risk_labels = input_data.get("risk_labels", [])
    selected_candidate = selected_candidate or {}
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "chinese_term_candidate",
                "target_uid": str(selected_candidate.get("candidate_uid") or ""),
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "english_term": str(input_data.get("english_term") or "")[:240],
                    "course": str(input_data.get("course") or "")[:160],
                    "chapter": str(input_data.get("chapter") or "")[:160],
                    "limit": input_data.get("limit") or input_data.get("candidate_limit"),
                },
                "output_payload": {
                    "candidate_count": len(candidates),
                    "selected_chinese_term": str(selected_candidate.get("chinese_term") or "")[:120],
                    "selected_candidate_score": selected_candidate.get("score"),
                    "source_uids": [item.get("source_uid", "") for item in candidates[:10]],
                    "chunk_uids": [item.get("chunk_uid", "") for item in candidates[:10]],
                    "risk_labels": risk_labels,
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
                "latency_ms": latency_ms,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def record_alignment_verification_audit(
    event_type,
    *,
    input_data=None,
    run=None,
    output_data=None,
    card_uid="",
    error_code="",
    error_message="",
    latency_ms=None,
    audit_context=None,
    commit=True,
):
    input_data = dict(input_data or {})
    output_data = dict(output_data or {})
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "alignment_verification",
                "target_uid": getattr(run, "run_uid", "") if run is not None else "",
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "card_uid": str(card_uid or input_data.get("card_uid") or "")[:120],
                    "english_term": str(input_data.get("english_term") or getattr(run, "english_term", "") or "")[:240],
                    "chinese_term": str(input_data.get("chinese_term") or getattr(run, "chinese_term", "") or "")[:240],
                    "course": str(input_data.get("course") or getattr(run, "course", "") or "")[:160],
                    "chapter": str(input_data.get("chapter") or getattr(run, "chapter", "") or "")[:160],
                    "provider_name": str(input_data.get("provider") or input_data.get("provider_name") or getattr(run, "provider_name", "") or "")[:120],
                },
                "output_payload": {
                    "run_uid": getattr(run, "run_uid", ""),
                    "card_uid": getattr(run, "card_uid", card_uid or ""),
                    "provider_name": getattr(run, "provider_name", output_data.get("provider_name", "")),
                    "provider_type": getattr(run, "provider_type", output_data.get("provider_type", "")),
                    "prompt_version": getattr(run, "prompt_version", output_data.get("prompt_version", "")),
                    "output_schema_version": getattr(run, "output_schema_version", output_data.get("output_schema_version", "")),
                    "parser_version": getattr(run, "parser_version", output_data.get("parser_version", "")),
                    "provider_response_status": getattr(run, "provider_response_status", output_data.get("provider_response_status", "")),
                    "verification_status": getattr(run, "verification_status", output_data.get("verification_status", "")),
                    "alignment_decision": output_data.get("alignment_decision", ""),
                    "recommendation": getattr(run, "recommendation", output_data.get("recommendation", "")),
                    "alignment_confidence": getattr(run, "alignment_confidence", output_data.get("alignment_confidence")),
                    "estimated_cost": output_data.get("estimated_cost", {}),
                    "retry_count": output_data.get("retry_count", 0),
                    "risk_labels": safe_json_loads(getattr(run, "risk_labels", "[]"), output_data.get("risk_labels", [])) if run is not None else output_data.get("risk_labels", []),
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
                "latency_ms": latency_ms,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def record_provider_governance_audit(
    event_type,
    *,
    provider_name="",
    policy=None,
    usage_record=None,
    run=None,
    input_data=None,
    output_data=None,
    error_code="",
    error_message="",
    audit_context=None,
    commit=True,
):
    input_data = dict(input_data or {})
    output_data = dict(output_data or {})
    policy_data = provider_governance_service.serialize_provider_policy(policy)
    usage_data = provider_governance_service.serialize_provider_usage_record(usage_record) if usage_record is not None else {}
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "alignment_provider_policy",
                "target_uid": policy_data.get("policy_uid") or usage_data.get("usage_uid") or getattr(run, "run_uid", ""),
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "provider_name": str(provider_name or policy_data.get("provider_name") or input_data.get("provider") or "")[:120],
                    "course": str(input_data.get("course") or getattr(run, "course", "") or "")[:160],
                    "card_uid": str(input_data.get("card_uid") or getattr(run, "card_uid", "") or "")[:120],
                },
                "output_payload": {
                    "provider_name": str(provider_name or policy_data.get("provider_name") or usage_data.get("provider_name") or "")[:120],
                    "policy_uid": policy_data.get("policy_uid", ""),
                    "run_uid": getattr(run, "run_uid", output_data.get("run_uid", "")),
                    "usage_uid": usage_data.get("usage_uid", ""),
                    "card_uid": getattr(run, "card_uid", output_data.get("card_uid", "")),
                    "course": getattr(run, "course", output_data.get("course", "")),
                    "provider_response_status": output_data.get("provider_response_status") or usage_data.get("provider_response_status", ""),
                    "estimated_cost": output_data.get("estimated_cost") or {"estimated_cost": usage_data.get("estimated_cost", 0.0)},
                    "blocked_reason": error_code or output_data.get("error_code", ""),
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def record_provider_preflight_audit(
    event_type,
    *,
    provider_name="",
    preflight_run=None,
    course="",
    error_code="",
    error_message="",
    audit_context=None,
    commit=True,
):
    run_data = provider_preflight_service.serialize_preflight_run(preflight_run) if preflight_run is not None else {}
    try:
        audit_record_service.create_audit_record(
            db.session,
            AuditRecord,
            {
                "event_type": event_type,
                "target_type": "alignment_provider_preflight",
                "target_uid": run_data.get("preflight_uid", ""),
                "source": audit_context_service.normalize_audit_context(audit_context).get("source") or "api",
                "input_payload": {
                    "provider_name": str(provider_name or run_data.get("provider_name", ""))[:120],
                    "course": str(course or run_data.get("course", ""))[:160],
                },
                "output_payload": {
                    "preflight_uid": run_data.get("preflight_uid", ""),
                    "provider_name": str(provider_name or run_data.get("provider_name", ""))[:120],
                    "course": run_data.get("course", course),
                    "check_status": run_data.get("check_status", ""),
                    "overall_ready": bool(run_data.get("overall_ready", False)),
                    "blocking_reasons": run_data.get("blocking_reasons", []),
                    "warnings": run_data.get("warnings", []),
                },
                "changed_fields": [],
                "result": "error" if error_code else "success",
                "error_code": error_code,
                "error_message": error_message,
            },
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=commit,
        )
    except Exception:
        if commit:
            db.session.rollback()


def provider_type_for_name(provider_name):
    return provider_governance_service.provider_type_for(provider_name)


def get_serialized_provider_policy(provider_name):
    policy = provider_governance_service.get_effective_provider_policy(
        db.session,
        AlignmentProviderPolicy,
        provider_name,
    )
    if policy is None:
        data = provider_governance_service.default_policy_data(provider_name, provider_type_for_name(provider_name))
        data["policy_missing"] = True
        return data
    return provider_governance_service.serialize_provider_policy(policy)


def record_alignment_provider_usage(provider_name, run=None, input_data=None, output_data=None, audit_context=None, commit=True):
    input_data = dict(input_data or {})
    output_data = dict(output_data or {})
    record = provider_governance_service.record_provider_usage(
        db.session,
        AlignmentProviderUsageRecord,
        provider_name,
        run_uid=getattr(run, "run_uid", output_data.get("run_uid", "")),
        input_summary={
            "card_uid": input_data.get("card_uid", getattr(run, "card_uid", "")),
            "course": input_data.get("course", getattr(run, "course", "")),
            "chapter": input_data.get("chapter", getattr(run, "chapter", "")),
            "estimated_cost": output_data.get("estimated_cost", {}),
        },
        result_summary={
            "run_uid": getattr(run, "run_uid", output_data.get("run_uid", "")),
            "card_uid": getattr(run, "card_uid", input_data.get("card_uid", "")),
            "provider_type": getattr(run, "provider_type", output_data.get("provider_type", "")),
            "provider_response_status": getattr(run, "provider_response_status", output_data.get("provider_response_status", "")),
            "estimated_cost": output_data.get("estimated_cost", {}),
            "error_code": getattr(run, "error_code", output_data.get("error_code", "")),
            "error_message": getattr(run, "error_message", output_data.get("error_message", "")),
        },
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=commit,
    )
    record_provider_governance_audit(
        "provider_usage_recorded",
        provider_name=provider_name,
        usage_record=record,
        run=run,
        output_data=output_data,
        audit_context=audit_context,
        commit=commit,
    )
    return record


register_provider_policy_routes(
    app,
    core=route_core,
    models=ProviderPolicyModels(
        AlignmentProviderPolicy=AlignmentProviderPolicy,
    ),
    record_provider_governance_audit=record_provider_governance_audit,
)


register_provider_preflight_routes(
    app,
    core=route_core,
    models=ProviderPreflightModels(
        AlignmentProviderPreflightRun=AlignmentProviderPreflightRun,
        AlignmentProviderPolicy=AlignmentProviderPolicy,
    ),
    record_provider_preflight_audit=record_provider_preflight_audit,
)


def build_alignment_verification_execution_dependencies():
    return alignment_verification_execution_service.AlignmentVerificationExecutionDependencies(
        db=db,
        models=alignment_verification_execution_service.AlignmentVerificationExecutionModels(
            concept_alignment_card=ConceptAlignmentCard,
            provider_policy=AlignmentProviderPolicy,
            provider_usage_record=AlignmentProviderUsageRecord,
            verification_run=AlignmentVerificationRun,
        ),
        provider_registry_service=alignment_provider_service,
        provider_governance_service=provider_governance_service,
        verification_service=alignment_verification_service,
        concept_card_service=concept_card_service,
        current_time_text=current_time_text,
        record_alignment_verification_audit=record_alignment_verification_audit,
        record_alignment_provider_usage=record_alignment_provider_usage,
    )


register_alignment_verification_routes(
    app,
    core=route_core,
    execution_dependencies=build_alignment_verification_execution_dependencies,
)


def knowledge_ingestion_models():
    return knowledge_ingestion_service.KnowledgeIngestionModels(
        source_model=KnowledgeSource,
        chunk_model=KnowledgeChunk,
        version_model=KnowledgeVersion,
        audit_model=AuditRecord,
    )


def normalize_governed_source_type(value, scope_type="course"):
    raw = str(value or "").strip()
    if raw in knowledge_governance_service.VALID_SOURCE_TYPES:
        return raw or "unknown"
    lowered = raw.lower()
    if scope_type == "personal" or "student" in lowered or "学生" in raw:
        return "student_upload"
    if "teacher" in lowered or "教师" in raw:
        return "teacher_upload"
    if "textbook" in lowered or "教材" in raw:
        return "textbook"
    if "paper" in lowered or "论文" in raw:
        return "paper"
    if "course" in lowered or "课件" in raw or "课程" in raw:
        return "course_material"
    return "teacher_upload" if scope_type == "course" else "unknown"


def governed_source_role(language, scope_type="course"):
    if scope_type == "personal":
        return "student_private_material"
    if language == "en":
        return "english_course_material"
    if language == "zh":
        return "chinese_reference_material"
    if language == "mixed":
        return "bilingual_reference"
    return "unknown"


def governed_trust_level(scope_type, quality_status, source_type=""):
    if quality_status in knowledge_governance_service.REVIEW_PARSE_STATUSES or quality_status in knowledge_governance_service.BLOCKED_PARSE_STATUSES:
        return "low_quality"
    if scope_type == "personal" or source_type == "student_upload":
        return "student_uploaded"
    if source_type in {"course_material", "teacher_upload"}:
        return "teacher_verified"
    if source_type in {"textbook", "paper", "reference"}:
        return "reference_material"
    return "unknown"


def build_governed_ingestion_metadata(
    *,
    parse_record,
    title="",
    course=None,
    course_name="",
    chapter="",
    language="unknown",
    source_type="",
    scope_type="course",
    owner_user=None,
    owner_user_id=None,
    document_id=None,
    knowledge_base_type="",
    visibility="course",
    content_hash="",
    license_note="",
    extra_quality_flags=None,
):
    normalized_language = str(language or "unknown").strip() or "unknown"
    normalized_scope = str(scope_type or "course").strip() or "course"
    normalized_source_type = normalize_governed_source_type(source_type, normalized_scope)
    quality_status = getattr(parse_record, "quality_status", "")
    flags = set(safe_json_loads(getattr(parse_record, "quality_flags", "[]"), []))
    flags.update(extra_quality_flags or [])
    if quality_status:
        flags.add(quality_status)
    effective_owner_id = owner_user_id if owner_user_id not in (None, "") else getattr(owner_user, "id", None)
    owner_role = getattr(owner_user, "role", "") or ("student" if normalized_scope == "personal" else "teacher")
    return {
        "title": title or getattr(parse_record, "source_filename", "") or "Untitled knowledge source",
        "course": course_name or (course.name if course else ""),
        "course_id": course.id if course else None,
        "chapter": chapter or "",
        "language": normalized_language,
        "source_type": normalized_source_type,
        "source_role": governed_source_role(normalized_language, normalized_scope),
        "owner_type": owner_role if owner_role in {"system", "teacher", "student", "admin"} else "unknown",
        "owner_id": str(effective_owner_id or ""),
        "owner_user_id": effective_owner_id,
        "visibility": visibility,
        "trust_level": governed_trust_level(normalized_scope, quality_status, normalized_source_type),
        "document_id": document_id,
        "scope_type": normalized_scope,
        "knowledge_base_type": knowledge_base_type,
        "source_filename": getattr(parse_record, "source_filename", ""),
        "file_type": getattr(parse_record, "file_type", "unknown"),
        "content_hash": content_hash,
        "license_note": license_note,
        "quality_flags": sorted(flag for flag in flags if flag),
        "created_by": getattr(owner_user, "id", None),
        "access_method": "document_parse",
    }


def build_document_chunks_from_parse_blocks(document, parse_record, parse_blocks):
    chunk_records = []
    document_flags = set(safe_json_loads(getattr(document, "quality_flags_json", "[]"), []))
    document_flags.add(getattr(parse_record, "quality_status", ""))
    for flag in safe_json_loads(getattr(parse_record, "quality_flags", "[]"), []):
        document_flags.add(flag)
    for index, block in enumerate(parse_blocks, start=1):
        text = clean_text(getattr(block, "text", ""))
        if not text or contains_ocr_placeholder(text) or contains_formula_placeholder(text):
            continue
        confidence = getattr(block, "confidence", None)
        ocr_confidence = int(round(float(confidence) * 100)) if confidence is not None else 100
        block_flags = set(safe_json_loads(getattr(block, "quality_flags", "[]"), []))
        for flag in block_flags:
            document_flags.add(flag)
        if getattr(parse_record, "quality_status", "") == "partial_text":
            block_flags.add("partial_text")
            document_flags.add("partial_text")
        chunk = DocumentChunk(
            document_id=document.id,
            course_id=document.course_id,
            user_id=document.owner_user_id,
            owner_user_id=document.owner_user_id,
            chunk_index=index,
            parse_uid=parse_record.parse_uid,
            parse_block_uid=block.block_uid,
            language=document.language,
            page_number=block.page_number,
            slide_number=block.slide_number,
            section_title="",
            content=text,
            source_type=document.source_type,
            source_location=block.source_locator,
            ocr_confidence=max(0, min(ocr_confidence, 100)),
            ocr_provider="",
            ocr_status="not_required" if block.parser_type == "native" else getattr(parse_record, "quality_status", ""),
            ocr_error="",
            quality_flags_json=json.dumps(sorted(flag for flag in block_flags if flag), ensure_ascii=False),
            created_at=current_time_text(),
        )
        db.session.add(chunk)
        chunk_records.append(chunk)
    document.quality_flags_json = json.dumps(sorted(flag for flag in document_flags if flag), ensure_ascii=False)
    db.session.flush()
    return chunk_records


def parse_block_records_to_chunk_items(document, parse_record, parse_blocks):
    items = []
    record_flags = set(safe_json_loads(getattr(parse_record, "quality_flags", "[]"), []))
    if getattr(parse_record, "quality_status", ""):
        record_flags.add(getattr(parse_record, "quality_status", ""))
    for index, block in enumerate(parse_blocks or [], start=1):
        text = clean_text(getattr(block, "text", ""))
        if not text:
            continue
        confidence = getattr(block, "confidence", None)
        ocr_confidence = int(round(float(confidence) * 100)) if confidence is not None else 100
        parser_type = getattr(block, "parser_type", "native")
        block_flags = set(safe_json_loads(getattr(block, "quality_flags", "[]"), [])) | record_flags
        items.append({
            "language": getattr(document, "language", ""),
            "page_number": getattr(block, "page_number", None),
            "slide_number": getattr(block, "slide_number", None),
            "section_title": "",
            "content": text,
            "source_type": getattr(document, "source_type", ""),
            "source_location": getattr(block, "source_locator", "") or f"block {index}",
            "ocr_confidence": max(0, min(ocr_confidence, 100)),
            "ocr_provider": "",
            "ocr_status": "not_required" if parser_type == "native" else getattr(parse_record, "quality_status", ""),
            "ocr_error": "",
            "parse_uid": getattr(parse_record, "parse_uid", ""),
            "parse_block_uid": getattr(block, "block_uid", ""),
            "quality_flags": sorted(flag for flag in block_flags if flag),
        })
    return items


def create_parse_record_for_saved_file(save_path, document, filename, mime_type, audit_context=None):
    parse_result = document_parse_quality_service.parse_document_with_quality(
        save_path,
        filename=filename,
        mime_type=mime_type or "",
        now_fn=current_time_text,
    )
    parse_record, parse_blocks = persist_document_parse_result(
        parse_result,
        audit_context=audit_context,
        stored_path=getattr(document, "storage_key", "") or "",
    )
    document.parse_uid = parse_record.parse_uid
    document.parsed_text = parse_result.raw_text
    document.ocr_required = bool(parse_record.ocr_required)
    document.ocr_status = parse_record.quality_status if parse_record.ocr_required else "not_required"
    document.ocr_error = parse_record.error_message or ""
    document.quality_flags_json = parse_record.quality_flags
    return parse_result, parse_record, parse_blocks


def create_formula_block_for_quality_gate(document, file_path):
    ext = Path(str(file_path or "")).suffix.lower().lstrip(".")
    if ext not in {"jpg", "jpeg", "png"}:
        return []
    if not looks_like_formula_image(file_path, ""):
        return []
    formula_result = get_formula_ocr_provider(FORMULA_OCR_PROVIDER).recognize_formula(file_path)
    formula_block = FormulaBlock(
        document_id=document.id,
        course_id=document.course_id,
        owner_user_id=document.owner_user_id,
        scope_type=document.scope_type,
        page_number=1,
        slide_number=None,
        bbox_json=json.dumps(formula_result.bbox or {}, ensure_ascii=False),
        image_path="",
        latex=formula_result.latex or "",
        plain_text=formula_result.plain_text or "",
        provider=formula_result.provider or FORMULA_OCR_PROVIDER or "none",
        confidence=float(formula_result.confidence or 0),
        status=formula_result.status,
        error=formula_result.error or "",
        quality_flags_json=json.dumps(formula_result.quality_flags or ["formula_ocr_required"], ensure_ascii=False),
        created_at=current_time_text(),
    )
    db.session.add(formula_block)
    db.session.flush()
    return [formula_block]


def block_document_by_quality_gate(document, job, parse_record, audit_context=None):
    blocked_reason = parse_record.error_message or f"Blocked by parse quality gate: {parse_record.quality_status}"
    document.parsing_status = "blocked_by_quality_gate"
    document.error_message = blocked_reason
    if job is not None:
        job.status = "failed"
        job.error_message = blocked_reason
        job.finished_at = current_time_text()
    record_document_ingestion_audit(
        document,
        parse_record,
        "blocked",
        blocked_reason=blocked_reason,
        audit_context=audit_context,
        commit=False,
    )
    db.session.commit()
    return {
        **parse_quality_summary(parse_record),
        "ingestion_status": "blocked",
        "blocked_reason": blocked_reason,
    }


def document_parse_error_response(message, status_code, audit_context=None, details=None):
    return api_error_with_audit_context("DOCUMENT_PARSE_ERROR", message, status_code, audit_context, details or {})


@app.route("/api/document-parses", methods=["GET"])
def list_document_parses_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    page = max(1, int(request.args.get("page", 1) or 1))
    per_page = max(1, min(int(request.args.get("per_page", 20) or 20), 100))
    query = DocumentParseRecord.query
    for field in ("quality_status", "parse_status", "file_type"):
        value = str(request.args.get(field, "") or "").strip()
        if value:
            query = query.filter_by(**{field: value})
    q = str(request.args.get("q", "") or "").strip()
    if q:
        like = f"%{q}%"
        query = query.filter(DocumentParseRecord.source_filename.ilike(like))
    total = query.count()
    records = query.order_by(DocumentParseRecord.id.desc()).offset((page - 1) * per_page).limit(per_page).all()
    return api_success_with_audit_context({
        "items": [document_parse_quality_service.serialize_parse_record(record) for record in records],
        "pagination": {
            "page": page,
            "per_page": per_page,
            "total": total,
            "has_next": page * per_page < total,
        },
    }, audit_context=audit_context)


@app.route("/api/document-parses/<parse_uid>", methods=["GET"])
def get_document_parse_api(parse_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    record = DocumentParseRecord.query.filter_by(parse_uid=parse_uid).first()
    if record is None:
        return api_error_with_audit_context(
            "RESOURCE_NOT_FOUND",
            "DocumentParseRecord not found.",
            404,
            audit_context,
            {"parse_uid": parse_uid},
        )
    blocks = DocumentParseBlock.query.filter_by(parse_uid=record.parse_uid).order_by(DocumentParseBlock.block_index.asc()).limit(200).all()
    return api_success_with_audit_context({
        "parse_record": document_parse_quality_service.serialize_parse_record(record),
        "blocks": [document_parse_quality_service.serialize_parse_block(block) for block in blocks],
    }, audit_context=audit_context)


@app.route("/api/document-parses/test", methods=["POST"])
def test_document_parse_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    upload = request.files.get("file")
    if upload is None or not upload.filename:
        return document_parse_error_response("file is required.", 400, audit_context)
    filename = secure_filename(upload.filename) or f"parse-test-{uuid.uuid4().hex}.bin"
    suffix = Path(filename).suffix
    with tempfile.TemporaryDirectory(prefix="lexibridge-parse-test-") as tmp_dir:
        temp_path = os.path.join(tmp_dir, f"upload{suffix}")
        upload.save(temp_path)
        parse_result = document_parse_quality_service.parse_document_with_quality(
            temp_path,
            filename=filename,
            mime_type=upload.mimetype or "",
            now_fn=current_time_text,
        )
    record, blocks = persist_document_parse_result(parse_result, audit_context=audit_context)
    return api_success_with_audit_context({
        "parse_record": document_parse_quality_service.serialize_parse_record(record),
        "blocks": [document_parse_quality_service.serialize_parse_block(block) for block in blocks[:20]],
        "raw_text_preview": (parse_result.raw_text or "")[:500],
    }, "Document parse quality recorded.", audit_context)


def knowledge_governance_error_response(exc, audit_context=None):
    if isinstance(exc, knowledge_governance_service.KnowledgeSourceNotFoundError):
        return api_error_with_audit_context("RESOURCE_NOT_FOUND", str(exc), 404, audit_context)
    if isinstance(exc, knowledge_governance_service.KnowledgeChunkNotFoundError):
        return api_error_with_audit_context("RESOURCE_NOT_FOUND", str(exc), 404, audit_context)
    if isinstance(exc, knowledge_governance_service.KnowledgeIngestionBlockedError):
        return api_error_with_audit_context("KNOWLEDGE_INGESTION_BLOCKED", str(exc), 422, audit_context)
    return api_error_with_audit_context("VALIDATION_ERROR", str(exc), 400, audit_context, {"reason": str(exc)})


def build_knowledge_governance_metadata(data, user=None):
    course = get_course_by_id_or_name(data.get("course_id"), data.get("course", data.get("course_name", "")))
    metadata = {
        "title": str(data.get("title") or "").strip(),
        "course": str(data.get("course") or data.get("course_name") or (course.name if course else "") or "").strip(),
        "course_id": course.id if course else data.get("course_id"),
        "chapter": str(data.get("chapter") or "").strip(),
        "language": str(data.get("language") or "unknown").strip(),
        "source_type": str(data.get("source_type") or "unknown").strip(),
        "source_role": str(data.get("source_role") or "unknown").strip(),
        "owner_type": str(data.get("owner_type") or getattr(user, "role", "") or "unknown").strip(),
        "owner_id": str(data.get("owner_id") or (getattr(user, "id", "") if user else "") or "").strip(),
        "owner_user_id": getattr(user, "id", None),
        "visibility": str(data.get("visibility") or "course").strip(),
        "trust_level": str(data.get("trust_level") or "unknown").strip(),
        "license_note": str(data.get("license_note") or "").strip(),
        "created_by": getattr(user, "id", None),
        "scope_type": str(data.get("scope_type") or "course").strip(),
        "knowledge_base_type": str(data.get("knowledge_base_type") or "").strip(),
        "access_method": "document_parse",
    }
    return metadata, course


@app.route("/api/knowledge-sources", methods=["GET"])
def list_governed_knowledge_sources_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        result = knowledge_governance_service.list_knowledge_sources(
            db.session,
            KnowledgeSource,
            request.args.to_dict(),
        )
    except knowledge_governance_service.KnowledgeGovernanceError as exc:
        return knowledge_governance_error_response(exc, audit_context)
    return api_success_with_audit_context({
        "items": [knowledge_governance_service.serialize_knowledge_source(source) for source in result.items],
        "pagination": result.pagination,
    }, audit_context=audit_context)


@app.route("/api/knowledge-sources/<source_uid>", methods=["GET"])
def get_governed_knowledge_source_api(source_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        source = knowledge_governance_service.get_knowledge_source(db.session, KnowledgeSource, source_uid)
    except knowledge_governance_service.KnowledgeGovernanceError as exc:
        return knowledge_governance_error_response(exc, audit_context)
    chunk_count = KnowledgeChunk.query.filter(db.or_(
        KnowledgeChunk.source_uid == getattr(source, "source_uid", ""),
        KnowledgeChunk.knowledge_source_id == source.id,
    )).count()
    return api_success_with_audit_context({
        "source": knowledge_governance_service.serialize_knowledge_source(source),
        "chunk_count": chunk_count,
    }, audit_context=audit_context)


@app.route("/api/knowledge-chunks", methods=["GET"])
def list_governed_knowledge_chunks_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        result = knowledge_governance_service.list_knowledge_chunks(
            db.session,
            KnowledgeChunk,
            request.args.to_dict(),
        )
    except knowledge_governance_service.KnowledgeGovernanceError as exc:
        return knowledge_governance_error_response(exc, audit_context)
    return api_success_with_audit_context({
        "items": [knowledge_governance_service.serialize_knowledge_chunk(chunk) for chunk in result.items],
        "pagination": result.pagination,
    }, audit_context=audit_context)


@app.route("/api/knowledge-chunks/<chunk_uid>", methods=["GET"])
def get_governed_knowledge_chunk_api(chunk_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    try:
        chunk = knowledge_governance_service.get_knowledge_chunk(db.session, KnowledgeChunk, chunk_uid)
    except knowledge_governance_service.KnowledgeGovernanceError as exc:
        return knowledge_governance_error_response(exc, audit_context)
    return api_success_with_audit_context({
        "chunk": knowledge_governance_service.serialize_knowledge_chunk(chunk),
    }, audit_context=audit_context)


@app.route("/api/knowledge-sources/from-parse/<parse_uid>", methods=["POST"])
def create_governed_knowledge_from_parse_api(parse_uid):
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    data = request.get_json() or {}
    parse_record = DocumentParseRecord.query.filter_by(parse_uid=parse_uid).first()
    if parse_record is None:
        return api_error_with_audit_context(
            "RESOURCE_NOT_FOUND",
            "DocumentParseRecord not found.",
            404,
            audit_context,
            {"parse_uid": parse_uid},
        )
    metadata, course = build_knowledge_governance_metadata(data, user)
    if course is not None and user.role == "teacher" and not can_manage_course(user, course):
        return api_error_with_audit_context("PERMISSION_DENIED", "无权管理该课程知识源。", 403, audit_context)
    parse_blocks = DocumentParseBlock.query.filter_by(parse_uid=parse_record.parse_uid).order_by(DocumentParseBlock.block_index.asc()).all()
    try:
        governed_ingestion = knowledge_ingestion_service.ingest_parse_record_to_governed_knowledge(
            db.session,
            knowledge_ingestion_models(),
            parse_record,
            parse_blocks,
            metadata,
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=True,
        )
    except knowledge_ingestion_service.KnowledgeIngestionBlockedError as exc:
        return api_error_with_audit_context(
            "KNOWLEDGE_INGESTION_BLOCKED",
            "解析记录未通过知识治理入库门禁，未创建 active KnowledgeChunk。",
            422,
            audit_context,
            {
                "parse_uid": parse_record.parse_uid,
                "quality_status": parse_record.quality_status,
                "blocked_by_quality_gate": True,
                "blocked_reason": exc.blocked_reason,
            },
        )
    source = governed_ingestion.source
    chunks = governed_ingestion.chunks
    return api_success_with_audit_context({
        "source": knowledge_governance_service.serialize_knowledge_source(source),
        "chunks": [knowledge_governance_service.serialize_knowledge_chunk(chunk) for chunk in chunks[:20]],
        "chunk_count": len(chunks),
        "source_uid": getattr(source, "source_uid", ""),
        "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in chunks[:20] if getattr(chunk, "chunk_uid", "")],
        "parse_uid": parse_record.parse_uid,
        "quality_status": parse_record.quality_status,
    }, "Governed knowledge source created from parse record.", audit_context)


@app.route("/api/evidence/search", methods=["POST"])
def search_governed_evidence_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    started_at = datetime.now()
    data = request.get_json(silent=True) or {}
    query_text = str(data.get("query") or data.get("q") or "").strip()
    filters = {key: value for key, value in data.items() if key not in {"query", "q"}}
    limit = data.get("limit", evidence_retrieval_service.DEFAULT_LIMIT)
    if not query_text:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_evidence_retrieval_audit(
            "evidence_retrieval_failed",
            query_text=query_text,
            filters=filters,
            error_code="missing_query",
            error_message="query is required.",
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            "query is required.",
            400,
            audit_context,
            {"audit_error_code": "missing_query"},
        )
    try:
        result = evidence_retrieval_service.search_evidence(
            db.session,
            KnowledgeChunk,
            KnowledgeSource,
            query_text,
            filters=filters,
            limit=limit,
            audit_context=audit_context,
        )
    except evidence_retrieval_service.EvidenceRetrievalError as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_evidence_retrieval_audit(
            "evidence_retrieval_failed",
            query_text=query_text,
            filters=filters,
            error_code="validation_error",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            str(exc),
            400,
            audit_context,
            {"audit_error_code": "validation_error"},
        )
    except Exception as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        db.session.rollback()
        record_evidence_retrieval_audit(
            "evidence_retrieval_failed",
            query_text=query_text,
            filters=filters,
            error_code="retrieval_failed",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "INTERNAL_ERROR",
            "Evidence retrieval failed.",
            500,
            audit_context,
            {"audit_error_code": "retrieval_failed"},
        )
    latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
    candidates = [
        evidence_retrieval_service.serialize_evidence_candidate(candidate)
        for candidate in result.candidates
    ]
    record_evidence_retrieval_audit(
        "evidence_retrieval_completed",
        query_text=result.query,
        filters=result.filters,
        candidates=candidates,
        result_count=result.total,
        latency_ms=latency_ms,
        audit_context=audit_context,
    )
    return api_success_with_audit_context({
        "query": result.query,
        "filters": result.filters,
        "total": result.total,
        "candidates": candidates,
    }, "Evidence retrieval completed.", audit_context)


@app.route("/api/terms/chinese-candidates", methods=["POST"])
def generate_chinese_term_candidates_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    started_at = datetime.now()
    data = request.get_json(silent=True) or {}
    english_term = str(data.get("english_term") or data.get("query") or data.get("q") or "").strip()
    if not english_term:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_chinese_candidate_audit(
            "chinese_term_candidate_generation_failed",
            input_data=data,
            error_code="missing_english_term",
            error_message="english_term is required.",
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            "english_term is required.",
            400,
            audit_context,
            {"audit_error_code": "missing_english_term"},
        )
    try:
        result = chinese_term_candidate_service.generate_chinese_term_candidates(
            db.session,
            concept_card_model=ConceptAlignmentCard,
            term_model=Term,
            terminology_card_model=TerminologyCard,
            chunk_model=KnowledgeChunk,
            source_model=KnowledgeSource,
            english_term=english_term,
            course=data.get("course", ""),
            chapter=data.get("chapter", ""),
            limit=data.get("limit", chinese_term_candidate_service.DEFAULT_CANDIDATE_LIMIT),
            filters=data.get("filters", {}) if isinstance(data.get("filters", {}), dict) else {},
            audit_context=audit_context,
        )
    except chinese_term_candidate_service.ChineseTermCandidateError as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_chinese_candidate_audit(
            "chinese_term_candidate_generation_failed",
            input_data=data,
            error_code="validation_error",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            str(exc),
            400,
            audit_context,
            {"audit_error_code": "validation_error"},
        )
    except Exception as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        db.session.rollback()
        record_chinese_candidate_audit(
            "chinese_term_candidate_generation_failed",
            input_data=data,
            error_code="candidate_generation_failed",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "INTERNAL_ERROR",
            "Chinese term candidate generation failed.",
            500,
            audit_context,
            {"audit_error_code": "candidate_generation_failed"},
        )
    latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
    event_type = "chinese_term_candidates_generated" if result.candidates else "chinese_term_candidates_not_found"
    record_chinese_candidate_audit(
        event_type,
        input_data=data,
        result=result,
        selected_candidate=result.candidates[0] if result.candidates else None,
        latency_ms=latency_ms,
        audit_context=audit_context,
    )
    return api_success_with_audit_context(
        {
            "english_term": result.english_term,
            "course": result.course,
            "chapter": result.chapter,
            "total": result.total,
            "candidates": [
                chinese_term_candidate_service.serialize_chinese_term_candidate(candidate)
                for candidate in result.candidates
            ],
            "risk_labels": result.risk_labels,
        },
        "Chinese term candidates generated.",
        audit_context,
    )


@app.route("/api/evidence/bilingual", methods=["POST"])
def retrieve_bilingual_evidence_api():
    audit_context = get_route_audit_context()
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return attach_request_id_to_response(error_response, audit_context)
    audit_context = get_route_audit_context(user)
    started_at = datetime.now()
    data = request.get_json(silent=True) or {}
    try:
        query_data = bilingual_evidence_service.build_bilingual_evidence_query(data)
    except bilingual_evidence_service.BilingualEvidenceWorkflowError as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_bilingual_evidence_audit(
            "bilingual_evidence_retrieval_failed",
            input_data=data,
            error_code="missing_english_term",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            str(exc),
            400,
            audit_context,
            {"audit_error_code": "missing_english_term"},
        )
    try:
        result = bilingual_evidence_service.retrieve_bilingual_evidence(
            db.session,
            KnowledgeChunk,
            KnowledgeSource,
            query_data["english_term"],
            chinese_term=query_data["chinese_term"],
            course=query_data["course"],
            chapter=query_data["chapter"],
            concept_scope=query_data["concept_scope"],
            limit=query_data["limit"],
            filters=query_data["filters"],
            auto_generate_chinese_candidates=query_data["auto_generate_chinese_candidates"],
            candidate_limit=query_data["candidate_limit"],
            selected_chinese_candidate_uid=query_data["selected_chinese_candidate_uid"],
            concept_card_model=ConceptAlignmentCard,
            term_model=Term,
            terminology_card_model=TerminologyCard,
            audit_context=audit_context,
        )
    except bilingual_evidence_service.BilingualEvidenceWorkflowError as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        record_bilingual_evidence_audit(
            "bilingual_evidence_retrieval_failed",
            input_data=query_data,
            error_code="validation_error",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "VALIDATION_ERROR",
            str(exc),
            400,
            audit_context,
            {"audit_error_code": "validation_error"},
        )
    except Exception as exc:
        latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
        db.session.rollback()
        if data.get("auto_generate_chinese_candidates") and not str(data.get("chinese_term") or "").strip():
            record_chinese_candidate_audit(
                "chinese_term_candidate_generation_failed",
                input_data=data,
                error_code="candidate_generation_failed",
                error_message=str(exc),
                latency_ms=latency_ms,
                audit_context=audit_context,
            )
        record_bilingual_evidence_audit(
            "bilingual_evidence_retrieval_failed",
            input_data=query_data,
            error_code="bilingual_retrieval_failed",
            error_message=str(exc),
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        return api_error_with_audit_context(
            "INTERNAL_ERROR",
            "Bilingual evidence retrieval failed.",
            500,
            audit_context,
            {"audit_error_code": "bilingual_retrieval_failed"},
        )
    latency_ms = int((datetime.now() - started_at).total_seconds() * 1000)
    if query_data.get("auto_generate_chinese_candidates") and not str(data.get("chinese_term") or "").strip():
        candidate_event = "chinese_term_candidates_generated" if result.chinese_term_candidates else "chinese_term_candidates_not_found"
        record_chinese_candidate_audit(
            candidate_event,
            input_data=query_data,
            candidates=result.chinese_term_candidates,
            selected_candidate=result.selected_chinese_candidate,
            latency_ms=latency_ms,
            audit_context=audit_context,
        )
        if result.selected_chinese_candidate:
            record_chinese_candidate_audit(
                "chinese_candidate_selected_for_draft",
                input_data=query_data,
                candidates=result.chinese_term_candidates,
                selected_candidate=result.selected_chinese_candidate,
                latency_ms=latency_ms,
                audit_context=audit_context,
            )
    record_bilingual_evidence_audit(
        "bilingual_evidence_retrieval_completed",
        input_data=query_data,
        result=result,
        latency_ms=latency_ms,
        audit_context=audit_context,
    )
    record_bilingual_evidence_audit(
        "concept_card_draft_payload_created",
        input_data=query_data,
        result=result,
        latency_ms=latency_ms,
        audit_context=audit_context,
    )
    return api_success_with_audit_context(
        bilingual_evidence_service.serialize_bilingual_evidence_result(result),
        "Bilingual evidence retrieval completed.",
        audit_context,
    )


def filter_cards_for_user(query, user):
    if user.role == "admin":
        return query
    if user.role == "teacher":
        manageable_ids = [
            course.id for course in Course.query.all()
            if can_manage_course(user, course)
        ]
        return query.filter(db.or_(
            TerminologyCard.owner_user_id == user.id,
            db.and_(
                TerminologyCard.scope_type == "course",
                TerminologyCard.course_id.in_(manageable_ids or [-1])
            )
        ))
    course_ids = [member.course_id for member in CourseMember.query.filter_by(user_id=user.id).all()]
    return query.filter(db.or_(
        db.and_(
            TerminologyCard.scope_type == "personal",
            TerminologyCard.owner_user_id == user.id
        ),
        db.and_(
            TerminologyCard.scope_type == "course",
            TerminologyCard.course_id.in_(course_ids or [-1]),
            TerminologyCard.status.in_(["auto_approved", "approved"])
        )
    ))


@app.route("/api/terminology/cards", methods=["GET"])
def terminology_cards():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    q = request.args.get("q", "").strip().lower()
    scope_type = request.args.get("scope_type", "").strip()
    status = request.args.get("status", "").strip()
    course_id = request.args.get("course_id", "").strip()
    page = max(1, int(request.args.get("page", "1") or 1))
    page_size = max(1, min(int(request.args.get("page_size", "20") or 20), 100))

    query = filter_cards_for_user(TerminologyCard.query, user)
    if scope_type:
        query = query.filter_by(scope_type=scope_type)
    if status:
        query = query.filter_by(status=status)
    if course_id:
        query = query.filter_by(course_id=int(course_id))

    total = query.count()
    cards = query.order_by(TerminologyCard.id.desc()).offset((page - 1) * page_size).limit(page_size).all()
    if q:
        cards = [
            card for card in cards
            if q in card.english_term.lower()
            or q in card.final_chinese_term.lower()
            or q in card.ai_translation_candidate.lower()
            or q in card.concept_explanation.lower()
            or q in card.alignment_reason.lower()
        ]
        if user.role == "student":
            record_usage(user.id, "term_search", 1)
            db.session.commit()

    return jsonify({
        "status": "success",
        "count": len(cards),
        "data": {
            "items": [serialize_terminology_card(card) for card in cards],
            "pagination": {
                "page": page,
                "page_size": page_size,
                "total": total,
                "has_next": page * page_size < total,
            }
        },
        "cards": [serialize_terminology_card(card) for card in cards]
    })


@app.route("/api/terminology/cards/export", methods=["GET"])
def export_terminology_cards():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    if user.role == "student":
        plan, subscription = get_effective_plan(user.id)
        if not plan.export_enabled:
            return api_error(
                "QUOTA_EXCEEDED",
                "当前会员套餐不支持 PDF 导出，请升级到 Basic 或 Pro。",
                402,
                {"plan": serialize_subscription_plan(plan)}
            )

    q = request.args.get("q", "").strip().lower()
    course_id = request.args.get("course_id", "").strip()
    favorite_only = request.args.get("favorite_only", "false").strip().lower() in {"1", "true", "yes"}
    mastered_only = request.args.get("mastered_only", "false").strip().lower() in {"1", "true", "yes"}

    query = filter_cards_for_user(TerminologyCard.query, user)
    if course_id:
        query = query.filter_by(course_id=int(course_id))
    cards = query.order_by(TerminologyCard.id.desc()).limit(300).all()

    if q:
        cards = [
            card for card in cards
            if q in card.english_term.lower()
            or q in card.final_chinese_term.lower()
            or q in card.concept_explanation.lower()
            or q in card.alignment_reason.lower()
        ]

    if favorite_only or mastered_only:
        records = StudentTermRecord.query.filter_by(user_id=user.id).all()
        favorite_ids = {record.term_id for record in records if record.is_favorite}
        mastered_ids = {record.term_id for record in records if record.is_mastered}
        if favorite_only:
            cards = [card for card in cards if card.id in favorite_ids]
        if mastered_only:
            cards = [card for card in cards if card.id in mastered_ids]

    if not cards:
        return api_error("RESOURCE_NOT_FOUND", "当前筛选条件下没有可导出的术语卡片。", 404)

    try:
        pdf_buffer = build_cards_pdf(cards)
    except Exception as exc:
        if "中文字体" in str(exc) or "PDF_FONT_PATH" in str(exc):
            return api_error("PDF_FONT_UNAVAILABLE", "PDF 导出失败：当前环境缺少可用中文字体。", 422, {"error": str(exc)})
        return api_error("INTERNAL_ERROR", "PDF 导出失败。", 500)

    filename = f"lexibridge-terminology-cards-{datetime.now().strftime('%Y%m%d%H%M%S')}.pdf"
    return send_file(
        pdf_buffer,
        mimetype="application/pdf",
        as_attachment=True,
        download_name=filename
    )


def get_authorized_card(card_id, user):
    card = db.session.get(TerminologyCard, card_id)
    if card is None:
        return None, (jsonify({"status": "error", "message": "术语卡片不存在。"}), 404)
    scoped = filter_cards_for_user(TerminologyCard.query.filter_by(id=card.id), user).first()
    if scoped is None:
        return None, (jsonify({"status": "error", "message": "无权访问该术语卡片。"}), 403)
    return card, None


def create_pilot_feedback_for_card(card, user, data, source="student_card_detail"):
    feedback_type = normalize_choice(data.get("feedback_type"), FEEDBACK_TYPES, "other")
    severity = normalize_choice(data.get("severity"), SEVERITIES, "medium")
    feedback_source = normalize_choice(data.get("feedback_source") or source, FEEDBACK_SOURCES, source)
    reported_issue = str(data.get("reported_issue") or data.get("feedback_content") or "").strip()
    if not reported_issue:
        return None, api_error("VALIDATION_ERROR", "反馈内容不能为空。", 400)
    classification, root_cause = classify_feedback(feedback_type, reported_issue)
    course = db.session.get(Course, card.course_id) if card.course_id else None
    priority = map_feedback_to_priority(feedback_type, severity, classification)
    feedback = Feedback(
        term_id=card.id,
        terminology_card_id=card.id,
        user_id=user.id,
        user_role=user.role,
        course_id=card.course_id,
        document_id=getattr(card, "source_document_id", None),
        alignment_run_id=getattr(card, "source_alignment_run_id", None) or getattr(card, "alignment_run_id", None),
        course=course.name if course else "",
        english_term=card.english_term,
        chinese_term=card.final_chinese_term,
        feedback_type=feedback_type,
        feedback_source=feedback_source,
        severity=severity,
        priority=priority,
        feedback_content=reported_issue,
        reported_issue=reported_issue,
        expected_result=str(data.get("expected_result") or "").strip(),
        actual_result=str(data.get("actual_result") or card.final_chinese_term or "").strip(),
        evidence_comment=str(data.get("evidence_comment") or "").strip(),
        classification=classification if classification in CLASSIFICATIONS else "teacher_review_needed",
        root_cause=root_cause if root_cause in ROOT_CAUSES else "unknown",
        status="submitted",
        created_at=current_time_text(),
        updated_at=current_time_text(),
    )
    db.session.add(feedback)
    db.session.flush()

    if feedback_type in {"translation_error", "evidence_error"}:
        open_statuses = ["submitted", "triaged", "in_review", "needs_more_evidence", "converted_to_backlog", "converted_to_evaluation_item"]
        open_count = Feedback.query.filter(
            Feedback.terminology_card_id == card.id,
            Feedback.feedback_type.in_(["translation_error", "evidence_error"]),
            Feedback.status.in_(open_statuses)
        ).count()
        card.feedback_count = open_count
        high_count = Feedback.query.filter(
            Feedback.terminology_card_id == card.id,
            Feedback.severity.in_(["high", "critical"]),
            Feedback.status.in_(open_statuses)
        ).count()
        if should_escalate_card(feedback_type, severity, high_count):
            card.status = "pending_quality_control"
            card.risk_note = ((card.risk_note or "") + " Pilot feedback escalated this card to QC.").strip()
    if severity == "critical":
        add_system_log("warning", "pilot_feedback", f"Critical feedback {feedback.id} submitted for card {card.id}.")
    return feedback, None


@app.route("/api/terminology/cards/<int:card_id>", methods=["GET"])
def terminology_card_detail(card_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    card, error = get_authorized_card(card_id, user)
    if error:
        return error
    if user.role == "student":
        record = StudentTermRecord.query.filter_by(user_id=user.id, term_id=card.id).first()
        if record is None:
            record = StudentTermRecord(user_id=user.id, term_id=card.id)
            db.session.add(record)
        record.last_viewed_at = current_time_text()
        db.session.commit()
    return jsonify({"status": "success", "card": serialize_terminology_card(card)})


@app.route("/api/terminology/cards/<int:card_id>/favorite", methods=["POST"])
def terminology_card_favorite(card_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    card, error = get_authorized_card(card_id, user)
    if error:
        return error
    data = request.get_json() or {}
    value = bool(data.get("is_favorite", True))
    record = StudentTermRecord.query.filter_by(user_id=user.id, term_id=card.id).first()
    if record is None:
        record = StudentTermRecord(user_id=user.id, term_id=card.id, last_viewed_at=current_time_text())
        db.session.add(record)
    record.is_favorite = value
    record.last_viewed_at = current_time_text()
    db.session.commit()
    return jsonify({"status": "success", "is_favorite": record.is_favorite, "card": serialize_terminology_card(card)})


@app.route("/api/terminology/cards/<int:card_id>/mastered", methods=["POST"])
def terminology_card_mastered(card_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    card, error = get_authorized_card(card_id, user)
    if error:
        return error
    data = request.get_json() or {}
    value = bool(data.get("is_mastered", True))
    record = StudentTermRecord.query.filter_by(user_id=user.id, term_id=card.id).first()
    if record is None:
        record = StudentTermRecord(user_id=user.id, term_id=card.id, last_viewed_at=current_time_text())
        db.session.add(record)
    record.is_mastered = value
    record.last_viewed_at = current_time_text()
    db.session.commit()
    return jsonify({"status": "success", "is_mastered": record.is_mastered, "card": serialize_terminology_card(card)})


@app.route("/api/terminology/cards/<int:card_id>/feedback", methods=["POST"])
def terminology_card_feedback(card_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    card, error = get_authorized_card(card_id, user)
    if error:
        return error
    feedback, feedback_error = create_pilot_feedback_for_card(card, user, request.get_json() or {}, "student_card_detail")
    if feedback_error:
        return feedback_error
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "Feedback submitted.",
        "data": {
            "feedback_id": feedback.id,
            "feedback_status": feedback.status
        },
        "feedback": serialize_feedback(feedback),
        "card": serialize_terminology_card(card)
    })


@app.route("/api/quality-control", methods=["GET"])
def quality_control_queue():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    requested_status = str(request.args.get("status", "")).strip()
    statuses = [
        "pending_quality_control",
        "needs_more_evidence",
        "conflict_detected",
        "rejected",
        "auto_approved",
    ]
    if requested_status:
        statuses = [requested_status]
    query = filter_cards_for_user(TerminologyCard.query.filter(TerminologyCard.status.in_(statuses)), user)
    cards = query.order_by(TerminologyCard.id.desc()).all()
    return jsonify({
        "status": "success",
        "count": len(cards),
        "cards": [serialize_terminology_card(card) for card in cards]
    })


def update_qc_card(card_id, user, status, data=None):
    card, error = get_authorized_card(card_id, user)
    if error:
        return None, error
    course = db.session.get(Course, card.course_id) if card.course_id else None
    if user.role != "admin" and not can_manage_course(user, course):
        return None, (jsonify({"status": "error", "message": "无权处理该质量控制项目。"}), 403)
    data = data or {}
    if "final_chinese_term" in data:
        card.final_chinese_term = str(data.get("final_chinese_term", "")).strip() or card.final_chinese_term
    if "concept_explanation" in data:
        card.concept_explanation = str(data.get("concept_explanation", "")).strip() or card.concept_explanation
    note = str(data.get("note", "")).strip()
    if note:
        card.reviewer_note = note
        card.alignment_reason = f"{card.alignment_reason}\nQC note: {note}".strip()
    if not validate_card_status_transition(card.status, status, user.role, system_action=False):
        return None, (jsonify({
            "status": "error",
            "message": f"非法状态流转：{card.status} -> {status}。"
        }), 400)
    card.status = status
    if status == "approved":
        card.alignment_status = "accepted_translation"
        card.approved_by = user.id
        card.approved_at = current_time_text()
    if status == "rejected":
        card.rejected_reason = str(data.get("rejected_reason") or data.get("reason") or note or "Rejected in Quality Control.").strip()
    card.updated_at = current_time_text()
    db.session.commit()
    return card, None


@app.route("/api/quality-control/<int:card_id>/approve", methods=["POST"])
def quality_control_approve(card_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    card, error = update_qc_card(card_id, user, "approved", request.get_json() or {})
    if error:
        return error
    return jsonify({"status": "success", "message": "术语卡片已通过质量控制。", "card": serialize_terminology_card(card)})


@app.route("/api/quality-control/<int:card_id>/edit", methods=["POST"])
def quality_control_edit(card_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    card, error = update_qc_card(card_id, user, "approved", request.get_json() or {})
    if error:
        return error
    return jsonify({"status": "success", "message": "术语卡片已编辑并通过。", "card": serialize_terminology_card(card)})


@app.route("/api/quality-control/<int:card_id>/needs-more-evidence", methods=["POST"])
def quality_control_needs_more_evidence(card_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    card, error = update_qc_card(card_id, user, "needs_more_evidence", request.get_json() or {})
    if error:
        return error
    return jsonify({"status": "success", "message": "术语卡片已标记为需要更多证据。", "card": serialize_terminology_card(card)})


@app.route("/api/quality-control/<int:card_id>/reject", methods=["POST"])
def quality_control_reject(card_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    card, error = update_qc_card(card_id, user, "rejected", request.get_json() or {})
    if error:
        return error
    return jsonify({"status": "success", "message": "术语卡片已驳回。", "card": serialize_terminology_card(card)})


@app.route("/api/subscription/plans", methods=["GET"])
def subscription_plans():
    plans = SubscriptionPlan.query.filter_by(is_active=True).order_by(SubscriptionPlan.price_monthly.asc()).all()
    return jsonify({
        "status": "success",
        "plans": [serialize_subscription_plan(plan) for plan in plans]
    })


@app.route("/api/subscription/me", methods=["GET"])
def subscription_me():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    plan, subscription = get_effective_plan(user.id)
    totals = get_usage_totals(user.id)
    return jsonify({
        "status": "success",
        "plan": serialize_subscription_plan(plan),
        "subscription": serialize_subscription(subscription) if subscription else None,
        "usage": {
            "period_start": totals["period_start"],
            "pages_used": totals["pages_used"],
            "pages_remaining": max(0, plan.monthly_pages - totals["pages_used"]),
            "ai_calls_used": totals["ai_calls_used"],
            "ai_calls_remaining": max(0, plan.monthly_ai_calls - totals["ai_calls_used"])
        }
    })


@app.route("/api/subscription/mock-payment", methods=["POST"])
def subscription_mock_payment():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    data = request.get_json() or {}
    plan_id = data.get("plan_id")
    plan_name = str(data.get("plan_name", "")).strip()
    plan = db.session.get(SubscriptionPlan, int(plan_id)) if plan_id else None
    if plan is None and plan_name:
        plan = SubscriptionPlan.query.filter_by(name=plan_name).first()
    if plan is None or not plan.is_active:
        return jsonify({"status": "error", "message": "会员套餐不存在。"}), 404

    active = get_active_subscription(user.id)
    if active:
        active.status = "replaced"
    subscription = UserSubscription(
        user_id=user.id,
        plan_id=plan.id,
        start_date=current_time_text(),
        end_date=future_time_text(60 * 24 * 30),
        status="active",
        auto_renew=False
    )
    billing = BillingRecord(
        user_id=user.id,
        plan_id=plan.id,
        amount=plan.price_monthly,
        payment_method="mock_payment",
        payment_status="paid",
        created_at=current_time_text()
    )
    db.session.add(subscription)
    db.session.add(billing)
    add_system_log("info", "mock_payment", f"User {user.id} activated {plan.name} plan by mock payment.")
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "Mock payment completed. Subscription activated.",
        "subscription": serialize_subscription(subscription),
        "billing_record": {
            "id": billing.id,
            "amount": billing.amount,
            "payment_status": billing.payment_status
        }
    })


@app.route("/api/usage/me", methods=["GET"])
def usage_me():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    totals = get_usage_totals(user.id)
    return jsonify({
        "status": "success",
        "usage": {
            "period_start": totals["period_start"],
            "pages_used": totals["pages_used"],
            "ai_calls_used": totals["ai_calls_used"],
            "records": [serialize_usage(record) for record in totals["records"]]
        }
    })


@app.route("/api/admin/users", methods=["GET"])
def admin_users():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    users = User.query.order_by(User.id.asc()).all()
    return jsonify({"status": "success", "users": [serialize_user(item) for item in users]})


@app.route("/api/admin/users/<int:user_id>/role", methods=["POST"])
def admin_update_user_role(user_id):
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    target = db.session.get(User, user_id)
    if target is None:
        return jsonify({"status": "error", "message": "用户不存在。"}), 404
    role = normalize_role((request.get_json() or {}).get("role", "student"))
    target.role = role
    db.session.commit()
    return jsonify({"status": "success", "user": serialize_user(target)})


@app.route("/api/admin/usage", methods=["GET"])
def admin_usage():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    records = UsageRecord.query.order_by(UsageRecord.id.desc()).limit(300).all()
    return jsonify({"status": "success", "records": [serialize_usage(record) for record in records]})


@app.route("/api/admin/billing", methods=["GET"])
def admin_billing():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    records = BillingRecord.query.order_by(BillingRecord.id.desc()).limit(300).all()
    return jsonify({
        "status": "success",
        "records": [
            {
                "id": record.id,
                "user_id": record.user_id,
                "plan_id": record.plan_id,
                "amount": record.amount,
                "payment_method": record.payment_method,
                "payment_status": record.payment_status,
                "created_at": record.created_at
            }
            for record in records
        ]
    })


@app.route("/api/admin/logs", methods=["GET"])
def admin_logs():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    logs = SystemLog.query.order_by(SystemLog.id.desc()).limit(300).all()
    return jsonify({
        "status": "success",
        "logs": [
            {
                "id": log.id,
                "level": log.level,
                "module": log.module,
                "message": log.message,
                "created_at": log.created_at
            }
            for log in logs
        ]
    })


@app.route("/api/admin/ingestion-jobs", methods=["GET"])
def admin_ingestion_jobs():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    jobs = IngestionJob.query.order_by(IngestionJob.id.desc()).limit(300).all()
    return jsonify({
        "status": "success",
        "jobs": [
            {
                "id": job.id,
                "source_id": job.source_id,
                "document_id": job.document_id,
                "status": job.status,
                "started_at": job.started_at,
                "finished_at": job.finished_at,
                "error_message": job.error_message,
                "processed_pages": job.processed_pages,
                "created_by": job.created_by
            }
            for job in jobs
        ]
    })


register_admin_alignment_run_routes(
    app,
    core=route_core,
    models=AdminAlignmentRunModels(
        AlignmentRun=AlignmentRun,
    ),
    serialize_alignment_run=serialize_alignment_run,
)


register_legacy_provider_admin_observability_routes(
    app,
    core=route_core,
    models=LegacyProviderAdminObservabilityModels(
        AICallLog=AICallLog,
        AIProviderConfig=AIProviderConfig,
    ),
    serializers=LegacyProviderAdminObservabilitySerializers(
        api_success=api_success,
        serialize_ai_call_log=serialize_ai_call_log,
        serialize_ai_provider_config=serialize_ai_provider_config,
        summarize_ai_calls=summarize_ai_calls,
    ),
    registry_seed_service=ensure_ai_registry_seed,
)


@app.route("/api/admin/personal-access-audit", methods=["GET"])
def admin_personal_access_audit():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    audits = PersonalAccessAudit.query.order_by(PersonalAccessAudit.id.desc()).limit(300).all()
    return jsonify({
        "status": "success",
        "audits": [serialize_personal_access_audit(audit) for audit in audits]
    })


@app.route("/api/admin/model-registry", methods=["GET", "POST"])
def admin_model_registry():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response

    if request.method == "POST":
        data = request.get_json() or {}
        registry_id = data.get("id")
        registry = db.session.get(ModelPromptRegistry, int(registry_id)) if registry_id else None
        if registry is None:
            registry = ModelPromptRegistry(
                created_at=current_time_text()
            )
            db.session.add(registry)
        registry.provider = str(data.get("provider", registry.provider or AI_PROVIDER)).strip()
        registry.model_name = str(data.get("model_name", registry.model_name or DEEPSEEK_MODEL)).strip()
        registry.model_version = str(data.get("model_version", registry.model_version or "local-mvp-v1")).strip()
        registry.prompt_version = str(data.get("prompt_version", registry.prompt_version or ALIGNMENT_PROMPT_VERSION)).strip()
        registry.retrieval_version = str(data.get("retrieval_version", registry.retrieval_version or RETRIEVAL_VERSION)).strip()
        registry.enabled = bool(data.get("enabled", registry.enabled))
        workflows = data.get("allowed_workflows", ["term_extraction", "bilingual_alignment", "student_answer"])
        registry.allowed_workflows = json.dumps(workflows, ensure_ascii=False) if isinstance(workflows, list) else str(workflows)
        registry.known_risks = str(data.get("known_risks", registry.known_risks or "")).strip()
        registry.rollback_target = str(data.get("rollback_target", registry.rollback_target or "local_heuristic")).strip()
        registry.owner = str(data.get("owner", registry.owner or user.email)).strip()
        registry.updated_at = current_time_text()
        db.session.commit()
        return jsonify({
            "status": "success",
            "registry": serialize_model_registry(registry)
        })

    ensure_model_registry_seed(owner_user_id=user.id)
    registries = ModelPromptRegistry.query.order_by(ModelPromptRegistry.id.desc()).all()
    return jsonify({
        "status": "success",
        "provider_status": current_provider_metadata(),
        "registries": [serialize_model_registry(registry) for registry in registries]
    })


def admin_ai_prompts_post_handler(user):
    data = request.get_json() or {}
    mutation_request = LegacyPromptMutationRequest.from_payload(data, actor_user_id=user.id)
    result = execute_legacy_prompt_mutation(
        request=mutation_request,
        dependencies=legacy_prompt_mutation_dependencies(),
    )
    if result.outcome == "validation_error":
        return api_error(result.error_code, result.message, 400)
    if result.outcome == "persistence_error":
        return api_error(result.error_code, result.message, 500)
    return api_success(serialize_prompt_template(result.prompt), result.message)


register_legacy_provider_admin_configuration_routes(
    app,
    core=route_core,
    models=LegacyProviderAdminConfigurationModels(
        AIProviderConfig=AIProviderConfig,
        AIModelRegistry=AIModelRegistry,
        PromptTemplate=PromptTemplate,
    ),
    serializers=LegacyProviderAdminConfigurationSerializers(
        api_success=api_success,
        serialize_ai_provider_config=serialize_ai_provider_config,
        serialize_ai_model_registry=serialize_ai_model_registry,
        serialize_prompt_template=serialize_prompt_template,
        current_provider_metadata=current_provider_metadata,
    ),
    registry_seed_service=ensure_legacy_provider_registry_seed,
    seed_models=LegacyProviderRegistrySeedModels(
        AIProviderConfig=AIProviderConfig,
        AIModelRegistry=AIModelRegistry,
        PromptTemplate=PromptTemplate,
    ),
    provider_selection_factory=lambda: env_provider_selection(os.environ),
    default_prompts=DEFAULT_PROMPTS,
    model_version_factory=lambda: os.environ.get("MODEL_VERSION", "local-mvp-v1"),
    prompt_post_handler=admin_ai_prompts_post_handler,
)


register_legacy_provider_admin_healthcheck_routes(
    app,
    core=route_core,
    models=LegacyProviderAdminHealthcheckModels(
        AIProviderConfig=AIProviderConfig,
    ),
    serializers=LegacyProviderAdminHealthcheckSerializers(
        api_success=api_success,
    ),
    registry_seed_service=ensure_legacy_provider_registry_seed,
    seed_models=LegacyProviderRegistrySeedModels(
        AIProviderConfig=AIProviderConfig,
        AIModelRegistry=AIModelRegistry,
        PromptTemplate=PromptTemplate,
    ),
    provider_selection_factory=lambda: env_provider_selection(os.environ),
    default_prompts=DEFAULT_PROMPTS,
    model_version_factory=lambda: os.environ.get("MODEL_VERSION", "local-mvp-v1"),
    local_readiness_service=legacy_provider_local_readiness_service,
    credential_presence_resolver=legacy_provider_config_credential_present,
)


@app.route("/api/evaluation/sets", methods=["GET", "POST"])
def evaluation_sets():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    if request.method == "POST":
        data = request.get_json() or {}
        name = str(data.get("name", "")).strip()
        if not name:
            return api_error("VALIDATION_ERROR", "evaluation set name 不能为空。", 400)
        course = get_course_by_id_or_name(data.get("course_id"), data.get("course_name", ""))
        if user.role != "admin" and course is not None and not can_manage_course(user, course):
            return api_error("PERMISSION_DENIED", "无权为该课程创建评估集。", 403)
        evaluation_set = EvaluationSet(
            name=name,
            course_id=course.id if course else None,
            discipline=str(data.get("discipline", "")).strip(),
            description=str(data.get("description", "")).strip(),
            split=str(data.get("split", "test")).strip() or "test",
            locked=bool(data.get("locked", False)),
            is_locked=bool(data.get("is_locked", data.get("locked", False))),
            created_by=user.id,
            created_at=current_time_text(),
            updated_at=current_time_text()
        )
        db.session.add(evaluation_set)
        db.session.flush()
        items = data.get("items", [])
        if isinstance(items, list):
            for raw_item in items:
                item = normalize_evaluation_record(raw_item)
                if not item["english_term"]:
                    continue
                item_course = get_course_by_id_or_name(item.get("course_id"), raw_item.get("course_name", "")) or course
                if user.role != "admin" and item_course is not None and not can_manage_course(user, item_course):
                    continue
                db.session.add(EvaluationItem(
                    set_id=evaluation_set.id,
                    evaluation_set_id=evaluation_set.id,
                    item_id=item["item_id"],
                    split=item["split"],
                    discipline=item["discipline"] or evaluation_set.discipline,
                    course_id=item_course.id if item_course else None,
                    english_term=item["english_term"],
                    expected_chinese_term=item["expected_chinese_term"],
                    expected_alignment_status=item["expected_alignment_status"],
                    english_context=item["english_context"],
                    english_evidence=item["expected_english_evidence"],
                    chinese_evidence=item["expected_chinese_evidence"],
                    expected_english_evidence=item["expected_english_evidence"],
                    expected_chinese_evidence=item["expected_chinese_evidence"],
                    negative_english_evidence=item["negative_english_evidence"],
                    negative_chinese_evidence=item["negative_chinese_evidence"],
                    difficulty=item["difficulty"],
                    tags_json=json.dumps(item["tags"], ensure_ascii=False),
                    annotator=item["annotator"],
                    reviewed_by=item["reviewed_by"],
                    disagreement_note=item["disagreement_note"],
                    version=item["version"],
                    created_at=current_time_text()
                ))
        db.session.commit()
        return jsonify({
            "status": "success",
            "data": {"evaluation_set_id": evaluation_set.id},
            "set": serialize_evaluation_set(evaluation_set),
            "items": [serialize_evaluation_item(item) for item in EvaluationItem.query.filter_by(set_id=evaluation_set.id).all()]
        })

    query = EvaluationSet.query
    if user.role != "admin":
        manageable_ids = [
            course.id for course in Course.query.all()
            if can_manage_course(user, course)
        ]
        query = query.filter(db.or_(
            EvaluationSet.created_by == user.id,
            EvaluationSet.course_id.in_(manageable_ids or [-1])
        ))
    sets = query.order_by(EvaluationSet.id.desc()).all()
    return jsonify({
        "status": "success",
        "sets": [serialize_evaluation_set(item) for item in sets]
    })


@app.route("/api/evaluation/items/import", methods=["POST"])
def evaluation_items_import():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    data = request.get_json() or {}
    try:
        result = import_evaluation_items(data.get("file_path"), data.get("evaluation_set_id"), user)
        db.session.commit()
    except PermissionError as exc:
        db.session.rollback()
        return api_error("PERMISSION_DENIED", str(exc), 403)
    except FileNotFoundError as exc:
        db.session.rollback()
        return api_error("RESOURCE_NOT_FOUND", str(exc), 404)
    except Exception as exc:
        db.session.rollback()
        return api_error("VALIDATION_ERROR", str(exc), 400)
    return jsonify({"status": "success", "data": result})


@app.route("/api/evaluation/items", methods=["GET"])
def evaluation_items():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    set_id = request.args.get("evaluation_set_id") or request.args.get("set_id")
    split = str(request.args.get("split", "")).strip()
    query = EvaluationItem.query
    if set_id:
        evaluation_set = db.session.get(EvaluationSet, int(set_id))
        if evaluation_set is None:
            return api_error("RESOURCE_NOT_FOUND", "Evaluation set 不存在。", 404)
        if not can_manage_evaluation_set(user, evaluation_set):
            return api_error("PERMISSION_DENIED", "无权查看该评估集。", 403)
        query = query.filter_by(set_id=evaluation_set.id)
    elif user.role != "admin":
        owned_set_ids = [
            item.id for item in EvaluationSet.query.all()
            if can_manage_evaluation_set(user, item)
        ]
        query = query.filter(EvaluationItem.set_id.in_(owned_set_ids or [-1]))
    if split:
        query = query.filter_by(split=split)
    items = query.order_by(EvaluationItem.id.asc()).limit(1000).all()
    return jsonify({"status": "success", "items": [serialize_evaluation_item(item) for item in items]})


@app.route("/api/evaluation/sets/<int:set_id>/items", methods=["GET"])
def evaluation_set_items(set_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    evaluation_set = db.session.get(EvaluationSet, set_id)
    if evaluation_set is None:
        return api_error("RESOURCE_NOT_FOUND", "Evaluation set 不存在。", 404)
    if not can_manage_evaluation_set(user, evaluation_set):
        return api_error("PERMISSION_DENIED", "无权查看该评估集。", 403)
    items = EvaluationItem.query.filter_by(set_id=set_id).order_by(EvaluationItem.id.asc()).all()
    return jsonify({
        "status": "success",
        "set": serialize_evaluation_set(evaluation_set),
        "items": [serialize_evaluation_item(item) for item in items]
    })


@app.route("/api/evaluation/run", methods=["POST"])
def evaluation_run():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    sync_requested = str(request.args.get("sync", "")).strip().lower() in {"1", "true", "yes"}
    data = request.get_json() or {}
    set_id = data.get("evaluation_set_id")
    evaluation_set = db.session.get(EvaluationSet, int(set_id)) if set_id else None
    if evaluation_set is None:
        return api_error("RESOURCE_NOT_FOUND", "evaluation_set_id 无效。", 404)
    if not can_manage_evaluation_set(user, evaluation_set):
        return api_error("PERMISSION_DENIED", "无权运行该 EvaluationSet。", 403)
    if not sync_requested:
        meta = current_provider_metadata()
        now = current_time_text()
        run = EvaluationRun(
            evaluation_set_id=evaluation_set.id,
            triggered_by=user.id,
            created_by=user.id,
            provider=meta["provider"],
            provider_name=meta["provider"],
            provider_mode=meta.get("provider_mode", ""),
            model_name=meta["model_name"],
            model_version=str(data.get("model_version", "")).strip() or meta["model_name"],
            prompt_key="term_alignment",
            prompt_version=str(data.get("prompt_version", "")).strip() or "v1",
            retrieval_version=str(data.get("retrieval_version", "")).strip() or meta["retrieval_version"],
            alignment_version=ALIGNMENT_PROMPT_VERSION,
            commit_hash=os.environ.get("COMMIT_HASH", "local"),
            split=str(data.get("split", "test")).strip() or "test",
            status="queued",
            created_at=now,
        )
        db.session.add(run)
        db.session.flush()
        background_job = create_background_job(
            "evaluation_run",
            user,
            course_id=evaluation_set.course_id,
            evaluation_run_id=run.id,
            scope_type="course" if evaluation_set.course_id else "global",
            owner_user_id=None,
            input_data={
                "evaluation_set_id": evaluation_set.id,
                "split": run.split,
                "model_version": run.model_version,
                "prompt_version": run.prompt_version,
                "retrieval_version": run.retrieval_version,
            }
        )
        db.session.commit()
        return api_success({
            "evaluation_run_id": run.id,
            "job_id": background_job.id,
            "job_type": background_job.job_type,
            "job_status": background_job.status,
            "run": serialize_evaluation_run(run),
            "job": serialize_background_job(background_job),
        }, "EvaluationRun 已进入后台队列。")
    try:
        run = run_evaluation_set(
            evaluation_set,
            user,
            split=data.get("split", "test"),
            model_version=str(data.get("model_version", "")).strip(),
            prompt_version=str(data.get("prompt_version", "")).strip(),
            retrieval_version=str(data.get("retrieval_version", "")).strip(),
        )
        registry = ensure_model_registry_seed(owner_user_id=user.id)
        registry.last_evaluation_run_id = run.id
        registry.updated_at = current_time_text()
        db.session.commit()
    except PermissionError as exc:
        db.session.rollback()
        return api_error("PERMISSION_DENIED", str(exc), 403)
    except Exception as exc:
        db.session.rollback()
        return api_error("VALIDATION_ERROR", str(exc), 400)
    return jsonify({
        "status": "success",
        "data": {
            "evaluation_run_id": run.id,
            "status": run.status,
            "metrics": safe_json_loads(run.metrics_json, {}),
        },
        "run": serialize_evaluation_run(run)
    })


@app.route("/api/evaluation/runs", methods=["GET"])
def evaluation_runs():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    query = EvaluationRun.query
    if user.role != "admin":
        allowed_set_ids = [
            item.id for item in EvaluationSet.query.all()
            if can_manage_evaluation_set(user, item)
        ]
        query = query.filter(EvaluationRun.evaluation_set_id.in_(allowed_set_ids or [-1]))
    runs = query.order_by(EvaluationRun.id.desc()).limit(100).all()
    return jsonify({
        "status": "success",
        "runs": [serialize_evaluation_run(run) for run in runs]
    })


@app.route("/api/evaluation/runs/<int:run_id>", methods=["GET"])
def evaluation_run_detail(run_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    run = db.session.get(EvaluationRun, run_id)
    if run is None:
        return api_error("RESOURCE_NOT_FOUND", "EvaluationRun 不存在。", 404)
    evaluation_set = db.session.get(EvaluationSet, run.evaluation_set_id) if run.evaluation_set_id else None
    if evaluation_set is not None and not can_manage_evaluation_set(user, evaluation_set):
        return api_error("PERMISSION_DENIED", "无权查看该 EvaluationRun。", 403)
    return jsonify({"status": "success", "data": serialize_evaluation_run(run)})


@app.route("/api/jobs", methods=["GET"])
def list_background_jobs():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    page = max(1, int(request.args.get("page", "1") or 1))
    page_size = max(1, min(int(request.args.get("page_size", "20") or 20), 100))
    query = visible_jobs_query(user)
    status = str(request.args.get("status", "")).strip()
    job_type = str(request.args.get("job_type", "")).strip()
    if status:
        query = query.filter_by(status=status)
    if job_type:
        query = query.filter_by(job_type=job_type)
    total = query.count()
    jobs = query.order_by(BackgroundJob.id.desc()).offset((page - 1) * page_size).limit(page_size).all()
    return api_success({
        "items": [serialize_background_job(job) for job in jobs],
        "pagination": {
            "page": page,
            "page_size": page_size,
            "total": total,
            "has_next": page * page_size < total,
        }
    })


@app.route("/api/jobs/<int:job_id>", methods=["GET"])
def background_job_detail(job_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    job = db.session.get(BackgroundJob, job_id)
    if job is None:
        return api_error("RESOURCE_NOT_FOUND", "BackgroundJob 不存在。", 404)
    if not can_view_job(user, job):
        return api_error("PERMISSION_DENIED", "无权查看该后台任务。", 403)
    return api_success({"job": serialize_background_job(job)})


@app.route("/api/jobs/<int:job_id>/events", methods=["GET"])
def background_job_events(job_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    job = db.session.get(BackgroundJob, job_id)
    if job is None:
        return api_error("RESOURCE_NOT_FOUND", "BackgroundJob 不存在。", 404)
    if not can_view_job(user, job):
        return api_error("PERMISSION_DENIED", "无权查看该后台任务。", 403)
    events = BackgroundJobEvent.query.filter_by(job_id=job.id).order_by(BackgroundJobEvent.id.asc()).all()
    return api_success({"items": [serialize_background_job_event(event) for event in events]})


@app.route("/api/jobs/<int:job_id>/cancel", methods=["POST"])
def cancel_background_job(job_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    job = db.session.get(BackgroundJob, job_id)
    if job is None:
        return api_error("RESOURCE_NOT_FOUND", "BackgroundJob 不存在。", 404)
    if not can_mutate_job(user, job):
        return api_error("PERMISSION_DENIED", "无权取消该后台任务。", 403)
    if job.status in {"completed", "failed", "canceled"}:
        return api_error("VALIDATION_ERROR", "该任务已结束，不能取消。", 400, {"status": job.status})
    job.status = "canceled"
    job.canceled_at = current_time_text()
    job.finished_at = job.finished_at or job.canceled_at
    job.progress_message = "Canceled"
    job.updated_at = job.canceled_at
    add_job_event(job, "canceled", f"Job canceled by user {user.id}.")
    db.session.commit()
    return api_success({"job": serialize_background_job(job)}, "后台任务已取消。")


@app.route("/api/jobs/<int:job_id>/retry", methods=["POST"])
def retry_background_job(job_id):
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response
    job = db.session.get(BackgroundJob, job_id)
    if job is None:
        return api_error("RESOURCE_NOT_FOUND", "BackgroundJob 不存在。", 404)
    if not can_mutate_job(user, job):
        return api_error("PERMISSION_DENIED", "无权重试该后台任务。", 403)
    if job.status != "failed":
        return api_error("VALIDATION_ERROR", "只有 failed 任务可以手动重试。", 400, {"status": job.status})
    job.status = "queued"
    job.error_code = ""
    job.error_message = ""
    job.progress_current = 0
    job.progress_message = "Retry queued"
    job.finished_at = ""
    job.canceled_at = ""
    job.updated_at = current_time_text()
    add_job_event(job, "retry_queued", f"Job retry queued by user {user.id}.")
    db.session.commit()
    return api_success({"job": serialize_background_job(job)}, "后台任务已重新排队。")


@app.route("/api/knowledge/documents", methods=["GET"])
def get_knowledge_documents():
    """
    教师端：读取课程知识库文档列表。
    可选参数：
    - course
    """
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    course = request.args.get("course", "").strip()
    knowledge_base_type = request.args.get("knowledge_base_type", "").strip()
    language = request.args.get("language", "").strip()
    owner_user_id = request.args.get("owner_user_id", "").strip()

    query = KnowledgeDocument.query
    if user.role == "student":
        joined_course_names = [
            course_obj.name for course_obj in Course.query.join(CourseMember, Course.id == CourseMember.course_id)
            .filter(CourseMember.user_id == user.id).all()
        ]
        query = query.filter(db.or_(
            KnowledgeDocument.visibility == "global",
            db.and_(KnowledgeDocument.visibility == "course", KnowledgeDocument.course.in_(joined_course_names or ["__none__"])),
            db.and_(KnowledgeDocument.visibility == "private", KnowledgeDocument.owner_user_id == str(user.id))
        ))
    elif user.role == "teacher":
        manageable_names = [
            course_obj.name for course_obj in Course.query.all()
            if can_manage_course(user, course_obj)
        ]
        query = query.filter(db.or_(
            KnowledgeDocument.visibility == "global",
            KnowledgeDocument.course.in_(manageable_names or ["__none__"]),
            KnowledgeDocument.owner_user_id == str(user.id)
        ))

    if course:
        if user.role == "student" and not can_view_course_name(user, course):
            return jsonify({"status": "error", "message": "无权查看该课程知识库文档。"}), 403
        if user.role == "teacher" and not can_manage_course_name(user, course):
            return jsonify({"status": "error", "message": "无权管理该课程知识库文档。"}), 403
        query = query.filter_by(course=course)

    if knowledge_base_type:
        query = query.filter_by(knowledge_base_type=knowledge_base_type)

    if language:
        query = query.filter_by(language=language)

    if owner_user_id:
        if user.role != "admin" and str(owner_user_id) != str(user.id):
            return jsonify({"status": "error", "message": "无权查看其他用户的个人知识库文档。"}), 403
        query = query.filter_by(owner_user_id=owner_user_id)

    documents = query.order_by(KnowledgeDocument.id.desc()).all()

    return jsonify({
        "status": "success",
        "count": len(documents),
        "documents": [serialize_knowledge_document(doc) for doc in documents]
    })


@app.route("/api/knowledge/entries", methods=["GET"])
def get_knowledge_entries():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    course = request.args.get("course", "").strip()
    language = request.args.get("language", "").strip()
    q = request.args.get("q", "").strip()

    query = KnowledgeBaseEntry.query
    if course:
        if user.role == "student" and not can_view_course_name(user, course):
            return jsonify({"status": "error", "message": "无权查看该课程知识条目。"}), 403
        if user.role == "teacher" and not can_manage_course_name(user, course):
            return jsonify({"status": "error", "message": "无权管理该课程知识条目。"}), 403
        query = query.filter_by(course=course)
    elif user.role == "student":
        joined_course_names = [
            course_obj.name for course_obj in Course.query.join(CourseMember, Course.id == CourseMember.course_id)
            .filter(CourseMember.user_id == user.id).all()
        ]
        query = query.filter(db.or_(
            KnowledgeBaseEntry.course == "",
            KnowledgeBaseEntry.course == "Global",
            KnowledgeBaseEntry.course.in_(joined_course_names or ["__none__"])
        ))
    elif user.role == "teacher":
        manageable_names = [
            course_obj.name for course_obj in Course.query.all()
            if can_manage_course(user, course_obj)
        ]
        query = query.filter(db.or_(
            KnowledgeBaseEntry.course == "",
            KnowledgeBaseEntry.course == "Global",
            KnowledgeBaseEntry.course.in_(manageable_names or ["__none__"])
        ))
    if language:
        query = query.filter_by(language=language)

    entries = query.order_by(KnowledgeBaseEntry.language.asc(), KnowledgeBaseEntry.term.asc()).all()
    if q:
        ranked = search_structured_kb(q, course=course, language=language, limit=50)
        terms = {item["term"] for item in ranked}
        entries = [entry for entry in entries if entry.term in terms]

    return jsonify({
        "status": "success",
        "count": len(entries),
        "entries": [
            {
                "id": entry.id,
                "language": entry.language,
                "course": entry.course,
                "chapter": entry.chapter,
                "term": entry.term,
                "definition": entry.definition,
                "source": entry.source,
                "keywords": entry.keywords
            }
            for entry in entries
        ]
    })


@app.route("/api/knowledge/seed-demo", methods=["POST"])
def seed_demo_knowledge():
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response
    created = seed_demo_knowledge_base()
    return jsonify({
        "status": "success",
        "message": f"Demo knowledge base ready. Created {created} entries.",
        "created": created
    })

@app.route("/api/knowledge/upload", methods=["POST"])
def upload_knowledge_document():
    """
    教师端上传课程知识库资料：
    1. 接收 PDF / DOCX / PPTX
    2. 解析文本
    3. 切分为知识片段
    4. 写入 KnowledgeDocument 和 KnowledgeChunk
    """
    audit_context = get_route_audit_context()
    file = request.files.get("file")

    if file is None:
        return jsonify({
            "status": "error",
            "message": "没有收到知识库文件。"
        }), 400

    if file.filename == "":
        return jsonify({
            "status": "error",
            "message": "文件名为空。"
        }), 400

    if not allowed_file(file.filename):
        return jsonify({
            "status": "error",
            "message": "文件格式不支持，目前只支持 PDF、DOCX、PPTX。"
        }), 400

    course = request.form.get("course", "").strip()
    title = request.form.get("title", "").strip()
    language = request.form.get("language", "zh").strip()
    source_type = request.form.get("source_type", "教师上传资料").strip()
    knowledge_base_type = normalize_knowledge_base_type(
        request.form.get("knowledge_base_type", ""),
        language
    )
    owner_user_id = request.form.get("owner_user_id", "").strip()
    visibility = visibility_for_kb_type(knowledge_base_type)

    current_user = get_current_user()
    if AUTH_REQUIRED:
        if knowledge_base_type == "student_personal_kb":
            user, error_response = require_current_user({"student", "teacher", "admin"})
        else:
            user, error_response = require_current_user({"teacher", "admin"})

        if error_response:
            return attach_request_id_to_response(error_response, audit_context)

        current_user = user
        audit_context = get_route_audit_context(current_user)

    if current_user is not None and knowledge_base_type == "student_personal_kb":
        owner_user_id = str(current_user.id)

    if not course:
        return jsonify({
            "status": "error",
            "message": "课程不能为空。"
        }), 400

    course_obj = Course.query.filter_by(name=course).first()
    if knowledge_base_type != "student_personal_kb":
        if course_obj is None:
            return jsonify({"status": "error", "message": "课程知识库上传必须绑定已存在课程。"}), 400
        if current_user is not None and not can_manage_course(current_user, course_obj):
            return jsonify({"status": "error", "message": "无权管理该课程知识库。"}), 403
    elif current_user is not None and current_user.role == "student" and course_obj is not None and not is_course_member(current_user, course_obj.id):
        return jsonify({"status": "error", "message": "个人知识库只能关联已加入课程。"}), 403

    original_filename = secure_filename(file.filename)

    if not title:
        title = file.filename

    ext = file.filename.rsplit(".", 1)[1].lower()

    unique_prefix = datetime.now().strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:8]
    saved_filename = f"knowledge_{unique_prefix}_{original_filename}"
    save_path = os.path.join(UPLOAD_FOLDER, saved_filename)

    file.save(save_path)
    magic_ok, magic_error = validate_upload_magic(original_filename, save_path)
    if not magic_ok:
        try:
            os.remove(save_path)
        except OSError:
            pass
        return jsonify({
            "status": "error",
            "message": "知识库文件校验失败，文件内容类型与扩展名不一致。",
            "error": magic_error
        }), 415

    parse_result = document_parse_quality_service.parse_document_with_quality(
        save_path,
        filename=original_filename,
        mime_type=file.mimetype or "",
        now_fn=current_time_text,
    )
    parse_record, parse_blocks = persist_document_parse_result(
        parse_result,
        audit_context=audit_context,
        stored_path=saved_filename,
    )
    parse_summary = parse_quality_summary(parse_record)
    if not document_parse_quality_service.should_allow_term_extraction(parse_record):
        blocked_reason = parse_record.error_message or f"Blocked by parse quality gate: {parse_record.quality_status}"
        record_document_ingestion_audit(
            parse_record,
            parse_record,
            "blocked",
            blocked_reason=blocked_reason,
            audit_context=audit_context,
            commit=True,
            target_type="document_parse_record",
            target_uid=parse_record.parse_uid,
        )
        return api_error_with_audit_context(
            quality_gate_error_code(parse_record),
            "知识库文件已解析，但被解析质量门禁阻止，未进入知识库入库。",
            422,
            audit_context,
            {
                **parse_summary,
                "ingestion_status": "blocked",
                "blocked_by_quality_gate": True,
                "blocked_reason": blocked_reason,
                "chunks_created": 0,
            },
        )

    extracted_text = parse_result.raw_text
    valid_parse_blocks = [
        block for block in parse_blocks
        if clean_text(getattr(block, "text", ""))
        and not contains_ocr_placeholder(clean_text(getattr(block, "text", "")))
        and not contains_formula_placeholder(clean_text(getattr(block, "text", "")))
    ]

    if len(valid_parse_blocks) == 0:
        blocked_reason = "文件解析成功，但没有得到有效知识片段；OCR 占位符、公式占位符或空文本不会进入知识库。"
        record_document_ingestion_audit(
            parse_record,
            parse_record,
            "blocked",
            blocked_reason=blocked_reason,
            audit_context=audit_context,
            commit=True,
            target_type="document_parse_record",
            target_uid=parse_record.parse_uid,
        )
        return api_error_with_audit_context(
            "PARSING_FAILED",
            blocked_reason,
            422,
            audit_context,
            {
                **parse_summary,
                "ingestion_status": "blocked",
                "blocked_by_quality_gate": True,
                "blocked_reason": blocked_reason,
                "chunks_created": 0,
            },
        )

    if knowledge_base_type == "student_personal_kb" and current_user is not None:
        allowed, reasons, plan, subscription, totals = check_usage_quota(
            current_user,
            page_units=max(1, len(valid_parse_blocks)),
            ai_units=0
        )
        if not allowed:
            return jsonify({
                "status": "error",
                "message": "个人知识库上传额度不足，请升级会员。",
                "reasons": reasons
            }), 402

    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    document = KnowledgeDocument(
        course=course,
        title=title,
        filename=file.filename,
        saved_filename=saved_filename,
        file_type=ext,
        language=language or "zh",
        source_type=source_type or "教师上传资料",
        knowledge_base_type=knowledge_base_type,
        owner_user_id=owner_user_id,
        visibility=visibility,
        parse_uid=parse_record.parse_uid,
        text_length=len(extracted_text),
        chunk_count=0,
        created_at=now_text
    )

    db.session.add(document)
    db.session.flush()

    knowledge_metadata = build_governed_ingestion_metadata(
        parse_record=parse_record,
        title=title,
        course=course_obj,
        course_name=course,
        chapter=str(request.form.get("chapter", "")).strip(),
        language=language or "unknown",
        source_type=source_type,
        scope_type="personal" if knowledge_base_type == "student_personal_kb" else "course",
        owner_user=current_user,
        owner_user_id=owner_user_id,
        document_id=document.id,
        knowledge_base_type=knowledge_base_type,
        visibility=visibility,
        content_hash="",
    )
    governed_ingestion = knowledge_ingestion_service.ingest_parse_record_to_governed_knowledge(
        db.session,
        knowledge_ingestion_models(),
        parse_record,
        valid_parse_blocks,
        knowledge_metadata,
        audit_context=audit_context,
        now_fn=current_time_text,
        commit=False,
    )
    source = governed_ingestion.source
    chunk_records = governed_ingestion.chunks
    document.chunk_count = len(chunk_records)
    document.text_length = sum(len(getattr(chunk, "content", "") or "") for chunk in chunk_records)

    if knowledge_base_type == "student_personal_kb" and current_user is not None:
        record_usage(current_user.id, "document_parse_page", max(1, len(chunk_records)), related_document_id=document.id)

    embedded_count = 0
    job = TaskJob(
        job_type="knowledge_parse_and_embed",
        status="completed",
        payload_json=json.dumps({
            "document_id": document.id,
            "knowledge_base_type": knowledge_base_type,
            "chunk_count": len(chunk_records)
        }, ensure_ascii=False),
        result_json=json.dumps({
            "embedded_count": embedded_count,
            "embedding_status": "not_generated_in_governance_gate"
        }, ensure_ascii=False),
        created_at=now_text,
        updated_at=current_time_text()
    )
    db.session.add(job)
    ingestion_status = "partial" if parse_record.quality_status == "partial_text" else "ingested"
    record_document_ingestion_audit(
        document,
        parse_record,
        ingestion_status,
        audit_context=audit_context,
        commit=False,
        target_type="knowledge_document",
        target_uid=document.id,
    )
    db.session.commit()

    return jsonify({
        "status": "success",
        "request_id": audit_context.get("request_id", ""),
        "message": f"知识库资料上传成功，已解析 {len(extracted_text)} 个字符，生成 {len(chunk_records)} 个受治理知识片段。本任务不生成 embedding。",
        **parse_summary,
        "ingestion_status": ingestion_status,
        "source_uid": getattr(source, "source_uid", ""),
        "chunk_count": len(chunk_records),
        "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in chunk_records[:20] if getattr(chunk, "chunk_uid", "")],
        "document": serialize_knowledge_document(document),
        "embedding_count": embedded_count,
        "job_id": job.id,
        "preview": extracted_text[:800],
        "knowledge_source": serialize_knowledge_source(source),
        "sample_chunks": [
            serialize_knowledge_chunk(chunk) for chunk in chunk_records[:3]
        ]
    })

@app.route("/api/knowledge/search", methods=["GET"])
def search_knowledge_chunks():
    """
    知识库片段检索接口。
    用于后续 AI 翻译前，从课程知识库中找中文依据。

    参数：
    - q: 搜索关键词，必填
    - course: 课程，可选
    - limit: 返回数量，默认 8
    """
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    q = request.args.get("q", "").strip()
    course = request.args.get("course", "").strip()
    course_id_raw = request.args.get("course_id", "").strip()
    knowledge_base_type = request.args.get("knowledge_base_type", "").strip()
    language = request.args.get("language", "").strip()
    scope_type = request.args.get("scope_type", "course").strip().lower() or "course"
    owner_user_id = request.args.get("owner_user_id", "").strip()
    discipline = request.args.get("discipline", "").strip() or None
    kb_version_id_raw = request.args.get("knowledge_base_version_id", request.args.get("kb_version_id", "")).strip()
    retrieval_backend = request.args.get("retrieval_backend", RETRIEVAL_BACKEND).strip().lower() or "lexical"
    include_debug = request.args.get("include_debug", "false").strip().lower() == "true"
    limit_raw = request.args.get("limit", "8").strip()

    if not q:
        return api_error("VALIDATION_ERROR", "搜索关键词 q 不能为空。", 400)
    if retrieval_backend not in VALID_RETRIEVAL_BACKENDS:
        return api_error("VALIDATION_ERROR", "retrieval_backend must be lexical, vector, hybrid, or hybrid_rerank.", 400)
    if include_debug and user.role == "student":
        return api_error("PERMISSION_DENIED", "Student accounts cannot request retrieval debug details.", 403)

    try:
        limit = int(limit_raw)
    except ValueError:
        limit = 8

    limit = max(1, min(limit, 30))

    course_obj = get_course_by_id_or_name(course_id_raw, course)
    if scope_type not in {"course", "personal", "global"}:
        return api_error("VALIDATION_ERROR", "scope_type 只能是 course、personal 或 global。", 400)
    if scope_type == "course":
        if course_obj is None:
            return api_error("VALIDATION_ERROR", "course scope 检索必须提供有效 course_id 或 course。", 400)
        if user.role == "student" and not is_course_member(user, course_obj.id):
            return api_error("PERMISSION_DENIED", "无权检索该课程知识库。", 403)
        if user.role == "teacher" and not can_manage_course(user, course_obj):
            return api_error("PERMISSION_DENIED", "无权检索该课程知识库。", 403)
    if scope_type == "personal":
        target_owner = owner_user_id or str(user.id)
        if user.role != "admin" and str(target_owner) != str(user.id):
            return api_error("PERMISSION_DENIED", "无权检索其他用户的个人知识库。", 403)
        owner_user_id = target_owner
        knowledge_base_type = "student_personal_kb"
    elif owner_user_id:
        return api_error("VALIDATION_ERROR", "owner_user_id 只能用于 personal scope。", 400)

    if not language:
        if knowledge_base_type == "en_course_kb":
            language = "en"
        elif knowledge_base_type == "zh_course_kb":
            language = "zh"
        elif knowledge_base_type == "student_personal_kb":
            language = "bilingual"
    if not knowledge_base_type:
        knowledge_base_type = "en_course_kb" if language == "en" else "zh_course_kb" if language == "zh" else ""

    top_results = retrieve_evidence_results(
        q,
        course_id=course_obj.id if course_obj else None,
        course_name=course_obj.name if course_obj else course,
        language=language,
        scope_type=scope_type,
        owner_user_id=owner_user_id or (user.id if scope_type == "personal" else None),
        limit=limit,
        knowledge_base_type=knowledge_base_type,
        discipline=discipline,
        knowledge_base_version_id=int(kb_version_id_raw) if kb_version_id_raw.isdigit() else None,
        retrieval_backend=retrieval_backend
    )
    if user.role == "admin":
        audited_users = set()
        for item in top_results:
            if item.get("visibility") == "private" and item.get("owner_user_id") and str(item.get("owner_user_id")) != str(user.id):
                audited_users.add(str(item.get("owner_user_id")))
        for target_id in audited_users:
            record_personal_access(user, int(target_id), "knowledge_search", None, f"admin searched personal KB for query={q[:80]}")

    if user.role == "student":
        record_usage(user.id, "knowledge_search", 1)
        db.session.commit()
    elif user.role == "admin" and audited_users:
        db.session.commit()

    return jsonify({
        "status": "success",
        "query": q,
        "course": course_obj.name if course_obj else course,
        "course_id": course_obj.id if course_obj else None,
        "knowledge_base_type": knowledge_base_type,
        "language": language,
        "retrieval_backend": retrieval_backend,
        "count": len(top_results),
        "retrieval_version": RETRIEVAL_VERSION,
        "data": {
            "items": top_results,
            "message": "" if top_results else "No evidence passed the relevance threshold."
        },
        "results": top_results
    })


@app.route("/api/knowledge/rebuild-embeddings", methods=["POST"])
def rebuild_knowledge_embeddings():
    if AUTH_REQUIRED:
        user, error_response = require_current_user({"teacher", "admin"})
        if error_response:
            return error_response

    data = request.get_json() or {}
    course = str(data.get("course", "")).strip()
    knowledge_base_type = str(data.get("knowledge_base_type", "")).strip()

    query = KnowledgeChunk.query
    if course:
        query = query.filter_by(course=course)
    if knowledge_base_type:
        query = query.filter_by(knowledge_base_type=knowledge_base_type)

    chunks = query.all()
    count = rebuild_embeddings_for_chunks(chunks)

    job = TaskJob(
        job_type="rebuild_embeddings",
        status="completed",
        payload_json=json.dumps({
            "course": course,
            "knowledge_base_type": knowledge_base_type
        }, ensure_ascii=False),
        result_json=json.dumps({
            "embedded_count": count
        }, ensure_ascii=False),
        created_at=current_time_text(),
        updated_at=current_time_text()
    )
    db.session.add(job)
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": f"已重建 {count} 个知识片段的 embedding。",
        "embedded_count": count,
        "job_id": job.id
    })


@app.route("/api/jobs", methods=["GET"])
def list_jobs():
    if AUTH_REQUIRED:
        user, error_response = require_current_user({"teacher", "admin"})
        if error_response:
            return error_response

    jobs = TaskJob.query.order_by(TaskJob.id.desc()).limit(50).all()

    return jsonify({
        "status": "success",
        "count": len(jobs),
        "jobs": [
            {
                "id": job.id,
                "job_type": job.job_type,
                "status": job.status,
                "payload": json.loads(job.payload_json or "{}"),
                "result": json.loads(job.result_json or "{}"),
                "error": job.error,
                "created_at": job.created_at,
                "updated_at": job.updated_at
            }
            for job in jobs
        ]
    })


@app.route("/api/ai/term-suggestion", methods=["POST"])
def ai_term_suggestion():
    """
    AI 证据对齐：
    基于英文教材知识库、中文教材知识库和当前术语上下文生成可追溯的术语对齐结果。
    API Key 只在后端环境变量中读取，不暴露给前端。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    data = request.get_json() or {}

    term_id = data.get("term_id")
    save_to_term = bool(data.get("save_to_term", True))
    auto_publish = bool(data.get("auto_publish", False))

    term = None
    if term_id:
        term = db.session.get(Term, int(term_id))

        if term is None:
            return jsonify({
                "status": "error",
                "message": "术语不存在。"
            }), 404
        if not can_manage_course_name(user, term.course):
            return jsonify({"status": "error", "message": "无权处理该课程术语。"}), 403

    english_term = str(data.get("english_term") or (term.english_term if term else "")).strip()
    context = str(data.get("context") or ((term.courseware_sentence or term.context) if term else "")).strip()
    course = str(data.get("course") or (term.course if term else "")).strip()
    chapter = str(data.get("chapter") or (term.chapter if term else "")).strip()

    if not english_term:
        return jsonify({
            "status": "error",
            "message": "英文术语不能为空。"
        }), 400

    suggestion = generate_alignment_result(
        english_term=english_term,
        courseware_sentence=context,
        course=course,
        chapter=chapter
    )

    if term is not None and save_to_term:
        apply_ai_suggestion_to_term(term, suggestion, auto_publish=auto_publish)
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": "AI 证据对齐结果已生成。",
        "suggestion": suggestion,
        "term": serialize_term(term) if term is not None else None
    })


@app.route("/api/ai/term-suggestions/batch", methods=["POST"])
def batch_ai_term_suggestions():
    """
    批量生成 AI 证据对齐结果。
    为避免一次上传后产生过多 API 调用，当前本地演示版每次最多处理 10 条。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    data = request.get_json() or {}
    term_ids = data.get("term_ids", [])
    auto_publish = bool(data.get("auto_publish", False))

    if not isinstance(term_ids, list) or len(term_ids) == 0:
        return jsonify({
            "status": "error",
            "message": "没有收到需要生成 AI 证据对齐结果的术语。"
        }), 400

    limited_ids = term_ids[:10]
    results = []
    failed = []

    for term_id in limited_ids:
        term = db.session.get(Term, int(term_id))

        if term is None:
            failed.append({
                "id": term_id,
                "reason": "术语不存在"
            })
            continue
        if not can_manage_course_name(user, term.course):
            failed.append({
                "id": term_id,
                "english_term": term.english_term,
                "reason": "无权处理该课程术语"
            })
            continue

        try:
            suggestion = generate_alignment_result(
                english_term=term.english_term,
                courseware_sentence=term.courseware_sentence or term.context,
                course=term.course,
                chapter=term.chapter
            )
            apply_ai_suggestion_to_term(term, suggestion, auto_publish=auto_publish)
            results.append({
                "id": term.id,
                "english_term": term.english_term,
                "chinese_term": suggestion["final_chinese_term"],
                "confidence": suggestion["confidence_score"],
                "review_status": suggestion["review_status"]
            })
        except Exception as exc:
            failed.append({
                "id": term.id,
                "english_term": term.english_term,
                "reason": str(exc)
            })

    db.session.commit()

    return jsonify({
        "status": "success",
        "message": f"批量 AI 证据对齐完成，成功 {len(results)} 条，失败 {len(failed)} 条。",
        "processed_count": len(results),
        "failed_count": len(failed),
        "skipped_count": max(0, len(term_ids) - len(limited_ids)),
        "results": results,
        "failed": failed
    })


@app.route("/api/ai/student-answer", methods=["POST"])
def ai_student_answer():
    """
    学生侧 AI 问答：
    综合英文课程知识库、中文课程知识库和学生个人知识库回答问题。
    """
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    data = request.get_json() or {}

    question = str(data.get("question", "")).strip()
    course = str(data.get("course", "")).strip()
    owner_user_id = str(user.id)
    if course and user.role == "student" and not can_view_course_name(user, course):
        return jsonify({"status": "error", "message": "无权访问该课程学习资料。"}), 403

    if not question:
        return jsonify({
            "status": "error",
            "message": "问题不能为空。"
        }), 400

    try:
        answer = generate_student_learning_answer(
            question=question,
            course=course,
            owner_user_id=owner_user_id
        )
    except Exception as exc:
        return jsonify({
            "status": "error",
            "message": "AI 学习回答生成失败。",
            "error": str(exc)
        }), 502

    return jsonify({
        "status": "success",
        "message": "AI 学习回答已生成。",
        "answer": answer
    })


@app.route("/api/upload", methods=["POST"])
def upload_file():
    """
    教师上传文件：
    1. 接收 PDF / DOCX / PPTX
    2. 保存到本地上传目录
    3. 解析文字
    4. 抽取候选术语
    5. 检索英文教材证据、生成中文候选、检索中文教材证据
    6. 写入 SQLite 数据库并返回 AI 证据对齐结果
    """
    audit_context = get_route_audit_context()
    user = None
    if AUTH_REQUIRED:
        user, error_response = require_current_user({"teacher", "admin"})
        if error_response:
            return attach_request_id_to_response(error_response, audit_context)
        audit_context = get_route_audit_context(user)

    file = request.files.get("file")

    if file is None:
        return jsonify({
            "status": "error",
            "message": "没有收到文件。"
        }), 400

    if file.filename == "":
        return jsonify({
            "status": "error",
            "message": "文件名为空。"
        }), 400

    if not allowed_file(file.filename):
        return jsonify({
            "status": "error",
            "message": "文件格式不支持，目前只支持 PDF、DOCX、PPTX。"
        }), 400

    course = request.form.get("course", "Data Structures and Algorithms").strip()
    chapter = request.form.get("chapter", "Chapter 4 - Hashing").strip()

    if not course:
        course = "Data Structures and Algorithms"

    if not chapter:
        chapter = "未指定章节"

    if AUTH_REQUIRED:
        course_obj = Course.query.filter_by(name=course).first()
        if course_obj is None:
            return jsonify({"status": "error", "message": "旧上传接口必须绑定已存在课程。"}), 400
        if not can_manage_course(user, course_obj):
            return jsonify({"status": "error", "message": "无权向该课程上传课件。"}), 403

    original_filename = secure_filename(file.filename)
    unique_prefix = datetime.now().strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:8]
    filename = f"{unique_prefix}_{original_filename}"
    save_path = os.path.join(UPLOAD_FOLDER, filename)

    file.save(save_path)
    magic_ok, magic_error = validate_upload_magic(original_filename, save_path)
    if not magic_ok:
        try:
            os.remove(save_path)
        except OSError:
            pass
        return jsonify({
            "status": "error",
            "message": "课件文件校验失败，文件内容类型与扩展名不一致。",
            "error": magic_error
        }), 415

    parse_result = document_parse_quality_service.parse_document_with_quality(
        save_path,
        filename=original_filename,
        mime_type=file.mimetype or "",
        now_fn=current_time_text,
    )
    parse_record, parse_blocks = persist_document_parse_result(
        parse_result,
        audit_context=audit_context,
        stored_path=filename,
    )
    parse_summary = parse_quality_summary(parse_record)
    if not document_parse_quality_service.should_allow_term_extraction(parse_record):
        blocked_reason = parse_record.error_message or f"Blocked by parse quality gate: {parse_record.quality_status}"
        record_document_ingestion_audit(
            parse_record,
            parse_record,
            "blocked",
            blocked_reason=blocked_reason,
            audit_context=audit_context,
            commit=True,
            target_type="courseware_upload",
            target_uid=parse_record.parse_uid,
        )
        return api_error_with_audit_context(
            quality_gate_error_code(parse_record),
            "课件已解析，但被解析质量门禁阻止，未进入术语抽取。",
            422,
            audit_context,
            {
                **parse_summary,
                "ingestion_status": "blocked",
                "blocked_by_quality_gate": True,
                "blocked_reason": blocked_reason,
                "terms": [],
            },
        )

    extracted_text = parse_result.raw_text
    preview = extracted_text[:800]
    terms = extract_terms_from_text(extracted_text)
    partial_quality = parse_record.quality_status == "partial_text"
    parse_quality_metadata = parse_quality_risk_service.build_parse_quality_metadata(parse_record)
    input_risk_labels = parse_quality_risk_service.parse_quality_to_risk_labels(parse_quality_metadata)

    upload_record = CoursewareUpload(
        filename=original_filename,
        saved_filename=filename,
        parse_uid=parse_record.parse_uid,
        course=course,
        chapter=chapter,
        uploaded_by="demo_teacher",
        upload_time=current_time_text(),
        parsed_text=extracted_text
    )
    db.session.add(upload_record)
    db.session.commit()

    course_record = Course.query.filter_by(name=course).first()
    valid_parse_blocks = [
        block for block in parse_blocks
        if clean_text(getattr(block, "text", ""))
        and not contains_ocr_placeholder(clean_text(getattr(block, "text", "")))
        and not contains_formula_placeholder(clean_text(getattr(block, "text", "")))
    ]
    try:
        knowledge_metadata = build_governed_ingestion_metadata(
            parse_record=parse_record,
            title=original_filename,
            course=course_record,
            course_name=course,
            chapter=chapter,
            language="en",
            source_type="course_material",
            scope_type="course",
            owner_user=user,
            owner_user_id=getattr(user, "id", None),
            document_id=0,
            knowledge_base_type="en_course_kb",
            visibility="course",
            content_hash="",
        )
        governed_ingestion = knowledge_ingestion_service.ingest_parse_record_to_governed_knowledge(
            db.session,
            knowledge_ingestion_models(),
            parse_record,
            valid_parse_blocks,
            knowledge_metadata,
            audit_context=audit_context,
            now_fn=current_time_text,
            commit=False,
        )
        governed_source = governed_ingestion.source
        governed_chunks = governed_ingestion.chunks
    except knowledge_ingestion_service.KnowledgeIngestionBlockedError as exc:
        db.session.rollback()
        blocked_reason = exc.blocked_reason
        return api_error_with_audit_context(
            "PARSING_FAILED",
            "课件已解析，但没有可进入知识库的有效治理知识块。",
            422,
            audit_context,
            {
                **parse_summary,
                "ingestion_status": "blocked",
                "blocked_by_quality_gate": True,
                "blocked_reason": blocked_reason,
                "terms": [],
            },
        )

    first_governed_chunk = governed_chunks[0] if governed_chunks else None

    if len(terms) == 0:
        ingestion_status = "partial" if partial_quality else "ingested"
        record_document_ingestion_audit(
            upload_record,
            parse_record,
            ingestion_status,
            audit_context=audit_context,
            commit=True,
            target_type="courseware_upload",
            target_uid=upload_record.id,
        )
        return jsonify({
            "status": "success",
            "request_id": audit_context.get("request_id", ""),
            "message": f"文件上传并解析成功，共提取 {len(extracted_text)} 个字符，但没有抽取到候选术语。旧的待审核对齐结果不会被清空。",
            **parse_summary,
            "ingestion_status": ingestion_status,
            "risk_labels": input_risk_labels,
            "source_uid": getattr(governed_source, "source_uid", ""),
            "chunk_count": len(governed_chunks),
            "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in governed_chunks[:20] if getattr(chunk, "chunk_uid", "")],
            "filename": original_filename,
            "saved_filename": filename,
            "upload_id": upload_record.id,
            "text_length": len(extracted_text),
            "preview": preview,
            "terms": [],
            "upload_folder": UPLOAD_FOLDER
        })

    # 只清空同一课程、同一章节下的 pending，避免误删其他章节待审核内容
    Term.query.filter_by(course=course, chapter=chapter, status="pending").delete()

    saved_terms = []

    for index, item in enumerate(terms):
        # 已通过质量控制的同名术语不再重复进入 pending
        approved_existing = Term.query.filter_by(
            course=course,
            chapter=chapter,
            english_term=item["english_term"],
            status="approved"
        ).first()

        if approved_existing:
            continue

        courseware_sentence = item.get("context", "")
        alignment = generate_alignment_result(
            english_term=item["english_term"],
            courseware_sentence=courseware_sentence,
            course=course,
            chapter=chapter
        ) if index < 12 else {}

        risk_note = alignment.get("risk_note", "")
        if input_risk_labels:
            risk_note = "; ".join(part for part in [risk_note, f"parse_quality: {parse_record.quality_status}"] if part)
        confidence = int(alignment.get("confidence_score", item.get("confidence", 60)))
        if input_risk_labels:
            confidence = min(confidence, 79)
        term_record = Term(
            course=course,
            chapter=chapter,
            english_term=item["english_term"],
            chinese_term=alignment.get("final_chinese_term") or item.get("chinese_term", "待质量控制"),
            explanation=alignment.get("explanation") or item.get("explanation", "待质量控制：系统已完成候选术语抽取，但尚未得到足够教材证据。"),
            context=courseware_sentence,
            courseware_sentence=courseware_sentence,
            english_kb_evidence=alignment.get("english_kb_evidence", ""),
            ai_translation_candidate=alignment.get("ai_translation_candidate", ""),
            chinese_kb_evidence=alignment.get("chinese_kb_evidence", ""),
            final_chinese_term=alignment.get("final_chinese_term", ""),
            alignment_reason=alignment.get("alignment_reason", ""),
            review_status=alignment.get("review_status", "pending"),
            confidence=confidence,
            status="pending",
            ai_status=ai_status_from_confidence(confidence),
            risk_note=risk_note,
            parse_uid=parse_record.parse_uid,
            parse_block_uid=parse_blocks[0].block_uid if parse_blocks else "",
            parse_quality_status=parse_record.quality_status,
            parse_quality_flags=parse_record.quality_flags,
            input_risk_labels=json.dumps(input_risk_labels, ensure_ascii=False),
            source_uid=getattr(governed_source, "source_uid", ""),
            chunk_uid=getattr(first_governed_chunk, "chunk_uid", ""),
            ai_model=alignment.get("ai_model", "")
        )

        db.session.add(term_record)
        saved_terms.append(term_record)

    db.session.flush()
    if input_risk_labels:
        for term_record in saved_terms:
            record_parse_quality_risk_audit(
                "parse_quality_risk_propagated",
                "term",
                term_record.id,
                parse_uid=parse_record.parse_uid,
                quality_status=parse_record.quality_status,
                risk_labels=input_risk_labels,
                forced_status=term_record.status,
                audit_context=audit_context,
                commit=False,
            )
    ingestion_status = "partial" if partial_quality else "ingested"
    record_document_ingestion_audit(
        upload_record,
        parse_record,
        ingestion_status,
        audit_context=audit_context,
        commit=False,
        target_type="courseware_upload",
        target_uid=upload_record.id,
    )
    db.session.commit()

    course_obj = course_record or Course.query.filter_by(name=course).first()
    if course_obj is not None:
        for term in saved_terms:
            sync_term_to_card(term, course_id=course_obj.id)
        db.session.commit()

    saved_terms_json = [serialize_term(term) for term in saved_terms]

    return jsonify({
        "status": "success",
        "request_id": audit_context.get("request_id", ""),
        "message": f"文件上传、解析并完成证据对齐，共提取 {len(extracted_text)} 个字符，生成 {len(saved_terms_json)} 条待审核对齐结果。",
        **parse_summary,
        "ingestion_status": ingestion_status,
        "risk_labels": input_risk_labels,
        "source_uid": getattr(governed_source, "source_uid", ""),
        "chunk_count": len(governed_chunks),
        "chunk_uids": [getattr(chunk, "chunk_uid", "") for chunk in governed_chunks[:20] if getattr(chunk, "chunk_uid", "")],
        "filename": original_filename,
        "saved_filename": filename,
        "upload_id": upload_record.id,
        "text_length": len(extracted_text),
        "preview": preview,
        "terms": saved_terms_json,
        "upload_folder": UPLOAD_FOLDER
    })


@app.route("/api/terms/pending", methods=["GET"])
def get_pending_terms():
    """
    教师端：读取待审核 AI 证据对齐结果。
    可选参数：
    - course
    - chapter
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    course = request.args.get("course", "").strip()
    chapter = request.args.get("chapter", "").strip()

    query = Term.query.filter_by(status="pending")

    if course:
        query = query.filter_by(course=course)

    if chapter:
        query = query.filter_by(chapter=chapter)

    terms = query.order_by(Term.id.desc()).all()
    if user.role != "admin":
        terms = [term for term in terms if can_manage_course_name(user, term.course)]

    return jsonify({
        "status": "success",
        "count": len(terms),
        "terms": [serialize_term(term) for term in terms]
    })


@app.route("/api/terms/<int:term_id>/approve", methods=["POST"])
def approve_term(term_id):
    """
    质量控制通过术语：
    1. 根据 ID 找到候选术语
    2. 更新中文译名和解释
    3. 将状态改为 approved
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    term = db.session.get(Term, term_id)
    if term is None:
        return jsonify({"status": "error", "message": "术语不存在。"}), 404
    if not can_manage_course_name(user, term.course):
        return jsonify({"status": "error", "message": "无权处理该课程术语。"}), 403

    data = request.get_json() or {}
    chinese_term = str(data.get("chinese_term", "")).strip()
    explanation = str(data.get("explanation", "")).strip()

    if not chinese_term:
        return jsonify({"status": "error", "message": "中文译名不能为空。"}), 400
    if not explanation:
        return jsonify({"status": "error", "message": "解释不能为空。"}), 400

    term.chinese_term = chinese_term
    term.explanation = explanation
    term.status = "approved"
    term.review_status = "reviewed"
    if not term.final_chinese_term:
        term.final_chinese_term = chinese_term

    db.session.commit()
    course_obj = Course.query.filter_by(name=term.course).first()
    if course_obj is not None:
        sync_term_to_card(term, course_id=course_obj.id)
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": "术语已通过质量控制。",
        "term": serialize_term(term)
    })

@app.route("/api/terms/batch-approve", methods=["POST"])
def batch_approve_terms():
    """
    教师端批量审核术语：
    接收多个术语 ID，以及每个术语对应的中文译名和解释。
    只有中文译名和解释都不为空的术语才允许批量通过。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    data = request.get_json() or {}
    items = data.get("terms", [])

    if not isinstance(items, list) or len(items) == 0:
        return jsonify({
            "status": "error",
            "message": "没有收到需要批量审核的术语。"
        }), 400

    approved_terms = []
    failed_items = []

    for item in items:
        term_id = item.get("id")
        chinese_term = str(item.get("chinese_term", "")).strip()
        explanation = str(item.get("explanation", "")).strip()

        if not term_id:
            failed_items.append({
                "id": term_id,
                "reason": "缺少术语 ID"
            })
            continue

        term = db.session.get(Term, int(term_id))

        if term is None:
            failed_items.append({
                "id": term_id,
                "reason": "术语不存在"
            })
            continue

        if not can_manage_course_name(user, term.course):
            failed_items.append({
                "id": term_id,
                "english_term": term.english_term,
                "reason": "无权处理该课程术语"
            })
            continue

        if not chinese_term or chinese_term == "待质量控制":
            failed_items.append({
                "id": term_id,
                "english_term": term.english_term,
                "reason": "中文译名为空或仍为待质量控制"
            })
            continue

        if not explanation or explanation.startswith("待质量控制"):
            failed_items.append({
                "id": term_id,
                "english_term": term.english_term,
                "reason": "解释为空或仍为待质量控制"
            })
            continue

        term.chinese_term = chinese_term
        term.explanation = explanation
        term.status = "approved"
        term.review_status = "reviewed"
        if not term.final_chinese_term:
            term.final_chinese_term = chinese_term
        approved_terms.append(term)

    db.session.commit()
    for term in approved_terms:
        course_obj = Course.query.filter_by(name=term.course).first()
        if course_obj is not None:
            sync_term_to_card(term, course_id=course_obj.id)
    if approved_terms:
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": f"批量审核完成，成功审核 {len(approved_terms)} 条，跳过 {len(failed_items)} 条。",
        "approved_count": len(approved_terms),
        "failed_count": len(failed_items),
        "failed_items": failed_items,
        "terms": [serialize_term(term) for term in approved_terms]
    })

@app.route("/api/terms/<int:term_id>", methods=["GET"])
def get_term_detail(term_id):
    """
    读取单个术语详情：
    用于教师根据学生反馈修改术语前，获取当前术语内容。
    """
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    term = db.session.get(Term, term_id)

    if term is None:
        return jsonify({
            "status": "error",
            "message": "术语不存在。"
        }), 404

    if user.role == "student":
        if term.status != "approved" or not can_view_course_name(user, term.course):
            return jsonify({"status": "error", "message": "无权查看该术语。"}), 403
    elif not can_manage_course_name(user, term.course):
        return jsonify({"status": "error", "message": "无权查看该课程术语。"}), 403

    return jsonify({
        "status": "success",
        "term": serialize_term(term)
    })

@app.route("/api/terms/<int:term_id>/update", methods=["POST"])
def update_term(term_id):
    """
    教师端修改术语：
    用于根据学生反馈修正已发布术语的中文译名和专业解释。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    term = db.session.get(Term, term_id)

    if term is None:
        return jsonify({
            "status": "error",
            "message": "术语不存在。"
        }), 404

    if not can_manage_course_name(user, term.course):
        return jsonify({"status": "error", "message": "无权修改该课程术语。"}), 403

    data = request.get_json() or {}

    chinese_term = str(data.get("chinese_term", "")).strip()
    explanation = str(data.get("explanation", "")).strip()

    if not chinese_term:
        return jsonify({
            "status": "error",
            "message": "中文译名不能为空。"
        }), 400

    if not explanation:
        return jsonify({
            "status": "error",
            "message": "专业解释不能为空。"
        }), 400

    term.chinese_term = chinese_term
    term.explanation = explanation
    term.status = "approved"
    term.review_status = "reviewed"
    if not term.final_chinese_term:
        term.final_chinese_term = chinese_term

    db.session.commit()
    course_obj = Course.query.filter_by(name=term.course).first()
    if course_obj is not None:
        sync_term_to_card(term, course_id=course_obj.id)
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": "术语已修改并保持为已发布状态。",
        "term": serialize_term(term)
    })

@app.route("/api/terms/<int:term_id>/reject", methods=["POST"])
def reject_term(term_id):
    """
    教师驳回术语：
    保留记录和证据，标记为 rejected，学生端不可见。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    term = db.session.get(Term, term_id)

    if term is None:
        return jsonify({
            "status": "error",
            "message": "术语不存在。"
        }), 404

    if not can_manage_course_name(user, term.course):
        return jsonify({"status": "error", "message": "无权处理该课程术语。"}), 403

    data = request.get_json() or {}
    reason = str(data.get("reason", "")).strip()

    term.status = "rejected"
    term.review_status = "rejected"
    if reason:
        term.risk_note = reason
        if term.alignment_reason:
            term.alignment_reason = f"{term.alignment_reason}\n教师驳回原因：{reason}"
        else:
            term.alignment_reason = f"教师驳回原因：{reason}"

    db.session.commit()
    course_obj = Course.query.filter_by(name=term.course).first()
    if course_obj is not None:
        sync_term_to_card(term, course_id=course_obj.id)
        db.session.commit()

    return jsonify({
        "status": "success",
        "message": "术语已驳回，学生端不会展示该术语。",
        "term": serialize_term(term)
    })

@app.route("/api/terms/<int:term_id>/delete", methods=["DELETE"])
def delete_term(term_id):
    """
    删除术语：
    可删除误抽取、误审核或测试产生的垃圾数据。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    term = db.session.get(Term, term_id)

    if term is None:
        return jsonify({
            "status": "error",
            "message": "术语不存在。"
        }), 404

    if not can_manage_course_name(user, term.course):
        return jsonify({"status": "error", "message": "无权删除该课程术语。"}), 403

    db.session.delete(term)
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": "术语已删除。",
        "deleted_id": term_id
    })


@app.route("/api/glossary", methods=["GET"])
def get_glossary():
    """
    学生端：读取质量控制通过的词汇表。
    可选参数：
    - course
    - chapter
    - q 搜索英文术语 / 中文译名 / 解释
    """
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    course = request.args.get("course", "").strip()
    chapter = request.args.get("chapter", "").strip()
    q = request.args.get("q", "").strip().lower()

    query = Term.query.filter_by(status="approved")

    if course:
        query = query.filter_by(course=course)

    if chapter:
        query = query.filter_by(chapter=chapter)

    terms = query.order_by(Term.id.desc()).all()
    if user.role == "student":
        terms = [term for term in terms if can_view_course_name(user, term.course)]
    elif user.role == "teacher":
        terms = [term for term in terms if can_manage_course_name(user, term.course)]

    if q:
        terms = [
            term for term in terms
            if q in term.english_term.lower()
            or q in term.chinese_term.lower()
            or q in term.explanation.lower()
            or q in (term.final_chinese_term or "").lower()
            or q in (term.ai_translation_candidate or "").lower()
            or q in (term.english_kb_evidence or "").lower()
            or q in (term.chinese_kb_evidence or "").lower()
            or q in (term.alignment_reason or "").lower()
        ]

    return jsonify({
        "status": "success",
        "count": len(terms),
        "terms": [serialize_term(term) for term in terms]
    })

@app.route("/api/feedback", methods=["POST"])
def submit_feedback():
    user, error_response = require_current_user({"student", "teacher", "admin"})
    if error_response:
        return error_response

    data = request.get_json() or {}
    card_id = data.get("terminology_card_id") or data.get("card_id") or data.get("term_id")
    card = None
    if card_id:
        try:
            card = db.session.get(TerminologyCard, int(card_id))
        except (TypeError, ValueError):
            card = None
    if card is not None:
        scoped = filter_cards_for_user(TerminologyCard.query.filter_by(id=card.id), user).first()
        if scoped is None:
            return api_error("PERMISSION_DENIED", "无权反馈该术语卡。", 403)
        feedback, feedback_error = create_pilot_feedback_for_card(card, user, data, str(data.get("feedback_source") or "student_card_detail"))
        if feedback_error:
            return feedback_error
        db.session.commit()
        return jsonify({
            "status": "success",
            "message": "Feedback submitted.",
            "data": {"feedback_id": feedback.id, "feedback_status": feedback.status},
            "feedback": serialize_feedback(feedback)
        })

    # Legacy Term fallback for old course glossary pages.
    term_id = data.get("term_id")
    content = str(data.get("reported_issue") or data.get("feedback_content") or "").strip()
    if not term_id:
        return api_error("VALIDATION_ERROR", "缺少术语卡片 ID。", 400)
    if not content:
        return api_error("VALIDATION_ERROR", "反馈内容不能为空。", 400)
    term = db.session.get(Term, int(term_id))
    if term is None:
        return api_error("RESOURCE_NOT_FOUND", "反馈对应的术语不存在。", 404)
    if user.role == "student" and not can_view_course_name(user, term.course):
        return api_error("PERMISSION_DENIED", "无权反馈该课程术语。", 403)
    feedback_type = normalize_choice(data.get("feedback_type"), FEEDBACK_TYPES, "other")
    severity = normalize_choice(data.get("severity"), SEVERITIES, "medium")
    classification, root_cause = classify_feedback(feedback_type, content)
    course_obj = Course.query.filter_by(name=term.course).first()
    feedback = Feedback(
        term_id=term.id,
        user_id=user.id,
        user_role=user.role,
        course_id=course_obj.id if course_obj else None,
        course=term.course,
        chapter=term.chapter,
        english_term=term.english_term,
        chinese_term=term.chinese_term,
        feedback_type=feedback_type,
        feedback_source=normalize_choice(data.get("feedback_source"), FEEDBACK_SOURCES, "student_search_result"),
        severity=severity,
        priority=map_feedback_to_priority(feedback_type, severity, classification),
        feedback_content=content,
        reported_issue=content,
        expected_result=str(data.get("expected_result") or "").strip(),
        actual_result=str(data.get("actual_result") or term.chinese_term or "").strip(),
        evidence_comment=str(data.get("evidence_comment") or "").strip(),
        classification=classification if classification in CLASSIFICATIONS else "teacher_review_needed",
        root_cause=root_cause if root_cause in ROOT_CAUSES else "unknown",
        status="submitted",
        created_at=current_time_text(),
        updated_at=current_time_text()
    )
    db.session.add(feedback)
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "Feedback submitted.",
        "data": {"feedback_id": feedback.id, "feedback_status": feedback.status},
        "feedback": serialize_feedback(feedback)
    })

@app.route("/api/feedback", methods=["GET"])
def get_feedback_list():
    """
    教师端读取学生反馈列表。
    可选参数：
    - status=open
    - status=resolved
    - course=...
    - chapter=...
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    status = request.args.get("status", "").strip()
    feedback_type = request.args.get("feedback_type", "").strip()
    severity = request.args.get("severity", "").strip()
    course = request.args.get("course", "").strip()
    course_id = request.args.get("course_id", "").strip()
    chapter = request.args.get("chapter", "").strip()

    query = Feedback.query

    if status:
        query = query.filter_by(status=status)
    if feedback_type:
        query = query.filter_by(feedback_type=feedback_type)
    if severity:
        query = query.filter_by(severity=severity)
    if course_id:
        query = query.filter_by(course_id=int(course_id))

    if course:
        query = query.filter_by(course=course)

    if chapter:
        query = query.filter_by(chapter=chapter)

    feedbacks = query.order_by(Feedback.id.desc()).all()
    if user.role != "admin":
        feedbacks = [feedback for feedback in feedbacks if can_manage_course_name(user, feedback.course)]

    return jsonify({
        "status": "success",
        "data": {
            "items": [serialize_feedback(feedback) for feedback in feedbacks],
            "count": len(feedbacks)
        },
        "count": len(feedbacks),
        "feedbacks": [serialize_feedback(feedback) for feedback in feedbacks]
    })


def get_authorized_feedback(feedback_id, user):
    feedback = db.session.get(Feedback, feedback_id)
    if feedback is None:
        return None, api_error("RESOURCE_NOT_FOUND", "反馈记录不存在。", 404)
    if user.role == "admin":
        return feedback, None
    if not can_manage_course_name(user, feedback.course):
        return None, api_error("PERMISSION_DENIED", "无权处理该课程反馈。", 403)
    return feedback, None


@app.route("/api/feedback/<int:feedback_id>", methods=["GET"])
def get_feedback_detail(feedback_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    return jsonify({"status": "success", "data": serialize_feedback(feedback), "feedback": serialize_feedback(feedback)})


@app.route("/api/feedback/<int:feedback_id>/triage", methods=["POST"])
def triage_feedback(feedback_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    data = request.get_json() or {}
    new_status = normalize_choice(data.get("status"), FEEDBACK_STATUSES, "triaged")
    if new_status not in {"triaged", "in_review", "needs_more_evidence"}:
        return api_error("VALIDATION_ERROR", "triage status 必须是 triaged/in_review/needs_more_evidence。", 400)
    feedback.status = new_status
    if data.get("classification"):
        feedback.classification = normalize_choice(data.get("classification"), CLASSIFICATIONS, feedback.classification or "teacher_review_needed")
    if data.get("root_cause"):
        feedback.root_cause = normalize_choice(data.get("root_cause"), ROOT_CAUSES, feedback.root_cause or "unknown")
    feedback.updated_at = current_time_text()
    db.session.commit()
    return jsonify({"status": "success", "message": "Feedback triaged.", "feedback": serialize_feedback(feedback)})

@app.route("/api/feedback/<int:feedback_id>/resolve", methods=["POST"])
def resolve_feedback(feedback_id):
    """
    教师端将学生反馈标记为已处理。
    """
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response

    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    data = request.get_json() or {}
    resolution_note = str(data.get("resolution_note") or data.get("note") or "").strip()
    if not resolution_note:
        return api_error("VALIDATION_ERROR", "resolution_note 不能为空。", 400)
    feedback.resolution_action = str(data.get("resolution_action") or "no_action_needed").strip()
    feedback.resolution_note = resolution_note
    feedback.status = "resolved"
    feedback.resolved_by = user.id
    feedback.resolved_at = current_time_text()
    feedback.updated_at = current_time_text()
    if feedback.resolution_action == "card_updated" and feedback.terminology_card_id:
        card = db.session.get(TerminologyCard, feedback.terminology_card_id)
        if card:
            card.status = str(data.get("card_status") or "pending_quality_control").strip()
            card.reviewer_note = ((card.reviewer_note or "") + f"\nFeedback {feedback.id}: {resolution_note}").strip()
            card.updated_at = current_time_text()
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": "反馈已标记为已处理。",
        "feedback": serialize_feedback(feedback)
    })


@app.route("/api/feedback/<int:feedback_id>/reject", methods=["POST"])
def reject_feedback(feedback_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    data = request.get_json() or {}
    feedback.status = "rejected"
    feedback.resolution_action = "no_action_needed"
    feedback.resolution_note = str(data.get("resolution_note") or data.get("reason") or "Rejected by reviewer.").strip()
    feedback.resolved_by = user.id
    feedback.resolved_at = current_time_text()
    feedback.updated_at = current_time_text()
    db.session.commit()
    return jsonify({"status": "success", "message": "Feedback rejected.", "feedback": serialize_feedback(feedback)})


def create_evaluation_item_from_feedback(feedback, evaluation_set, actor):
    if feedback.converted_to_evaluation_item_id:
        existing = db.session.get(EvaluationItem, feedback.converted_to_evaluation_item_id)
        if existing:
            return existing
    card = db.session.get(TerminologyCard, feedback.terminology_card_id) if feedback.terminology_card_id else None
    english_term = (feedback.english_term or (card.english_term if card else "")).strip()
    expected_chinese = (feedback.expected_result or feedback.chinese_term or (card.final_chinese_term if card else "")).strip()
    if not english_term or not expected_chinese:
        raise ValueError("feedback 缺少 english_term 或 expected_result/chinese_term，无法转 EvaluationItem。")
    item = EvaluationItem(
        set_id=evaluation_set.id,
        evaluation_set_id=evaluation_set.id,
        item_id=f"PF-{feedback.id:05d}",
        split="test",
        discipline=evaluation_set.discipline or "pilot_feedback",
        course_id=feedback.course_id,
        english_term=english_term,
        expected_chinese_term=expected_chinese,
        expected_alignment_status=(card.alignment_status if card else "exact_match") or "exact_match",
        english_context=(card.courseware_sentence if card else "") or feedback.actual_result or feedback.reported_issue,
        english_evidence=(card.english_evidence_snapshot if card else "") or "",
        chinese_evidence=(card.chinese_evidence_snapshot if card else "") or feedback.evidence_comment,
        expected_english_evidence=(card.english_evidence_snapshot if card else "") or "",
        expected_chinese_evidence=(card.chinese_evidence_snapshot if card else "") or feedback.evidence_comment,
        negative_english_evidence="",
        negative_chinese_evidence="",
        difficulty="medium" if feedback.severity in {"low", "medium"} else "hard",
        tags_json=json.dumps(["pilot_feedback", f"feedback_type:{feedback.feedback_type}", f"severity:{feedback.severity}"], ensure_ascii=False),
        annotator=f"feedback:{feedback.id}",
        reviewed_by=f"user:{actor.id}",
        version="pilot_feedback_v1",
        created_at=current_time_text(),
    )
    db.session.add(item)
    db.session.flush()
    feedback.converted_to_evaluation_item_id = item.id
    feedback.status = "converted_to_evaluation_item"
    feedback.resolution_action = "converted_to_evaluation_item"
    feedback.updated_at = current_time_text()
    return item


@app.route("/api/feedback/<int:feedback_id>/convert-to-evaluation", methods=["POST"])
def convert_feedback_to_evaluation(feedback_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    data = request.get_json() or {}
    set_id = data.get("evaluation_set_id")
    evaluation_set = db.session.get(EvaluationSet, int(set_id)) if set_id else None
    if evaluation_set is None:
        return api_error("VALIDATION_ERROR", "evaluation_set_id 必填。", 400)
    if not can_manage_evaluation_set(user, evaluation_set):
        return api_error("PERMISSION_DENIED", "无权使用该 EvaluationSet。", 403)
    try:
        item = create_evaluation_item_from_feedback(feedback, evaluation_set, user)
    except ValueError as exc:
        return api_error("VALIDATION_ERROR", str(exc), 400)
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "Feedback converted to EvaluationItem.",
        "data": {"evaluation_item_id": item.id},
        "feedback": serialize_feedback(feedback),
        "evaluation_item": serialize_evaluation_item(item),
    })


def create_backlog_item_from_feedback(feedback):
    if feedback.linked_backlog_item_id:
        existing = db.session.get(IterationBacklogItem, feedback.linked_backlog_item_id)
        if existing:
            return existing
    category = map_feedback_to_category(feedback.feedback_type, feedback.classification)
    priority = map_feedback_to_priority(feedback.feedback_type, feedback.severity, feedback.classification)
    title = f"{priority} {category}: {feedback.english_term or 'pilot feedback'}"
    item = IterationBacklogItem(
        title=title[:220],
        description=(feedback.reported_issue or feedback.feedback_content or "")[:2000],
        source_type="feedback",
        source_feedback_id=feedback.id,
        course_id=feedback.course_id,
        severity=feedback.severity or "medium",
        priority=priority,
        category=category,
        status="open",
        owner="unassigned",
        target_pr="next",
        acceptance_criteria=default_acceptance_criteria(feedback.feedback_type, category),
        created_at=current_time_text(),
        updated_at=current_time_text(),
    )
    db.session.add(item)
    db.session.flush()
    feedback.linked_backlog_item_id = item.id
    feedback.status = "converted_to_backlog"
    feedback.resolution_action = "converted_to_backlog"
    feedback.updated_at = current_time_text()
    return item


@app.route("/api/feedback/<int:feedback_id>/convert-to-backlog", methods=["POST"])
def convert_feedback_to_backlog(feedback_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    feedback, error = get_authorized_feedback(feedback_id, user)
    if error:
        return error
    item = create_backlog_item_from_feedback(feedback)
    db.session.commit()
    return jsonify({
        "status": "success",
        "message": "Feedback converted to backlog item.",
        "data": {"backlog_item_id": item.id},
        "feedback": serialize_feedback(feedback),
        "backlog_item": serialize_backlog_item(item),
    })


@app.route("/api/backlog", methods=["GET"])
def backlog_list():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    query = IterationBacklogItem.query
    if user.role != "admin":
        manageable_ids = [course.id for course in Course.query.all() if can_manage_course(user, course)]
        query = query.filter(IterationBacklogItem.course_id.in_(manageable_ids or [-1]))
    priority = request.args.get("priority", "").strip()
    status = request.args.get("status", "").strip()
    if priority:
        query = query.filter_by(priority=priority)
    if status:
        query = query.filter_by(status=status)
    items = query.order_by(IterationBacklogItem.id.desc()).limit(500).all()
    return jsonify({"status": "success", "data": {"items": [serialize_backlog_item(item) for item in items]}, "items": [serialize_backlog_item(item) for item in items]})


@app.route("/api/backlog/<int:item_id>", methods=["GET"])
def backlog_detail(item_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    item = db.session.get(IterationBacklogItem, item_id)
    if item is None:
        return api_error("RESOURCE_NOT_FOUND", "Backlog item 不存在。", 404)
    course = db.session.get(Course, item.course_id) if item.course_id else None
    if user.role != "admin" and course and not can_manage_course(user, course):
        return api_error("PERMISSION_DENIED", "无权查看该 backlog item。", 403)
    return jsonify({"status": "success", "data": serialize_backlog_item(item), "item": serialize_backlog_item(item)})


@app.route("/api/backlog/<int:item_id>/update-status", methods=["POST"])
def backlog_update_status(item_id):
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    item = db.session.get(IterationBacklogItem, item_id)
    if item is None:
        return api_error("RESOURCE_NOT_FOUND", "Backlog item 不存在。", 404)
    course = db.session.get(Course, item.course_id) if item.course_id else None
    if user.role != "admin" and course and not can_manage_course(user, course):
        return api_error("PERMISSION_DENIED", "无权修改该 backlog item。", 403)
    data = request.get_json() or {}
    new_status = normalize_choice(data.get("status"), BACKLOG_STATUSES, "")
    if not new_status:
        return api_error("VALIDATION_ERROR", "status 不合法。", 400)
    item.status = new_status
    item.updated_at = current_time_text()
    if new_status in {"done", "wont_fix", "duplicate"}:
        item.closed_at = current_time_text()
    db.session.commit()
    return jsonify({"status": "success", "message": "Backlog status updated.", "backlog_item": serialize_backlog_item(item)})


@app.route("/api/pilot/report", methods=["GET"])
def pilot_report_api():
    user, error_response = require_current_user({"teacher", "admin"})
    if error_response:
        return error_response
    course_id = request.args.get("course_id", "").strip()
    course = db.session.get(Course, int(course_id)) if course_id else None
    if user.role != "admin" and course and not can_manage_course(user, course):
        return api_error("PERMISSION_DENIED", "无权生成该课程试点报告。", 403)
    if course:
        feedbacks = Feedback.query.filter_by(course_id=course.id).all()
        cards = TerminologyCard.query.filter_by(course_id=course.id).all()
        jobs = BackgroundJob.query.filter_by(course_id=course.id).all()
        backlog_items = IterationBacklogItem.query.filter_by(course_id=course.id).all()
        eval_run = EvaluationRun.query.order_by(EvaluationRun.id.desc()).first()
    else:
        feedbacks = Feedback.query.all()
        cards = TerminologyCard.query.all()
        jobs = BackgroundJob.query.all()
        backlog_items = IterationBacklogItem.query.all()
        eval_run = EvaluationRun.query.order_by(EvaluationRun.id.desc()).first()
    usage_summary = {
        "documents": Document.query.filter_by(course_id=course.id).count() if course else Document.query.count(),
        "evaluation_runs": EvaluationRun.query.count(),
        "active_students": len({feedback.user_id for feedback in feedbacks if feedback.user_role == "student"}),
        "active_teachers": len({feedback.user_id for feedback in feedbacks if feedback.user_role == "teacher"}),
        "searches": UsageRecord.query.filter_by(action_type="knowledge_search").count(),
        "favorites": StudentTermRecord.query.filter_by(is_favorite=True).count(),
        "mastered": StudentTermRecord.query.filter_by(is_mastered=True).count(),
        "exports": UsageRecord.query.filter_by(action_type="pdf_export").count(),
    }
    markdown = generate_pilot_report_markdown(course, feedbacks, cards, jobs, eval_run, backlog_items, usage_summary)
    return jsonify({"status": "success", "data": {"report_markdown": markdown}, "report_markdown": markdown})

@app.route("/api/terms/clear-pending", methods=["DELETE"])
def clear_pending_terms():
    """
    开发辅助接口：清空 pending 候选术语。
    默认只建议本地开发使用。
    """
    user, error_response = require_current_user({"admin"})
    if error_response:
        return error_response

    course = request.args.get("course", "").strip()
    chapter = request.args.get("chapter", "").strip()

    query = Term.query.filter_by(status="pending")

    if course:
        query = query.filter_by(course=course)

    if chapter:
        query = query.filter_by(chapter=chapter)

    deleted_count = query.delete()
    db.session.commit()

    return jsonify({
        "status": "success",
        "message": f"已清空 {deleted_count} 条待审核术语。",
        "deleted_count": deleted_count
    })


if __name__ == "__main__":
    with app.app_context():
        db.create_all()
        ensure_schema_columns()
        seed_demo_knowledge_base()

    host = os.environ.get("BACKEND_HOST", "127.0.0.1")
    port = int(os.environ.get("BACKEND_PORT", "5000"))
    debug = os.environ.get("FLASK_DEBUG", "false").strip().lower() == "true"
    app.run(host=host, port=port, debug=debug, use_reloader=False)
